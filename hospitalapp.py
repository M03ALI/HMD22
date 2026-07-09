
"""
Hospital Weekly Dashboard
=========================

Two modes, one shared SQLite database:

  • DATA ENTRY (analysts, password protected) — pick a day from a weekly
    calendar strip and enter that day's data:
        Patients   : currently admitted, new admissions, discharges, ER visits,
                     ICU patients, surgeries, births, stillbirths, deaths,
                     referrals out
        Capacity   : total beds, beds available, ICU beds available
        Staff      : doctors, nurses, support staff, specialists on call
        Ambulances : available, fleet total, calls responded, avg ER wait (mins)
        Supplies   : oxygen supply level (%), blood-bank units by blood type
        Departments / Medications / Tests (editable tables)

  • PUBLIC DASHBOARD (everyone, read-only) with two view types:
        - Day : each day has its own dashboard, chosen from a calendar.
        - Weekly     : Mon–Sun roll-up across the week.

Run:
    pip install -r requirements.txt
    streamlit run hospital_dashboard.py

Analyst password comes from env var HOSPITAL_ADMIN_PASSWORD (default "changeme").
"""

import os
from datetime import datetime, date, timedelta

# ── Hospital-local time ─────────────────────────────────────────
# Cloud servers run on UTC, so a naive datetime.now() is hours off the wall
# clock in Gaborone. Every user-visible timestamp (report "Generated" lines,
# dashboard "updated" badges, observation stamps, audit fields) goes through
# now_local() instead. Change APP_TIMEZONE if the hospital ever moves zones.
APP_TIMEZONE = "Africa/Gaborone"          # UTC+2, no daylight saving
try:
    from zoneinfo import ZoneInfo as _ZoneInfo
    _APP_TZ = _ZoneInfo(APP_TIMEZONE)
except Exception:                          # tzdata unavailable -> fixed UTC+2
    from datetime import timezone as _dt_tz
    _APP_TZ = _dt_tz(timedelta(hours=2))


def now_local():
    """The current wall-clock time at the hospital, as a naive datetime (so
    it formats, stores and round-trips exactly like the old server time)."""
    return datetime.now(_APP_TZ).replace(tzinfo=None)
import sqlite3
import hmac
import hashlib
import time
import base64
import struct
import json
import secrets as pysecrets
import logging

import pandas as pd
import streamlit as st
import plotly.express as px
import plotly.graph_objects as go

# ──────────────────────────────────────────────
# CONFIG
# ──────────────────────────────────────────────
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s %(name)s: %(message)s")
log = logging.getLogger("hospital_dashboard")


def _maybe_cache(**kw):
    """Use st.cache_data when available (real Streamlit), else a no-op passthrough
    so the module still imports under the headless test stubs."""
    cd = getattr(st, "cache_data", None)
    if cd is None:
        return lambda f: f
    return cd(**kw)


DB_PATH = os.environ.get("HOSPITAL_DB_PATH", "hospital_dashboard.db")

# ── Desktop-software downloads ──
# When the app runs as installed Windows software (in a native window), the
# in-window browser can't hand files to the OS, so a normal download button
# appears to do nothing. In that mode we ALSO write the file straight to the
# machine's Downloads folder. On the web this is a plain download button.
IS_DESKTOP = os.environ.get("HDM_DESKTOP") == "1"


def _downloads_dir():
    d = os.path.join(os.path.expanduser("~"), "Downloads")
    return d if os.path.isdir(d) else os.path.expanduser("~")


def save_to_downloads(data, file_name):
    """Write bytes/text to the user's Downloads folder without overwriting an
    existing file. Returns the full path written."""
    folder = _downloads_dir()
    base, ext = os.path.splitext(file_name)
    path = os.path.join(folder, file_name)
    i = 1
    while os.path.exists(path):
        path = os.path.join(folder, f"{base} ({i}){ext}")
        i += 1
    payload = data if isinstance(data, (bytes, bytearray)) else str(data).encode("utf-8")
    with open(path, "wb") as f:
        f.write(payload)
    return path


def dl_button(container, label, data, *, file_name, mime=None, **kwargs):
    """A download button that also saves the file to the Downloads folder when
    running as installed desktop software."""
    clicked = container.download_button(label, data, file_name=file_name,
                                        mime=mime, **kwargs)
    if clicked and IS_DESKTOP:
        try:
            path = save_to_downloads(data, file_name)
            st.success(f"Saved to {path}")
        except Exception as _e:
            st.error(f"Could not save the file to Downloads: {_e}")
    return clicked


# ── Offline charts for desktop software ──
# The slideshow and the chart-colour preview render plotly.js. If a copy of
# plotly.min.js is bundled with the app we inline it so those work with no
# internet; otherwise we fall back to the CDN exactly as before.
_PLOTLY_JS_CACHE = None


def plotly_script_tag():
    global _PLOTLY_JS_CACHE
    if _PLOTLY_JS_CACHE is None:
        import sys as _sys
        _PLOTLY_JS_CACHE = ""
        _here = globals().get("__file__", "")
        _here_dir = os.path.dirname(os.path.abspath(_here)) if _here else ""
        for _base in (getattr(_sys, "_MEIPASS", ""), _here_dir, os.getcwd()):
            if not _base:
                continue
            _p = os.path.join(_base, "plotly.min.js")
            try:
                if os.path.isfile(_p):
                    with open(_p, "r", encoding="utf-8") as _f:
                        _PLOTLY_JS_CACHE = _f.read()
                    break
            except Exception:
                pass
    if _PLOTLY_JS_CACHE:
        return "<script>" + _PLOTLY_JS_CACHE + "</script>"
    return '<script src="https://cdn.plot.ly/plotly-2.35.2.min.js"></script>'


# ── Authentication / security ──
SESSION_TIMEOUT = 30 * 60     # auto sign-out after 30 min in a session
MAX_FAILS = 5                 # failed attempts before a temporary lock
LOCK_SECONDS = 60            # lock duration after too many failures
PBKDF2_ITERS = 600_000       # OWASP-recommended work factor for PBKDF2-HMAC-SHA256
THROTTLE_BASE = 0.4          # seconds added per failed attempt (brute-force slow-down)


def _get_secret(name, default=None):
    """Prefer Streamlit secrets, then environment variables."""
    try:
        if name in st.secrets:
            return str(st.secrets[name])
    except Exception:
        pass
    return os.environ.get(name, default)


def _pbkdf2(pw, salt, iters=PBKDF2_ITERS):
    return hashlib.pbkdf2_hmac("sha256", (pw or "").encode("utf-8"), salt, iters).hex()


def make_password_hash(pw, iters=PBKDF2_ITERS):
    """Create a salted PBKDF2 hash string to store instead of a plaintext password."""
    salt = pysecrets.token_bytes(16)
    return f"pbkdf2_sha256${iters}${salt.hex()}${_pbkdf2(pw, salt, iters)}"


ADMIN_HASH = _get_secret("HOSPITAL_ADMIN_PASSWORD_HASH")
ADMIN_PASSWORD = _get_secret("HOSPITAL_ADMIN_PASSWORD", "changeme")
TOTP_SECRET = _get_secret("HOSPITAL_ADMIN_TOTP_SECRET")   # optional 2FA (base32)
# True when the app is running on the built-in default password with no hash and
# no configured accounts — i.e. effectively unprotected. Surfaced as a warning
# on the login screen so it can't be missed at deploy time.
USING_DEFAULT_PASSWORD = (not ADMIN_HASH) and ADMIN_PASSWORD == "changeme"
# Optional: restrict the master/bootstrap admin password to one email. When unset,
# the admin password works with any email (handy as a master/recovery login).
ADMIN_EMAIL = (_get_secret("HOSPITAL_ADMIN_EMAIL") or "").strip().lower()


def _verify_pbkdf2(pw, stored):
    """Constant-time check of a plaintext password against a stored
    'pbkdf2_sha256$iters$salt$hash' string."""
    try:
        _scheme, iters, salt_hex, hash_hex = (stored or "").split("$")
        calc = _pbkdf2(pw, bytes.fromhex(salt_hex), int(iters))
        return hmac.compare_digest(calc, hash_hex)
    except Exception:
        return False


def verify_password(pw):
    """Constant-time check against a stored hash (preferred) or plaintext fallback."""
    if ADMIN_HASH:
        return _verify_pbkdf2(pw, ADMIN_HASH)
    return hmac.compare_digest(pw or "", ADMIN_PASSWORD)


def _valid_email(e):
    """Lightweight email sanity check (no external deps)."""
    e = (e or "").strip()
    if not e or " " in e or e.count("@") != 1:
        return False
    local, _, domain = e.partition("@")
    return bool(local) and "." in domain and not domain.startswith(".") \
        and not domain.endswith(".") and len(e) >= 5


def _b32_secret(nbytes=20):
    """A fresh base32 secret for setting up an authenticator app."""
    return base64.b32encode(pysecrets.token_bytes(nbytes)).decode("ascii").rstrip("=")


def _totp_at(secret_b32, when, step=30, digits=6):
    """RFC 6238 TOTP using only the standard library (no extra dependencies)."""
    s = (secret_b32 or "").strip().replace(" ", "").upper()
    s += "=" * ((8 - len(s) % 8) % 8)            # restore base32 padding
    key = base64.b32decode(s, casefold=True)
    counter = int(when // step)
    digest = hmac.new(key, struct.pack(">Q", counter), hashlib.sha1).digest()
    off = digest[-1] & 0x0F
    code = (struct.unpack(">I", digest[off:off + 4])[0] & 0x7FFFFFFF) % (10 ** digits)
    return str(code).zfill(digits)


def totp_verify(secret_b32, code, window=1):
    """Verify a 6-digit code with a ±1 step (±30s) tolerance, constant-time."""
    code = (code or "").strip()
    if not (secret_b32 and code.isdigit()):
        return False
    now = time.time()
    ok = False
    for w in range(-window, window + 1):
        try:
            if hmac.compare_digest(_totp_at(secret_b32, now + w * 30), code):
                ok = True
        except Exception:
            return False
    return ok


def two_factor_enabled():
    return bool(TOTP_SECRET)

DAILY_FIELDS = [
    ("current_inpatients", "Currently admitted (inpatients)", "Patients"),
    ("admitted", "New admissions today", "Patients"),
    ("readmissions", "Re-admissions", "Patients"),
    ("discharged", "Discharged today", "Patients"),
    ("er_visits", "Emergency / ER visits", "Patients"),
    ("icu_patients", "ICU patients", "Patients"),
    ("icu_recovered", "ICU patients recovered", "Patients"),
    ("surgeries", "Surgeries performed", "Patients"),
    ("births", "Births", "Patients"),
    ("stillbirths", "Stillbirths", "Patients"),
    ("deaths", "Mortality", "Patients"),
    ("referrals_out", "Referrals out", "Patients"),
    ("referrals_back", "Referrals back", "Patients"),
    ("beds_total", "Total beds", "Capacity"),
    ("beds_available", "Beds available", "Capacity"),
    ("icu_beds_available", "ICU beds available", "Capacity"),
    ("doctors", "Doctors on duty", "Staff"),
    ("nurses", "Nurses on duty", "Staff"),
    ("support_staff", "Support staff on duty", "Staff"),
    ("specialists_on_call", "Specialists on call", "Staff"),
    ("ambulances_available", "Ambulances available", "Ambulances & Emergency"),
    ("ambulances_total", "Ambulances (fleet total)", "Ambulances & Emergency"),
    ("ambulance_calls", "Ambulance calls responded", "Ambulances & Emergency"),
    ("avg_er_wait_min", "Avg ER wait (minutes)", "Ambulances & Emergency"),
    ("oxygen_pct", "Oxygen supply level (%)", "Critical Supplies"),
]
FIELD_KEYS = [f[0] for f in DAILY_FIELDS]
FIELD_GROUPS = ["Patients", "Capacity", "Staff", "Ambulances & Emergency", "Critical Supplies"]

BLOOD_TYPES = ["O+", "O-", "A+", "A-", "B+", "B-", "AB+", "AB-"]
DEFAULT_DEPARTMENTS = ["Cardiology", "Emergency", "Delivery / Maternity", "Surgery",
                       "Pediatrics", "Radiology", "ICU", "Outpatient (OPD)"]
DEFAULT_MEDICATIONS = [("Paracetamol", "Available"), ("Amoxicillin", "Available"),
                       ("Insulin", "Available"), ("IV Fluids (Saline)", "Available"),
                       ("Adrenaline", "Available"), ("ORS Sachets", "Available")]
DEFAULT_TESTS = ["Full Blood Count", "Malaria RDT", "Blood Glucose", "X-Ray",
                 "Ultrasound", "COVID-19 PCR", "Urinalysis", "HIV Test", "ECG"]
DEPT_STATUSES = ["Operational", "Limited", "Closed"]
MED_STATUSES = ["Available", "Limited availability", "Not available"]
SPECIALTY_AREAS = ["Cardiology", "Surgery", "Obstetrics & Gynaecology", "Paediatrics",
                   "Anaesthesiology", "Radiology", "Orthopaedics", "Internal Medicine",
                   "Emergency Medicine", "Psychiatry", "Pathology", "Ophthalmology",
                   "ENT", "Dermatology", "Urology", "Neurology", "Oncology", "Dentistry"]
DEFAULT_WARDS = ["Medical Ward", "Surgical Ward", "Maternity Ward", "Paediatric Ward",
                 "ICU", "Emergency", "Male Ward", "Female Ward", "Neonatal Unit"]

PRIMARY, TEAL2, ACCENT, INK = "#006868", "#02A6A6", "#60D8F8", "#06343A"
LIGHT_BG, GRID, WARN, DANGER, OK_GREEN = "#F0FAFA", "#CFE6E6", "#F4A340", "#F85050", "#2E9E5B"
STATUS_COLOR = {"Operational": OK_GREEN, "Limited": WARN, "Closed": DANGER}
STATUS_SCORE = {"Operational": 2, "Limited": 1, "Closed": 0}
MED_STATUS_COLOR = {"Available": OK_GREEN, "Limited availability": WARN, "Not available": DANGER}
MED_STATUS_SCORE = {"Available": 2, "Limited availability": 1, "Not available": 0}

HDM_LOGO_B64 = "iVBORw0KGgoAAAANSUhEUgAAATMAAAFaCAYAAAByo7rMAAEAAElEQVR42uy9d7glV3Xm/Vt776o66cbOauWcE9mYbHLyYKLHnvkINsYevg+MA9hjoo0TNmCDGQcGbIIxSWBMNEYESSCBiAKhgEIrdL7xpKra4fuj6qTbt1uhW0hArec5fW7fe2LVrne/K71LQghU9pNrg7MnYe0vxh5T/uzH/jwwXd7bLBBFAgGCs4g2YDPQ5SNET75++eR+DkSjtwvlzZUfQQBV3ntAB1BSvK/4AC4DpcoP6MB5iCNQQt/liI4J5XMJgVgMGnA5GDP6noOP5dd8r4nPWz7Wjf1KTR6qyn6CTSow+ykBM9ZBqjUXultzZatQXMwAWebRSqENeBcQCYhWhBDwQXAEAoIvX1MK3MMDOQQ3Bia59yiliMY+xgDUwghsJAYMEHxOcJ4oKp6R5zkhBKI4ppP1SeI6Un4Z7x2R0gjQ66XU68nEt/JrQWqdYxLWe1xlFZhVdt8BtMkL2g8v1uIxAj5McBIFBOsRoyCA8x5E4XXB5gLQ8yFoJWIh9Mv3EqAbYHlllZtvv53lTptut0vuLGma0u228R6iSJMkCY1Gg5mpaeZnZ9i0YRMbpqcLEAM0gTqCLtlUUnLAYAtwFQHnLN57lNEYpXHBo0Ud5Hsz8b1lHdRfh6xVVoFZZfcVc0wyDvCoEiAEwAshhDF3UA3d09xZosTggNRDRghKCZ3y+W3g9qVlrr7uWr79vau4+vpr2blnN71+holi+rnFeY/WGjR47xERtNbEcUy33UZCoNFogPd46zj26GM486zTecCFF3DaSSeyNaqjgRpQB1EeEgFxIKpwTREIweGcwxizBszUnQazg3jjlVVgVtm9TsukADM3jFn5oWsXDd0tVTCzIJOxNA1ZCKTOBotCGY0F9maWG3ft5IMf/zhX33ADt+7ZiyQRKo7oZH2sc8S1Oi7o4WuFEAgiiAg+2AJ0RKGUIjiPtRYVIIoi4ihCJNDprFIzmm3zG3jk/R7Ekx/xKM7Ytpl6ydJiIO878DnNRg3vLSKCSCDNUpI4GcLXeitZDoC3yiowq+w+7V8GgRzwJT+ZALNQgtlYksBJwCpIJYRQspoecMvSIp/70pf59MVfYs/yElJrsJplpM7jjOApnquMoVar0e700TrCCzjnQAnGKEQEay3eFiwq1oYQAnmWkec5kdLoWKOUoLyDzBJbx6yOOGnLVp7w8Efw5Ec9kgjYUDMSAdZ66kYhhDHo8msOysFAqwKzCswqu4+DmR/GgQbsbABmGtADEAsjmpILZEBfijjYCo7rbr6Bj3/qM1z6jW/Szi1ORTgT005TokYDMZrcBVABFRmyrE/W6RE1ZxApEgTe+8KVlTJlCeBD4RI6h80yjIloNpt472m324AnijR1ZdDeI2mGyjISEaYiw+Me/gie9dQnc/yGDUQBWoJoX4B04Xr64hiEg1EyNelWV6umArPK7sNgxsjVHMsWIoy5liVVSwXSMpjfA76x43re/q53ct2NN6GTBkudHkQ1TK1BO8tI6k16eRGA15EB8UUwXhXuos3BBo+IQinBe4/1DhHBGFPE5PIcEaEWxTjnSPt9lChqzQTnXPF6BDSC9g4DKOcg71NTith6HvWgB/Li//6rnDQ7AzbQMiLiIVaMgZlfZ4VPglkFaBWYVXZfwrDynMngAg4lE5Kxy9TL8LJ1mUfVFF0PuSL0gF29Lu++6EN8+PP/Sao1AYUThRODE40TRUDhx1+TMAyzSxgVNwQp3ViZBBO/xrmTsaWmAsMyjyAeGWBu+Z1UUQiC5DnTUYxOM6TT5wXPfDYvetYziIEGSBJAe4dWGrwtGajC9/uoeo2ALxMjUsbz1KgGrQK2Cswquw+B2ThCeI+zFlBoVRSe9voZcTMmBTIIC8Anvnwx7/7gv7FzeYW8lpApQxCFRwoQEzWEkyCqgILyvST4EtDK3w3hSo3Blx973oGxKhUGeOuH8b4h/jKeafREWpGudGgpzaZmi3xhidOPPY4X/cr/5OFnnUyDIlHQXl4JG2amZehSWwtagICXSTAT1sl2VlaBWWX3HpgBRXxqvPrf+oKhlZX7VsFSGkIWCx2B1735LXzxyq+jphosdDqo2hRW6YK9DNlYWZfGePKzYEoDkAE/BKX1wGzyfu2iW/N9pHicLxne6D09Sori3kSEzuISGxtTqDwnRvHLT3sKL3j6U5jTijqIAbJ2n6lGrXzrUGZYwTPKtErhgFdAVoFZZfc1MPPOoZQeFZqJAgV9D11dVOdfcfMOXvumN3HL4n58rY41Quog6Kh0HsvYksiY+1f4fkFGQMYEoLEG0A5lB8sm+vL9VPl6Be8LYwF8n2ckcYxyDoOQiGD7KSbr8+CTT+K1v/PbbGpNMQ3UQGzP0qoZ8swWyQfxE2Cmh/0ElVVgVtl9BsyQgPM5RkVFyUVmCSrGR9CFsAh85POf52//+d2kUYSNE0IUsdrrkdQaWOuHQDZiSuNuoBqLaa0JsIs/gGXdYQnEuNs5yESWz1OBwi0dY2c2BESKZEOW9cHmJJFGISTB08hSZoziTa9/A2dt3Y52MKcRE8pWqTD+XiAiqPE2gArVfmqsKr75SQY1KaBgFGsqYl5EkAJLAT7z1a/yjn99P91aTNao0wmBheUVdNKgn/kyRsbwNmBdgi/5mkeHwf2ofk2CKuNh4zEzdXDwKkFqYukNnz8CMhWKNisJxTsZUcRxQtrtoCNDVG/Q954QaVadpV+P2J2n/K9XvYrLr7sOr2HVE/IywbA2ZCeDmNr4rbKKmVV27zKzQSzIBUcsEaGs/bfANXv3hYu/dSVvfs8/sx+HtKboi8KhUcrQSlpFnZeWIfMaQNIoU7lezGus1EEO9bgxMBN/IKpQBOfHM6DiCyAjSBm/KxhVnvZJpluknVVEG1pTDVaXFmk26nR7K9QRkn7OUY1p/vjlr+D8k05h3kAUEFO+zSALq0VNgphU7KxiZpUdGWA6BDkId8DKADQKhS76KhG6EHZ5Hz7/zW/wln/5Z7pGE83MkonGeiG4QK3WYHnvXkRk+DrjruVYx9MwJqZCEc8a3EsowOdgy0i8Gv093PEyGzx2xDT9EOiiRp2016fRmkJrzerCIs3mFN1en6gxRaoNZnqG/Wmf33rVH/C9m25ihaIo2DIq+VC+ArGKmVV2jwGZKxmDQR3QnjTOd2QosqNBwJb/M9aB1nQFOhBWgPd+/nP8n/e9FxvV6QtYJXjMBLAMco1O+UPscEVpxaj+6+Cgete//IF1aevG1cbiXWvNi8fhIATqXtFwEKU5Ne/4qz9+Hfc7/jgSPHMoifOscF21Lr54pMabIiqrmFllh2v+DuiYP8jfhoxOa3JX/JwB//7Vy/inD3+QtjKsSCDVEU4MjMWjhhxI/LoLwI8tj3HmFo7klX8wIBv8bfx2CEAUE6F0RFRvsZrl5FGEbzT4rVe9kqt27wQUKSE4QglkYSgG6auQWQVmld0bJ6Cs8A9FD+TgYux7yA2sQLj0B9/nHe98JyvdHlG9gTJmUBE2ShaMuXJFycJdJFTCAQmDe5XduqJ1qtvtUm9NoeKElTQlmZ7l9179an64fxcZgo8SvMvB2VJKqFp7FZhVduR9/UO6omMqZQEIZSFo+d9ME1IIe3p93vhXb2Z/p8vU/Cb61mJ9IJTZSj8OQsoXQPbT4GMFUEohWtHudul7j2k2WEr77Ot2efPf/z09oE0IPk7wxgz3hTxLKzezArPKjiSQqXEXsgxM+7Vgt7ZqPgRyYDlAB3jVn/wJS84yt20be5eX6fVTTK0+pFMDQPPicQJurNRKrXNbz/Vc73Zvm1EamxUy23GjSQbsa3eJZ2cJ9QZf+94PeM9nPkmO0CYEa4qYo3OO2lALrbIKzCo7MkB2cNIx0k8dxI5Eyt8KGQQn8Pf//ilu2L2XhX7K7sVlWrOzNGfn6Pezsr9y0kUc/PzTsHi9c4Vum0noZylSi6nPzNALgcU0Q8/O8H/e+14uv/F6uggphbqI9/Y+AseVVWD2U3QKZHAawnqAJ2NuZvEArwRXapJdt7DMuz76EZa9pzG3AasNS+0evTQfve4dLAAJo9tPmkVKU4tieisrNKZa9BeXyCXQ7mc0N2xkKctpC7z1n/6JleDxFOUaUZSQpWm1/Cowq+zHw9z8JJAJoIreyRxCBrzxb/+G1ESsWkcvgEVotqbxaVY0Va8BqcHP8lMSAPfWkacZcatOp9uBeq2QLYoilns9pFYnak3zvR/dyPs/9lH6FGUtnkCcVG5mBWaVHbaFMpDvgx/1WkoBWn4NRfPBg2icd0M12dUs48Of/RzX33IbuTFIvUYeQEWGfr+PrtVGrUdj7GvIxPjJr7GSQCHoWE6e0kojxuBdUZMnKiIPikwZMmX414//O7esrpBC8AjOhaE45OA2OC+DW2UVmFV2Zy/IsSr8SfZUMDLnLEoUeZ6hTKFLtup8cCbmI5/+DL0QyEpG5hTDglNNoRAx8F+Lyv3RSR///082ex0HZY8QRjMPKBSRermjtWEDHet409++jT6Fiz4ccFxZBWaV3ZMXaSH0rLUGFErMsMhTa8VHP/sprttxM7mAlTAx4FchSAio4W3won4sPjbKW96X6sbuMsNFDSWxR6wzoELxXbU2eKXp+YCPY775w6vZsbhAF0LOpJx2ZRWYVXZ3L8SxQtZxn28IODBU3tKRoe8CFsIy8L6PfozahjmsKkFIlQxPPIHSdQp2CGLrVdJL+Ok6lsWC9mgPOgR0KNx5pSNW+n36StMl8E/vfy9Z6a676gqowKyyI8nCZOQzrTHvSoUMQGkhAz75hYtZylNSKVxLLwdWgAVV/lxW+Q9u67GvO1M/prhz9Wg/9s1ABoA0+vRSyhZJCERKkzmLJBE+jsijiM9dehm3ra6yEkLwVdVsBWaVHWl36SAnSGu8h9wXTCIHLvqPT6MaTbrWDl3PcTdTRFCqqIz367CyID8d8TIoin8H7Vl+oI4RRoAGgeA9mAirFNRr5Ebzif/6PE6EUPUAVGBW2ZEDsYNO4g6hUF0NBXezgfCDG2/h5p23s9LvQxyDEkQVsyoL+R6PL8pChxwrUKrEToCaulPSPOu5pveVurSCaZbikmMcU8r5BAqKso0kIXeWnk1Z6nTpB88nPv+fWKpG8wrMKjuCLuakdxlkUp7apilaC1ZBJvDZr3wZX0vwOiZ3Hq01WhQieuKUhiCEMZ9S1gGCnw47kHUOtdko2KkxBqzFRDVMrU40PcPOlRUu/f73ySGEg8iGj1PmAzadSqW2ArPKJlmOCsLI2fEggSCeENxwGrgDVgNhX4BPXPYVloInVxqlE7wtpv+KF1Q4MKqlSnlrVd6G16IaxNXuWizsvqSaMcpe+rHsbDH7M1fFHFDRiqzXx5gYgyagWeykyNxGPv7lS+hRTLByuUeJApFSsdeNmGcYTYsPQyDz5a1axxWYVTYBavhBoN+t+aMUF5LAVTfcQNs7cq3xslZ/f737n/6TrMZcymKyU9HN6kQNM5VKKcR58jTDmJhaa4aFbo/Lr/ouqxQdAVFswHty74aCmesNNR7+Pvx08dsKzCo7Im7mkPWMV5wPSy0UrrxYv/b1K8hsXrQ0VYfuzjmh3qOUQmuNcw5rLSKC9552u81Nt95CFnwIAr7UipNBwbFIxbwqMKvsiABd2R3gS/bw7e9+D4sgeiS2GKpk3B2aK9uboigqpJPyHK01relZLrvi8oLFASiNMsVUzRDCwRUcgyp3oeryqcCssju2UE4XL0em7e/0uGXn7aAEVw7qrYDsTixwpbDW4pxDa40xZsjM0jzj8u9+Bwd0vQ9BgSpVTGRNgP8AuaZqIEoFZpUd2uWUSWpW1E9BuOGWm+lmKUGbISMb9V5WdijTWg/dzDzPyzo8RZbn7Nh1O20gL9V7J1z+g7mZsuZWWQVmla0hY2VdGRRZtUExqAWuvu5alNEluAnBy0ShbGUHP6YD8AKw1uJ9UdKiIkM7zbhlYS9e6VEcMnjU+EYxUBs54NKpLp8KzCpb56pzI0AbXUMECjC7ZecuMAbnPY4w0u8P1Sm8Izczz3OstRhjiOMYpRTOOdLckgG37N41ISMeQkCUOih7rqwCs8rWuzgG7qJShZhiyQgcMlHHf8NNN2K9AyUkpahg5WbeOWYWxzEignOu0JArM5wohY8MV99wAwpwA1m5CVbmhhvNxOtSJTorMKvszl+IjGY77tq3H680uQukWYb3fpilq+zgNi64ONpAik1DtCL3jsWV5aK1aaj3VmYzx+JiBwyeoaoyq8CssgMAa/w/o0p2Gf5dgP2LC+g4KurPQhEDqpjZYR57gaA0u/btL45zOU9zdFz9+n5lRckqMKvs7pkDchdAdFELVQKZrk7fHS/wNe77wPUMIeBdQKKYfYtLCKBlVPQ/IWUuo01lHMiqo1+BWWUHPRUHng4PoR8gqdfIsmzoJuEqjfojYSKapcXlAwiYiOAJ6w9LDlUioAKzyu6yecDaQK1Wx1o7vNBCCD9VSrH32PE7RMxMKYV1jna7O2Rew2JZkTGP0g9jl5VVYFbZXbThNTXmGo3cJVXq21cDOQ7bhXdFe9MI6AoxjMFGctDJ7RU7q8CssjXuzNprZExBexzIBtX/4gVvXfm4iivc4QJfEzPzQjniL2CDnzgPagBmw1IMhVrr/lcIdp81Ux2Ce5mBBVe4jYUqFyIBQTB4QoCaKAneBkkS8lKEUSN47bHikJ95QDuY3JFnHO7DGBD5IVCVWmgyGrdsFIgqoE2XT1SAH/ifAzATNcGcK6uY2c+4+RH7KjXsJy+3yb87KUo3BrGyIFWl093aQO5Mk344MHvpB/9WCFaBWWWVVVZZBWaVVVZZZRWYVVZZZRWYVVZZZZVVYFZZZZVVVoFZZZVVVlkFZvcVW9s/OWitGS/kHDxm0AsdKIZpFPVnB3/dqjfzyJkceICrg1KBWWWVVVZZBWaVVVZZZRWYVVZZZRWYVVZZZZVVYFZZZZVVVoFZZZVVVlkFZj+9tl5hgFrzy0q8obLKKjD7sdmdmaA0/gi/5udwkJNWnbjKKqvA7D5PwfydYGmVVVZZBWaVVVZZBWaVVVZZZRWYVVZZZZVVYFZZZZVVVoFZZZVVVoFZZZVVVlkFZpVVmmOVVVaBWWWVVVZZBWaVVVZZBWaVVVZZZRWYVVZZZZVVYFZZZZVVVoFZZZVVVoFZZZVVVlkFZpUd0u6MnllllVVWgVlllVVWWQVmlVVWWQVmlVVWWWUVmFVWWWWVVWBWWWWVVVaBWWWVVVaBWWWVVVZZBWY/u1bpmVVWWQVmlVVWWWUVmFVWWWUVmFVWWWWVVWBWWWWVVVaBWWWVVVZZBWaVVVbZz7KZw36F8coDWf/Xa/4EeAhqvT/cSfOE4VPV6A3lXrg/zL0kiCsOQwA9fFENuOF3E0CCQgSC+PLbe1Q4fHkhfwR2OFmv+iQogvJH9LOpifcr/hpEHfbyHbyGGlu4w0Mrg9MsgBBQw9P/U2GH+iJyD7y2TP5J7jNgFoDg13xQVV5qaxeiH33w4MdAQa37JSe+aADEj72tx5eLTxh7jcC9c3/I9SATjxu/8IOAE12eCAtegS+vHtFoAQuI0xhnEAk4lWMloBwl6Hm8HNn1dlcgSHtQg+M/DoNB8L4ANM/k+Rs+vrz3UqyP8Xf3qBJQRq/pg0cHUMEj5eOsFI+9I0BeD3C9KJTI8I8yvnC1w6NQCIJgAjiBvPxzdIQvxHsCi+54bR7iBWTtdafWJynhrn+4IKPDrI7gcTx8N1PGbod84bHfiBp7nj/0twlMLPzBa92n/GO5a/fjhMoBgVACvGN8J7CDrz8GBE55gngkKLRXw+M6AIi7ej86LyULHPu7Cuvfj5+PIB6nLF7Ai19/lYs/ECzHgSwMnuKHT1P4CQY2YGFBiucUG8HhnzofQvF6a9abLzfN0aZdHKHAgY/9iTY51Hr2d4nR31lUuqc2AXO4ByKUe6pag22EgjdMsi41OhDl7zR+5DKMX1QHAzQZcy+lfDXx9+p6CFC4gHfmXo0/SxGNu6sqQHmxOnFYIEeHXDucsgSR4jWCQgWDDmp4lYlXqLtxPzimw2PrJ/+/3v3geV48VocCXMs9SXtVfr4CkHwYrQ1fPn8ciKTc/SX4ETsL4+RA4fEFkAFWTW6Oyt+5jS3IekyiWDseV4CXDKiCL4FMyrWpIEgRChi8Dh5CKE7s8PP4H/u9HNbzy+8sB0KV3AHBUhPXth96ZHdWWFnd19zMUDILPxZCUoeirzJ67Ohwjp4taw705DbIQZxtNdpB76WN7VBELIQwks4OYcJDlQEwBxmugiC+jMuE4aUUxJc3RRChcLKkcK+CLwAI7vY9Y2xLrfn/hDs49v8BKIRxJhakdBlHO/dEWG+cDQYp14srnUp/ANhofIkXqoCbEtCGF86ROH8hMHhVvyY8EAagNdw4xzboMETiNVzlJ+0efOlQD+4HMelDOW/l6T7o5ck4sK3zoHtCTf6wEwD+jrhkuBOQGAYLxpUX9jpxsLWMcILr34vRi3An/i4yukDKDy5SfonB1V76WkECgmMAe4NvpoND8HiJQCDXCh9UAQaHwUwlHPr7jMfjBh/VA2pwcftBwkKV50wNmffa/ETxXqp4TlAIfjJ6MMHfB2/oh4yvSH6oSaa1ZgXeUU5k/O/ag5aADmOegaznDskEW5xY+IMd/N5KQB3mfbElhon7u7KRh7Eo5/rsS60PavfAZXvYYKbuzDce+1GvYXF4zdB3HLgyYzvhxBuMJQn8ml/eGyT/zvBkWfMAWZvylVEcMRAIZcZswAFc6VZGDmp4gvWkolDeY4IUi088KoziT3f2/s4EcocAVq7cMMZOijOnhhDkS14+iGuNc7RhLKx01waxNxcCiMeJGrqcB2a6Ch9WgirWj4xQ1ou6+2AuHvEBXW4wYUCSx9h+GHA0kckreBjqUPf+Rnq3E1d+Yk3KAYtZHXzTHjtJ6iAEx4/h/LpM7QgD2mGBmZQvENZ8lUP6zkENQa1YPWoYD/EHANWIAYxymMVPYXi4o8nkCz/+qgx1R8RsPQo+cbUWjxrsj74EMQvBAviIVg4hg5p4UgK4gIgDsSWbO7yqEnUQRjZcd+UJPdB9HLCxYl93ak0oaW1iyE+6KQCuBIQC39RY3GYsIRBKQLyDnL6E9RnaenCnAeUC+OLvWYllCoVFUBRusx77Pq5c21qVYCx3MWZ6BO8Pm4gENXQr16VcE4vFr7va5aCMbOTCjift7kkf6ojUmckBzMmvH94av3ZLFjYIBodyYdjisIXBus2GryO4shLLlWmDtTGre9PuiMmNV46Nl5UMooVS8p1BHHIQi4wAHRTzvYw5m+IywSKIMwQlhesl4bCZ2QHhgzXuZUAOADkvCisJVgxOFLkWMgVWKayCXPthXMVPlFnoCT9NyqQAQcqAdEEBw+BzyqAUxw+TFOsH9O/q2lWID4QQyIF+ufYGzu3QOVAFBys32lBusjLmZoV7cd3d7eWvBeIxds06JGJd0BxnduFAZiwTKa4DmZrcQ4V6dwrMDqbJJWEsqo8DVdBz63O0KJRWZexnbJv2fnQZhwBKCCI4IC3IB/3iZ3Kg66CbQ7sPndzRs45+CNhQfi4/WnSDQPuPe1bl4TA876WI3ZQxsYHL5hToYElcm6mFvTz/IQ9htrtM5PoYcWgPIhpbsoe7A2bjgDUAN39A4bMM64JCCFjv8N6Te0cWIq7duciudp/d3S4rytOrGbqxphNrUhHaPkeSmNSBVxqlIpwPiBicczTrDdrdLkprTFyj1+2gawne5YhSGBzBr9katOAHjE38YQFaMAaL0AWyAEl5gXeCpyGK8dXqJzfocF8oD5I7AQuyJr41XpodMQC00e9N6TwPHmNdoKYLyAvW4fKUKKmBL70wpXBZhjYRuByMwdscFZmxshdX/qvQyD2SAJC7Ih643mMlyIGRUfHgAz44lNZkLid4MFFtlL9yRbgsg9ArwatjYamfsb/TZWElZamfkXmFlQgrBqs0VkV4pfFSXMTB2cmSph83kN1RzOwOPo9HowJE3k/EbpyC2GdM5yts37+T6LOf5bj2Ai3Xph4sgkNE4V3hpg93yLt4H0QOyjGLJKuUGVUZBhFcKNmMiknjFh0dkyWGXhKz16XcsLif6xf2savbpq0DeZLQEcElNfIoIUWBielbh3MOJQZlIrpZho4S4jgmz3Nc3ieONOJHSQ6PAlWwU1DosD6Y3Sk3M3jqIpgsI06Ln7EpcRLRz/sAxNpMsFU/ZITqgFhgwZB/fPdDX10C4oWgDrzHlfV/5b0RQ9JIaCQNklpMw8Rs27yJE044gZNOPJ4TjjmWTXMz1Abg5gNTSnDeEyslScmAbJ4SRcno2NiABIcoPXZQCsxwMr4JqGHdgtwbzOxQF2buPEYrQvC43KIloKIyjuUDaI3WEU4P2VawQKahC/yoAzvbffYuLtHPLRLFeB2R5hF9K+g4QVRcxtWKPgIbBBcCPgS00qNAzN0Ak3s6mXlHQOiCHzpdqjy9AYWlYAxJiMBqWjls6Hla1lMLFh+ycnFEI55wN8BsUDS6vsM89h1V6QYrGX5npx1p3iUNAasULoo5VsecbSLy7ceRJjErSnHj6jI/2LeHHy2tcotdwNdrMDtN5gImMuQ+xQcHwVFvTNPvZdh+SpLEaF9eJGGQ6fX4oJGDQNTBNhe17vFXZCFAkpC5lJAkYDU9FXCmifP5ZGRExpmrGsXxSnAbFDP/uO7H2apSRd3f2nsRAa0xYor/e8F7h+p28L1VtAhX7r4FfdW3ibQm76XUjObMk0/hfmefzQPOP4cTjz6GWW3wEFKgBuJVEfLJXSDRRbw3uECsCx5o+11MvXbAhnJPxs0Oj5mJDKsLFOCdAx/QWkO5o/d8QLTQh7ASIJcC1G7b3+eanXvYLw3SKMErRZo7erkFUcS1OrUkIeuXJ4TyvcJYAFSP16rdsVt8T9hhtRIJ2DIalVg1ViKgcCgSB3NumVN372DLf3ycMxb30rIdkpBhQ1peUnrIEu7W58cNC3HXbPbDdqKBi1lczEVNnIjgVMCqPk4cggEVE9Dk3uBEk5mIrJawM8/pTDdJN8zznf27uWLHzezxjtXYsKo8vlaHqM5Kz0JUx4kmhECSJLg8R0JeuuChKBwu3esBSK1tEbtLcSOt6S8vMTO3kfbqMonREALGaFzwSPniXkZxwvF43SA7e6/FzMIdX7MiRUnQ+E1Kz0ZHxbH21iEhkBhNXQxic3y/j+QZNeAh97sfT37MY3jAuRfQ0kINsAHqgjgfQl2JGMCnHiMgkQJvCUqGzMyXaQA1jBDfS2C23uOCQPAyqqgoT7At+9iKvjnCKrAMrOSwY3+Pm3btpWM9Mj1DzzToofEIQQmiiosnSx15bkmiuCghDSWhKE+E1gVZCI6hCzL+GSfque5FO9TxDUAwBRAZV4CZLmvufFAkHqbTBU7ZfSPbP/ExzljaS9N3icWRK1ssCjvemnTXoWyYwRnW97GmZ1bwBEIoGFko45ygcMqi64EsZEV21Rsk6NL9iwhiSFFkScyqViwoYamekM9MsyfN+MrN13Ft6HKrzQmNafT0PAuZpx2EzAeo18rP6IpWp2G6RBPCyM28awt+8v/WO2q1Gr00J4o0sTHkaVbk4XwBZgWIhxLEZKJTBX/vgtkdnmHnhpvPeEx5cNORop92wXuMiYqYWJpRU4aZZhPX6xGFQE0p0pUVts5v5OE/93M8+QmP58SN8xigBaRpTk2EpjESK3C5LY6SkXXdTH0PlGbcKTBbCxLjAOGD4H2xv4suKtTzsqwgBzrA7j5cv2eZWxZW6EpM1Jqhj2Kln+K0BhOBMlg8uSsWjdGGKCrW8nD3LdOc3k9mZNb7CveVQSKHAlNfFtuFsmShqC7LyyybJvGBmWyRU3bewNH//lFOX9pLw/aJlMOKQwPGy2H6un7kdq7dqYRhm0LwMspkCgWgiMeSgirWhFIGEU1wgZAVYFBLGvS9J1MaV6+zhGdVBN9ssVfD4tZ5Lr7uGr6/czcrUZ29GPz8RlYiRR4Z8mGhSnHyVSgvBa8Z8ca7D2YEhWhF7h1REpN2OuA89WaLPOujzKiZP5THaajUEdR9IpN+Z5iZUmrI0Lz3OOeKCydPiZr1gqH2+ygR6rUmWEevvUojToiUkHf7xCKEPCNSwkxritOOP5YXPefZXHDccTTKmFUCEgN5P6NWi4etYiPsKhMA94BsxmExszKxXu5wgVwVma8UwgqwmMN1uxe5ee8SC1aIZjfi4ib72xmpV7SmIsQ6vM1xASQy6MSAgtxBluXgwyieWJ4QkeJ9is+khmUDP3agCodfHhDGm8/FgWRIAHF16i5nU1jhtF03sf1jH+XMxb00sxSDJ3hbugr+ML+DP8jvZQha4yA26K2kjPFJUHgCGR6risRFEEHpQtLIZylaKUIImCghV4rUQ9CGVTS3pRnm6GPpTs3z1V07+cqtt7NDKZYaCf2pGm3lsDKQPCo6ARSCcgWYWe3vlIzReudIBTAqor1vH42tW7ChAODERJggtFdXiBIzDGOMl5eEe7xq6siEObTWOOdw1hYsQATRGmMMRgQdPM7leD/2WqIL1iaKLMuIYk3eT6nFCYZA8J4sy9jYqOH37eF+p5/Ci375V/i5006nDuJzRyPSOJtjjFoDZpRgNjiO4bASaBOPtSEcEHcYv8jW1HWWsRM/DFZjHUEbrAjdMkC4DNy4lPGj3fvZ2+1jpufpS8T+bk4WIuJWA28g60IDhxIHSuNDwAaPU6DEIBq0KHzwBdWjkGURVQT/QwiIig9xINZnbfcEoB1uRnS0kPplr3NMw2bM9Jc4bdfNHP2Jj3D24j6mrMNQ0lMV8GXKe7A0/ESQutgBZbzZ+2AX/tCtnASHgFmjVKHIFKRG8JiiaFkUmQg9gb4WrAIjoMQROUtNBNNLixIAAWctsZgii5lMsavbp11vsTTVYm9zii/tuJmvL+xhcbrOUj2ibaTYKFHooDBBo22REMi1xyk/dLV9+Y0Gcb9hkH4QipiId5XdFVFE6ixpnqO1RiG4LGeq1SQts5pri0b9RKv0vdeTFOTQf/c+IAJKaZQaxJ49znmCc8Sq2HgICusdLgha6yJr7Rw6jsn7PZJGY1iTZ7OcWj2GfpctSUJ75062zc/xqIc8hJc+//m0iiQBCYielJco1uN4c74f/y6HALM74ZJKPwQMoJ2f6DNyZVW+CmMqPWogDuiHdSMaTeY1PSWhC9zi4Ie7c25cXGHFByRJyKWg5kXRqymCqIwuMplobF3bCnOw+1FG6l51Iw+TlQ0KSscvOilZSMPmzPd7nLrzJrZ+6kOctbSHpk2JgyN4P3mifemH6yJT7EplizgHHQbZzjJTE8Z8dq1xqUcT4a1FTTcIPiezaRnDjMm8xwWPqif0lKdTT1icrvP163/Env37+cE1N9Lp5vSU4fsLe9DSIAspW+uzbGo0mHKWBxx7LMdv2sD555yO66wyhyHOc+oK8jyn4wJqdo49qWd5epprbc7nb7uZHyjL8vwM+0STKoORiJA5asGAOFzsigSUK+N4orBq0IQfiPygdKO4WJwq/q+DL4PnagjU4662rNHQO2jW9G5mke8r9wcjMeNu+MGOg3FgrGcqScizHr7X4fhtW/iN//krPP6c86kB0yAJgBusOUswelgUHnvDsJhv6Or4UbVuWCO+uqYwf7zg95Bg5ko3JBI1EKHC+hzRZSYToe88omOWIfxoBa7e2+PWbkpbx7ikRhZcCU7lnhlUKR9TKif8DAt3F0zHT6T5danhJdgSzHJO3XkTWz79Ic5c3kUr7xEHT/BhWAOmQrlYZARmVhWso5aDeF0sCiWgCqVaRVELmDtLZBIwDYhjuou7iBsxqc9RjSbetFjs5yRz8+zP+nzv9pv5j29cwSW372Qh0nSDwumIC85/IMv9lDA1wxkXXMjO2/dy+803UUfo3HYb/Z03shmNZ5Uz5jfzyDPO4tHnnEU97bFtdor9O3fSSpqoKKYfxeyODLckCf+1bxeXLSyytHGe5ahGxzqipIlLM0RDXzK8QOzK9iIMVo/q4nSQ8thOglnkB66jIgiV3Z2N3Ct0EPJ+ShQrpps1VvbuYUMj5jmPfzwvfeZzaATQvT7z9Zr0VpZDfWZKPK5sF1MYrybBbBzI1gOzkgyvC2Y2uNLNVBPlzYMMRCg7Bk0YVUMhkOHpo4IHFoGbFxzX7V3mlk5Ox9TJkzo+isoArh/r4RqpK/ysWxjKYB8azE7ZVTCztWBWlCmoQgus3GxQDGNXEoQ4FwhlIlwJXhe1ZYInBI+OIbU5kUS0V5eZmZ/GhpwewgoR+0MTP7+diy79Gp/8xje4HctqXOeUB/4cpz/skWw9634ce/rpBO8wBqzPSYNDGnU6aYZWMTNRjaVrruPaK6/g+quu4PuXfIH5lWXm8pRHnXMOjzrnbE7fMEu8sozptmk16oR6nd3Oszg1z3UI7/nOt9k7P8vS7BQLHiTSZEBWZlYT6zFlN4ST4hh4pZEQDnC1B8wMSg+kArO7DWaNuIb3nm53lUYtJtaC765Qc44ztm7hz1/1h5w8M4fLMmbiWLK0j9ZCZCK6vS7NpDV2PZTaeGNC5cMM2WSgnvV6RSUENwYyY2A28F7weG+JgqBVBAGch56GXAhLwDV7Ldfu2suy02S1KXo6olfukIVufShhb1yrrCzCDBWYDfTBRhdcAWZ1l7Ohd2gw8yW1Nb7ITHqlhu6rANqp4X+8KoL0QQbdBp4+OVoLRgzBe9K0i5qqs2pi2s05vvCj3fzTZ77M9blFbzuWRz/3OTzkCU+muXE7+1LoN+dYdZakqZC8yzTCH/7hq3jx776M6W1H4ZI6C/sWOW56njjro9IVav1lvvHxj3HFJz7BTd/7BsdpzdMvOJfnPvwhbEq7zPqM5Z27mN+8jSUfc5MI7RNO5H1XfZtvpj0Wp+q0CaRJndTEEAyR94VLKQEvAasLUPOlKzlQyR3EykZB/QrMDifEkvUyWq0WeZ5ibUZkFDo4TN5jQ1KjYS1v/P0/5CGnnEpSupw2S2nEyeDkTCTBfNkVO5QPGo/Zj1B0DZgNujF8XsrPqDV13+Dww673iEJB1OVlz2RM6ABX7nXcuNhmX6ePT1q4pEnbBTICJla40v1ZC2ahbD647ye372lA88ON5FBgtvlTH+LMlV1MlWDmXVHv5cbArOi3VBPZutE598PmbwhDBui0YK2lIQlaRaQq0GvV+fy1P+QD3/wOX15oszK3nef8zqu44Bcej603yXSMI6GbBzouEM/UabsOM1qx82vf4j2vfBXHPfWJvPB3X8GN3RWS1gzSz1F5Tj145pIE2ivMCVx3+aV85r3v5LpL/4szEsX/+8TH8nNbN7Ox36OeOXKnCRs2clWWs7B1M5euLHDJrttZnZ1iRy8lbc1jxWCCx3iLDoEQHE4VYOa0LmJqXg+zoVDGFCs7bKtFdfr9Pl4cURQhRkh7baQ81xsbDRq5569f/RqOn9/AlnqdJCB1AZ/lSGTKa2DgbcqwhGOcd8lEbmN9MFMDNVfPgRI1ZYXN8I2yACEGF8NiBlftcVy9a4E9DrL6NH0TDYEMpUopmLJMbrxwr0z9ItWWeEeJhDvLGgJqDZCNM+1y6ImMzrIq3axarUbmoUfCompwazLD333xCt588eV8ZdXx8y/+bf7qM5/nhMc/gfbcHJf/4IfsXemw6By22WB6W4vlNKXvPanNyTs9aE4zFdfIOx3mWtMoo3GxwUxPY+M6e1b73LB7kdv6nvnz78+L//rtPPv1f8YP67O87mOf4m8++0X2RC1WVEJcS2B5hROU4riVDg/QdZ5xwqnM3rqbk03CbGape4cZVLpL0YBvgsYgxQAUXw5BKQF8whWtlt5hWT9L8QR0FNHPM3qdHuiY2uwceZxwy/IKu/pdfuMPXsn1+/aQU3h0nb5Domii5GvSvTyUrUUrJomRYo3sc1kYHpVswQG5gi6ERQg3LDi+f+tuloKiF9XIooQeilxARwqtwVpf4mZZbYxG0CWoVVg2TFWXrs/BWlMOVUukBgpoY6KGBZD5dReAHsvk6QCd5TYmmaLTmOf6aIo3fPly/vIH17H77Afx2x/8JA94/m9ye9wkTM8RvPDh3/8D3ve3b2O6pulkbW7a28U0G8xPzdIgZro1BbtvB6Op1Zv4viVd7mBzcDZQVxH5/kXe+eLf4N/f+15kag6/7VhOe8qzef1Fn2PDo3+RD9x4Oy9/94e5KZlmj0QQx2xwgdlde3kwhodazcsvfCgbb7iNozodZrMeic+LHJdSOGVADDrEaC9EvmCuw3YxoXTHKxfzcM35HIdD64g4qWOSOt55epkjFcP0tu341jSLzvOyV7+Wy390IxlATWMn1Ij9UN5zcJOSZPmxzP/6w0HGo2cyYmITAngBtJhikASKHMJOB9/fk3P94iqrKiJL6uTKkCuFLyl9CGWbkQ8EXwbZbKmU4X1RoOcmK/kru3tAqIYB7vEky9h9SckHksh6OK6tSMSYWo1+bLhFhLde/F+897rrOfP5v85v/uO/0D/hDFY3bqfT2sD+3NL3ng33vz97L7sUsRlJpJiabuAcrCyt4hCs0bBhHtVssNLtIhbmWtNERoF10O/z6Ys+Ckp44lOfwqrSXL/SZimZwm85gV/7i7/jF3/3jVyeB17zgY9yi05YFk3Wzzm63qC2cw9HLa6wbfde/uAJT+To5RU2ttu0+j2MtVgPFo0NBhdM2WJVaMcOjsFw2YlUiajDtLhWQylFv98n6/YwxiA6LudFaBZXV1nNLUxN0VXC777+9Xz31p2sFGo5wUlRUC1Bla1OkxqFdyVVqF/zmteMUpx+nMH54c9BafoQ9gE3rnquWVhid99hWy16ItiyTEMphVKFt6sDGFWkbqWUkkH0UCmp7B2fnAnxsxQrK+n1oJBxcIVJSYsFTxQ8deuZby/RvP4HbE7bJN6iKQBMTWi4y1ipThhLc/syRRrKmsFQuPkmomstHa1YnZ3lLz71KT664yYe+JLf4CkvewXLqkknarLPKVwtQdcTsDmm3eFHV3wNvX0rJ5x6Ct1Oj+mpGpnL0EYzs3GOmVNP5UEPfyQhTggimCSmn3ZJ8Gz0woff8haYafCEF/5P9hiNb06hiem1c6Kkzslnns2Go7bx8c99ih3XXcf9zjmHeggY62hoRQ1Hw+aorM9JJxzP7r176OQZWaRJdQ1fa5G1+8SNKazLi0OBK9zQoay3Hm0IhAqV7jYz86XXpdBKE2wZ+RqUAuWOpFkvJMGsI+D42uVXcO7592NuukUEr4tE8M4X10IAl2UoTAFyjFRa/MBlHCr/TkKfKopO1SitOIytUM63VKS2EExc8nDtwiI7uj16SUKqo0K7XaSUOA5oHwq2MJwfVBYnjtWLHIk2oMoYU9sb24UGACZ+6L+6PAdTFM3muUclDXpBYWfnWJ3bwJ9++CN8/KbrOOfpT+Npz38hHSfkSiNRjdZMTN96ut0+3nse/NCfg7kZLvmvz+P27WGryYlW9zMbchoGesFy3s8/mNRZ4hCYVYr+bbezPYmp9Tvs/NF1sLCXC5/8eBZtBrU6vczjnaLebGGjGjvTjPOe+lRe9MY/55L2An/20YvYIcJKEtPDIXmPDVi2dFc4ob3Kc047jdONYb6fkuQ5tt1heus2UudBIlB6srjah+LmQrUGj3C4RPsiyaJ9MQawNTNHv5uSOo9PYjIT8aO9u3jd37yFJWDV+eARxBjy3JH3M3RUAzWcvHjHhb0DZvbq1752xJCGYCbDWxDoasLtFq7e3+OG1S6rOoakSRrGqODA3ZHi4lIlqjKWZ0AmpbTXym3/TC6CcocZAvwwYFmUGgyY2dS1P2BTycxMCMV+M+jnZM0uNTaEJBCQUscKG9C1Fn2EBQ/76i3e8PFP8KnbdnPSE57Kr7/6DXzzm9/nsi9eynlnnI1zjr4r1IFbClohMB2DSQxPfuwvsKlV56Zvf43rvnEZut1h88w0UovYt7jA0Zu3UM8cS9dcR+emm7jl6u8yLYGjN83Trsc8+AmPpr5xI6udnKm4RpJZVJbiszZXXfUNvnzFl3nkYx/FlqOP5mOf/iTX33QjZ597DrOtOv3lBZraE2V9ojwnVprmxg3sXFpiKfekJiYPYK0FPRhOHIaHRkLBHNQgZ1Yh2mHtpyUrKuPuMuEppP0etaSG864QkVSaxnSLm2+7nV27d/OEBz8I73mts/51oEjiQgux324HU0teN44PYVj7OvRDJopmzSiLEEaMTI0mGnYdpBpuWfZct2uBXqOOqU/Rz0JZw0s5wLV0TcOo2M2PtYiMj2SfvJqrBXHQnM2acW3jYoDrpQLGV5hXRQFiCKW+nBVcEILSrKBZbho+9K3v8Lndy2x8zJP59Tf9DW0Ul1x6Bdd/6RL271vi+S99KThNrVYj3bsf21niq9+5gt7Sfq78whf54kc/XGzF1kHq4ejt/PfXvJoTTzsDWVnl79/8Vm78zy9Atw0JYDQXPOzhbD7uaPbdcB279+zhlDMehF3tEgdPI9F86eJL+PSb/gxOPIpnPvsZnPHkp/BsES56wxv4i49exB889Umct2UbYWkPibOoTpeo3mIlMjz9tNNxO27japexZHv0Iz1UYfFla8tAg96UkkZ+TStdZXfdxmNaI/l1P1yvSg0wQNOzGR5NsmGOL1xxORd96Xye/oiH0lLFuen2UmqRodZqSciL0o31QeLASJpRa6iblFu+Q5EBS0LYncGO1Q5LXrCSkOWQZ4F6YsD6iSycnxidBl6xrkTNoRqff6Z2NhmpmA6xfc2UlrDedjjsVRvrVR2O5J48qs4FjDIoFbE3C/Q3TPO9xd38n8u/jL/gMbz0H97DDb0+1npe/Duv5A1L/5sbP/dJ3tpf4Xde9rvkt9/KpR/4AF/+wHuhVrqwKdSPP4UXvOTldJ3n9r27+OwH3sP73vRm/uKtf8OXPvlZbvzCFzjraU/i3AvOpxbF/OdHPsS3PvEZiAVUF3RC7fjz+aPX/xn16Rof/cjHueR9/wzbj+J3Xv/nLKUQNWa48Jeew20/uoUr3/8eLrnuJo4550w26zqJExIlrNx+OyeeeAJL+xZ47NZNdK69Hp92iKZnaaOwSgN6uKGq8clP1WZ6ZIBMCiLjlJqYvJUkMWmvjzKK3OXUanV6/S7dzLNtZpa//qd/5H7nXcC22UZoAI16UlwR3iFRURwmpWCiGgOq9Soh9Gtf/dphs3OQctAsgkXRQ0Jb4Du3tblhuYttzmBVTOZAK1OQvKFUTKFNPpDk8WX6m5KcjV+sg7SrlK0LP8vp8SKUMzZkdqxHzQRPY+BmXlckAOKyMHRU1V+2NQ2GrZbT0QfDSEQVTeYutVBrkjWmuSHP+IcvfoEfiuHX3/L3LM5swDZncVGC2MBjHvNorlnYw65LvsSWrVv4j/f+M9/+wud58MMfxC+96Fe5fXmRlV2LvO7P38wxZ1/IlpNO4bhTT+TGfbtZ+Mbl/PwjHsk//NHr2HDKybzg917BplNOYHZ2jkc99GF87j8+QYLnla/5A+Y3buD7n/g03/7Otzjn/DN579vfCvPzvPzP/5z6/BbyuEFPGUzS5Lxzz+VrX/kK1195JY+834XI6grT2qDwNOfnWbp9J1s2bUBpRX2qwS37dkNicBSlGgFdVPurckpWKActi6oA7Ygs5MF+GoobxeQwa10holouWBNp4lqdfp7jrUNbz46bbuTBD3kITSUYeF0IHqMEvC3PTxgLn8gwEnYgsJZUcDDezOGLCUt4+sAt7cCt7TYrHoKpoSQikohYFJQUHhnMTVSFWONAJaMsTnQyVisyZG7lULmqbPGwVpAXhVVqTPVhspctBEEZg/OK1AtpZPjcld/g8ltv4mkv/jWOOv0kutbSyy2p1/SimI6J+JVfexEPfu6z2TY/ww1f+wpHnbiJ+z/6AZz3yAdxe38FanXmjz6Gdl1xO206iefM884AZ+nt3glZypmnnECmUlYkxU8lpM2IaMtG0iThlLMu4AEPeBCPf9aTWfj+ZXzrykt5yC8+ld9+65uJth3Nkm7Siabo2CYrecw+qfGMl76cm43mLz74QeKjj2eP9aRK6C/sY1uSUNuzh63LSzxodopTY2Gq26ZuMyJXaudTjMTzZSeEU35YfV7Z3Q/8Dz2AwfWuPKE8tmIoRFudpxYndLt9VpdWqDVahMhg44jLv38VX7/qB7QpZoQ4UWQ2Bz2p2CFhVANLONBlUZMTxwWNwRORoUMX+MYNN7KqNboxTWaL+jHlwFtPpGVNhmEkpxIOieDVAlqbBTrUQimC1escxzU+qR9mMcf63ZzH2UAyM03bKH6wtMQnvncVctKZ/Pyzf5VdqWVqw0b6zuGV4KKE5V6PudmN/PKznsfHL7oIYsPtt97G2/7yr3nBi/8XT/2VF0JS4w/+6k3sx+KnGvTxzG3eCElMZ3GpUDCNNM3Zafres2It/37xF8n2LXK/Jz+Fv/2XD/C//+DVfPbTn4R6wqVf+gJPfsJjac3Osdy35LpoJFeNmMwk+Kl5jn/Qz/OgZz+PSxZ3cvG115E3m1gxaNEY0Uy5wBbnqe/aybPPPZ/taY+ptE/kMtTYhCdBkEEXSlVndoQXtB+uxYFajnN5oRmXptTrdUyS4L3HukDHeVSryZ+97a2sAO3yDEUmKepTw50PsysEcueIEWIMkmo0ii5w1WLGYtxkRddwEhUqn64orjVKMVDedYzGcInIEEF1WWVuylKngQflZVDCONmC8zMZM1vTaT9irwrxgvJFS84k6DmcKkZ4EVTRb+iLI+qUwxk/1KuPlMEYw96sw+pckw9+/9t8G8WzX/smdje34OY2s2tlBR2B0h4XLJGOyZf6/M1fvJXrv3El0bHH84zf+j0e9aJXkRzzAC54zLN4xO//Pg981tNJTY1eJgRVZ7WTg3Mk9Ri0Jw2OXt9jomkcTU4690Gc8YIX8pTf+C1WZ+Y55om/yP/467fDGeeweN31vP4Vv8PK3iWMiogB5TyZ7+Frni6KTtziSb/+W+htJ/N/P/dJ0laNJKmhnMZ6hfOGWmY5xik23baTZ556BhtdikpXadQF7/tosWjvqEtMyBTi1UEmYPo1t7UnTt1HJmfe+zHf8fiZCmriFkIhp+5CkWrJSwFMyuJ5lSSsZJZVFfjHD32g6DSCEGwo8pNejen1Te7fa8NTxksgijSBQLCgtKKXwUoMO5ZW6ZqITEVEzqDCCCgHbuPBWm3WZRsTiYJqIYxnf9wh/q7CZCbTDxR0y+ne5XAvBgKEQYoxYGWtLO00hZkZ9vici666kpMf8RQ2nH0B+1WdxdU+yVSrmEJkPVlmmYpiVm7byfVfupRjHvggfvv3fxdlZlh0iof98hzLWc75j3wM9ekWq51isIonjOZ3Ur65ErwThAQRYXrjJp7233+VdrfNc37t14kkgO3zl+e9g0vf/y987H0f4FuXX8bDnvHf6WSBTtpDT2l6eYZYqKkYMzXPOY97Itf889u4et9+mh42RxG6VsenPVynTSvUOTap41zg3KkZlvp9bl7aR2NmBmsLiWiHQYmuFuA9vb5LyXQAYwpxCV+2/qjI4HKLNwofcj5/2aW86JnPpCmmCIr5QXDMrwsla0fXqUICaEANhaChrwi3LTn2rrTH5IFHWTdX6mX5Kt5wj5tfd8NQYzGEAXOQYdRg1OepCs3GuEZmIq67fTcpikc+/omkDnRUo1ZrAAprCyUTpYo4xRVXXAauz1N+6ReReo1cgyQRefCIifBesbRvZQhe62m1iwhxHJNlGWhFq9VkeXkZHRlya+l2uwRypqbrPPeXnwPL+7nhqu+S7t+L5DmzzSbKC7GuEdVq2Dghb0zzsF96Lrulzmeuupp8Zoo0Erp5jzhS1OOIxFnqacaGnuPhm7ZzbCdnexzhV1ZoRhFKKXquj6/pSj3jnl6/5fCUAYB578nzvKgBBCSKQSskitm1bx+f+OxnCEDf+rGlvqaH9iBlXkYk4HyOqASvoWsJ1sANe3ZjI1O2fYxqxUbBfH9wBlbZXUsChVFBxdgw6EOyubAeE5YDmbHSET6KWCHw71/5Cqoxw8kXXEjanGJfr09IDL1eDxEhMgZxOcFmnHb80VzsM979f99JXyBInS2nncVDH/ckjjrxJObnp1jupli/BmLL2QzDexR5njPVarGw0GZ2ZgaXdVltd/jkRz7Ijh98E5b3cfr8DMzOsbnZYL4Wsy/PsSbGieBwxKKLQShOOOqYE9lw6jlc8sPv8xsPfTDTkcJ1OkxHNeoocufwaUaLLsfFTZ5w0sm885rvs237Fpa7fYIC3WhinS9KKqs1fI8ys4FazuDnwUYnWpGlOc5Z4kjRrNX5+Gc+x/Of8BQirQJ6sIzDsOj+UIlnJYAjlDlM6CjYZWFflhNqzVEfexixBF8K/A1ulR2ZoOnBYhLjA0Xu2JUfCT0SFBbBmZhd3R7f2ruL8x71GOqbt5ITFQoH4kkSjUggVhC7PpsSxc3f/Tr0lul991uE5QWwXXZ/4dN89F3voNmKuW3nHpwNEyxs4DCH4IZg5pyj2WyytNKj2WySOUuaOz79mc+x4z+/ALuWoeP54RXfg9U+t/7wGhKbEZtiN46TmMwHetaSh4CXiJ41/MJTnsmCDXzr2htItUI3IpzLsN1VojynKTBrHVtWVnlAo8V5Uy1m2h0imxXqs4XWeLX2fgzMbBAbzvMc59wQ1PI8BxOhkgQnQt97du1f4NLvfR9roMdA8VruVMqwmLGjNDZ4ctFkCq6+bYGsViMVTQga5dXExeXLAJyqdrTDPdUTdGptfbGfKJwdF0dxa15jEAQYiTipUAxfdSGQiuLmxUX6CA943OPpqJiFbo/alpjcJuQ+p1lLsN0OU5Hhlhuv59P/8A7mTj2Zl/3xG3BT08xs2851O/fSDuCVMDu/mSz3xeSsNagaQhgGbJVS9PopRjfIckfA0mjUefyTnsjq/e/PGduPYoMJrO74Ef/49rdx/Te/xcc/9mEe/N9+FSnixDRqhuDAZRaJYnyc8NAnPImP/NWr+eFtt/HwE7aRKEGLwwhgLUoMLSP4LGNp904ef8qp3Pitr9Oe30xPGVa6fYjiyrP4MYDZIG42ALLx/5u4NugdJ4qa9NIlPvypT/Loc86i50JQRoYNSgcGWiaZmhroCAUpxs8uAzs7PXqiKYZsmVKeowSvqrTiniFn4cBY2Xigc1jyUrbmDHrhDnydtSmeiNU8Y9dqh2UMWWOKNorW/AyLy5B2A2m7eP80L4aTfOZLX4b5Ddz/cU/Az29AbzmKW3o5zaOOYf6o41hpO7ptwCmCLzJWEyJ7ZV+vUopOr0e9XidJNNOzBmU0u/bs47gTt7H1xBNYaibsrsf0t23mF3/zxdCI+NLFn6eYIOhJ+71iDotk+NAnCzkuUsiGKfoqYm+7h1fFqDjEQ5SA0mTWYkOOyTts1cL84gqPOPYE6r0+UW5piEHZkRbcOk7Lmts6bLrySu7QoigqZnQaQ5IkRGXMUmtNlCSkaZ8st+TWY+I6GZrvXncdt7X7aCNrMMxPhGfkwDPmsN6iAAvs64GNaiz1LI7ogPTz+GCIyg7fwiGO7agzQw7e2zp+qa2jRmKtB52wp90jMwknnnseodlkteeJYqjXhJmpiOX9lg3TdVxqC/pP4Bee9CRaG7ey3LfEU7N0c08/95g4IYqLQfThAN9XUMWgRkSEer1ejJLrZSwuWlCGDZs2cuvtyxBFhCRhd68PUzMcddqpZTC4qHYUmzJTTwh5Ttbvk9RiJFL0VGDB9tl2+qnctryEqjeKREM/hSwDY5BIk2Y96nUN7WWO8o5zWnNscIGmc8S5Iw5SMbN72JxzZFmGtXbIxtI0Jcuy4WDoKKkBir0Li9RnZlju9rjkym+w6sYa2ddu+OsWzTofjIqwQAbsXFqlEwRTn0JUTAgyVJMZsDMJlYt52ExsICE+cBDXUd4VUTgPURyT5hmiVZmtFoIfZ0OB8SR1CCPRAKU1aM2PdtyKj2tkOqIbHDpR+AB5Dj6DLTMG37b49gotrSCztFc6dFY6KDQht4VOsJEhQ8xyEFMMjR02tIdQZESdwzlXDI0OAWNiosTgETLvMbU6wYPkjshBTdfJOjmYmKYIqr2CSTvUxIHNiBNDL0txRnCRwTXrMDvFYpbTdYG0Z6lHzWE2JXOWKDGkaY9WpJjLcjb1Uk6pT7FVInQvRVlfTNcel3Sv7Iivc6310LUUEYwxwzVj80JvTmuN0hGZ9+h6k89c/EWUHguiuCIrSjkicThxKYyDmVIylOsA2pmj5yAPggsHjqAfNOhK2TojVeHgPcTYStURrfBaCFqXY9HUIWmeWpM86PV6BVjGMVu3H4M3Gq80rtDNI2mAUp723kWmyDh6psm137oS8pS3/Nmf4dsdNiQxLQRJOzRig8279CwQ3/H30GXWyoVAZgO5d2gTo4wmWEcj9xw3Pcvem3fw1j/7S+j2ae/cSWfnrcyIJ/Q7hTquGKIoxlmh28/Ig2LrMSewL+9gtSLWERKKFnJElYWZgdhoQr9P3OlwTBTxsOOOJ15YYkZHJNpUC+1eNAVorXDeFvE0U/TPdp3j1r17Wckh8wQ/FPMZaPyHddOaatAYboF+gJVeRuaKnV6PuS5FDsAX4mtDQlAB2WEHSOUOXFAlBCn6K32Zog4ySncXzKc8QeOJGgVBBaJYo5Ui6+XMTM+VFdmq1N105DbH2SWO31jD7dvBy3/5mbS/dSUbjjmW1at/yOtf8hI+9Pf/wGte+ltc943L8b1lkrrHNC0ZjnCAct7kf6215e48UBrQBKXp5ZaVxSXe8JKX8G9veztvf/2fsvO/LmbzCSfDbTt502++hNVbb8bYrMBMq1EuKZQyXEziYo7aegxdHDYyZRZXCFGECwETNLFTRGhioCkBs7LEMQGOV5rEWbJ+b6xOb+3BV5O3dZM3Vczs8AHND7PvgsZSKFfvXljkhltvKWZXKzBGofShpGRKMHMFmIV2Br3cEiTCEBe1TzIqkA2DWQGeISOrBkIcGUAbBfdHPw9/p0BHBhcCbuIsjlWwhzC88KQsnUE8UaQxovC2cBMVQrCOYANJotFGSBSofpe//+s3wb49PP3XXsRfvu3t/I9X/gEE+OZHL4JbbsOIQQLkeY7ogMcdvHC6LA0Jw3qzsdhfCMVYiOCh2+HrH/oAvX27edCv/gr/+41/yivf9jZQ8JcveykqTalHMcppQlq0xk2ZBOMNIfNEEoMKWOsJQYOJsMGjXEC7QEhzjIdmrIn6HZpL+3nUyacw5RytJBpOqars3oupGVW4mbl3xdDwKCJowze/+71hsD+EQPAD7bn1lU4MxaASUmC567AolDJoUVg3uNgCQRUzAYpBqmqyEqCyu+9OyqELj633xcVp9FBGcHyuYJARcDCWAQ3l5pOlfXzN0owibKdDwxiCz8ldlyRp0m/nTOsGobvI7m9+D44+gYf/0nPYieG8ZzyT1z/tl0hzh+32cVrTI0ZMTJYGRHTJWsIBiYeBRVFECIL1AS+C9cVOG9cT4i0b+P1/exdpdwXtNTURlr1nyxmnc/oD7scPv/pVdt56E1umNxIrii4FB0YHlHMs7NmNCY7gHMrE2F5vVBIiCqxHjCbYFKKM2SSh3+9y2oYtNLq9Ig5Xi6tFeK8ys4D3dphI8kqTYZlt1fn6d75F/vQnkQPGFfMDROmCUK3nZgYplPpzYKnTIagIpBRYs6O6Ml9O31bj47oOQvcqu+uA5g/C2ArXUqGjCFf2Xo5FV9f9OYRQMiaPMYqQZzSjiP07dxLynFYtAXE46wvxACvs2nErOM/p557P/tzRbk5zQy/ndutYQuHrU2zcNg+qjg8RIrqIS623QMfGSQwTEhQgppTCek8IoJMabefx9Sau2YK5DSxgcM0pzr/f/UAC11//Q1Y7K4im0MXyAZv3McazsOc2Ihxic5J6jVwK1hgpPVrtHiJtyNIOknaZzlKipf2cNDOD7qXlZPfK7iWfhCjSWJvhnENHBlGK1OZYUdx4645Ro57RiFKjmP96YOaLuTXBAUvtDkEE70cj5wowK7Sfqj62e/rUHvjzYCPRkQG1JuMWwuTvwjg4Fi5pEhuMErbOb2RleZF9u3ZilBAbwFkatQTtoRklkGZYhPqGDbSjhF6jDtNTuLhGVK9x7Y8WUEZjIilaoMKaGNM6RTt+FL1FmQJ0PUJuA1nfkVDH9wTjG+zb00E35ljNIK43imEsiSI1oagGj4Ak4OihTc7enTdxwvw8UWaLLhYlmMHAMudAK4J1xeaMx6ddpnDMOMsDTzqZrbUGuvIy792NvMxwKqUKMcfgISrktTtpn3Q4JE4QpYfdAENAkwPArHhCt9cDpXHlbiUSRoOeyhiZH5/eVNlhMjI1ERsrDnoh6WxCwZo0miARWdwg0zW8GET0ULtMwtqAaMHIBv2zvV6HhtacfPR2WggrO28ntDuILxZRlmegYGbrJtiyieu/exUqtZD3UT4n63RoxpruUpvjjp7HZ55e1zI93SJ3WaldNyaCvoboGKWLOF0AlxeyPnUPTS00Beq5Y4Opo/opGxtNdK+HznMu+eolEBzHnngCzalWsagV6FgVwp55h8Vbb+SM7UcXE5n6PYzWJCYGF7A2KyZcRELwjsgkJCYiEoXpd9nU7zC9tEDTZhg/GdKfXN0yscUUOl0DvYZqdz9cy7KMKIqI45jgPOQ5tVqTzDuC1iwuLdFO0+AZac/JQdxBFdDlkN8i++VCQCKNlWKOnQmBKAi6DC57VZzMQYFnVXR4mGzM+yJoPwA0sRgsiYUoA5zBmxZudiuLKsGrGtZ6RALBu7J1SEaX4aDdrCibL4pJO6scPz/LHI4d3/oudQcmFOKEQSwuCSxrx8wpJ8H1N/IPf/hqtoec6dV9bHQ9mp0VttUN37nsazSNpR55ev1Vgil652zZ2ztoMB8fPuxtwKio0K4CataT9LtsVtDdcTPR8hJ6cR+btCdp72dr3uGSi/6NG774X7BlI6eceCppJyMqY4uZzWg2alxz5eVE/S5bawn1zDFjYkzuICsGrJgkLnZ6HEErhBh8Qtd5GkazaWUfj9o8y9bckfQLCWerFE5CyWodoZQcJww6nj1euSK5Eg6lhVbZnY2Y6SjGuaJvU6HRyuCzQvMs84Edt91GLUnEls/QwRMR1p0gbtaGvQoXZTAazpdTsCGEUoto+LIVkB0JG6pdrGFWOhTzB5VSZKLIG1N0TJ2cfjnMq5QlHmg/ySQzCqXcr89yZqZn2eQ0m6OEb37xCzzi136DuolppznGB3INSaPJ//Pil/DWy7/BzV+7lE+9+52ccb8L0bUG37r2Zv7z//4LmIgX/sPfwfQ0rakpci/kBzrGk8kNCYTgyPOcRiPBBIWycOstt/B//vSNsOs2LnjMozj93HM546QT+MynP8XnPvZvUI957steRj8XGq0WbVtm0hWs7NvL9y75CluM4cxjj2MqjglZivYBtB5KkXhc2bDnUU4hKkK0wQTLRg8n1mvE+5ZIooRUGwYrW3SZPPCDYuTJAMCo1q+S3DiSwAZFpYRi5K3003ww861Y5L5Qh12vwLmqGryXgUyCQg28ljA4qQWRdgqC0aRpwE81SWNDqorBMcVlFMoTG1B+pAukvRrK72iEfpoylTQ554RjuPSqb9Hbv5dcxyS1GRpa01vu0KjX2HLM0fz2372Fv37Nq/nyP72TL//Le4qoPQLbT+DJL3wRcxvmSVWC7UCsFULAiB8b1lq6X6pUPBCP9TmNqQarmScKKfPzLTbMtXjKC1/Ap97zLr71sY/xrf/4j2IkXXcVZlu84M//jKPPPpe2ivFecK4YJaaVpaUD37v4i5yRNDhx6zaiXh+f5uVAC11krgZlbSoU9Wfeo8QTSaGVFSnNpulpJLsJ7XJiVafrbXlZSSFWiS8HLFd2T14D65d3FW7lSrtbDKBhTQHFHYPZhG5jdaR/HPtREcNnMHBJlaUOxVAYj1XQJZC3mvh6nWw54BTDwLUvM5yxjGJnOsjw1GoVYTttWo0G9zv5BJrX/oDrv3E5c49+IrXWPO1Vz/xsk+X9S0jDsOXs0zj75x/AVR/Zyct+6zcJU0303DzHnn8hSxKxP3eYRoQOkKUBY0CrUYZ7krF7UIHc5UQGsB60YqHdIeSeU847lzee9yaWbrqBXTt2UFPCP/3d26nPz3DquReyEDQSt+j1PUkCPs/BdlHdFbo338iDzz6bhgv4tE8sIMYUgf/gQEArhfOulHIv2mCUD4jP0RKoB5hLEozL0d6C88XwjdJd1iKVAOmPBdDUsKZ1UGo0yFp2el38BOCpO+B2EyR6nUB1dbzvMSAbao8NMjNBAaacugSZBLoKOo2EfLpFqjVepFSiY5TZUWMNsx60N2hnwAfqSjEVLGds2cAJSY1LPnERm2oRPkupNRX72pCZiBWl2YPDz06DzTn9lJO539lnE0eKr1xxGas+I6rX6OcOFyCOpHSHPV55nBrTafflVzIaHWvaKx2MUkRK421gam4K4pjvXXcN1954HQ998AO44Kwzse0VmrPTLHZzTHMOLzUIEUkMETmJ6/K1T3+CeXIee94FmH5KyCxmsHeHURgkhCKmqEQV4/ecR3lPjCKyATo9jt+0kXpwKGtReLQU8xTwoWRnYRR+kaLGUqpK8SPuXo7YWFEqNsjip2m6pgJpEMfkzrmZoRBGKd9KUbVt3LOnUgFOBsHzUgiTwSwHT6qF1ShidsNG+jtuIBOhXjhDSAgEAl7KnGIolYFDWQPmcozydFf2s33jFi7Yvpn3fuUL+H17aGxr0QkQYhCp4yLHYrrCife7gB/E/8oHPvIxvnzFpfhuF7Yfxe+//xH0nUMpQ6/viLUmorjIh15uGNQielQI9Nodoul5WjoheKGeGFxP6K9mtGLD+//tfYRvfI2PvvGNPOsFz4f2MscefxKNmQ2seEOWCcYIeR+aGpL2Khd/+F85P2pywsw0jW6XuJyGXbAyD1GR8fQ+MEo8Fo8pM//UAJNnnLB5I63VFVaCI9YGpwLeOZSEMmZWbef39IbuUeu6jYMSI6FIUCqKQUpFYfTYCDg5KGdT615y1em8pyj2wDezB6j2BhRWAi6OaOuYMLeBfpyQiUbQxQScQYvSxHPH+jQDaCVIr8O07fPIM09jK44vfuD91Ptd0l4HnYAoxVK7jTSmOP7cC9l4/wfxxf/8L3wQnvL/vYxXvf3/0Kw1ydMU53KSloYErAhOxsoUZFQJJEAjqaGtQuWKtJ2CLWradC8lsZ7f+73f5TmvfCV60zwfesc7kKO289inPI2g6nR6RXGtUWCzHjXf45rLvkz3huv5hfPPxXQ7xGUcrJiobYdZXO998bNWBYJJ0YCO9ygfiIHYBzbXa8R5H+M9EgI4jwoQG4MPhUbaZP9sJRt0zzGzydY+KGW3y+hlWC+sNv5Ka3UxVDi4XzoaNFvZkTRfyo978cOuiqGWGeCNoWsi7NQMvbhOT0c4MUgQTJCyuIZR/V9QRR9kkKIIsVljuh5De4kHHncMDz3mGP7z3e8i7NnJXE2Tpz2ctTSbMyhdR5km93/Iw2F2A4973i/z+Cc8ic3TG1i8bRfaBvAW513RS6fAikaFIm1O8ARxw43T574oAwkw3arR6+Wk3R5bpqe49dpraMUx97/wQv78j98ISnH2hRdS3zBPz0Gc1BBV4FHDeEx/ha9++IMcI55HX3AuNW8hTQneMtatDMGTD+SZtcK5vLgUtCpaZ/IcQ6CGp2EtUb+H8RblA946RDRG6aLuaZ2A9UQApoqpHd7aX3eH92MAVYKZ93cYxVfeQ1ySvCTSOJ8Xg4TLLPcALYsYRDhgzmNlh8nMytiXC6NWpcEhHmilp87TE0VozWKnZrCtaToOEF0ozg6nNJUXmCqkrb0xhEjje1181qelAq1el6c/8P5s7K9wyQfej1lZoO5y6giSWhIiso6j3pyGXsZpZ5yDWLj92pt482+8jL99/esLXbNgwQdyARWDSx0NE0Mtopv3QHtCYtAqodfJimEUUTFvotGo8Q9v+xve9Tsv54ovfomGMjRrddCarhZcvUYuathy6q2npS03fP0Sdn3tMp527vls0pqad+jyOOFd2S9VjL3TcVLMkA0BMbo4Ls6iddmG5RyxC8R5Ss1Z6qqYuh2pCAngshwjqkgcFMUdY10OqgKxIwVm3k8MPFEhEGmDzXK0FqammkV9olIyFjjDZQcefxVJob2gKXrYFIL4ArSsPQg7q+yIWlhzXINMMmltIqxEdLTBbdrCgkTYqIEth/+KD2XsofQwyyZzi8MGX4gjakPsAlvimLPmN3D/+U1c9q/vY7rXJu52qRnBaM3K7kVaUZ3zH/BAUIoPvOd9fPnzX+Avf+/3wOY891f/BzO1OiHNaDUMFNqJNOt1ut0urK7Qmp2ByNBPM/p5Rmu2QdCwd3+KkoBWgV947CNg0wY+8Td/y8fe96988IMfBKU48/73Z18/RdUiUuto1CFr7yPq7OeL7/1ntuQZv3j+hcx6T0MrdPBI8MMY4/jgFz+eASuVRQdTtgtHPRDbnE31Bto5ooG7PDY+L8jBEzeVHQEnU6lh610IoRhNl6VEuugPbzYaDEbPBMaGZmt1gKupZCwI3awnGCVF8DMMbgeqylaAdiRdzEMc2wDBCyiD05qOjpCjjmV/vUU3rpGJKotjR4x56JqKxStflF2VAYXu8grdPfvYGDzPetD92NRd5WN/87dsNZqs3yFNe5yweY6Q5rhGxAXPewZ7r76KD/7lXyAb5/ijf3oHp5x0Ir7fpx4ZOgtLzDVAWU+326Y+3YJGjT0Li6AiWo0mkkSs5jkruWPztoSWEUK/zbZjN/PHf/sXHHX6Gfznu/6Zi//90zROOo0HPfqxNDZuZKXfQ4yjvbjI8bM1brz0Ym677Cv88rkXcpKKSLp9VJZNSPiE4U2NjmtQozkKw41dQDxxCMTWcdTsDGI9RgZKJDIRahmfMq/K6Vdr3aHK7qbp0drVCPgiAZNEMeIsm+fmhukBzaRowQExM0KOphANnWo2h6H+YujA2iutipfd44HQsaG/ChBf7FZeGfomoje3kd7GrSybmDSukUvAhzAxrLm4yAJSUnibOfCKRnOaDa0Wc1q4cPsW7r9llis/9AEWvv8dZo2npoXuYp9anJDWNM/99edz4gUXgFE871efx9T8FMrl3Hb9NfzV6/6IKO3ASp/EptRrMe1OB8QwNTsPqWd53xJGeazrE0eezlKXd775TXz7yxdTSzRJPeaVv/PbECXIpk28/q/fwmI/ZzWz9LIeifJsn4ro3XQN7/6jV3Gy9zzzfven1e0zpRTKuTEVlzCxGYxXqozPvxppykvR3ZJlbJ2dQ9lCn6EI+5WzHcdYbmX3bMx4sA2JChiBSCDPUrZu2YIe21qkbARYbxi9IhTZJw1MN2pIcHiblzU3B7/wqhN8hGNnw42inDI0SHKKxrkAypCamL31OvboY9gfRfQjQ6ZV6VqNS2p7JOQIlghFEsegEgiG3kqb7u49bDHwkqc+kaP8Mv/6xtdRW12g7lJIUyyBVQnc2llh1dmi5KEekeZ9ZmPDFz76YfjGpcz6lBnfZz5WGJtjjAEd0Uk99D3HbNkGWcqU8cwlsDEK7L/0K3zmr/4SY1OcTYliDc06oduhpzX1mXl6vZQtG+ap+Yzayj7+421vZtPCXn7/qU/jWGWYDeC6/WLeAAOt5APX6LAmLEwyrUH8VwfQ1jLbaBCsQ0JZZiIHD1MXnooalp9UdnjmypkRg/YkIwqjNc4W/bhzU80SzMpJWkwy8TWoVDxIA816GUj2g0kqdxDrqQDtyLGygejl2AizIglTZCWD0mRRzD4dkW3eRDo9RUdrcq0Ja6SBpBRK0cEj3hEyR7bcAQv11gwbZ6Yx3VWOTRS/+dhfYOmqr/KZf/x75r1jpt6gn2XUpmfQU9Oc85AHw/QUX//6N2goyBf3sXD9D0HD0TWI+kvobJXY5uzZsRNUTGNmA4hm/+59tKyjmfVRC7tppascf/x20KA7XXyny549u6C9wux55+OSGkurfVr1aWy3R1hZ5OYrvsa3P/ZRnnvW+TxodgPNdjETINIKUWqdTWEysTkM2YdJIBMfwDsMgbrWRRY2CFJeOiGMpVSGJU2y7rSyyu4mIyuPrWiF92XkLBTRMZ/mbN+8lZiyrMyHwVxzGMqUHnAVeaQUcWmUSQBjDCIy0qKq7B51Lge7/egU+wMSBI5AXwnLRtOfnsZs3kjPCD5SBKOKM1jKVIsP6JJ5iPOouEHcmAFdgyzD9fuEzgpJZ4mHHbOVX9iynS/8/Ttof/8aesvLSBRjuwFRdR78pKfCCSdw3Re+wpcv+hSytMALnvok2HE9b3nVK5hKHCrr0Nu5my++79+gNccZFz6Akx75aL722f+kc/MtbHeeTXmfT//LO7np0i9y/7POJF5sc8OV3+M1r3k1zLZ43otfhEsS4qhB3knRec6cCO98/Rs4u97kORfcn+0uUAvQX14qlnMUDQ/QCKyKKnIdBOVlQsN/kI0fZuW9J0JhEDRSamuNnJrJxL1UAHaPQZrH+2JWRJGfKWL15517LrqMmmgJEFwhXFBeD561oxpdG5RmkVpYBj5+/SJ7VYNuMPRcQIkpMmOqWAyK4k3xZr3zXNldAbIwarR1pW6/YAtw8/EwoO9cTqyKcgjjUo7vdzn52qtJvnoxZ7T3Md9fQoeAcsUJV6FwlfQgMxd0MRcuMuQuJWpEeB3ohEA3afK9nuPl7/sot2w6nld99D9Ymt9AG4+JFFPKk+/bz1vf8EZ63/5Wkbq0Hci7xUI89XTimc1kP7oVak2e+Qe/yznnnE22aw9/8so/gN27mT5uO7qhWPzuN4t5hs1pCsV2BUdt5ddf82o2nXwaPSK80zQ0TOcd3vHbv8W+z1/EXzzlaTx1ZoZjXA7tVZhqkgF58NRD0WkwjDWW7oIa+iEjIBvPdAYUq7Fmz3SLq+fmeePV13LD9Bz9pEU/BJKQgXgyPd7GNFB1KC/A0r2v7O5bKMsy8iylXksIaUoSPA3veeXzX8j/eOjP0fRIpMB7W4g4AjYU0ujjeQAJYRWPkNJkP4Sr23DJj3bSTqYI9Sb9zJcjzsBJEVBWgHFqeCH6CswO84QylAQY7PrDhJmGLO3RMIYaAdXts0VgbmGJ1rev4MzvXsaZ+Qr93gr1WMj6PYyJIU7Is0AcTMFaggdlcSoQxJeChJpUxeyOEr662uW3L/og7vxH8rJ3/gvLU3N0gcR7Ihy232f3LTfzncsvx+7fw1EKPve+98K27aipaeaOP5bnvOQ3ibcfgwuaGedh/wJf+ND7uPKSi0lvvYH4xBP5xef9P1y3Yw9dL5x7/wdz+rnn4moJXW9xzjKlhObqEv/+V3/Cdz/wbl76gAv51XPO4Ni0y1SvC85CHGPjGIsQu0IuyY8xLz2YgSkypFchhKFbM4hFrkYRu+oNbti6jdd+/ypunJmjHxlSRj2AxSYDgh0WlIsvwgFVmIVh4kkdED/0B4kJDxZ98X+jFGmeIzVBlCekKXVn2ULEv77lbzm9WZeGL5l3SbR96WYWYQQ1BDND8EFJJBqoAxtq0JJA6nM6WQ+IhrubCmqYNZJBJ3GlrnH4QMZIpnyc7Q7i1qKLzcS7wn1MVcTC9Bxq+wms3Hgte/d1mYrqWNfDqKLK3ecRSqJiYhEeVKGFHpQt38uAV0S5YwbL/TfO8buPezKv+tzn+Phf/znPeNVrCSoidZZcCbrZ4qizz+G4s86mkWWw63Y+96nPceyxx/OKV/8Rnek6y/Umy3GdNIcsTdm+cRO/9CvPA9vmsotu5tV/+he42a1st8LMxm0st3us6ph+PyeOFU0DM7bL5R//AFd84N286MLzefpZZ3B8PaGWd5FYQCV46wjWFgmHNZeNOoQzc+BlGErGNXLxVXBle5guNxmN4MrX9WPJr8FQnyoUczghFnzROpYHS97rUa9p4hS2TE1xVLM+6fSFseTOOmdUuVKATpwlAuYNbGzUiXHg3XDKD6HoAdRBofxosEDFyg7PZKwAc60QsxdwAURHhHJwrzWGnkAeG2TTRla3b2dHrUG/1sQ7QYsqZG6cJxI1mnEqHi+2rGgHLxpUEUJoKWGDczz2pON57smn8c33vZOvvuddTGermJDTmpti3/IqXafZ385Y8oI66ig2PvBCdlx/LZkWVGuG/Z0eK90+phbTmJqik2aYKOHb3/8+bN4C8xvpJgmto7dx68oivlHDaSFSQkMs9c4St1/+FT76Z6/j4Zs38bTzzuaYRg27skKv16Ob2TKnVei0ReM1X4OvWU0l/4mKlvXTFEIRr6zV6oRuSl0UD3vQg4ohwQy1NocRrbE+jMlWTKWMDC4qE5AGsH1uhsjn1EvhP1Qo+nTLmMFQS6taN0d2l1rHnCs2E0coZJ1jQx9oe0en2aBz3HHs37yVfcqQRwmIoAVi0cML3UvASShjcgUb8WKwohGt8VmG6Swx31nmxQ//OX7phBO56M9fy7c/9kE2aI/tdpiensU6QUyLXBL25o6TLrw/5Ck/+OHVLHc6TM/MU4sTXL9gMtZ6FvbvpXvjTZz8sEfS0xGpMexpd5BGQi9kuGBpJRqzskD3mqv4u5e/hDNjeM3znsmJytPodohyS71eJ+gI50HKUYghyyq/4Cd7Kyeq1YsEli1kmExukV7Gkx/zC9QGYDauyhhGallr4UcVsw+Lqv9IimnlR80paiHHhBxkMtagS//VS1VrdqSSAOO6ZoPZCuNMwxdpF3IRrNb0NHSAxSRh+aijaR9/EruiBt24gUODmELp1TmcKprYgyqasYdZv6DxKHrOYuKCGU2nK5wVe174wHN4zMZZPvLHr+V7n/0EprMKvT7NusY7hZgY6lNsOPpoCIGdt9xC3s8JDqaSBPp98k7OXLPJzh03AYGTzzgXH9XBRDgcST0Bn6PIibJV5PabeetLX8xp/TZ/+JTHcrLtsN2mNNMOSXAYU8PoGMEgokB0IX80lnI8gJVVfcT3+RBLLkUrkxaF5I4tU7OcuHUbx81vJAbR653CQdPumjicwim8DyilMYByMK9gLomI8hQd7DrBvcl4T2X3INApQaQQaR0AWm40NonoxjH7p2bpHX0iKxs3s1Jv0FVxUbpOwPmcIEWoQMoexkHsU5XlIC54JBKUeBo2pdVe5LRYeNXTn8hZ4vjgG/8EdfutbKvFdPZ1adSg3fcsrLR5wAMfDMawZ9dO6vUmWS/DpRk1HWFEkfX63HzDjyDPOfqEU+h6TeaFOKnTXVmmaaBuu2wOKW//vZfT2ruDP/xvT+eB0y1q+3Yx41ISn+OzjHR5BbGFO+BtuT1Hler7Tzic4fMM54te2aaO6e1f5L897gnUAZe7cvMNY4HldYBtCGah6DcbcLZIBUmAEzfO0QiWyBeV5BV43UNxAzl0LG14KzN0riygdcbQM5oFXaO3aSvuxJPZNTXDatIgNzG5D6AcXlkQV/ZZj6IMKhTj7Bq1Ou3VLkUricW1Fzm6lXBsnvKaZ/8SW/ft5O0v/U3qKwvMKI/RECWKZr2FywNYx9VXX0MSNxAd0+/0Aag1NRLF3HTrbaAUc0cdQ9yaInWBfrfPppk5kjRjpr/Me//kNXDND3jl457GhfUm251jNnjwORJHxI1GoUGmEwxlc7uzZMEdsM0O2Fml7vLjjHz5gwROFAcdyVfqxKl6vRwg4/C9HnVRPPrBDyYBqUelerCUySuZdDfXeUczbHRyOIwqCmiPna8xp6HhLbF3CGXsZqwUoyoePHKANqGUMT7GrxxqGulCIoVhA4+i72AlaFYbM2THn8BtrRn2tQpA66PQWqPKdjU1jBWokQvmPQaD9qPfmVaT/p6dbAmO81stXvCAB6Kuv5qL/uZNzEmftL2K9xlQqHFMnXoWWc/inUIpQ7MxRQhCuxPIgxRgtmETUp9iz1JKZOrM1Fq0d+1jQ4AbL/4cl7//XTzppBN53MmnsNk64tU+UVIjpDlpuwNKFZ9RdCEJLgoiPSybqOwneO37HGdzWnFMZB3PeMIT2dpskfd7wYzpxQUZA81xUJMJZlb4n5aALWtqI5AacO7xx5CkPerBEfI+WhcZ/dQVgyykajw/InGD9X4eAJouEjOILW66nOiEL7ptQxSxrBW7m1PMPPjnuKUxxUJzhixKyHOHQRUj2ChOXvC6CCUNXiyz1E0Ncg2qDllOTSckec5cr8OvPPD+PPGYo7jsvf/Ijq9/iabqEkU5Qob3jn4/I4rqxSxPL/iyB04nglcaqzTMbaCXO5IkIQ6CX+6y2dSQ22/nPX/4Kp549DZe8IiHMd9PaeRFqIPUE6IEXWvi8hwVx5Cl4B2iFM57UHo4C2EtO1sbPxv0/1XZzh9ziH+sk0hEcKVopjEGl/XBOeaadTpLC0wnCc9/7vMwwHStLrbUIAuqrKiQsfG/62QAhhVnloDD47EooEVRc7a9VSfJujR00cWuNHgckQKf2+ps3cPu5rDdqSwVG2STB0kYHxmWg2dlapbb61Oos87jlrio9zJTs/g8FO0FJY0RyqZe5QghLRfEQIZAAxEBhfFQt5ap5UVe+uQncCKWj//dX9HI2mxoaPLOCo04Io4N/aUFvEtJtCfPOjQbhizLCcrhfAbNhEiDOIf0uswnijkV+PS738XWbIVnnXcWW71j2nuU9YipgYrIvSqGtzDZszqo96rsJ4F5FedMaz0EszRNCSHQbNSpGWFlz242NRs8/hEPZyY2Q7FYpYoJZWslnML6XmbZ9SETITk0nhrIHHD61s000i4zxqFtF2cz6g1Dt+uIY3NAv9sB/W+VHXrnCge6m2GN26nD6DYeRwPwqujZXI0S9tansKeeycJRx7CvMcuyaZBKDEQgMWEgpiKCVZ6+yrGKQsNfFF4ZcmXIdFEwGnnHlM3Y4nOeeObJ7P7mV/naRf9GdtsOjp5rkXWWSfM27L+VZuLw+QKNKMNnS0R0qMc59Bagt4iWLjO1lPlGwC7tZOXWa7j4A+/iUduP5mHHHc+Mc4jN8QjoiCBJIWFbOBtlgdGg+HcA6qoKddzrC9ivXzhcZhuVGIIXBiwriiK01libkac9fLvNMXNz1EPgBc99DnUA74qk4xiQrZ8+OAiYKRQaVU5jchgKjbNjG3DShjlkdYkpE1DeQqkIacyBwo2VHVmGtragdm04NXcOVaux6hVLcYNdU7PY085mYct2fuQjVmtT5KaONxEeMyxXEBFEK3JdzOYsbkKudAFuZcazqQTVXuGXH/cLnDY9w5c//CFaeYrpdWipQBwctOro7gIt12VzAo2sQ8v10N0VVK2GihRxb5FGfxW/cBtbo8DH/uHtTNs2z3nUI2mmfWo+J+RZsXt7sM5jxIwdFz8CtHLosK4UX38i3MyBq1lghiGO48Ltt55po7FLi7zwWc9iS5KQgNSUJrc5IoM6CjX278HTDsaXXRmFkyG4UDxU8JigpK4I5x87x45vL4Cv4bVlpd1mqjVFmg4+bHXS7q4NR13K+oCmCq4y3Ib8OrpcQQSPoa/q3Boyjj3mBKTdJU1zkj27IChark8stqQ4ReGsUsX5dsqPWKGfDEX0ul02bNvMjSttHnrOeXz3siv537/5v8iTFvVWi94tO2F5P3/5ilfQE+g5i+gEHQyxBPyNN0Pa5a9e8Qoya9k+M8/qzTew9P0f8KhtR7N1ehq7skhDJ2gbMCJ458F7EIWSsR16rE5ODzoZqiV0nzbnHFEUDdmZ9344hyFJEur9LsfOzPCrj388Ug7X0rqQip/swASZkNn0E0BXXEuqGGGvB8zMCdortEBNoOaRWeCc444lzvqoPKOZxDhXaPZVi+nwgGxYKBtGSYCwRsd+0DY2HCwz9vdIRfgsEInGqoQlXee2uMnSiSeRnX0Oe2Y3sJjU6fpAQAjl+HRxAWOFyAf0IGU60EHzBVio4NGRoZemdLtdvvHVy5lqtdjYmqJWi8AHmpu3sPH0c8iyDHGWWMX4NEcFyLsZM9uOp3XsSfgspxnF7N1zO61Gg1NOPYWbdu7ga1d+HRWZYtFrQZxFEzAqEHxWpvDXREkKMfgDduvK7ns2Hm4aDC3x3pOmKS7NaCnD//uCF1EDZg0SAz7LC30/Rtl3CYzroq9rxgESyh0vhEJNoSzvUCHQUEInh9PmIvb25lhd7KLiiJVVi45NBWaHC2iDHWyCQI9Aa62KxuQcR4UGvA3FQtGKvsTsDJ58ZgZOOJ6p1RWyH3lc2iZkfbzy4Dw6KyfiaI9RFqtdMTPAC8YHxGlQYCWwmmcspj2Wfcb/94rf4bhHPYZes0WqIrQ3xLHBSUY362Ma06yu9JiuTyPO4V1Ka7pOx3bp97toF5gF5tKUlz/ukXzv1lt5wvnnEmVZwcayHCKFqCJ2m4/RfnXAQh7vc6nsPrm+lcI5R/DFZKwoioYJgFpk+G+PfSyPOv88DIiypbY/grMWtbYo+g5Ot/F4xPsimyV6OHcQb/EuoKKE6QjpB8K5R7XYtdpj5+I+puvz9APkpUjaxHvKSKdL3cnFdrDYm5dDP/YnvT90sHcMSjHGW8fGv+Nad7Q4xh6ckChFnjlMorEKRCfc2k/R9QbbTj+LhTzHdlbJVmCKPjXXR/uxGBQBE3xR7jCYhF5uijquoWLNlJkmJfCBf30fp+7Zx4qOsHGdvjOFSqjvEUUGKxE6Tsi6GVoVo9py34cIkppGe490u2S33ML+bofmycdgEYwxqKwQNyhYl8eNrV4ZCZGNZGTXbAgHjRLLwLku1+X4Ql1n7RJG9wpfYeYh1u4otuvHBEbVcB17LYRQcH6tQlHilaZsnJrirK1befEvPYOEYsK8ACF3iNFoJVjvQOmDTsNae95NTIDyDaUclDq4rJTxgEWFwJREaOChJ2/iyhv2c/vyLlpTG1hwntpMk27Xo0qhOlEQXPHpfLkwh1NVGNT6KMohOZOB7jUXrVrnEJZTw8rHHF6P6L0ZQPYy8v3DGFD7dR97oFyzKl3BEEBHgsUiChweSWIWVOA6FbPhlNPZoCJWbvoRG/fcxhYUNZOCOFCeQOF2ahGsKDJVNKCjBK88wQcSHVMDbr35Zq5//7+CboCpjebmSCg2wcG3UXrgB0Pag1hDKAQiSfvQ7jJFoNWs470lT3O0SSCOio3UCKJN0Tc8riQrRTLAq0KNVA/69A4FZgEQGY4zCQTWC/QO3Hd/B2skDBjyfcDNPZz1OwxljH2Pg73e6LHjU6uK6V8+t2grxNqgnBRBi0iTiicYjbcpJEKeZTQDSNrhmFaLd7zqVcxTuJbiyy68cmBJIGCUxpfzHUZeSjnQZJ0jbwoGpUqd80GILZRPDBAssTL4ssNvs4Yzt20gCvvYsbrA/Obt3Lhrkc2b5+j1CkDzeUB0kcEomgukYH/Do6LW+E7rsy8/drCRUZAvDIey/nTsmEH8xAV5V5hnUIIMpogXVBslhcRPZhJuDxZ31NFgIrpa8ArM4l6S9jJJf5U8OBIlGFWM+dKh+DwOhfUeFyBpTrHnttuJEV776v/NlnPvx4atR9PJHWkoNqdaWf/mcVgJeCU4BWm3x9TUFKlNibSQpV1ib9lUb/KMB57DbTfvoPWAB/z/7P13vGzZdd8Hftfe+4SquumFft2N7kZoxAZARJIgIolAAhRJMYkSTckjWR6lkSXbypQsz1gjz0i2PCPNfKwxNR/ZHslBgRJJMCA2ciYBIpBIBAig0ej83rupqs45O6z5Y59T4b77Oj6gG0Tt/lSfd0PVrdpn77XX+q3f+i12YqC9fEiFQFWS1JM047erCZCFzhtpKWr5mOf/yvsh2KzWm1buifzBw4gXHtTgTJxS3bNq9E5+L8bsSYv0Ov4RgmbOapCU2ywZpSws3eEUk4Snn7/Af/fzf5udqHnt6Yq+v5H+7JFT38eDJRvdQvZXZRHOLLIGCxnihCNS4WQH9KlbkG44y5G/jzsuP8BN15/n4NIxW1tbdF2ictk4Tr2nsCWQSEayEF5vhBYe2KDQJmbhlayB37Ju99JglXVVLvexgfCP62IahP7koZzo05+7UDMZuGhxvW9ka+BiqeiFM1xXPINqbPFfFg7vSZy3yp7OKcI8Y1U+IS5R2Bongkah3Bpx0CW2d89SVhV//W//dW753u9DxjX7sxmmGGEV6qTYvoN4RHNPAiOM6zFHB4fUdY1vG7YnY+L0iMO77mUbywue8xzCrOHyxUNurLcy1OHbnIgobO959Quld+NtOmnkV8pcVj2N1cNuUJ2VIXkgfyAkrB5LVJJ6w736GumU1ZflyJcClmscSHGYQjAqpJQPy6EdnDUGjYF0fMRoMuHcaBt7+ZC/85f/M2698Qa2V4gQaw15HiU9wp3mAfRe+XIRpIQxQjF4RMAtO4Zwyw0c3XEvfj7nunHN/uE+o+0dsDBrA2VZoikvHLOWVr1yM4r2hacnN/kJQ7b6GjpAPt+heIaKIUq+I1FlEJPFqiBGiMbQitCmQChq7IULuMLQWsO8KGkfuBc9vITEDCeUhWQEViOiQmkM0+Mp3pWMt7aYtw1bkz0u71/k+L4OKpePcrXYmKW4NUn29AfspIukGNna3qZr5ly89x72RhUaY0bFuoA1hq3xON/V4yloxE6KK49h7VFFGdLwVxqx9RiMK9R7FwZAuDZu3R+YtbTEE+VhRASLxFR/sHiBEDqwuSZYVSmAMJ9zbmsbM51SJ+Wv/Lk/z8uf+XRSB7a8ts6EG4D6k2CzDtK0PRZiMJREoiZBnJ4H3A6kW67nM1+/nznCTuFofEPjBcoaY4Tk8+vZE9SDpbb6yVlaskcSq8JsZhl66nqc/+1N3H1wzO+hvM4FnbBv2CGaxQAGhCgmsNU2UzrubMHvnOW620bMdvfovvQl3Be+iGuVNDJs2YhNntQ02OSxbkStyvntbT56x9eZAX/zr/88z33ta3H1iIuzI9L2GFWlCgN+Jb0gQcZet6oRzazNfQesEH2Di56xT/zFN72B3/u9LyPPeQ7WWtK8w5RFrhX2HleVp0zIw9uYSR/at/1OL1If9hgnmiUbffDIYHVuVXuw35j8TWdQC2k+R4xhr7BUsynbSfgLP/cn+Okf+H66uef8qBCJwz3QNY/s0VYOOemDSlk52bKR6YHdlMA4hpZ0pQqluCwtBXrbLkzsdfzmF75CtbVL66e5to5ECGaBB0tvuQaAdWiEYtcmrseOFgat/1qH8HKZJUW//R2yk2oZjwUAXrjmJv9vYI5Zk33iGEtmKXDRCGlcU99UghtxrCMeuPMO9g++znY4ZhdlZAtsihBbymrEfUeX2d7e5ly5zd/7r/8B/p/9z9B1cHYPunm+GWFY5UMtZX/DkoHo+4UVYFLB5csY55jEGTfd8lxSShzNjqjdKCcIQiLG7K1LDz9kF3x17tLDDMY345HALLLa0WqFlbDq9KyLh2bytXGWjohxBksk+Y66LHDzlnGA/8tf/Wu8+kUvZhukHBU006A7YycxXbuN7DiJLbCKO2ThRjtkplLCGQsEJArbWJlYcGN0+7uexgd/7+vYyR774rg4m2KqrWUrrhOEz7Qqg6snAtChgasug4fcgccsDFl62MjSQ+MGj59v/2BWat17uLqHZtY+yTBHccDSDMQWigSFGZPwXPQdUk4INz0VV59levY81dfGbN/7NZjtUxUlRhsILV2YYUdjzp0/x33dETtPeR7P+UN/mMOiojNCVTpMVJxPve2SDP6LZkgieLa2tpgeH1BYZVRY4vE+93/5K9z7rvtxo4LRqGKskLrA7PAS1aik2h4xn88petXcwSvP6yetJKjkVANvHsFB8J06hn6tV4L9Zs0DW72uRlAquZA8JsUZIPhcFhkDxrdcNxmzV5T8vb/yN3jRk29lzyCuD752xk7aNuaO9ugS8H8M3plb5eEseE4rnpqYvp4vac+5EAgBUaEuhS54zrlKihJ91W238Nt37NPNGqQYcewbkq0XEttDZ+gFqH9axkRZJgMkc1ek5/rkRG12a801sObpGmXDHs8RyIRne+I8GubWeygFCtt7wVKQSouXkrt8R3fmOq4fbXF2d5vwpTHpzt/HH17mAo7KOGJqwAh3fu33GQF/5s/9WW75gTdwuSyJrqBM4HzCxlw9kASCE4LJtJxCDMG3OGcYl47ZwQPslSU7SfmP3vpmvvq5LzC99WkUEUxKjLZHgDKbzRCXIwNd8Jb6DSS61HtbOYzNQx1Oq+38NnDZFYY/YbgCdjrl0FeTFnNvrUFSoiST7OuYqIGyqHjyzi7/7d/5u1xfTxaGLLWJssh/saptZjycMGTDvx+pQXPo0scRWcdNcxKjt3fWLheEK/tu0YnCGRIdE0p5ikG3n7rH71+EL9z3AHcezzF7Z2hcSSuGiBCSQZXcRVoypjPEuUqOVgw5wrViSGEwaOs+4zDiY1mVvYLrExE7eTheg/aeEICN2vOxEqr5XsZe2Scp+N6blX6+jQjJjLhUG47FcHjzk7nx7A56/Q0cf/5zHN59Nze6lq04w5K4brzFTc7yv/y3/zXyb/4Vvt6iSZEqCZVYCgVNgY6EWkNwOVuoXcBKT/2IHS54tp1F9g94KsILnnwLhULtLMSAMZYYA0Xp8N7nXuOargiNZFGzut6jx5z0enXFusvKRmEV99HF5lFVjDXEmAY6+hPaw3usa9fHrDYcQsD7gKvK7NnHiHOOLgaqqmI+n0FKjLfGdF2HbxtGkxG+9UiKSIKRD0wMmKbhT/30T/Knf+pn2AWqFYEpV53Q7X+QOX6kWU2HrFQWq8kcpRMGbQg55dQwLeEwuByzynnQ8gxs1+d50qzjM3feidYTpKxpcYgUiC2xA9k75f4b0tPGUAhJ6WLEJiitWWRXF6HGEFLJY7+bTxRDtopNPObXIi2wyGCW9Z1x0ZF7OIENRxGKoiY5JVYFF+oJZ/eu5/LvfZHDr36Jm5gz37+f62+8mZ974w/zv7/r/Xztkx8l4LC9Z2gRsnmIxP7udMO5h8MTiHmxoSSOgae5LV5wy9N4xXOey5nRiHB4TBEiSA/+A+OiIPnUpzKkD3N0ceAafZCMwNqJfDWIQb4lh84TediiwGvCp0Q0/T5UJfRs46qumR8fYquKUVUyn8/R6KlGBV0zp8LgUHbLEvWe66sxf/dv/m1e/KxnsEcWnxp6kYqYK9CVa7n93MlXWw03rzRoy5A0e1KWqAkjOe61CSZGpDQwmqDnJiVnJzfx+Xvu52sH+9hqG6nPMO9aNJVUpdBzNQkxT2AUEGewxuEMhG6QeumVE0xaHrYCJhm+va1ZWngbiXWy8OnY2Pq9WmaJTY+TpRXZoPzvQZsRsjFzCcpM/mdkwDph7oU7knBQ7XH9k/fY2T5LdcONzO79KnqHMJ5Neemtz+UlT3oGo6pEQ4NqoguRylpqY3AmQwPBgBchoaSYT9gQslpC5SzpeErRdJwfj3HaUYSEbwOFM4gYXM/6Jwxy3/3nWcFa1/MByyTRqTMmV4aaf1CizMdaAZA5gREKwZqCBLmsDSVET2hmlJMRqsrx0T5GlaJ0WI04A1UMbFuHf+Ai//HP/iz/0U/8ODvZiFEoUomiKVcAJWMQY040bTbXzKC5tU3SM9FF1g3ayQBv1YO3UhBSh1FDYfsO0zFQS5LzptTxuOK6W2/mqYcNX75vnweO93GmJkqimwqjUohGCSJ0InhRQhJ8gE6V2lancl90lWOkS9rGI7k+AfKZi3lf1kWYU2WBrraQ3QrGEQ0kNb3iRf/6QwgvS3DX0FcekUgxEcUQnSXZLR5Qxywqe3sFZ0cj5IZzXHjWMwhf/Bz3ff3rXF+O4eiQYj5nuzCIE0xocJ1Hokc0oiZXAGh/0+pqTBM8TdNQuoKtsqIEZDrl2LfEwlIZl8OKEHGuBGOJTYu1Rb/g+iz3wn1dqQM8kSxZO4BPqaYYFEg2/QMgagCRvju8En1fSVI4RBQphNg2mBTZKQqcQOxaNIV8gDUdz3/ms/kb/7d/wK1nz7AFmKiMjGSwX3Pv8dSH8N/MKXe5ICpdadAAO8jPnJL9k54nYekTA8JSYcgmSpSKJGMMI9DdnZqnbN3AXUfwtYuH3H98TAdIMHhRGsBYi7MF0TnUOZJkATcZOkithJfKesZzMAaP5PpEgfA5LVs0bNeeeX1yww7GbCi3CYP3JdmgDfWMRjOZFZYZzijQWZOlfiQRtMOLYIsCI0LTCpdTRSu7HBbK/dsjdnZGnLnt2dzxla9jv/51zh6PmM6PcBqpJLLtCsYua5hJ7NAUIEYkJPCH7NUVbI9ouw7tpqSg+DZQ1RMsDuMcBE+bOpwzGS8F7GC0xGSZcBmg6lXDZK4w/kaXCaeTTnhcMIY2vV+dzQCBxpD7N0SoSpcJqCkQmobdakRsOwof2KoLptMpO5Mxr3jpS/mPf+aP8rSz11FmbAynSGkEN5w5p6jQmm9WmBnXXL1eIWAwWSte2smMBhJ7Fzdhe0DVa9/MQhwOSJrAd+yVtWwBI0V3d+HW3R3unsFX773IAwfHBCPU1uGlpDGWVgNNSKQIztaLzT6EEyIDOJvWtMB4pNerYFbfumvqK8bSwmNaORLWFR4exhi8zoE8m/qssCHhYtaoyzy/nlArYK3DiEFpmYcGSSWlWCghpYquOMf9R0oRE896ys246iw3PfP5NPfezf7Xf5906X5GzYytpmU7eCYpUhlDSUclibErmB0d4GODqQpiSpRFQeEcgsfWNb6LhBQRMYgrCEaJCdT2PUDFECXTPaK58p6dNPSrJXFXFImv8BzjdzhJTTSLYIpktoIkKJxFYiTOG1LXMXaOuuuYlBXd9Ihta/nZH/9xfvSNP8xNu3u4EDjT2yQHhOBJAm3IooxmKKHuFWflRNLlmhrmoZvZemo/rW166WsHzVqitlcA1dgXKzmIIacDXM5YpqAUtoYIpYEzNvdxbUF3x/CMp53jgXiOe44T37h8mXuOprSzKcaWjMoJaio09p1ddEnkG8i2iyqFa3FjH5dr6t9/rrJYJSWmqx5dSy6eAsEuHeohzb5ez5qNYtlji1GWz1OBWZMoSoOrXG5pEwMxRkQdQRUflNHO9XTVWX4vgowKLo5r3M717Dz1WWwdPEC4/z6O7rmH+y5eZHx8wHbbsKOBnZRou45q+zqcgS60NLHDe8lFyOIwQQkpUYjJ3b+sJWogRsVYi+0libwxJAPegNF8eFrVhXrFSa/1ZOH0qpEbDFri219C6rHCJRojxoIk7Z0QQ+waXIIzO7vMDw6oQ+Qlz38Bf+zHf4zvec6zGMqpd0Eq50jBU7gCHzzOCM4acJamaxjZamHIFpnLb1K4Kb5/4UXoNci46ArZ9UTd5spSyPmpFFEc1jgUN7RkxA5V7pH8vz4lmXrgsVNHI+gcaIAjYN/DfUeR+/aPOZy2ubcdZlEIP5Q/pBXe2iPFytav8hif/+iviwana9nZpd55FIhi6YzDmwIvbklC1vUNbNK6B60rNJusHLv0ZKLksFQlHzJdAC9zKHJYlxKYaHMRsQHvI6PS4pvMVyOAMYFJ7Lium7Izm7LVzNmez6kPD9D778Pfdw9ycMBW6JiIIjFgBEpnib6FLrBVV8xmM8rSUVlL13UEiflkjbmjj+lj4yj5PqV+97rUe+Urh1l6yJTJ+tzMbcF+Mebu82f4x5/6Tb4+qZmXjmAtMSqWYsFzVNEFapyQJ1wa89Gsv5z1Djg0t4BDcAgSEtefv47nP/NZvOE1r+FFt93GeZcl1iuBuqdaSFJcShgrqKa+6kPoUsCYvp+IprUm1oMxy9wY4VrWcEhSXXcAdImbXT2wXVFDHTRSV4TZhidpLxslK+Q3lVV3X+h6W5fFhjLPLLDE6iKZI+V93nRdSASfCJrwmn8exVw1nLM9YZekV/w89b7dYwoXV5TKH+o6KDdof0NVc78kMXFxw8GgqjgNBLV01Yi7Wvj0vZcJ423EOmJvA1WX4P4Ve2ulqsPoeli2uqkXbr+EE69zsvORWX8tEi4FRtFTxdwoug6BOkZGsWMcPHWIhEsPULQdzKfIfE5qGkzbYWLAJqVpj0H9GkFSRZdijFdJPC8ztrosej6lGcwq+TL3S1jyl7wtuegt5voz/Ov3/grHY+GyKKEocRQ4tZgwYJLSmzLto5F+jVp9nL0786DGTJKiJmeHY0+jKeqKrdGYSVVyZjTmwtmzXH/hAuf2znD9dRd46s038aQL1zO2GfqwOYQcmv6tPbQ/VFZd4GU5JGsNdK62Pq9d2PwY28E92PMfTlysrHtGQ9h7MivV1xrq6s8UHrQV1VXgsW95Wl5O+ft64mfuxGfuHSAuAV+awofunHJYjogqmQ6jNpOOr1VG9VF6A1aHTkkJ0YTTSJESRYqUKTAGyuApfIfzkTJEXEo4TTgS1gxE3yvvykOtHXmEZX16YmMphmRH3PikszztaROCa/GyxZxEUqUWu5ZgSSc2cnoY6+/xHroSXOUYIKsWWnL3NcgqrwbwqpiklLYnWykZyF8NZ2UpsvpQ+/+R3MtrRM14nEFIdBHi2n6CzFUWiR3c2xVDsKhTvYo7ZJ6AmIiueBKh/9KuO1sUIB1oDZQue8ApeKLJW2moLHv8d5KD3kNGEsEkWpPI9R5FDi9Li9WSHoTAaSbZOhXo4lVxn4drzB4ObnRy04kIVhXXtRwd7fO0Yg/NzfiYYHCSuVJDoOFXNoxbESf6dqmMksV7zXGS6gorYeWzFdaIY51TusYvZVl3+ETLBLvHf2PrckWYAbSVPi2//HY84ZkO8kAjyb93NaT9m5E1uRanZY8iUqyoHKWVf5sloqbBQwoBLUrEZFVORdFvQXzz4J63WVMB1p4+kQX8stKKOrfIetO3iltgXUkpCof0JNkk+UA7GSY9FE/w0c5CmSLX7e4iVWQKbFEt2i66PtlnZBGOq4AU/c/Nyrp8og/VNeMlPbS1iIIWaKBkj214jtETzoCuZgVlQZ5/OGvlO8KYDdlIWfVYBkFQyXyj1TlbK4UQEO0eMtTVb5qr+8gAzJPv06hBBkUSjRiB2C8nS6YOFIM3YHpZfWNIKmhSkkasPH63MJ3C41pVXcj1vm5xb+VkYwqrdJoBhJyFlUX4tiBtG137+uR1TYHlkRqzCNP5EbvbO0zps3OACeD6mlZj6Eux8hIshjWrwwb6NuF3rJZ16ckDPxeJGyO9A6GYQYs/pvUN9xCh5ONp0J4AxmzwSczS0vcGSDShCzzFLMDbYdFDWuOunPr65lux2MzDXk9XB/Dy4nEsAf3eQ8D2Xa6MQkraJ4JMTnV/k9fOgx8A2gMFDw7JJV1edSWrigjJ2h7wN9BfVbJ8qUrCiOslf8yp1+wNrujgnbhqEsToqT9X4xlPtuiMowJCl3Tb9Ui/B1MBMSAm5TyfGEQjq+IMSHzwhipPKIO2AtL3JUbWmuX+ipGoisZE6qkUZvW5qk9og+aeGBNs18CkNQO1ks4dCs6VvCsSiohdMNtP3276oN+Xx5xSMQ9prBU99QYbyC3fBjqMgFL2xiyAdYvkMimhMWVFz35NGTFXVA18q4chrHlqSz5gvpZuacBWqreWnp0Iamw2Lj15eDB0mXpjFt7eyWvGzUxf4XAym75yjVnkcaGXrYKqYNViLHjvKbDUGFmk0k02ZLjU60xqphr3hkzNei/TJ26Iqf1+OlFxbwcQX/MWGzLjSo4WZF2LXK9ysOkp2crVf38rDdvjb8xUVk5MXVOaPYmYL3SPeiNnFhH+I/SG+NaWMwmykkFbkSdeoNhm1TnrF5dZ68VhxGHFZe8GUBX08TZkuiT/njTrqecqtl1CxK5gXLIivTPYgkVVam/4TN86c8WAcfpVevJ03xzx6teeCtTXjfTPE9QHqlEFwKgA5hncpu7B2D4uzsogy/ZASy2PJ7ba7UkDdMUhvuD3xFwNAH27SVDN924RMK1RktZD1ycCNu2u9WQ9/BBl1aoMNHa7/vVDWJ4hi2n5Nh2yAP5YWKmh/KNfMWYAcPuwwPQSdAOI+9gB6MeyFdeNWNIVAYDBONt8k9IVKCkLlWGbHmSC9GHgdrKyJ69yTbrsDbCstBBsBJN8FrHSHhRDgNATZVkkoxh6dPYeY1pgaN9eS279/vUwj+EKEuvCm5Wr70V5GGHnd45n9ni4Sk+YGIC17Kuu8NktJ1gXKrnYXrO+/4KWwONXkpNW2gMOTjZXvB9zqgHsAU0eNuP4qkzklT4Ij/BqNCEErLplkNC/3VyMrj1HyyzjymSuVCr4ttUTMg8dIp9M8py4PpG2rXsi7OfHbP8er8Uk8JiIXgsDZnKx7+pieTCP90QN5uM5hi5MJ2+KnNrxZ/1nKllmO8qD5SuvfhUSJg31uY/8+YaAIZLE5npVUYpeRcJnNAnLkq6wytJ2fZbzD8YZbK7ZnvyONmbXxLF5vFbUSl3lo11EJ+WV0omlpSuh1BOtKFpXJKLMgxheTvzOsmxqtaT+UV7lhNLLI7ouOrLmsjiWhG1do27LlezYVQ7Wt7VFW5cA06v40ebbIJh6AlAzHpsnp4/VO3pCnIkJXcDVp8WgV37mJ1LvY3OV98gJ7zGcCBMNmWtnH6NasD7q56UTX8tK+Y9ZZDDN0HFerm1k8cRBd8zDNmLyhN9Lj+tIJx6P3DilNS/n8biaR3VdXSBLSaCrmrA+LOsr7J4A3progz9WPbO1UPMUEObR4l7pMTyWFI8lMduSFUas0oeYK/xH6T1BG9akh75ThpwMt59glvwJEmamh/j6aifDMqUvj+uVx6h0m3LF3Aq/blnQ1f/GwoglkiGjOU+AVnlW0zqIvKqeolzRuuzUuy/LRMYjvfJYsMM0GFlZqENIH30WA9esV1GO5GqAQX8uLEyiewJ4K4/WrJqH/Ik8lBv6TVC/+AOAmZlHfVPkiXAdyNGP4LrMypmFYJ2e6JR1GiaoK3/zcd9Ekhae0iCJtHYvh89wFZb8yQzoI70+9o7wcrr6hnJF85PVemFWD53HswLgsRB2H4kh+jaIp90Tw4g92PXbzAd/VNdMGNar+KCm7zw93K5TVUIep3uXThxF5mr3Tq7ulT3WUPdaPTex5G4vmO99OLmUls/ZU7u6TuWJsH8eS9z4MD7C1agoTyAQ7QnimX2bG7JrMdSsUD1O42bJY96838x7l+TaGKcnQsC2UGFZmeylJ2x6b9N85zZFf4Le4+/wlg6bsRmb8QdlbIzZZmzGZmyM2WZsxmZsxsaYbcZmbMZmbIzZZmzGZmzGxphtxmZsxsaYbcZmbMZmbIzZZmzGZmzGxphtxmZsxmZsjNlmbMZmbIzZZmzGZmzGxphtxmZsxmZsjNlmbMZmbMbGmG3GZmzGxphtxmZsxmZsjNlmbMZmbMbGmG3GZmzGZmyM2WZsxmZsjNlmbMZmbMbGmG3GZmzGZmyM2WZsxmZsxsaYbcZmbMZmbIzZZmzGZmyM2WZsxmZsxsaYbcZmbMZmbIzZZmzGZmzGxphtxmZsxsaYbcZmbMZmbIzZZmzGZmzGxph9xwyVK7+X+hskupmfJ8QmWblHm1vyxBtuMwWPtxVb3yRywrjJKSeO0c1m+mbcBj0515KQfqa1/4lI2kzWxjPbjEfqnZ3cXKtfb27cN2fepbdqp3nK6z7zZmyM2WZsxh+AIZsp2BizzdiMJ5xh2sTsG2O2GZvxbb8BHqkhG8C1zdgYs83YjM3YjI0x24zN+BaOjQO2MWabsRnf3kM2Rm1jzDbj4Z/6qqgut4mIIAioklJau0Fd1+GcwxhDSgljTP8am3l8rMNaSwhhMdcPe051k9HcGLPNeOSOgsia4duMa2zQxDykYdowyzbGbDOukfewMWbfPA/ZWrtycCzjShHZhJgbY7YZ12yzAc65K0LSzbhm1mxhzHSzMzbGbDO+2cYMUkprRmzjqF07z8w5x+CQaYYtEXkYwP/mHmyM2WY80jCTTZj5TTVmdv17JxzfdOK6GRtjthmP0jNbNWaquvDQNtHmtTFmfXIY3bhaG2O2GY9hSFo785OYHOqwTP1bBUkBpwGDILrxEhYLWJclSSrLx2k/P/k8qwmXPJUMxqz/R4on/DBz4v+mv2+bsTFmm7EIXJSIstw8ESGJAdFh2zAxMDER1zXUCtollASW72iyk9FcJC69vluU5WP4ORoxolgU1UhKAaMJI0qRIlvacuPOUggzaSJpB0QgkhZ34ZRtszFoG2O2GesGbTBkq3ZJelKmkG3WpIAidRgRRARjWCPXficbtLWw8SGMuwzz13tmRQxUmuc4ASKKM9J/pYtQf7NRNsZsMx7GLRAMqAHNIaVFF5vJgThgpx4hMeYsm9BXA2xmLz2U8cKSVIgIKgaMoAtjlbAGrCydXIPNICXkiozN2BizzXgkYwHaLEKnYViytvneaIRJkYSSJPXkzs1mW3hjsvRoB29tdR7TCrg/lJJJUiZVScVKxL7gZ8jJu3NC92yzdTbGbDNOmLGTVc2KaFokBgxQAGfqmsoYUgqEASRCNlyntcWcMCSEtBZ+am+gFt9KATRiVTmztUU5GLPECQKfPOgG2Uz9xphtxtr0nwCYT4j/yeCZ1bBVWDR1AITEglLwnRxiZgdqPSM8JAUGg5YGZ8vI0isjURjl3N42RTZm4oBF7H7C65U1Fw02+eSNMduMK0733pjplae+kgZjJlslbNcFRhIikZR0wzODE0YsIZq9spMLe3WurAYKUSonnJkYLD1VQ1YOk83kbozZZjwag9bfidP2T4o4oLawVRgcEWcUTXGz31Y8tNMywgBJlxCY7Q2ZBWoSI1EmJk+91dVTpJdY6k8Yc9IL27T6e0KOa9Y382S5zcBSv1ZlOI9XcfUq4/60713t8z3c96uaN6NREHpvS5b+hjU5t2kiPO2mHX7nC0eY2GHFAC5vOP3OLT5PJ3qOmhM6Y6o5Odkl0OQZFxYzn1Fo4uYnncsbYPDGEn3srtDLAi0SB6o5AXDKWrgW6/ax7h+5yvu62utcq/WySg9aVqbIo56fx/I+zbXY7Ke94at9/1oZy8fNi1opK7razXk0C2WB/0jqd+OApeXtNLKwa2DLwJYxObMZN/7BQ86pgZh6+oUBCR07ZYFtjrlxC8o+jLespkGz8eqpyVeePqroaT/7JoxhPT3Yujrt+9+qw+3ke3ssf/ex7nF3rf6APMIT61p5Sd9sz3L1756sj7yW72dwDOSEp5FICIYSZAv0pu0J04M5k2qHQ98itvrODtElG/zltKUrzmkxvYeWEhUJ1805N6q4zkEFFL2hW8SOKiD5MFbJr27pD5lFnZnpobVvrUd8tb932vr9Vry3B/sbD8cGnPY7j/Z9m2ttBJ5o3tTgIT7U48Ged5pHds3CZwaHQBclOSqG4T9JigXQRAE8/dwebnbMxBqsgNFNVo2BdMyVPTCXkaNC8NQo8WifZ994gTFQgBh00ckcIwvGi/bijIpe80jjWq3pR7tHH8lzH8m+eaxR2WMxwO7bbdme9Ia+GaHst/K0lZWK8kTCDOXMSXqgOXseJkUqa+SGCr15a4svHF2mGu/RpsR3dB7nKrd/oGMYgejBWcFpJLUN141KbtozC34Zq14XNrvHMsgChf7QSb2htMtT6BpiZt8sh+Cb7Ww8kaSpHvMueLCb9c0MCf8g6HvJCkRjUGTlvwGUtpLNW2ktJTBJ8LxbLhCOjiiNbrpxP9QCN+C9p7BQWIFmzjNvvokRMAIxJERjTnsOmBjrLDKhD9meYEmWhxNZfLNDz4faiycxv4fCAB/LvjbfjIn9Vhmaa/13Hs7NPgl0XuvPnVY9joG5njKz3WhiAvKkMdxwbg9N4UrawHdyqLmGpaUFldbkqnysKruTmhvPOmrNlRVr86e68NAGgF8fpD7zW1G7+UiB/2/9YWEW3cIey+c7GW09mv3kruapL146nWL1Zf3D5E2oJ5bWlSlhFa7ANVSW8isnv78izILo6vXhlSUuX0/Wfl9XfsGsEL51pd6xf/f5c/QpeSVmJYuhiE9kvfxFM9Ksy19an09dfrE+F9lYDUmARMbD8lNyazkl0XUdVVXjgOfdfJ6Dr92LN9tE00PU/UY8bWklSf3fXCgRroS4qQfSV95T/zoD9eHkfJ/0CE+y7dfuN+uM/OV74qpr61R1nRU+mekL86NZ4PUnXj/rwsUIVWkxXUMZWm4+v8sEmAhiNSEykI/7tdZnKvP6jStmsp+fHlMT+vKoqxmah3Fs5ek3Dzqvw2c+7c+oygmxzitxXRWzZijWEiWqSB82n1z/w5uRE3WuV9xbvdKwquqpn0NPnZ9VasdyUayXx6ZejGH9yXpins2qolZAaVMgaCKheO8RY/Iba1sGT9uY/AjJo6pETQQUD3hSn9KOqMaMNYgSNfQ/ybso+BYR6Z8jhN445MWlJBKB2Ct+RVJW+upnL/WvK4gOj0E2J63gG71nE/LPfJdP3iQwi6k/dXXBIYoqxP7nw4LNxkoghX6xaBbv0wQx5K9jgBDAe4gBye+Utt8OCZYVN/1DEZIIMXVAImpAUGJqSXjUJpIJRBnmwFJUNS1QA08fw61FZJsOfIOzfVF6Tw7VBowHm/KcqCSiSURJfZKhXyyDeqFR0kDc7V/DQE5AWEhG8RLxEklGqWsIIZBSwNo8v67XLAwJKGEaE6mCaddhTP65CQkTEqT8Xhek+wQ+RBQIMVJU0IWciRQgBEWDxzrFipLahJV81iYgxpgPAEnE6MH17yOBNJ4dlC3f8oIbRlQBKhJWW0RznSYm9ZMVEYkYDdikuDUjadCUK56y7lnCp9jvn4hXT2BY5ytR63DPYyKllHXViPk5scuVHiJEHxY3IHa5bE37OlJNHjT0L5QfwbeoKiFlnbZ+JlBC3i8G8rvLXmbQgWySQHvS9UoCKSFrmnCpD6+HA3fYdyor6zp60PyZQvIZ+Y0+39T+PhOUpBl9jP1j3nT9X43MpodK9Pnk6ToQ7edFyTPa3yONPTFTV3RlVoyZINjeYR6wgaEhrcudNLKRqUf951GaWQsYSlflExFZ96oW2SFdNFgVY0kI2pcDW1fmm9D/bdOfigQlxbhwHBUl9SdQGrpN9P9Og/tg+hPT9N7S2qNPdCmUpRAX7zEuCKeDNzV4iWuhnhhC06A9ObDpmpzHT3lxz5qZ4iwpRXAu/57q4tRID4JbK2CsXfxiQrO0T38odP0iVxU06ELfzIFcBzznhjPEwwe4+boJ7fSAFANtm7AFjEaD9xUR9aAeUY+SHwxF1cGToicFj8ZAigmCLh6VFWIXSCFQljY3ze069i/OGJcO9Q0udYT5MX7WMK4yJy542Nk2dE1ka7skpUTnPSF5isogJcy8xw+erAFXWIxEykLw88TYgXrwLUwqoSoNKbTE0OCcoeny86S//wklqZKsxTgoCqgEtqziZke85Om3cAaYSEJjRwwheyayMN1LETk1CELwieATMfUeoF2GVRoj1tgFzmkl/ztqJKZ4hd7c4P2uOT+93FBKCVcU2TB4nztGqSLGMO88GIuKofPd4gAuigJjBWtlma/oPTQjeR/GGBdeZYyKYPJnjivukGi/59fZJ3bVXVzDDNPSkNj+tBEhpYSPAetcfgVjwUewMmxFpl2rUZW6LvExoAjjyUQwNq85m11tay1t6Ho3C1Tj4AXkv3VqJJYUTQkx+U2FEBCRtV6C3kdEhBACdV0O9mRg2uTTQPJmHCy5LAidOYDyomAdCvjgqW2BCPj+FYqUxfIWKSgDXYoYY/upyzV0Ftassi7CtKXLKqTegpnedV52Ap83UybjUXbztc8eqpCMwfehg6M/pHvKRIoBUxZ0Keb31sdIxmTjKCRKIHURUxQAdJp/1/QBi03rYaY3yyLyIeQI6GJzWwyu/8yiEEPCOEPo82si8ADo2+865HOXjonVBKm2OPaJZt4xcSVGlNIoSeLCM1MxoC4vGDUM5j2XXhuMGkQNVgVRxceAlIaZBnwIOOfY26poDwKFVcoCfDtlUu+QEhwcHmGqAlOXzH1Hktwqb348ZXt7m6Aw7zzBWsaTkvk84awgIeE0QWgYWUtzPGdn+wxNq6hRSpeYtoeMt2qiWEIaYQo4mmejJRGMxD5SALEObTsmGtgJLTcVgR949nXsAeMQGaPinCFGXRoXNVfgUatwUIyDcQBXGIwYfNf1vTf7w3+N52aWKMSK1LkYXfxewmARYshd61PnMbYAY3JcIvlwm/tOC2fEKKQQGbsSHztMUeRDUGNfFaLrWIoUiy99myiL3rC1HlsVWUFEBtK2WXv3outRzrDHo5g+yw5OLG3nKctiEQFagBD707p/T9bQhEjhLJ3P9bNJEp6g1lopMEgSCmfofMxJZQOJgE0wHBngwCydktUqQNGk2b1bOW2igBiXk9JLiIXEiu5ThHKQ50z5j6eB2gBI0P5mJqQq6MiGK+QlIyZ7nyqSPY1awQSFEEmSoHQkY+iG5/TQgGHde1o5U6/AKbR3ug3LUyeljsoYRBPRe2xRoQhehMH3sYBLK0eUwLHvMGWJ79+/9k5/FlZWxlgpeuMUQiChWFcsPFWXlsCBmp5PBqjPJ7hUxcKwZwdYVX2gEiujwqycpNBFJYowd+jXgLd/4U4O1bLfBcx4G1eMGNmCbjrHJt8brBxqDkI5Ri2iYHupoSiGJAajko1Z6o+HFNDCoqXNvKsYUR/RJjAqDK5Q2tAiSfA+Uo0m2MJxNJ3iRhVJFVtXzLoWNRYta4JYWlV8SBiEQqBAqVKi9A0jZ4hNR1HUdFE4mh6zuzMiaMvx/BhbT4jVLm0q8ClRFBZJ2dvG5Oyvxp4c62DPH/Gyp97ATWNhFxiTpyQlxRZyUqhksYbTCmQ8rL3VzWOGjauavTCNa/Qh61z2qh8EQVMSRkxO5hhD8B5XFEQM8xQ5ajudjEZ9QAZlv3eIkdLaflnkMMwhmEG9UxVFiMnQdoF6XK5hTMmDK9ZDtXRKOkVWdptiiP11dXt0SbUyskBdLXkfO+tIMS6ij8Hpzx2x8ivPV+xBG1THLh8j2UAPck4pGzORfBBf1ZgNceiS7ERE6BSOo9dZjFBV+URP+U2OrCW1nkojZ7dGQoqoLEMrWbA/FQrHLEQOQ6eNEYLJap++i4zqEpcS+JatZNmpKjG95zNrWw6jVxnVeDUkzZ7QcFP7CJAUckZvwIaWp0j2GyksMUZi6BhZR0XCeI+Nge2dHVE1RDGEnhqRKx4FSTk2xximTUMaVSQRvb9r8DFxZjTmUtNS1CXbCM3smBvGW2I1M8qzB9r2p/Xg5RlEczIi9FM0nTYqriAVliYlWo0YYygEKrH4wyPOTrakNpIxCJddOu+VrhTuTOh9CT78uS8xtQ6pthht7zJ2BRXgAjhdLozBYxWV7M5Ii0ogGrNQ8BAF2wOx1lqa6OlIxJSoihoblRF5Xo+6KZPdHcRZjo6mVK6iLB2zaUc9Kbl8OMWMRrQId+0f0IrB25IkjqiCQ7GaqIA6BW4cOcbWUDrDvPW4UYVYiEEhdYzHFUcRPnPHJZjsElUQawi912RLRyFggmebjvH0Es8+N+bVN52jQrEEujaBrWhDNlKxP2BSf6MGZr+KUluh7ddPYSxVUVIM60+VSpHayLq444K3llMJgzEbanDXgH4Dw29lh6YDW7LfNJrKirkKXgzGwNHhnElZsFc75odTxpVjUlUiRJzGnoDdh4MqpATqLJePWqUsiSI0nceYnBQwDoIGkF7/TZe9DYaoRPo3PEAw2XuTrIGMYT6fsz0ZUyq0s2P2xmPGxvQ5stB7i+DbQFGVvYMBX/3Gfbp7wwUaA/MuQ1GFdVhVKhHwDWdGtYimHLFJWhqzUxMa4AZhTaPap6QzCHjYtXrv9Ii2LLMLah1VUSAejpuWLVeQfAYDRZYCnYNFy98zdCFx2DYaxzX3HVym2N6hTYkuRfaMYX54mV1jcK6mTImRsQN4rtOYaJqGuUJSm/WoMBnUFHot/Oxl2SWOvehuFAV867FWsMZiQkfRtlzY2oKZp0wRa9yVWRZdyiIkoBjXfON4qvsxUuxuczRvuO/giK2dbfZjx9cPD7lhtM03Do/0XD1iu3RigdIVi5h/UWcp2TAP2Z2jmGhjy/7hnHJrTKs5pJeYKBOMYqCOQUtbSCID0tnICAVwg4HCwLPObnPH4REHs30u3HievXF2FGrIISMWo/YEWdeijDOg23uLg9Kt7Y/poHChdESgjVDZXPROA+Ox5c4HPMZGtrYtN+xMSMBsBjuuZGcbrtueYB188YFIbKaYaivDIqrU1uX75ZsM2qvnxpvOct14AIorpin/uzZCbCucycf52BlaTYSQs3pWBBGTAfvOw/yQUancem6L5950BktEaJjOO+ZB2G8bqLZ7o5Wh4YHOoT2GpJIIXdfXdRocSprPMlYmwtg5RorulCWVMSIpq28UYnIyJyasE5ZiTisQlC4BqhgVY3P2RKxjqp0excDhUYOMd2lR2iYx3h5xaea5+MA+1+3tMAueCiixuTdE7xmiGWdNIrQKxyQOp0dEY0nOZfzKWXwIJOkFLdUsqklkmZdGU8oJMSMLEEcHv1VBioKu88TZjHOjEcnk+2JSzPCBBlJKFGVNjEoQYRaS2u0tLjYtRyrMUmQ8HlE6aA7nuOi5rh4x81EnzopKDsZ1hRZgVlPlgzFriFgEJ5nPhLEEgYPQ8I3DQ/7WP/y/cxgiXsC3HYVPXJjs8Bf+wz/Ja178QjrAifbZj0XmOsfECskZxI35ysX7+cs//ze51HXIeIQidMdTztWOn3rDG/gjb3gjO+evpwG6GHUugi8c/+Ff/kvMEGIaMo2G1BuIZISuN6IuDZuwT8mLycaMkIHS5HG+47qq4G/++f8Tz3vqkxmZWqueOGl6l9ouk+4ZaDQwU/RQlfd84uP8k3/+z9HCEYPiqpJZM+XcZIv/4Ed+lB971fdzoXQSgRg8pcuhnIqCpAXOsNpU9lLb8o4Pf5B/8Yv/Fu8swUBRFDmrM5vzv//CP2OCMhEQZxYbLUXFIkxAuqj6kpuux4eO6f2XqPF87Le+wP/4P/3/8km/oGakBTa0UJjQPHFRIJnlRhsOh4QQNS9oYwyp6zAhcePOWf7u3/l5ilHNmd2Sf/e2d/Av/pd/ibWWG8/fQHs0ZT6dY13J3/9v/hHVZMTOzg77bT5hTRBCyJ2TnDWM6wnOH+MDfOAjn+EX/ud/xgOH+9RnzxI0QTPn7Hibv/YX/jI33/o0jO2FFlNG+5y1SFKK4DHtDOdbrp+MeO6NZ7gesNpQiZNupHrvvZf52//g/8UdD+wjpliGV4YeV0wZ6ljwQ/JcGxGqouS6s+e49alP4+k33cT3PO+5nBtNOLu1o2PnqJJSgYwL0wPhAzJmTuem6QqNCYgYOoTGwDs/+lH+3//jv4BqQhOUra0trG+x0fOHf+gH+bkf/zHGWqqTgU9kF1QhtYJHmIPePTvif3vzr/Dej3+Cw66hSSknHoxdGNdBbUTSauepfs0YSDLsCpfpUWowJNrZETvO8aJbb+Xn/5O/xF59FiMGY2IO9/v3E0PAFBm6un96zP3NnD/zn/3nzIsaO5mAEfYvX2Rihde85KX8rT/35xhXZe60pYPgQh8trHIzVgya0z7DqH1ucwjP2qRcbudc7Doudi2jvbNIUeFC4oHZnFAVUJUZ5xGG3GDPnDFEETRTL3QWI+X2Fvsh0FYlh82ccjyh2Bpz7/E+jRGoKkL/9tTmkpLjWcdUlCMRvGZDFkSX7roRut4UFQPLQHrW1oBLuYJpyv0mR67AH07RrQk6mtCRiZODlpWVIXmQ3TuVHJsfpUQaj/jFt78NP6ppUqLcHnE8n2Pqmpkx/K+/8mbe9IYf5Bi07TxnykJSynSLFdJebzYTObxFi51tmqrk/q6hqLZpUNR7XEyMi4IZiWnynKGkMrKGBZYIqe04X5Vyb4j6XbfczGw2Y1sTtp3TzQ9Qa+gsxN5gDZwsm/J9yqbc9l5JDo8MKXs4ww53lhBjxiedpSpr7tu/xPaZPQ6P9jnsLInIzvk9Dg4OuPvgXqpoMU7p1HPuwogv3TOj6XIIVZiCkJSqtKQEUT3TrsG1c8p6h90zOxzM9hmf3+I4NIg1jOqCo+khIXiaJmd9VRVbVUgCS4TYMIqJkSTObtV8101nuGAy07+Qgnma6ywmtq87z1cv348fbdF6Xaz9gYuX+sxdGsLNmLlohRgkBO67825+9867qYzwC/9yzgue/Sxe/4pX8/IXvZinnDuTA5yElCZjqIOnM0AN9IZBeiqbMULSRJsi3op6MUxJvPn22ynOnefe/WMmZ87xjYuXOLs9xqD82gfex0//kZ9kFgKlEawx2F6YMnvc0ieVoD67x0EK3HW8j9napjVKlzKmZhZJn6F+fmnMVBIq6+ogRjMYY8RiEOrdXebNlIOmYbS9zTxEnDFDDg8xlhgVWzjaBFNU02TEv/hX/yvd9piGgoPZDFs76nNnmM2mvO/Tn+Rn7rsXOX+eamc7uxdiFqG4XIUca5bguGYgsmei2KpGi4K5CG5nl6OUaIxhhpCqkmgMXYq9p2wRHAa3IHpYI0PWVoLmMLUToTWCjsZMFaYKZjTCO0vscaThPbYxodYQq4qDFGicJY5HzK3QWEtnDNHYfEL0GbgBDxq+lh5pjz71ZduWUVkzP55TSrFIHjjJ75yY8qOnV2ANDegD0xkf++zn+MrFBzgywrExHKH4usLbguOkHPiOd33ko1xqOigLOnJaXYxZECN1AIr7NWNA5tEzSxGZbHEMNK6gcQVdWdI5Q6uKqQrCGiMIrO2zpMYhXeJ6Z+UmI/LyZzybybxh0naMfKREMGQuUiSBEbZ3d4gp5FPTjkl2RDL5EaNBbU3QgqgFKiUxGKypKYsxzTyQfGJrNKFrPJVxjIqs3DFt5jlx4wydS8wlkZzhsANbl0iRaRBdSBTGEvwge2SJIriiIiZog8fWJU3w2RiHQKsxe/nWYK3QdV2eRAtN9DmxQ4drjzhrA9/1pLM8ZQTbfaa8QChMJZiSqQ/5EEqBthDauuBYlNYaGiPMUiA6h7qC1K9bLWumCscR0mSLpijZV2j3dvnw73+Jf/yv/jf++j/8b3j7Jz7N3Y3ngdZr2yeJguYkEZKYzQ9VJBGiXySXYpfZBGILPIa7jw74+Oc+z1cvXuQgeOZWmIrC7hbHBu7v5hwY4a0f+hCHXZszt2IIScFkTCkpfaYdoiY6BTOqSIVlHiLJObQsCMbQGccMoSsqWldylAJhPKGtRsxtQeMsjTHEsqTFEG1+nlrHNHmCMaQyY6sDtasUSwqBFCPWWrqkeAOX2457jo94/2c+zcXYMSsdZm+PYxEOibRVxdxafvGtb6VzjuOgmgv+Ta9eImhMy4zkqjEb6A6rfRqXHnAGhXPI5gjkRac95UGGE6Y3gdoT+jQlNOpa8UDmALHyeha1rk/z9o1vWa1V7EmK1kBZ0iXluGnxYnH1iKIq8W2bgXKgJFFJpkiUwxWgaZgYy8Q6XAhU0hcYayJq6htZaG4Wa3paiLUEErOotECoCt7y/vfkcNsYip0tZiHQJaWJyqzzxKLi37/1N5iSmGfPViMZu2BRwG6X4KueNseWYIaH6TOM9C49YgdCcJ9BXWZchNgpVmHXwo3b25wtS0YKaTZDvKd0BVVZkTrP0eV9pAuMihKNKSdI2hYNLQgUNocQVWGJbUvlHCkGNEQmVY2QODzYZ1Q48IqJIAMjMvULz1goCoLJh1RIGRvSmNn7xmTRxAy+L8tZNGYCbIpATEhQnOY5UyM9zysfQFaUedNQ147CRCrt2JKO5z7pPE/ZgS1glJAqgVHbkyAyXyv2OOEsBOa5Eh11mUPmnCOFHlBPikMojWXkKqwtiCEz521Vs+89YbzFxZi4Y3rEz//Df8i/fPOvcoThgbmnUVUjBcezqQowGY9lIEkPi92WhlkbaEEPk8du7/DW936AWUzstw31mT2moaOJntYIZmvMxbblV25/O43JB+5B26lxRd5zqjhZZltLY9muK4yP+OMZBYrzAX94jEugvqOQ7OlSCJQV8+jx8xlqHc6Vmb4VcqKmSAGaGWF2TGlzk50YI0agsoYYOk0xU0WM5JaIaoRpUg2l4y0f/CD3tVN8VdFZw9Q3OWqpSnxhiVXJhz75CS63HftdS5vphgtHQK5SPuXyQulJMNKL2MmyvCgTYqUn2+WNv2Tb96xzNQv6lRFdEHBT7+ZmYL53ZTWn/61YlmdHpgK43gkHcBG1Ck3TgcmhTey5Z37eYFOiTErRNdg+ASArbdpsbwhqm7laaTaH2YxqNCJMj9HQMi5KsQt4zCxKPGKKqCtBhPuahi/dfQ8f/u1PIaMRIUVC1+NwCK7I7r0ofOXuu/jE57/A9z7rmeyMx7RK32B2UVLQI3K6YrjTMgOsQ/JfQXtQNhms9ty3FcBNFh2CzSI7Ouih+fmM+cE+NYlyZ4vLvuWwaZFUUgShEGHixohX2thQlQ4KxRWG0LX440NKDBqFsbEU4gh+ToxzShFGZYGODe1sHxcslUKlDhcyxojpGYimz472xscgPQ63UtLVq+QaMdghi0tBJQXOQEiWpH1WLfrMF5Sekxgjk7qkne+zbRLbLvKCW27iph1hD7AtUjgWUtjD318kjFAmkzFNyLixnzfUhbBdVjRHR9BERlVN6CKEOaREhaGoBB+VNkRc7WgN2NGIS61nsrPDv37b2zk+Pub/+Md+hiJAVRQUdSXaZ8xBESc5Fh3S8qUjARdnHb93//188otfZK7CeGeHw26GKyrEOoJGqqqi08RX7r2P3/rdzzF5/nO5oR4RgFJAk9C1nqoqaBTaSwdsRdjDYLtA6SaIcXRWSCHSpchoPOKB6RGhc7iqQMSQKof4FqeJLYBmhktQSU64lOOa/fkhhSZGBnwzx43GUliX+WvGZc+TzI88Tsp9sxnv+OD7mSrgHG1osWWJ2oJwPEXLkuQc9+4f8Gvvehd/+id+gsOYOGd73t4wh4vqhaVhc4s4XmOfkrUDKXhhhJxaTNJlaZOkBWt4ABCtXFl8axZcHFnydIYQcKhL641cXlxLIvNgnOq6hs73YYUlqmdSlrzh5S/n9S/7buqUcCms6IItIVeALvgF1aFEsW3D857xDOqUvUgjJrsHpFzS04PdnSSmGO2M5c1vfxudsSTjiClhUkJiBDWEzB7muGvYHo/5tXe9k+c98+kcAOclh+CLKooebNcFBUAWBeQn6ztlETIPjYFX6qL6+Rv6RcYVzk8NbNcl3/P857Lz5/8sU2OYGqEpHWY8IaVE7CIuJL76lTt48+3voI0dXdcSm4CVyJmtCf+Hn/lZLpy/nvlR5odhM88sdC1j58A3nD13jge+sZ8b6VrXH3iZtKqqBJO9q4FbtCitIhJSXuBilJg6TIrE6HtCthBTIiWDUyGqElX7Mru8qAqBKgZcFxmT2NE5L3jqTdxcCed6/o7TCMkuS3wHn3+FHOu9x4cOawq2qpKXPv95vPwlL8LGSKnKuKo4Pj7GGAtG+Prd9/C5z3+RO+68iyOFeUx0xzPqM+eIWIJxXDo84Fc/8D5uuHAdP/261+BIeq4oJaUuz5BkOfSYIq7HgkVgDvii4lfefjtztURrcUbQzuN7YntRFMzmxxTVCEH55Xe8jRc/9znsZe9X9wziBly0r0E9N6r1h17+Cm575rOIxqLG0rUejFCPa46jp7WGX3/3e/jsV7/CPHm8z7SKErjl3Fl+9g//OOdGI2oEP5/h5w0JcFsjrIFxhDPjMTElrQWx9IWx1pEMTEHnRvnIpz7JHfc9gK0nNCFiMBQ9/NJpQlVpU2J7d483v+1t/OQb3wjA3tY4V17oSSbpaqH5KRXbCiuco/ww/SliVgNHyQs3rXoLK4XKupJpMCd06jXJMrs2cJuWfknvbaW+/i9R2AJnSwiJ2Mw5NxrzXU+5lV0NVDFmQ7Faodwv2Koa4b2HpBQCqZ2zbQtS01DWuV9bjNmvsdZijCNh6DTpxdmUw6S858MfpdjbZf/ggGJUM3KOkS0oTAYvk4IpCpoU+NQXvsA9B/vUIbG7u6OFIBiD9EW/GQgeuEwO0T4FowNBcXkUmJWmwHIKs1NWDH9MCWszMuO6Vm/e3eGm67+HmQiXY+QB39JaA0WNc5bSwKX9Q1IbEVFqV2Csg9ASDhtecNsLGbkx5uY6wwMFqM1lqCkoW6XwwFHkkm/YNhBsHxrbnhIjCWOyxys2Gy1rMss/SSLaXvbIRExSEiGTTw3EArzLxMwCISB0RlEnxCLfWqORsXbE2Yybzm3z3FtuZk/gAiAaGFsnNvV1tANksbbWTfZ+RbHGYkWxwXPL2TO85kUvYq8sKUKuFRyPakIIHM5mlN/zvVR/rOL3vnoH7/nYx/iV976H0dYu83lDG5XWdIwmI5Jz/PN//4s886k38bJnPYNxRLetFUFQTTkclKW0dwt673TGpc7zkU99miYq1lUcHRxy5swuJKWqKmazWT7yk9L4wKe+8EW+cXkfk5Qnb23RJJgIlGVB6CKJRBUTz7nlZm679RlY5wiqhLnHlQUdgRnKoSR+8+Mf47PBI0Zw1lAZRxE6bNvx6he9kB0Rtm2B+JbSZKpUExN1XdMeH1GrkroWV4/6Up1MD+mAQx85SpG3vfd9qLFY66ANTCqLn8958lOewj33ZnzRt3M6V7DfNrzrgx/iJ1/3WuZK3kuri9/JCWMmK/R5MSRJpF6BIS3IxH2Zkiw7P0STiH2mTm2fJRvOPlkyqEOfeUuyVG8YTKKukluvojiQyX2W2EVUW5wVSnGMnaPSxE1lRXV1URiJoGLtooRP6hGFQCxKCuuIPqEiuKIgSSJoJIlBxRGKyK//+ltIxnI8nSNFSSmW+cERP/KmH2Zcjfn3v/EWPEoxrmmODzloG379ne/kr/yJP8mhj4wKu7DqMsiayCmfVXoO2il9fc1JRYgVRYHkc5mW60+ZGrhhZ1tCzPelBXWtcHDxEnddusg0JKrJFs94xtMoxGCtRbGo5MxrwtH5jq3JGeYzzx1f/QbYksMuotZQVRXN7JAzOzskP8Wp51zM5Vuxd/ulJx5aMTgshUIVEq4L2ZiZLE6gAi7EbK56btJyg/chay8SkFNVkDSQuo6yaxnbmnPbWzz3xj3OGZgARlsmYkVjg7gy1wb2cx5FCDGS+gM+CfgUcaXDxoR2niJEdq3jBmuorMFVBTVIcFZ9nSsOPIkXPPlGrj/zep7/3Nv4R//s/8tXDw45c911dM5yfHxMc3TAucLxb9/6Vl787P+UY++Z2DGiihGHHVBhm4YqF4IteNcHP8DFozmhcLiioHKBMJvxxte/jqqq+Ff/9hc5c/Y8l49nlLbAi/Br73wHf/GP/4kMNRjU+yils2j0jEc1rizE9xU3KBgjYmyJNXCUCq0H4m7KdZ3WuuyohMiWK7AxcqEeMwJ2AbEZqlBNeIQRImztYFAWrM3U17MWjg700Ld84c5v8NkvfBEdjdDQH6A+8OxbnsyP/NiP80/++/8+80FtybxtOTee8Jbbb+fHXvc6DqcNu1t1XxyfC/2NnNwnvQXSXkplKDkaSJSRfFoO0i7JCtEkkkkEk8/TDuh6lYiWDJp7za/VgQbNhi8O5Duj+TVWDaacEJZYIb6qy0Y218Ulou9IvqObTa8oDVHC4pFI6nJ1mjh6D3LofFSUfblWJt5qX87UacQDx+q5dHDAW97+dkbbOzS+Y1yPEGBSlPzo617PG1/9GmrjaOZz5vM5yWYv4j0f+iB33/8AbcyMuFWVIBFZ8IoG3a0FWXNF+kSFtQMgScrKDn3sHvrw0rqMSYkmfDdFY4Mj5FZqXqnmKjcVRl54ww285ElP5Tpb4e+7xCjCmcmEpvO0IrQY5knpjMVMJhyGyDRZ5hS0xRhf7xAmZ5hWW8zrM9w1S+yr40hjDieNLnpWugRlTFSqyLyhnCvjLrIbErtR2UkwSZGtGNiOkb0IOwojVcRD9C2hx8dsT0R1SA63Q8coJG4oa5579hyvevIZnmbhepCxJtkVK0ZDXugp9DG6olkxRJPN67Cz5OSExmw4U+xD/ogJLSNgAjKJQcYkdghyliR7BtklcN4Ybtoa8cKnPJn/7r/4L7lxa4tLd9/NbDZDCke5u0Mj8Fuf+Qyf+dLvob3H33ZpkFVYHOhN6nSmyv7xlLe/+71IUYMrs5hDUeA08bpXvpzXvuL7OFOPSE3m+llrkaLkne97H/deupxVSsh71Hcdzhnadg4pc74qQWqDVH1yLHmYZMeZFHqHpShpUTyGDqX1kbKo8T7fiwoYiZUKI1ti5LqilC3NB0kRI1YTmgIxgSky2brRDMPc/t5302luxGPVUmIoQ+THXv86XvuiF/DUCzdgQqC0jqKomHeer999Dx/7xG8Tk/aav8vDPJ5y6C820TqfJK0pLRkVpE+he+PwUtKJY1/hckIPfNJL3utFH/R+7/We4PX+GPQImJpMyRgydEMmdMiGpv7U1NO0oHrJD8FmYFmFLkSwhvHOLi0wR5mjNAhzDHMMDYY5wgzVw5j0OKFtUk294R1kj7DZA0iaSXmFqVGM3nvpgI9/9vPstx2Xj46oqzFh3lJ0gZc889ncfPYM50Y1z731KZwZjygKhxiHGW8xS/DLb387nRi8rlAqZKlYmnnX6ECkXdWHWu3IPRj7KOsFqMM/YwjQF9FXZUVhXS+/ksO2USlUEfYsctt1E7771qfy1HO7TASOLl/EVA43rjHOZWmbBHU1Zj7rsKVFi4rjmI1cg+GBaYvZ2kLLCnEjhhNebJGztUovWZMWoX1VCruTERfO7HDz+TM85cIut57f4mnntnnK2R2efN0eN5/b4cLOmMkInLOItTjnMq+QmAmxqtQqnCktzzy3x4tuMlxfIGcEcR6cD5iUcOKyubBmBSKTFeR2WTgt1mYpqpgFFkrnkKgU5I1bGQcxEJsGQalIlDFKjco5W8tN29ucKSw/84d+mPPbWxSqiCpN2zIPETMZ8e4PfZCjrsu1lWWB9uXQsVdHEePk7v3L/M6Xv8zX7ruXYjImBiX5gB7P+d7bnsfTr7vAM85fzytf+GKYN+xOdmiajoAw9Z53vv/9zFLIf9MYXFmCcVTViMLY/AdXkkgOGHhwBAht1zOShOgjBktdjpg1LfVki7ZtB3kpIUVSaDK0ExI0XX+QhZw87DP2CZhGdOYjB63nnR/4IOVkC0/CFobUteyOCl750pfgm4Y3vOLljFSwKR/wTQjoqOZXb78dX5TsNylnNdGcqDwBnZkh/Zcxq0xVKEGKCC4lnOSyHEIkeQh2xJyK4LY5SgV3Nh33hI675nMuRuV+hTuT8g1r+VJI3KFwUBZc9AmtJsxjWIQiQ/nFSttCHdLJAyO5jFAEMMaBOMSU2HJEa0uORPjKfMq9JL46n3FH23BH0/D7x8fcnRJfOT7mK7MZd/uOrx/t842jS9x3dKAOcjNdAh0tQXzOpmFRDNMQMVs7/Ju3vJVjzZkqQqTwgZEP/NTrXss4eHZL4cd/6LVoNyWFjnI8psEypeRDn/osdx8ccWnqwYIPDaTMxu9iytJCK6Zb+r6ZhrTAC00aWPimzwgtM6Euge3LZdT04Wt2nfPDFj0LWDEWRgK1Ik/dKeXVz3sGNh5RlYEoLfM4R/GURpgUFd1xw/nRmPagyzp0VlCrqCSqUUnXtJRJcD6hvqfciMl1sH2IWJQlIQVSIezPjql24cyTxpw977j+LDz5LNxyBq4/J1x3Bp583vGkG3cQl8vvQgjgLK2JhCLTBsRHahG2bOK264ueDJs9VVtAUZYgJSSLUCyIz/RZVYPFRYNLeW1bVSxZJ02sJRkhpLSQxFooRJkCW497g2iobUWhQomyC7JXlvzED/0QZ4qKURdxrafEUdY1RzHx2a9/naOkTNEM26WET1mvSzHsh0bNqOZf/+qbMVsTDtsW64SxtewAP/Z9r+JsB7tt5PXf/TK2bYmJSlWNOPYdOq759XffzqXZjP22y+Fmn9jTPjtUWLto+OEGzbqeOlCZ3CPBIngfKFxFEQTpElVZM53PcNZSSG8ExTJydaZdIFCXxBT7cyJn4o2V7J1ZuDTr+KW33k4abdHYjH+2qSXQ8trXvJyxTVxXV7zp+1/BiA78HGMEGZccGeXjX/sKv/31r3Jf1+SKHyRnMts2J+96RQB30hOyvXSR1exWhhCYz+c4UyFlxTS1aFL+yf/nF7Btw6SWnI1oulz1bwwdCcqSNgZIQlHWqLUczGfUu7s00ROaLhszu646erKxt5XcoSj1SlyCwcfE+z74QX7/i59DYoc1WXOrrmtQQ9e2VFWPpMWAiZ4tV/DKl7yQ73/xi7NHFiPOGgqGiv6ImKwKcGk658Nf+Dx37R8QS0frA1ujMTI95pYze7zwGc/gjCswzvDi257JrTfdwBfvv8h02hBVwJXcc/EyH/v07/DU17yCI+91r3SSYsgJjqCIHWR+zRWSzwPob1Y8CF1w0laSAZp6Fd20SNasu3h9VUdf8Fto3qQdilWPmNhLMGSysMYssmdizjyW1mEKx6AsM8AFZuB8JcH1mbi2bel8JFklxiwVZGwmcv7tv/9fUiZhenxM6aqsEoJZFOFHjbTdnJQSVVFSjUf4kIUu1SwrEmwvE+XUU1BR5tqFnuBilkeDSM4D9Ry9jOf2JGM1Q+a8l7gZpKJiH97LoJO8FNTp4RY7AAR9ttkiWIHdumTmA897+tP48Ge/iLiSo64jGCid4yt33Y0WjrkP7LgCaw2FLXLxPpb7p1M++bU7+cq99zCXgmQcezu7HN/1DZ5ydo+X3vYcLkxGOOClz72NJ50/z2fvuY+ucrh6RJNajtqW933kI/zcG97INKLO5HySkXUl4CszgRngXi2AH2pzZYVLanSpmDPIBmYeKPiYsG4ldA4BsRXRwL3HnsPO89FPfZqDZk4qLdYaom/Zq0te/T3fTRk8k1IJ1vKm1/wAv/a+9+JjwNQVEeFgOuffvO03ePFf+kscgzazKefHVWblrwDJuazrKvLFSaCejGnUkJLQNg1aKJO9bZqjY5wk5qGvZk+pz0QqQQZOEEiMFG3L/vSI0e4u0/kMbF7ItS0Q31xBq9ChljADmoSY8bWoiS54qrpg/+iA2fSQue96ZoXibJELbH1AjOL64uA4O2Ybw9NvfjLixtnjM1Uucu4lWHzyYOAoJt0PLe/80Ac46hqmJJwtaWbH7CK8/hWv5Gw1Ys8ZPFD4yPe/7GV88d/8O7b3RkwjhK4lGuVX3vYbvOnlL8FqZLuYLLK5ufRkoVmnQ2ndoCmmogscbUlwOZEhOYE1LhSdZWUH9n6cil3huQ0BV5+pjpkmMLxIlqwBn3KVQfK9Z2YMxlrUZiFKnxLzFLLn3v/9oijAWabdDBVLJOHqkrvuupudM2cypuhy0bX3HWE2pTCW0jowDuuyqkMzm1FtbdF2AWNtlsheyGX3Rq0P1h+rDr9ZVI7IgiZ0xZ6XQRNVsf3vDqVuAoysxXnPM269lY989vMQE1ay1HmMkel0yvHxjLhd5BNMs9EMGL0/zOhKx5vffTtzDURTIBYuX76E08hrXvMqdna3Fod8XTte9OLn8eV3P4AdlUxTRJPQauKd730/P/bqH6ASw86oxLh83pn+BDwpi72W3F3Uop7Sx9LoqTKjQ2IwRsW5XEIVY4exjpjrprVBuf2D7+POe+6mmNS4qkBjolDlu556Ky988jM4UxTUCBHLa176Pbzlne/EVQVdiCRjUAOf+PRn+PK996E729w4ntCQZYRMkjWZ96trLUnieDYlasAWBldaTOWY7l8ixJzdmgHHQFsWdKXDl47gXI+PQesDURN7Z85RlCUphoXX1M2bNSNmFuWtPU5kslZWFz2usozGI9RkTsrR9Jj94yNmMdGIwbuKAx849hFflMzFcBQCDQJlTZeUNkIIMRPVe/dbYr4ZxhXMQI9i4o5LD/DbX/g8sbRUkzGjrQkacw3cK1/2fRA8JcgI2K4qfvBV38+Z0XgoMs6KEFsTvnrvXXzss79D5ywzknYpa80UNutvpZgL4weuX1oB/pNc6WSdPGx0cfvMUmNKVvbmiiKnnkhE0G9iSYJNPaPe2l44TknBU5rEyCg1kSp5Cu8pgselXjvfQOhm+dDpWeAhRWJMhKR0IXI8a9g+fx61WWbp8v4lDttjOpeQSQnW0MVE0MS8aVEjeFVaHyjHo3XBRNbJ2vYx93IxC96ZrHjIcmJjqMb1eR+4fQvcJtNrdra2shpxM6OqC5xzdN6zs7PHpYv7OdTtn9/FbNQPvOfr+/v85qc+Da4gojhrSaHjSRcu8CNvfCOx6xa0pbos+EM//INsTypm02N8iqgt0KLi97/xDT7xmd/J+PSa+Vknusup1uzUtjrIaRZC0tqcOGeX7UIk8zFnoPtd4sA3vOMD76XanlCOakLn6WZTwvGcH3ntG9D5nLL/c5Oi4rZnPotnP+OZ1GVFM5uTQqCoKqK1/Ltf/zViUTIj89ZazFrfZnPqzV0J+8qqQo3QBk8XQ/5xYbKHJ4rYAnUOKbJsXOxPnegjkuDMzm4OVZspbTfPpTcpk+Oqqlqr2F8dQ1H0aDwmWWHetszbBlWlKC11WTAqqz6RkPlh1haYwmHLAleWFGWZM6AiuLLAVfl7hqUHk9UX8il7HBLHknjHhz/CA/MjGsl1gs10xmQ04kUveCFPetL1nN0e5SwtsFXW7G7v8Mrv/l7mB0do8NSl44GjfUJp+dV3385B9By0HckUK12Xeg18XVEl7W/OOj2j78x0stRsxdiZkw2d1kLT9WYzi5/3IefCUErqqTcJa7J2w25tqMOcra5h0s3Y9nMmoWGkHaPUUWlkUhWU/SbLvDKhrMaMR1uM6m1SSMxnLbPjGVVRMt7eoiodKfV0DMlYje0Bf2eLnhyp+Ji5eZmft/SczOqm1Gth0E73yoYWgEYU0bjeDFjWRRqNsiYbugjixC6b+ugKk7BXfEm25Fff+k5i/xlrW+CPp4yd42k33cTYZTpQDfjkMZp40plz3PaMWykGuWrnMnNALL/xnvdw0HZMFZ2nk01EMsfxylYgD3UqpKt+PwUWtCAVg5iCABx1nv3Y8fEvfp67jw64ND/mqJkRY2RnNOHmCzfw0ue/kJ2tXapeD27kslrvy172MoxmGpCIkFJWt/3gR3+Lw9Zz78Fh9u4RjSt7xayd1Kyf3Mqys4txFlcYUttRV3WmRzQNZYgUTUfRekzTYZqOog2UIVEnYNYwcQW1cxTWUY5Gi5u7qLXS06P5JND43AzDOsnF2Ra6pskNRlJklBJlM8ccHzKJnkn0cHyAnc8ovKciQDvHpoD6hraZ9k0z+i2uOZTtgLkId166zLs++iHcZEIUKIqKEALee37gda/lsGm5rHCx9VxOkQe6liiG73/daymLAkm5u0AwkOqKj3/x83zhrjtpbWZex5UPac2VHLK1bkiS1gyOWW3UckIb6zSV3fSgTTVNj0kqqrkxSVJP0hYxntoErhvDG59/C69+3o288Lo9nmQDO/NDJtMDdro5u3jKvnpAQ0uKHu9zk5umaZlPZ5zbO8vYFmwZh5nOqbuOSYhsh0jZtVQxUqLQdlhV/HxOChFXlqj3p75v0WvTIzHPnfRh7ImQa+XfMpRcad/UQ5aKX3GABwQOj2eoQDmqaduWEAKjqubg4IDz587lrGlYSmVNfeCuy5d594c+jCtzBrMsS1KM1GXFT/3UT9F5jysccyD08MA8tPzwG9/EeDzOdaYKx97DqOYTn/ssv3fHHTQKwWSK1aLTmV6lAxSne19Z3SWeasQWdqHHvENIi35WHWgsC7w1/Ks3/zKzlDB1ha1KXF2jKvzAD/wArfccty33Nh0HCg9E5fJ8zste/WpGO1uZjO5K2i7Q+cCs7Xj7u97FZGeH+RpHc0Wc8cFaY3VdRypLiD5PSmgwtuaVL/1unn7jkyAmUojUriCEkN29whBjpHIFB5cPs+hg6fjFt/4GnoTaTBqMMZ60/6pXMDM0KyxoDhFLk4t+X/GSl/Ky73o+k6LEpIhv+/BVhXnbYMuClCKjqiB1c7RpeebNN7G7VWXYwnqiZiZyAo41SxW98wPv5+LhEebsGWxSurZlezLh5vMXeNptz0Yd3B8ixkRi12GqGsRy5vobufVZz+ZLd9/NrJ0hoxHHyRNUecv738d3PfnpzG3SyhiRsARml4F1WrZSE12o5uZyprTI8F6tnd7y9bTXUBu2m6xonS6LjxeF3b1bEVKgkISq72XDZ+yZXRywa+DWG2raG2qOgaMGmibR+jltzF7DyAqjwqJqwTq8DziEOJ3zY699HWNj0HaG0Ug9WjY4UbLggBrBFiU+BO452Od9n/5kVv/VZcub1cYgaXjz17wLsmb9uV5JdgBkTB9xDC6vrBDDZzHiMXzhK7+fPSxrsgKyMRgLVT1ia2uLosxF16EXHz2YT3nvhz7KUZuNxtxHau8ZbW0x2d1j78IFGBXcfXyMC4GyLtGqJIjjxptv5vobb+LeO7/R30dLMMIsKu/44Ae57danUo5HWW1YTp+m9VZ/vccmp2Njp31fVl40RY9x1YIo36ryyc9/ni9/4y58VWLHY5r5jKaLVPWYl7/q1Ux2J1mJJGRNmHkIxHHNlhWedttt3PeZTzFvcx8JcVko4B3vuJ2fesMbIQTObFeyqtHokqyz9vUEO90YQ12UtG2kKEq0FOJ0yutf+Upe8bznU/mECYFROSLEDq8J4yxd12EKR1VWHLQtflLxy7/xa5ln4xzaZdIfMS1CIk4hwVkE7QGu2llKY0nHx9x0/QW+//texlbbsVsVzKdNzqAZR4iRajzKeF/0TCYjJAZSO0fbOakucUaxYhaekgcOfcevveudlLs7XJo12KpGxDGdzrnPXuY/+St/FUeichCjx5a5tubg6JidvfMczGYczeaYuqDznqJymNGId37gQ/zpn/xjVD6ytTVhKDG7wpSfcqQs9T2vFKNb9bAGA7ZcX2lRGiYrvQfN+mrMaiay7LIjgMSA+I4ytZwxFT3xOPOGQLsatDZ4JrRMuKwdZcqlPz753LFIYFyPiIcH/Mk/8tO4tmOvLLG+zQL8MYE1RFcwixFTOJoYsVXFxz/3Wd75ofeztbtH28aVz2sW3b2GUPpa9MHt+3ytkZXT2uuvKJ+oIfWdjHrOol5uO+bW8Lkvf5mA0PiczbUG8JEnP+lGiIGxyZw8r3DsG46D59feeTumGmHGI2TaMA8dRVHwwPFl/tb/9f/MyLe4EBgXjuPjY9qkjHZ3OE7KYewVUJ2lqEe0XUs9HvPuj3yYP/rjf5iRFSZVnXsDYNbDRX2w5p4rVTqrfW9PeU7SDiMG18eaISlqhAcOD3jre96NVlUWdp034EoKZ5i2HX//H/0j5pcvY132RmOMJGMzDq+JRqELMctFlRWCYdodczlGPvSRj/K67/le2r7EKYn2cvesSl73vSsBOxSRkxm9hQgmBkLomBhLGRM7CLdMCiqKvglBydBRgFGx2LCFVjwQEmVKeFVCzFwekiwWj6yQZlfZ8ppST/rLRcsmCRoStTi2xHLLeEKNF7uztdZQNQFpMiH2zkdrDUVZUyJiyVhNlwLOlERg/7jh9g++n0aVaQjUW9t0IaKxZTQZc3A8pbQGSRHrsnR3ajtCCFTViNnhISEqrq6Yeg8iWFdyOJtzXVXz73/jrfypH/1RPGhQlZHtmzLYnFFNGvKiMGQxwL7h4iIk73GuQf5ytVUxKZc4JFjolJXOElKv9575D3miiwLJvR8W4u5Jc+PLECOVccSupTRCERJVqVKRslQ2BmdFwvIA0DkwbTuMb3EoVVnQaeYSdr5hu3SEecPNOzXngLrMxtH0G6MF7cg4S1s4LnpPGQLb1tC2bZYLEpcZ+s5mszMcQmIyOVcS2n/OwWNafKH6oI15JSVcWdI0bZaQdnbBeQyo2JTFE1vf4YoKa4SQMj2lAz0C9pPy5ne/iyMfmceQIw9j8PMpW6Xluc96BoWFJiZKa/CoHseOD3z8Y1w8nsHODhdn08xnJOHV40Mk+QSzKSYEKuN6CWuDHk9pohJdxolj2+LFUBUlBwcH1JMRv/SW3+Av/vGfowWVpOKMgO97ZRqznJqHOAxWu42LXNmMJaWIsYIPLYUtwAj37h9w8eCYj33y0zAa4X2LjMdZfMBHisJx1+WLxLaloCR0s0V4oQJe+2oh6zK8ZR3H0ynXbW/TTue8+W1v43WvegWzAFsF+BiorFs09H5QXGG12NmlXB1QxEQVAyOQcYRRhFHKnW9GIGPN1xFIlWLW2x/IoFGzCsdVXd4VvsvQkzNmMXhJSuFyTeXRwQGu767pSDgCZSI/evmdOu9zGVGIy0rxBBIBAbG0HqaNqpaWd3/gAxy2LZQlTTMDYyjr0ULVNKniSTTB04aYMYyqZtZ1NL5DJdE1LdtbYxBD03QU5YhpSLznYx9jhnDXrEGt6Cx4zYq+niTZS3V9l54UfZ8lcn3Rdm6W3Bkh2ky29X2SpA2KOphrLhtR6/Cay6hM3+ZLVnoaXNlJ2vXFtS5Xk+Mw6haSTYasEkKKkAKFRgoyQ74EqVCqGChj6qWQMq1EzTJ8thpyiVOuHZU6wvDYVmSL/BiTex5UMddw5iYsrIHqqVediFmr4DGPGGPG/Pq+El6VDs1VImLwxhJwSDkiist1lEZogKMEdxw2PDBv+eW3v4ODtsEVFWVZUtmsBuFi5E2vfS0Tlxu1NCnpQTentfCW9747i0R2PuujGTKx2gpi4Hg6JRpLsBbvLJ0rCKYgGpeLtbE4H9iZTCh6rl+1vc2Rj/zm7/wu9xwecKlrESO50sD16idDAkMeHE99WF6tKDF5irKk1ZDxaDH80lvfQnKOo9mM0WSC+gAxUYihaZqsf1YW2bnp5zz0iitrFTBtg/ct9WjEpekRcyJ3XLyfD33yU1yez+lArS3oYnfluz4NTF72O84Ofl/uhkvaE+k6RDxWPBZPgacgA7slUPY1Wyb1gHaPB0mv3qhXdf8zfXFowDtolFlrsYWjHE/6mlCnLY4WR2egG2oX+5PfL3pAo7mzTNm30rVIAUfzOb/1yU/xxa99pc8wWXC5TVboOlLXIslTF4axK6iNoVIYiSB9J+3KWWpjMDFg2o7CCIXYXAKmhm88cIl3feyjzEWZAtYVEkCHEqeFokhSJOWMWUTQosBbywx0JqKz3KycWfaKtCuFOdAI2gnaAhQVx6HTTuNaen4VBF4qR1hESyQV2ZilKj+0QimyhJGYvteaLCSfh0O9ACnSogRyLRud628TUULuFL7o17wUnjMpMcqHIjVKFdOCnb+aHNET/MeT3sOjHcUgyFgUBIRWlVhUNLlni7aIHoEeqegUdD+pHgTVixG957jhS/fczd/4e3+Pu/cPGO/uYgpH8h3zw0P26pKXPOc53HT2HJJVmUWMYX8+57c+/3m+eOc3sgySGkaFy53KjmfUMTIGJs5hYjYCudg5QohICDgfKHyHmc2I+/tZoLSweIVL8yl3Xr7Ihz75KWb95/A9bWKxm9Npu02uTAReLUG4wkEbvGURxz3Tmd4/nfKhj3+CWQgETZkc3cwxPvfh2C4LJOZmLDYmipgdjzLlpGGNMFGhTomdUYaI2tBi6pLOCQep420f/gA6GXG5m/d1zUuxhaWEzlXIs1cw1Fl22F749ivpapEevBXbp7bX0yiqemoqXE6cDapKivTUixLfNRAjcx9oRbgrtBTB59sQJfO2Fp2jbRYuTLnjUQoel5u9sV2PF1mqUBb80tveilYlpihoQ6IqS9rZnBHCxBi8al44nceKUorBqSIxYpxDYouqcqaqmM9m+HbOaHePSAa2C1fxK+94Jz/4yley3zZMqjo3hC5cpmSE3N3dlTXqLOoDjQ/MgEvB06pBYiAXXTntBSYQUQ3JE9BMM2lbPbM1ws9mnN+ZZMxRlqKAV73JYlaoOS4rpzCEXHbBr9dT+iuaXnQzs+2XbPxshOLiIFywsqRn3g9HzMKzFjEkzSq6y/Zssuq9r2khLVuhPSbMTARKR2w8s5SYAQfAUdOhIWBtVpHx3mNdQTmq+NId3+BdH/4Iv/re93OpbSi3MxRhUfYmk1wY17T83E/8JBNnOVeNWPRaLSr+zVvfQqgq5j5QVhUkTzubcX48oTLC0f5hxtlGo6zNZlJvgPLsxpj7lE32djhsZhzMppTbO1AUuPGYWDje/M538vqXv4IHtOGWqiagvUzO1TGwNeOV5Uoe1HMTLMZYPIY56JH3vPc3f4v9riU4y/b2hKPLB+xtbVP0ysem6xhVBV0zZ0KWvteBuyhLyaZkhFnT4kQw45K2bbBFQRMiH/z0J/nq/fdiJxXnyxHOFoMi0Lo7v1LCuSxyNku3dEEfGFqEGtPTBtMCOnQnmOsq6y3g3Ilkw9XWpHNZtdJrog1dRuZGIy43cz771a8ysoKkmJ/fE0AX7UhNFvbLPflysa3zgUlV8pxbb0UEDmcNd1x6gN/92u/T2swp1xDRrmFL4b/6q3+Vp11/gZ26Zm9UoUmxksunvPe5QW4vt934DjWWT3/+i/w//odf4NgHvC1Qybryd166yEc/82l+5GUv4zjkruNLSGepHzfgPd7Cxz//Bbad9A07cusWo4KG7FGIUULyuFFJ6SztwWVqge953m0cNl6LepTltnUdP0p9wX/PHF4oDsgCQDULlnsUciPXviGLrJuT3l83i96baxw47Uvq+8dAS8jtDQc5lmXnyL7xWS9tvQT605otTg+TH/XQI4RA1EyaVuu469I+H/7Up9lJgu07dR8eHlLWFVjD1+74Op/74u/x9Xvv5vJshrcFF4/2Obf7JMrgGRUFYTZD5lP+1M/9cZ59081cX48pQWYRvdi0fP3iZT75+19jZgsKU2ePtYucLyv++l/4izzrKU/m3PaYsXHM2zarh5iC2MOc1lpSzP0RLh8e8OW77+Yf/tN/yr0HB7gzZzB1jY+Br953Lx//3d/lB1/8YjxoiWRJUFmWMCV5OGGluXrqRA0qllZFLzYNXVHxK+9+F8ciyLgmpCykaWctr3/lK/izf+KPUyqk0FCXFanNRe0qvaQY2QFSDF5gv2v5C3/jb3C5adna3iLGyP7REYy2+KW3vYW/9h/8MRrQOjcvG+pYlgvDyCpXabmYTG/Uhn4CA2Et9cYvL+q0NgGDcmhalNb0r2dk0QF9tQLAnPAYZk2LGoctKsQaYrTMfMs73v8ebr/9Hbkn7hpXa0n6VTI/TWPqM3SebVPy0uc9n7/4Z/4M15/fRbdr/uW/+GWm1tJJdpmdCmWIPPP6G/iuG27i5r0xhebMlO+a3AjWGFLhsFXBrGkwNjciCUZ48TOfztmiYDZrEOOIYulSYlKNeOvt7+YHXvISdqXI8ihdbj9sTZbh9rHv5mwslML/8xf+BwqrmJQ7nzvKbPCSwYkj2URIni52OCJlCIxR/qd/8o8ZVRWxNyXmBBi+XI8BpFsYnuxeK6z1rO4pJLrujUWhx5DcsmfBSiNh6aV7nMqifMus5ssWGyn1fTLRKI78MItSrXw4hQWNZWnulqi0PEr+7Hg8ZtoFQlJSiHzidz/L7332C8jRjLrvB9D6BlyBjzlsKkZjgoHGGuq6YLs8w7xtmM9nuFCwYw0/9Lo38NM//CaKecMYaNuoWIvakn/7q7+GL0fMukAlgukaRiSu3znDS57xdPYKx5YKcTZlx1UkEYyxxJ7zZY1kCMIKe9edYzQacWFvj8PW06acJZ77xMSVvOX2d/LqF76QoxConcsy+MKpgP6pxuu0KoHhwFXNCiA9/BNdwfs//BG+cXmfMC5pQ6AsCoiJibP89A/9EDtRGQsUVZ37KxRlrvWVREjZexfJRrIVIaWW73/JS3jvpz7JvQeHuLpm++xZQuv54Mc+zp/+0R/lzJkSMVadPBjrV9Ki8UMyS+XhZPL3FsKKuvqZzRXpX12UJ7F4PUzPC+HKRirLGD1SjUaoD8xDl3EDK1TjESn4zGuzDlSWzVx7YxnIlI8QOsraUqhjJAWzeeBQPTKuuBgidxzv87aPf4RUT0gmq7tWCrum4A99zyu4sah5UlYKEldYohtRDFtSc1/MUBSosbTG6iGBcjTm9a96Jf/u7bdzOUXEOkxZ0QGf/dKX+M2P/zY/+b3fSxczQzyJyd4dJise5XwrXnJNZEda0CecWCRlBQgjlk4jOIOpx1hjCN2M2XSOL0qCWAKiDiNu2O4L+5SNSDLtEijV3KFLCJi8XZZC6P1hZ4d70xciR4MGEbyR3gD1OBwRl6BQyZhIEpwZns+apluSZVIjrEhFLZrnrFiroaGLPRG6PtrRtp6kmiW6J1ukEDiYNtTiEOdoVPGuwNWjrHUXAnPJ96Qcj9g/PmIyqui6jnNndgn7+/zkj/8EP/umH6ZOcMPWFg7ElZbjgN59/wO8+0Mfw1/YAwstidIKhbW8+lXfx1blOGMLdoBqPMlaY/09i8WyYXXKVGMOUJ0YeN0rX8k3fuM3OGznqFiKqqSwht/9whf5xKc/w+ue/zy8c0tx1pMx5mqN6ql78SrhZr+k5h4Om8Avve3t6GgEdYHOjumA7aLgtqfdynOfdANVG7hQOlFyAnrobRExuQZYtW9/mPFuRhNe+z0v432/+TEKzTptlatQjewfHXP7e9/HhR9+E1KOTovuhqAgv/kiBVyKWM0PlzJAu168cSV1JfUubEdvvFhRPYCMxkivgaS9BHfuRNCHNrn3JSllrMskinEJBNrmmK5rsUVO6UcD3mTWfTDQihKM0Ikio5JoLW2ItDExT55WFMY1xxp56/vegx2P6dQTu5YaYaRKGTxveOX3McmdnqTou4OPEUpyM5UyRooYGFmhzAkPmWDRdsobv//VmNhQasKlFokdmjyN73jPhz7E/dOGaEHKCjEZMyMGKoRacido71vE9aRSV6CugrJAC4t3gi/AjGtaSbSaaDVy3HZMzpzhcDZDinK10+BSIaG/F2bILpPvrcFjNWA1rDcNXj2kTiVfDhLfCSFg+9cxRER1heaT2+5lfG150K2WN2sP/IsOpTcBQ8BpxKWATTl5tKYiriuKDatcPF2n6gzRxIARW00UAuOyzFlYH9AYURHsuKKRRFdkLb5WIx0JUxVIaYjJE6NnZ1TjOs+2Ri6UFf/Vf/5X+aNvehN7xnKhcDiQpvVqBO67eJl3v/+9FHWBb5ssaJlaKpMoNfLaV70C2o5twMQgVc9qdzFSAkWKuBSzppp6ShIThLExvPZlL8P6ljPjcd5BKTCfz4iifODDHyIYmPsuFzOp9ty6pShorjbpZZF02QVMeiBhwb07UaUSeyji0vSYL331q9xx190cz2cE7ymqEZUm6hR502teRXd4xNnCUSlsg0wU6qQUCrVCnSIjZfkAmTh40bOfxS3nznKmKjlbjzm+fAlnLVVd8/bb38XcBzrI9KEIuVmHmnwyAq3vIAa2jcEcXmbHRMYS2SFh5y0SciPWhUFbyYImoBOYCxwL6i1oCpSq1N4j0yl1AYVEpJlSo0QibZ95yVk6l8M+3zK2kYIG8UeUJjCSRE3CErB9pX+mHyREc0V+mZRRUqo24ZpApQ4JUFdbBCxHXaJBePd7PkAKgTHCjjVMUgfTy7ziu5/P7k7FaOyImtuiLcipMZdiLfFCQxdi7taEsFtNOL+7zfe9+AXU2mDbQ/ZsxIWGqrD85qc/xVcO9vl6G+lKm3XAQsdOVeJSg7SHGD9jZBST+lb3vQaXeo8mjxEP2qKpwUiC0FEaoSoz12j3zFm6LizNUA9mat/paaTg5p4RFuMjJiTGtsBFj8aILUwvUdO3+epD26zgGBYqDAWI8Z5SIjZOGdsWwmVK11AUkaAtLZFGe7pDH9907TTjIz0WWwKpDYysozKKtjNGTpEwY2wj49RRhw7tugWu63GLVn6aeuTuCopP9hQjMFe0I2ErR2hmbDmyFFKYM9FA7VvKGKgcBOlIhRIkgAnE1FBIxKYOaaZMNLEVA6PZjO+6/kb+05/5Wf7p3/4v+JEXvIinVBNurCpGmcZGWRVykFSPxPPuj7wXK4Gt0LLdTNltZ+zEGd/3vGdzw2TC9eMJSZXKur41oyLW9NBDxp+MKiWGAmUEsleXnJ/U/MB3vxQ9usxYI3VKuRUgiQ989MPcf7jPcdfSqM/M/RhAPZ5ElIRqoiCxbaFOAcHjSNQG2qMjiirTUpr+vIgp++zRKBebmdpxzS+++ZcgKSPjqBrPqPVsB89NO2Ne/d0vZlwYRiZr0Pmp74sKhhNNcztANWsSWDsGytDxY9//Gia+w82OOTeqIDSE2HH56DIf+cTHuP/4cvY6Fc2gch/miEBlHXvViHEI7Cqk4xkkxQblwtYuVVRMDAhlb9XXuw0MfHYLbBUld1++zC4W0wZGVUUznZHaOU+aTBj5ns4wyR+gI9M2JqagaltGfav3YMBqzqy5mPlnau1CKkcX/QVSr4KQMNoXxfuENTC9vM/4qcKWM7znfR9CDqdUXSD6jlFVIqFlB+FNr3wFpmuoRuNF56m1Wquhu7gsKSSm517VoFti+aGXv4Lf+q2PURcFRwcHjKsK4xWNykc/+lHe8IY34KgpWs84BNLlS4yMZXdUgjGE6HscYRkWDlejvZR4gtoYNCbM8TF0nnOTLfbvuZenPO1WSvoEUc9gbn2gdaLN8RG7Vc2FckQ6OsZZwcWW9uiYC2d2mR8copPtrDnOUk9N+mzNAk9VGFtHGSJ7xnDYzUEj/rDFAFvliO7oiHL3bCbahoh1Qj0agUY63+FcicFwpnJy9P9v783DJTvK+/7PW1XnnF7uMotGGq3WviAQCElIQgLZgIBgbJaAY4fNBmwndoiJnfy8xXZiG7DjeMvPwQl2bBwSIAazGIOBnwGDWCRAIIQkJLRrtM6MZrlLd5+l6s0fVed09713JDnx8/wSoXme0X1mptV9zumqt97luxxe1/rQCrt6ffavrrDgDFk5ISsrelXDMQtL2KbNKGSmBH0EgI9XjBN6ghTW6P7DBzl+2zJ79u6nl/W7dslsxhFStmKMSUFE4oEslu27dnL6Kadz2okn8vSzzuGEHTvZtbyIqQP5pGZ7P5McqOpAPzOMvOoDDz7A9TfdwOGH97O+tkqxuMCOxQFNaKgPH+blz72SoUIG9EWkqQNZZmZuK8yDf1v7RvEMxLBsLd9zwQVcddUXWF7oszaexKl7ktv5zN98ihc/+9ns7B2t3qhYUaxxSYzVIGVFVle4yQSpSnpFj8xaTFVz0nedTL22ynBpmSyxvIL3qEZDZgPceeut3PXtbxNGJduXtzOuS3Q8YuewxzOfdC55U3PUYImqarSXOcl7WeQ9ZZukpaeCpQn6s5Q7fc6FF/HB972f1cMHWdp5FIfGI/p5zsA5PviBD/Ddz3o2olqjIY5YW0R+LXCgmugDayv4fkGwDrWWXlbgJxOY1Cy6Aucbdi8PxeE3QCpMEt0TvMBDBw9pvrjI/vURIc8oVagUFvoZ4dDDDMWwlPdZ6PWllyLHWoU+vL5KXWSMxFMm/wGX+lQmKBZLnazhZje7SvKc1ID3Sp5FQ5TcWQqTMRmPWFoa8PDDq5g8Y1xOOOWYbTy0NsKGQCZC3wkLLmM5y2QqNxO6QBZmtC6SuCZiDVVQghEeGq/rWOBwVZINhqyOxwz6ffrG4Ucl5eo6xx27i/0H1zDOUfuG4dICk7KhrCcR0uFchKbM3N9U2yuZz8xIkBcui+bGa+ssFT0WjLCryETqJrrhtPIzwP5yog+urNJkA9xggIYGfMNCv0c9HlOosmthSD+h9e1cINfuG/c17Du8qqvqCcOCA9WIYtsijoyyXMdNGo5ZXMCNS3YNh5IBVTmO6hp5gqUEaCqPyTIeXh/r2AoP1hPM0gK1gWoyomiUY4aL7LvjLk4/7ni29QuRMG1iK9GYxrY81E4exHQRoAywb7Suvl+wb3UVM+jhNbqszymVyIzOiCp1XWKCkmVZROLjo2a+hnj4DHr0o1o4GZHdISFCOUwM8DpuPHVuGXlFi4IKZVxX5EVGub7GjmLIgrEsimXBWskTX6rlic6W/K3/bIfpM8IDa2s6kYyHJiOWd25n7+F1Br0ehRFkMiasr3DMtmUK9fSMkcJGyac6GNaqRps856G1dRgMKE1Uia2rScSzBc+OrGBghaWskNwILu3T9XLCffv2aW/HURyuPAwXOiqSFdDJGkcNegxF2Zn3hKqhcFHaXWtFepK8O2ahWjET9al3Ngbdt7JCyAqkX/DQwRVcr0+RO2w1xjYVucsQDXUsHdRGCnwIeBHGqqw1lZYiaBbbreoDNJ6BzbG1JzNKv+/SDEs642CZ4XSJcUwaz+H1kXrryAZ9xk1gfTRm17ZhBKVORvRtxlJ/KG0WNKkCh8brSlFQEaikdbN2mGjMOKU8HYFFYBScM0zGJYXL0pAulqRNVVMUvYjNMtA0yurqCtu2L1E44eDefRy1Y4cMMotRjfcO4Mw0I8Dj2gmgGkITQaqSZxyelNo4yyQZuVibMRqNKKyLkIyyoakqbJZFXhqR02adw2aOyofo0bnBvVk2qotYQ+2bBNmKhP8MKNdW2T4csJhZ8eWELMvxhIgoN4aDk7GKzVHjaLxirVBXJRKUXuaoxyOWBkMGzkrrYdpCa9ohiwCZwqSEh0ermi0NWfdN0rGrWBosYCuPn4xZdBnLw6IboHltEv4tZtoSozIraxOdABNnGAvYXiSlZz5Qraxw0s4d4ic1gzxDZgKQT0oOklTeOp59wwzgF1YnNaWgWmSUIR52s+tnCsCN33fuHN772GM0hszEaWByDO8osYQoP93qjgXvsTauk8Nr67pWVQRjcYM+3tpon9fUZEXGEMvhlRWOX1qSPCUrofQRDmNaQ6EwN14zwXToVxVLLbBaNTrBslZV5IMC52BtdcQgz1jIHeXqYfp5xmKvnzykY669uj5irfLaWBfNbMYT8jzHWUOuSt8K9dqI5eFA8pnlOB6P6fcL9h44qG6wiHc5B8cVXgy9nkM9DBxUq4dZHvRYzHMxaVJYjyuyfp4QEdOu6Sy+FGDSeFxW8PDamnqxjL2nt7CAAmVVYTWwo99j74GDiPoUzMQlGYBoeeaT01BboRigVtVCRFwKFEaiwzhG04QplZrJ9RxVMFlSIcxQgbIGl8WlMmmC9pyRRhsKici0pox9nqyIk5vW5SkkfIxsmKsIm/vSs4WZTwbF7SJr6tSjcVCOPUVmEQOTCRT9+IZVpRTZLD7Fp8Uex95KqzKquOARH8C6+GHRJBK1sDapsb2M9brSQZZLO3vTsmGYO4JP1Z8qNhOaME0k9AhwA9mCEqxJGy3Myhp7yGyaQko7sAkEcXgMo1BrbiKntolkJtGZ1kD7WxudTonSxFhnJvbVpKbfz/ASvS3SgIuqRnMXS9zMpvVTNQSUosiiCmtSb42cjJZcGm+mkam6hCS+Rg5U6xWDQT4/KU/BjBTMhOnUtHMNrmooMkKI/ptlgKpptJdvHujPDbLqgHPRRLpFqPikJNsOUjUKLdP4Kpr24mkS2Na5jNLXOFsQMB1lx6RD2IOurK9z9DBOLqk0Zj0dXnUazDQNUcwM3U99jbicpvaYLGaYPj3GVjxyUjXaz51MqrEWLhNnHJoGIFo32Cy2Ndr+ednE/QGxYdpOocXEIO014LIsckV9jcsyfDBMmmQ/0X7fTQxmvmrIc0M9nkSBVpX4NWc2hbHQyS4Y5pMhDZEXW/SGVE1sFwBU0fKWANRVpYM8l5iZqZJ8V1Btxb/tVLMp+IhzqSuKLIegqPcRkJps4GIgCzNp8ExIMQZfNVibgbVMygrUUPRdwoc33ajeRCRMSu+jUWkIcZebRESdE5AI8T+BraExomDE4Jsa9YEiL0CVyWhEbzCEAOPDI+1vG8TqJKk5SNBI/jZTZJSaqH0/O33LUOqqInMZbcgsxyOKwYC68Uhmqb1nPB5rL8+knxddhtiOFpumxEj0YzZp0lD5ZHdi3aYeYMR8JT0yA3UzzeAMUwhDPRnjsoipCUaogse6olOfiCoHNZnJmNQTFRHpuyJOr3x03smsnXnYdkbldsa7M2UqqjHDNQjOTaNtXbYo+oDL43uMywkmt9175hiomwjutZamigdB0KmQZmg8ztopv772SG7ngpnOZGYWoPFT0S0M6hskc1RBsUYwqnPmy3PCJNLSe0LUEU+HcjwlAe/xBEyedWwHY2WGhh2ZF5OmQsSQ2YKyqnEuT1ix6cFsITrNu5nnHaIL2ZGC2awKrjGxdKvK+IyMiwZdxkHlfVK1SND4EB2cfFnFkg+Pb2KGZ2yW3g98E42lQ1VhXBz+BGMRF4NQ4xsKExdD7WeCUzpEop9OkrbyaXDj8k4ooGxiyWnaMnMmfmiKH5ImRKPRJD43l9M0DXnuGFeeIrdxUKJRPaFTy4iD8LaZHx1QfJIzMeLiJgpCVU8iBSOdADAVsIuSutPBeQgBYywm0WVCCIiJK70OPva3xEbvAPxUn0shM65L2UWTo1Oa7LS5ht2COr0JS1SX9LOCsiqjsUeWUVcVIoKzeWTzJ1BqUUTDDSMG1WZaLrf6Xx3wdwp38KGJstMS3XDquiYrYnnkNWBd1r6yC9qTlTXNskxMkSEmKWKIEDSk0lnmGv+thIvMDQKmVC2TsG9N7XGm5XomHSrRWGLOsDsyHD54glEaDeQSv5PxZISq18X+ooSmic+3Q1nOMENUCCYwkYDFUflSc7FSGIc2ynhtXYdLC6KtR4ST2BsLniyzqTSMWXeiu0cJl+6kiGsyNDUmy9K91V1G4L1PhiipdychfSfadhS7LL5zQ1CD14CYjEldYgzkRmYoeJsxVlXZROcim83UKGGuFxe3SoSy+BAIySVdEGpfx3VmMlq/zOC1Ey8YjSYsDAcxWFcVxmVRjXk8xg6LeD8yhUqZmUNaBUoixU1Dw2RtrNsWl4XU78p6PcaTdXr9fsqAhbrx0Ruj833whKZCXNbty6aJgdt1BtY+JTjRoAZjmQRPZiwZPk4iTaykRqOKPM/JXDx8TIpqvoySWopQ+oBx2YxsT5gZ4MyzO4Jq9BQIHmuiZqIxjqqqooERUI4nFL0eohoTfsUS8DQzDD1toquPRSInscg7qo0nUNY1PVvMLITQNStj4zoihK1YmhB9Ca1kXYApsiLxAoQ6kRlsCk+qnlxSqPLRgzFKs0Zg3RQv5LcsNadDR5MehI0LI22EPMsJGgiqiLGUdU2RFXiNJUJmovRI5IxJl/qqpisWQYxhQqRjoB4rLuGqkvCkiUGmrEqKvGBSjshdhrVC8LEnUlcVWZ5T13HxaABJWVbjp5nhLI919n5VFWNtt9h9U5O7LKb0RR4XoiHqZLVZmQ9kEt0uahqMNUyqCb0s70yKJ2Wkm6jOkOBF5oKZl0CFMKFWp8pAcmmqml6WDrqqwuY5dfJytFZomli2EaKreRP115MbdqSm+bLCFkX0erQWgqcOPj6nskn3NVtGphaAkRmJ8DDDDNDu0GuSUoazWZriR3CwyFT9bbbEl5lDRfFJhinMNeDbTCYm3O1UQijrCUXWayUfKcclg/4gHs5iYxCwNlYh4wlZrw/eR3ejXkZoaiRzKQM23RHWlpjBQI0y9iMd2EEE2JYNWVsjhoaG0AmQTqoJg3yAYGiqaAiCCcRxjk3eoRZrYnCMGKQ4mPGpP6TWERDGdZxKGo12lI2P+8SmLHg8GjEY9EBDKoUdGhRMNN/uhAqmOu4zApAdXyjGBt+QZzlro3VdGCyK94kN4xzluKLX7xPqOgazWeBi2KhI2iot6PQUamvr+BozbbZ07xG6RHsr5LDMqa+bjkEwS6CZNnE3CZl3fRkeXY7pkQRFQUJc3OnU0pnWfutiNFv2RheqlKl4H30d0112sNAw7TGJRKlvj27KHbupWwdGnSkf5QjXywYIgiqNNtikAR8ZBMkgtfaINRCS/ZtMH55Tg7RcNBu5YBpCVDHpfCan5UhIsj7xoJo6SUW36kaNWAqs2LTl1Yep5Zsxc8TzzlQjPbQGjxgbg0AIXVapdYVkbjqx0xi4rYmZmaZDY9b1JWbM8y0HSe6U6n08nBDq4HHGdTS+WClJ+u4NtXoysV2Vot2ELXRBLWZZAQ2BPL1Xe2gqPjoWWZPKZKHt5kWeZcpMm6YL7B0if0a0IVYFZlPbZDYzWw9lJFkjsaQLGkVPCfim7krCKS3cxj03s9H9rBjjDLjKdAAYmLq6mi5OmEfZf6ZjaWwU15zZ/2GLwCBbx49whH55GzREg26JxJ42UcP8X3bBLDEEmHaON25Z3VD4zQex9sUb0vUNWJqpP+QGaQ3Zii8W/s4/m9SNaG3ydA4vnxZdXeOc60bikoQBQalNG9hlBv8zJ+Y6M2XVGTklmQbzjR3/IxENNwSy7qezNEClIZa7KUZpVUd2gVW8kS6YWY2UzJQKT4+mdoeEJg40TNpM6aSsUt8js1l3yXUK+p1LUdKCF2OSD6l2BHLTvnbmEG5SA7nxMSi3B4qkaVlRFLH3lAKABk0lUgq8YX6XT4PuVPY9yyxWXOzBOoeIMKlKRGw0uJ4pMZvQ8h9t6hVPt3ToPDrDDA7KIKnhn1k3bUf4VJUHSa2WGJTajN11fbzQYjA2Zb+BLVQr0nBNOkaDIqlv1VQlmc0J2qDqyfJsBudvpsFMzZwLiO+C9MbRUitA0SIMzSOS0sMWfzJze9/MKXRvykZk4y4JR6S7bym02QqOzjJe5jOBLT5MppLPj4UZt2kv6xHGdLPXsNVrN7xGt4jO/wu6CWkDpTN8Bt8hYrvFVTex95G1J3oTM7OwxcBBpsiRtuJgXhg2dKWHTU9wVva6k7+eCW6qs03f+VOaEFhvGuwgR4ByXOmSy6VrIhsImemWpk38ckJqWmjshcZLmarX1r5CrENb8Gj7T94TtIlltYlYuNg3mdXkV+rUwG6DmU29WKuzRsbxI9dHI4pB7O00TUWRZV0GMyonFEURS3zvyW0WR6chgLOP/L2LoE0VM6VivvxF42QSwDoTe1nI/A7cAEmZ3aqSDo+uBeGjgnHwUd/Y5FEhOfioktu+rw/xc2ILMszlHGEmysx5grZfepjxJwAwFl+WsSzvXhoQGwcSumEDm3aur1sHkXnnzGlmZToDZDNfqW1BrtZNoc1sSE5m1rVlzg7xSJN7sxW2lvnUTrTRI0bIIwWz2UzOPFoQ0w3vpVuMjWTz+z+iYuPfQSFBjpD0zI0NNEzHconoqgkcu7Y+ZrAw6IKPL6s0EZ1OrDAyR5zWDV6qIpvpjdqVszKD69/6ec4yD2aTM+M1burcMgkxTW7SLS1IolKkVDHYONFslc1cwtzUviErUuO38oS6weYZ1pm50r975hqwIgksbAiV7wYz0yAQy0CVaYmtratSymzboF+vjTRbHIhHKTU29dvvJ/gpCn92QdugU0yKlQ3eBxvWlxG0DkhmqKvo4drr5VsvoGZmPfoNVcDsupOUGBJd4G0a3WpIvbTZL9ErVep9OmcJKRh1U8/WBakd6HRZn+2y962C2Vy2oxFC4jVEccg07skzl/ph080mWwUz5vf1fChqVUrMtJLaKsHYFD9CV57Lxv2v01mKNxu9MObPEnOkhIitgpl/hLAgjxauZrJjHqnPc4S/kw1XzNYYq/8d1ap5EtD8z+lpEWbGlNoJEWkkRbBeVyoiUrgs9iQQrMTpmjHRL85gu+GEqsz1NrYO1ppUPuQIQXb2eegmMGEbfK0KZeXRgWUlqIoRMqCokL6ji25qobEhjSsMLgVIvyHDba+7VqIjVvBaZJEBWysqGiiMEQtUlWfQTrzasq5JWZuV1KOa3YRTeaCWUN6CWUc+YDLDOkFRoZ/qtTbzX68qFYWFIhdqT6ZxmCOF7YLZXOaQLqiqa/IiY1R7TB7HAJUqTVDNrZEscVUj5TSQ52Z6Uh8p42tLWtU4rFE/dT9JPa/g47PIcxvhlmbzpmzqGpuZDvgdZmKViJmzGZwGa51ZAzESeI2IBy8JXZCSy8jD9dO+d8rkZkvK6b+ZI6RFYb5jpWa+ub5lMhKOICe0Yf+3KjpblJCPGsS2uNbYM9vUYgszLzJHbqA/Ugdej9R0P/JpsKX34xHST/Mogc48xqxt7g3CfJBWA2VXCChNVepi0ZMQlExSVhXq+L20TfO2XTqj5HCk/oC28INHvf6Qukipt6FTNU5fgSkM91eV/utf/7eIzXj9q1/LRaedKnmdvAQloDbQmDbhsIkSF7Xhq6rCe9XcZdJmRmE66NHJTEXaJi2WqOcfh15xd5lMJGu1JVFCqKcT6c5SPW6I7ty2wmpZQy/jwcMH9dpvXMcxxxzL+WedjQMpK7SXIy2w16Vnk6XysLFT5LjRDcEsfYcjD2Nfq+RZwjXGa6wIDDEMQFrOoG0Dwawd4KZSY3afpGGQGHCWugnUhC7DbAPU9Pqiqck00/Zd+u7nemTSAX/ns06dwVrKjKZgdIqqNZHRQUyAXhIrmM1UmXEls13Pzmw9WZOwGYquR5rAhS2CodmiuAszoGezQQHx72juPPP5zgud4I9swHjMBxMzF0hkw1SQDWDOzTX0ZrrRpnp4K0dp2SzTPPtaC0du8B8xJ5vJOM38SdF6f7TaaE1oXdiEftGTMniK1Bwm+EhcTOJyoWvtS5QATpF6rl8gsyl//Hf7iNfPFF/WAhDT6RwEmszRgH711pv4+l230wTPZQ/cw3mnnaomQ1ytHZxm9vttn0CGkOU5YEQ1mdom74SRwgMHDnL9t7/Ft267jdFkjDGG4XDIaaecyjmnnaEn7txJ30nSpELLBPLJEclNnrIW5lG2aRzgDawHpepZPVSN+Be//At8+447+d4Xfi/nnHU2OWiWw+FJo3nPMQ6Nalmx3B/IfGqwoYJIH9MYKBXGErTJMz5+1Wf4xs03Y/oF47pk1/adLPnA8dt3apHlnHnqKXri8SegdcMgy6UFanSWBZ0hQdi8UG1U5mhcnNiNQcdNlYKakoklS6PR4D2ZNXMYyYjoaNfOTAK4RbNc03S5Fsc4qZrc8+BDVE3NSSecEAGtjddlZ6UzGGyvuT3c0j7Q9viVLSoo2ZxYQJz+bkx05jOo+enoxj5cVJPdAG6e3a9Hap5t6O1tjAWuu5lHzQ0eJdtpI7hsUJrdNOGcfz95DIHXPFIfbtPnhpkH+mg/Z3S1ZPNcIgDfuvNu/Z3/9/fJsozf/s3fwPqANVbERC6r0iTwn00T0TCbwHfBuF0w0l6byHzqfcT7mP0yZ3VJ0mDCoAeBD37qb5g4SwjC3tUVaqJDVSOeCEgw3bCh6/Z5wQdPlmV4DdQaAc37x2P9+Oeu4mN/80lu37MnOedIxBllcULmvcdUNScsLHLa8Sdw0gnHR8NlX3PU8jbOf/K5+uQzz5Qo9rg15EQEXNoZeT7gxrvvZtsxx7B3PGKUHstNd9zJe+JjbQAAKH1JREFUH/3RH3HPPfeAhWdedhlveN1rdXdvKBqizlf3rcpMvGldrgVGq2u8+ed+jjpznHPJBVz1lS+z99BBtKy44klPoTrxBHYsLnPqaSfjMOKddJixCEieTqC7oULrLh9MDGTxAKAGrYBb7r2ba772Za764hc46fgTuOSCi7no6efr0b0FnE3fcvC4Dr+WDhzZ8JhkJl5rUgoRoRbDBPS6O+/gj/7sXdxwww2oQH9pyMtf8lJ+6PtewigEzVvajJi5UsjOAHE3ZVSPMkwLW0QFMxsUpwzSI/6/LQTczgXA2YPCbFmu6YYdPDuAdHYm45q/ADMXSOzGgCKPraiTI7zCbpUpytYDBzlSV/8xF5RHgjikrCGl3iqRJ+kMmCCMa7j2lm9z0wP7CFXDAysjjt02ZIWgBYZKA1YcPRAtociis/T6ZKT9XiFeDVaEoJENUzhwabqowXeIazAdB7b9gqzOKBCaJCguUPsQJ3pZThU8EwsHtOHqb97MerA4m7H/4UPkJAFGl/BzAUQF2xLzU3BxcezFRD1qHId8qf/k536aO/buBWfJC+GEo3ZzyslnsuPYExhnBe/+wAcwLqNnLMs7t7GycoAvfeEOjCiFsZx47HGcetwJqYSQKN5oJOnN0bETjEgUAxC45cEHcNt3sW/SMMkKRsCnPvcZ/uSP3sEVl13OD73uNfzpu/4b7/vMp7nxwQf4g1/5VS3UJI+tCBXyQGbocFtGo//rMb2BvPftb2cddB9wy213sPfhNYzCK17ySp755Ce1CihigUwcrivkTDfAUA0z6LVWO75CTJ/VWtVnwgh41wc/zDs/8J4ovmkN39q/n7+6+hqOXt7BH/7u73Fc0de8rNmeZxJLAUNowdFhdr1H8YE2PoSyweSOJkBp0H0+8Eu/+ds8fHiFF7zwRRx7/G7+5H+8i3f95V+wf99D/OIbfjx6LSd8oMwooksKaHFtpQ/YMpCZLZKKrbCj5jElI3bT68yRR5hsHQvMEf5tA8t2c0T9e7Am/Lu9j/xvvu4x/P/zWLbQQQSjz23y7hMhZOiNd96F3b6D1b17OVDVLANrKyu844//C/v27eO5z30uL3rOlboz2j5RA7Y3kDGNrqyvIPTIewUDh4wDLBhBJCQfogbT5RbTqY6f9sW7k7QJkaRde0+RJWqMNVTAVVd/hYPrI7LhkLr2TKoIBiYBgo1qOgElBZMN8A+iP6lHcdZx0kknkW3fziXPfCZXXnE5i/mQBck4DPzVV28kHyxQNRNUlN/6zd9iGQihZmgi5We8PmJpMIwagBKPxSCJzCM6k54rzkQ5qGAspcLICLU13HjX7fzhO/4z3/v85/LDr3oNPdvjB37wh/jKv3sL37z9Nq779u08+8zTkjFulHcXYmwwiY0i6eDq+IMRDaGj0YQQlMy4blJsQNx01seRO5kbsJMuwyuMnbAO/Pyvvo2v3XA92WLB8Sceyymnncbtt9/Onrv2cGA84kd/6k38wdt+g6dtP0rqsumoXXOTb532xkLC+DlsFFWU2DhcUXjvX36I+1cO4VzGy1/2CnbvXOQzX/0iN33rRj75t3/Dq1/+Ms7cfnQcbiTkjWyJFDCPYf+Y/+VYII/W0Za/n1ji+E7/FRGeM2gfQ5M8Hmvglvvv4eFqRL48ZOHobRyoJ/zsv/45ti/v4N69D/D2d/4Xdh1zLJef+2TNBFYJ3H7/bbzj7X/Ird++E8kXeP6VL+BNr3uN9gxSacD6BmezueVhNmSrLYSi/U+cDNLxPBuvrKvX0jk+8KEPUvR7EThbVxw8fCjei03cto2YjpkF4cuALQyFNYy0ZEEKedvP/LxO0pp3M9fmgbtvu5VqtIbLLYv9HoKnh0VNRpY6c/3hsMNetWgpg265CMtGyZ0wmVTU3rOwfZnb99zNO//sz9i9ezev/YEfYoftYYALnvJkloYLjMdjvvqFL/HsM0/rFElbFQoRSUN46TB4pGfXPoFDK2vUGihcEjE40on/WH4FGPsGCscfvvs9fOO2m1nYvsyJu47ibf/m37INy0GtuOa6b/Bbv/Pb7B+t87nPX8U5/+Alupw7IbEZJGjygZh+5yo6T5FyyToxQTj+6pMfJ+SORmFpeYEC+IEXfz+/f+cd1OM1vnb9dZx8xZVqjIphHijbwm7M32PC8v/3L/OdHMfi0MtsalJ6jZnNOnD/w/s5VI2xSwMOlBN+6S2/imSOt/7aW/n+l7+MsYF3vPtdeAMrwG0HH+AX3vrrHH/aafyzn/0ZdKnPp79yDZ+55quMg2oQiyQEfbQOm24mu2HS3Za9tY9zrknd4BUqJeq7Ocf9e/dy5917uPCCZ2CzHMlyHty/P5oZp/IiyMyRvGH12iKWF6GpyBWGGIYgMilZAlkEGUYbVAxwx7dvQX0NEjhqx3YGWFxU1xVS47msGjUQ/Q4DnQxy2+BmBmbinESHp6aJTknjCXsevJ8bb/4W/+RHf4yjewsIESaSA8ceeyyqymS0Hg2ITbLcS8postGqOwWIMFO1r49HnVhBx7Wd60LoPATmiOvHgHWYwvHRa7/MR6/6NKu+YnX1MD/6qlezG0vuAydKzneffyEve+GLGPQKvvSlL0QkWNA5hzI72xadHbW3rQFrqFQZe6+IcODwAWoaBssL7HSx0Hv2+Reh4zG5CA/tfSDh/8O0tkzvXSXUTngc7efv6GA221SftcCVhMk8tLbKmq8wg4Lhzu185fqvc+vtt/PP/vmbqLTmkksvo84sN917LzfvOwzAr//Ob7Ow6yh+7Ed+gksvuIzLn/dcbr/3br5w9ZfomQiDnDQlde0j2LIlKqYN34n7tX6lVggmDtBdFl8fkkNWA/zFBz/AqaefxmlnnkVZNeRFnwMrq+n0JjqSk6R2NgnkR2UDJGCdJTMmotobz3aTy0DBjCcUCoUgTQ333ntPjGw0nHji8bGcC1FQIAMGzsowdx24vTPF0NlG77wVTgH4smK8ukYIgaWlJU47+RSeetaT8CHQB+kJMvFRaSUI9IpiRkHCPMpwKpbbPpnwTpoaL1Cr756QmTlItgpks16gs3+3VjWUoO//+Md4cG2FfNjn/PPP56JzzmYR5GhrZKjIToRXvfQl0e+ibqLQaAL8xkZZ2IT83QhZCCnzbBLxTqzFZoZ9+x5iBPSBTJXzTj2dXQuLnLR7d+tsO/0tIUG9wyaw6v/tv77Dy8yklKDa6XGJgE2+evsPHGTiayYa+2j//V3/jTNOO40nf9dZDIDejl0EY8m3L3D9bd/mhls9d9y5h3/x5p+gpmKBHi943gv40Ps+zD133tNJnheun0pM2Xw2bnAtiuKJ0eqq8gGbGcYBtQb2TUo+f/U1vPYn/hn76obxeIJbXGBtUrZ4KpU5qYsZ5EzLhsktZTXGuoBL2mnUgaywaBUJ1zGLhNvvupODKwfRfsZ6NWEw6LGAicqoXpGkB99ypm0blWdJ8qJz01kL1I0Sypphr0+J4icVL/tH34dR6BsjZeXV5lawMZhaUU488YQIm0lO2GoTd1JlS1hBKyGwVk2idSKCT6RsmUNBaKdYe+Q1E78zLwZTGG7Z9zA33n47g4UFQu259KIL6UNkR1hD1jQsFJl813BZf/ZN/5xQN4QQqIzijE4HeK3G3UxWFrX+DIoSNAqn9q3jcKjIrGUlSeHsufd+BruPY+iEt7/l3zEOJYVxZARpe4Ayk/dLIso/nvKZJzIzEik4TBVabSqpDux/GKOwtLTE/ffcy8q+h/mRH/hBHE3Xtr/kmZfz8OHDXHPddfzVR/+a3dt38dwLL2c7vbjJ18bkjdKsjQmlEhKRWRCa4Dfhbzb+qYZO2tu6WBI3CgeD55qvX0tvMODJT34yTgxFluMVJnXFmg8dQHSW0OJnYAxIYOLryOuzjjo0qd411JMKySzkGUEDlaA33vZttHBkwz4mz9i2c0d3IvacRLXbEJH57X3OQWgkzKFRIOqvDZ1IOZ6gIbDUG9CMJjzl9LNYSgEpz62sK3r3/fdRliWZdVx62SUxWELHKIDZnxJ/t9CBBBheWV/r1DxUwDm3qbzvoDOz5arKpgzNAyPQv7nmS3iFelIxNBmXnn8hAP08yk6JBsy4JquDPO+CS3nRpc+msE6CGIIIwc7ILKXDVVOvs/WPF2T2rJNMDCcddwI9W2DF8anP/i19NwX8bjeFFLWXIkAezJQPm+7TpYP18RQAnhgAtE3iVnjQg5EIgBytrmGMY2VllWVbcMrxJ3LBk85jW1wI0oAOix550eer136dsLLCS7/3eSyTdWj5BVewkPVw4siMUBhL00QtYknloyFsmjIlSWUaVPfs2cP1117PjTd/i/Wqhn6Pwa7tfOrzV5HnPW6/5RYO7T+AleQH4ANrk5Jy2E8gjCAhSEebaSW/BXAuo6YmqKcwjrqpybI+NrnKiQSqpAp010P30xhDldD32xaXkqp45ASoxia2EUFDVJJwznUcna4/ZKQL6C6zlClIBA8H9z/MMdu2cfbuo6N6aWJLlRJ4/wfex3K/z1PPfQrbFxbwCoU1iE4t5kSissbGcUfb8B6Nx1EtOOHbrLVzbUS7tSX1NKDJFF2tYlgF7tn7EEWvT7M+xqpy3PISLukfW60ht7jG4KyBJojJhCqq6SFJSW0ONjQTRC0SHawyixPDejlRW/RkII6nnHk2Nz+wFzXCp666in/6qn/MEtBziPGQuzyxExIhTqYMA/s43MbfkcFMROZOcpNlaD2F29jUjzq8tpoMU6FvM84760n0iXZyWURxcMLOo9GyJkj0/Hv1y15Jnl5TKlpgEB/I8xxrkWg/lkVfgxCwJkrKZAnHpCESriXLGVVj/eM/fSef+cxn2L64g/POfzonHXccB8sxdzx4HwcOH+Kc08/mYx/6EDfsuZ9hv8fB8Qgxhp//xV/g5KUFFlAKVf3hV76K8857koySqF5gCs2wZIgoDYp1WfREtG2PJup7TYAHDxwAl0WJuSpw9M6jY0fORh+FJhW1dYhkdHFJ2XXDNHUqrhkI3lBbtHZCpZ7FxUV27dhJAIoEkS9B77jzTq77yrUMXcbPvPHHGUKUG8oMjffx+YsQgmKI9nCaNOjqJkQhe+Dw6gpeA0qGbxp2795N2QTtOzOFlcr8kEI3AyvnDob79+/HYCmwnLjrGHopv2rwok6i6Y2NqeHAGkaVp5dbmYRGaxMp5cFHMxp80pFzUXU4hCbRoKJSxqDoyQQ0By546tN576f/Fqxj3+oaX7/9Li477eSObFGWPqnkWiajEf1B0fUZtGmQVttcnghmj48ScxbVm1Zzo3Hac2h9FeMcToVytM7lFz+DXnpooVKKXOSopW2aqWCsYeeObSz1BuRtQ1lArGV9PKIYDph4dMkZCRpoGnBZLBvH47EOi544YxFjyLGM1HPzzTczmUx485vfzOWXPDPKEitMBN75/vdz1PYd/N7b3ooA1z+wn5/6pV/BOkPhHKeffRbfd9llHN3LGYpw7ulnyqjyZLmTCjChoWeyuFGk7SklN3GZlqIeCGK1Ag6ursUAoUKGY+dwKeqOTSb0Bj1EYHUy1kGvH0cEvsYZGzFzM8j8MFOmiYUxkC8MaAyE0NDv9zuoxFqIaddv/MZvsGNhiZ//lz/HyQuL0ASKLMqUO+fmgkwbi7z3yXtg2p1cX1/vzHGMMRR5jjUxkKlG6W3ZELRgMzSrDXICHDhwgLIsEVWW+sM2PkirfVZrQyYOLSuknzPIkoNSnkudtJ0zZ5hMJvRcLyoQJxK6tTm+mURturzA14E8M7LSoBef/zR6/SGH6hJnDR/71Ke46LQ3MAHtG0QKS9UEegJ5r09d1WS5xTdNp/TxuMFlfMdDMyQ2gmvA28SfNVAbmGjQ1XJMIDaJMzFceN652A7kEOUpj9m+E+ujfd2JJ57IQh7BAZGnCCG3+NwSCgdZ7IFVTUOeGcZVqd7XDPtDMcYwmozTKD16Dl5w3tPkzT/5Jp5xwYU0TdQ3TSplfPnqa7j04kviQAE4+9ij8ON1XKIaLSwtctFTnsLpZ5zJmaefGVUvUC1Bv/DVL+vHPvHJSA5XcCHiwGwiPZmEDp/l4XrgwIFDEAQXDAPN2L1tV5yT5gUTr0yCcnhtlRe9/CXaENkQKmw+LWZK6VX1OkYpJVA7wVtBc8daHd2z79v7ID/y429gebjAW37xX3PuiSdSr0xYdEbGa6PY/FedI0vMBqKQZM/bzzt8+DCkktsYQy/P26Qt9frCkdfKDDdydvPU6+MOM5j1ihnJpgjEFpPjQ4P0pwDpfp7TaM2h0Zr++Yf/QhXI816HwhcbYSZ1VWFtjmQFVVnisvieQ4fUAc558lPwJmOsgWu++Q3WEqRoHXQkaJmJ7hutqbGSfCbAZg4vMG6ax8znfiIz+z+/WZaAAyE5f86M4TH4EJjUE0zWY/euoxm2C1gDWe4ogZ1L23AajTfOOOP0rlfWCrmOQ4Pp5YlWFCV1hlkuHujlRbJy8Xgf6PV6sfEfPLnJaQgU1oiV6GK1HlQnCHvuu4/77tnDP33j61lMQ4ICeMrZZ3HdPXfQ1BUPP7wPDxSp+V96pZc7OaCNvv8vP0wvBL7/hS/CpmjjWu37tqRJpkZJ1ksC6OrqKiaC3+hj2T4saEALF8EaE+/1C9dczQUXXhgBt7IFYj4NV1oPh+uuu453fuQvGS0uxrLNCnc9eB9fvv4bfP6vP8YDD93LFVdcwSte/DKOX1qUQkEWegQPw6QzFzRQpiwsykcnCWtro96/SGREqDJaG3cKJ9YZnJlugmnDfXPv6giID6zC9uGQg3VJ2Xiqpp66ViUdjNY/2mZQl5G/JS72Zv/igx/gvj33Ur/Y02/ltH1AkktXlvWAgDYNea9gHA+WyLY0cPmll/O3117LcHGJfWurvOfDH+UZZ53Ng3vuZu3Qwzxw3x5c8Pz8m/4FvSJjMqkoevmc0ssTwezx0j/bLNXYnfJ1XSNBkeA585RTqLUdcrVQRKO9fo4VJYSG0884tZsU+dQHWxut0qD0+30mtaefWfHAeFLpQi+aohpxNFLTJGKzM9F+y2KoG0+tHm+N5lZwwOc/+1mO3raN888+hyz5SY4UPeOEE/jSjdextH2Jffv2dVI9NaqFFamBPQ88yHU33MiPvfrV+A3AKZljh8aILGbKUZZkOSgh0HMZvRSrWuXz3Fr5/Jeu1mdfcQXjptS+yyVLyHPZQCCOiiiGi85/OqsKf/65z7HQ6/NQOeKwEf7kfe/h1S/+Pq64/GK246JpbBqKVpMGM3R4lHp9Qr/fJ7dm3m4w6Z1ZmWroG5FYekrU3S+yfNMGkJlJTOsypjJVwm/BLq1wYq5w0u7dfGvv/XijrJfjFvCDquAkAxryPCc04HpCFeIBdM9DD+iHP/oxfvPX3oJ68HXA5gZrDaUGHCa6gnuQzDEJgdoYnSjccfc9etXXruO9n76KBVewPh5jfM0HP/oRDj/0ANuLjKO3L3HppZdw0VOfmvrEgaKXc3htVRcWFmRWiuiJMvPx8ADUY9BWA6qT7FE8VTmm5zJsgDNPPpVeZyYslE0VN6SLTlXee47eeVSH3OmpSI9kWmFg4ByDROtvUPq9XD76sU/o17/2dW2qupOIjvLE0dczlBNyZxnklp4VyeJQgc995m951sWXMgS0qbFEZdnTvuskcheDzT177qLCR2jHNEzpJz/1acgsV77gBfhkAjIrS9TVY2EzgLQY9DtHqMFgEEvQtILKEHTf6mG9/pvf5KKLLqJwRatktkmwtaMgAX1EnvX0C3jNK3+QZjQh6xUcGK1y1/69nH/pM4i5Fmzr59JM4mdlQ8cI1Z9/y7/Ru++/Vyd11WHEvI9foLVxuopM0VV2JvNq76Fj+GxgKDw6AyACdQsDZ5xyMpP1EarKWjmmTgMAK62ngWmNjjqsX4nqf3rnf+H8C57OOaedJb3MYhPWxAew1tCk8tmLUgUlGMM3b72FN/7kT/KOP/lTFhaW+L23vo1zTzqVheEQV+TUovz4j72RN7zudbzy+1/Gcy5+pgyLXIyJqsIeZXFhUZxK/KAngtnjJCvTqCzhVDoYlAmaJH0iNsdhsF459YQTGBD7S6jHJqqbyTNCAvHkWQxW1iuFh0JVFl1OXxzLWY+FRBv0wWulyqc+9SlOO+VURCOi24pltL6mBqFfDKLCgIdqbcJkXKkCt9x0M2uHDnPFpZeRgwydw4TYvzv79NPYthC5i1VVceDgAZI1sShweL3k81+8mvMvegbD4SINqt60bIJkfGHm1QDFKK0Tan9pQCPR7Wm4NIzDgUQTyo2Rr3z5q5x11lldQ94zlcdGTSe53AY0FyDz0AvI6buOJWsU5xxLu3ayUk/43f/0H7uJoSW6bAdgtfG6v1nnSzd8PWI30i8/0x8TjaYk6LzsuLU2bmJVFocLmzMT1bnJqzxCqdlyV5/21PNYGPZRAqurq6xUvssSvSoOhzR0xsg18MXrr+UrN13Pa1//Bg6ur0bTq9YjwCcCkkTGgxhLqZ7PfumL+itv/XUuuPRifulXf4WXPf+7OWXJcfFZT4puVhYOrx7i/n17sQH6aeqeeciMxWSGqDUcoGn+biKITwSz/9NLzDRynElMfGoA5wiDrKDXNBR1ze5tywhIHXz0uDQ5Hs+wyBhYy9AKeVzc4jshM2GxKKCqWFxcoAoxO3LGyTeuu153Hr2LxeUlsck302BYGCxKHKc2YHNQxfZyin4ulUevvvqL7NixjbPOOL3r8WRiMAHOOekElvt98A1Ff8g3bryJGlhrokTi9d/6FgcPH+LK5zwXC/TFxOCKn3fPMdKlTlZMh4UaFL3YIDdCsTDohBwrYLX2+olPfIKLL76EbQtLU2evrngPc1ptra6bMWBC9E/MgkfKktHhVUye8dmrr+a6W74VxQYB72AUAGf55Oc+w3Enn8yJp5xKUeRJO046IK16332nQWfsEY2LPakg9Hr9Dpzsk76XFyWoEFSioQ0zyvkzah9xqhsPivNOPoVjd+7EhsDh1UOsrK9Ng3CIpsgtMnesEWryZ3/+Hi669BKO234MO4eLUs9g47IsSotnCCYZ7jaIfuCv/hLF8PrXvJYFOrVfnnvZ5fjxGE0+qA/sf4ieiX3PzENPbAfEBUMVmmT0+fja/t+xwSy6r0eBRe8DGGgEmswySRy+QZ6xDVj0NcsLBQoanCFYh8GQI7KAxZVjwsphjt2+jTKgmgkTAxOLim8woWEiDVWsNvDARz/yMf7hP3wFwQqNQObyZECerqvIwULthNJFIT5j4ZqvfoXzznsyNpkS1UmwL1OkDzz76ReSBcDmfP3mWxkDmYMR8Kfv/u8cc9wxXPL0CxiAJEtgLILT6JoUEBpRahuVWptKKdJCGWR9chP9DyRZsK0CY9Db7rmHW267nZe++PtoJpOIMBcTaWIS8OJRE2IqkyKDJv6pd2hWwELPUpQTzjhmN8+67DJq43jr7/4H1oCHQddBJxk6Br54zVe5+BmXJmBxYiFIajBpjSQ/PU+UC19N0+VJ3WBNAZpR5MMIpA1Ba6LhS5CoLRZU5g89ncqBt9ma0RhMcuDlz7uS7b2coBVXXf35KAUFZHUqbQXWQmBi0A995uPcfMftvOENP4pLYdabQG0TX1JicLcKZdngjOWbt9zMzbfdwTMuuoheGvj0U5V/6glLuOAR31As9rn17jupmPEnDYIExSTRcWvyGR2zJ4LZ42igGZ2F2iCT2t8SgEIsYXWFoik5+ai2HyaomAhyxGB8zRDDzv4AU9UUBqmIeuwe2H3ULgprWC/XmRDVSL9587fVWssZp58iLSEcNsgkJ7J7iJIvul437Nm3j71793LFFVfQN1D6oE0UUMDFjSWXPP18jAYOrq5ww223IingfOG6G9iz90G+53u+J044K49W9XQBhFgGtrpu7SDEOSHUsADs2rEdCYHJ2iqT8ZgqvWYCvO+DH+DSSy9lsShkodeLkNimnk6HJeqOaUtvTq5EVR1wIHVTMchziiZQHTrMT/7wj7F9YYnJpOI/v/PdUessfd7Nd9/B9V/7Oi+68oUdF1Tmnp8mk5pokuxnZqrG2JT1GPpFEftb6RnPvc6ZCLbtAKvzv9sgp5UyBHn1C17EUp4zyHO+/vWvdc+mDU6lwCjUWgHv+R/v5Q2veR3LWZ8CxGnApkalmAjzaT+zVzgOjsa67ahdeIFDh1fbz5bGozauNY4/ehc5UJUlh1cOUgY/ZZUkfuwmXdkngtnjLkXDiOn4z61+lwNOPuF4rAZ2Lm+jBVUK4IxBU/N02ebya7/8S7zlV36Fo5YXo1JpKjclZSfro1WqakKtFY2veMc7/hOvf+OPUDatNVaYb9XMLLK19bFmIP3M8YWrPksvd5x71plMqirCNpxQqeIcTCqvF537ZM4962wWlhe44767uW09SlD/1/e9F7HC85/7HHpE3qBxM6CE1CXfKFGuTcxAFLj4gqdTAMcdtQvKMg0z4LbbbuPzn/scL3vpSztLjtB4skjijsI80lqcme6nqKFnTeqHZQwHPYZ5xq6FRRaBn/nRH8evrvH/ffQjfOvGWwjAmvf813f+GZdecBG7FrZhEhFNW35m9ztJDBnXYeWqAH48wtY1BcrAxP6ZleQE38JpyjKO+t3m7THLzRSNPdZQK03w/MTr34h4z1ev/TKH19c4HNC6B6sClQGbF/z73//3LOV9fvD5L2YnufRUyFXIAmQaJZM29uWGg764QY/VcswdD9zLKkoN2th4kBTAYtaDyYQFcZx/9jkUxkTV3S2s4MzjMJA9Ac1oF0wIiInqDx4o60bzzMmF5z1Nf/an/yX9fo8MxEp0FXKq2CyjqipClumFZ58rGVDXAWMNXpMoohWCNjznu7+HBx94gJ7k/Pbv/xZPe9rTOO6YXXJEl78Zo4TFQV8mIagYw1ev/jIXPu3pbCt6MiMLTx1qIGcxt+JBv+95V/K2P/wD+oMeX/vmtexZ3s5d99zJP3jO97Dc62FAysmEXl5sqZswF9As5OkvnnTGGUhVsbpvH8VaFQckAn/49j/gJS96EaefegpBFR8CmbNHPCmnemGC1kojyjA3LPf73HX/vejakF3Adz/tqTz8ylfy0Q//Jf/xP/wer3/9D3P3PXdSjyf87Jv+H3pEqAoSS7KtzmrfeO66f49+4drruPmePVx3260MrKHxDTd97Vp+9ZfvxIxG9Pt9rnzu8/SKy57JsNcTA9RVQ567R2iSRyhFT2BZrDz/gov1hiuv5GN//Qne8cd/zM/+1JtZRTUzwoF6nVtvvZXrr7+eX/zpf8U2k0m/y5pks/l2CkB1Qg4tLyzxghf9Az761x/nQx/9CK/83u+nn172tetu5PabbmL3cUezvDTkWec9o3WcojDuyNmLPBHMHmeRrNWUmBJwXQjigO86+mg5ZudOFRHGk4ku9HpC7bFZXOB5loMg47pRjJOhM1SThqJwOCuUIDReL7/kEn77d3+Ht//nP+Ch++/nl3/6X8l6WeqgKKZDLAkbDENjtmQFemKkVGXlwAH9x694BVVd6yDLRDVQ4cltxEuNak+eWV7xrCv40Ec+zJ59D/HBv3h/VPhoKl76gheyiMOh9Ho9fFOByTqF3fk1rhEgmxpSfYucumOnvuTKK/nYX38cMxlz83Xf5C/+/D1Uq6u87tWvwgGFRC4gAeqmIo/OT8yLQqQbDYJxwsAKq03NGSedxK2338aixCC1CPzoP3oFZx5/PFdd8wU+8uEPcu655/Bvf/EX2dlfYLDRlrQd5rSiAWoAz5677+X+e+5m+2CBFzzrWVTGUouwWGRsCw1DEYbDIU976nkMej2pG0/hLFk2pfxspK7P+cVqIARVsZafeN0bWch7vO/d7+XUE0/g+S+8kv6gz2e/9AX+7J1/wo+95vWcd/aTyOv5oNUFspk1oAJN42kyqxnwyn/4Mu6++27+/D3vZv+eeznrtDN54L4H+cQnP8lRgwEXnnMOP/NT/5weSlN7HeZO2pHyJmsNffxtZXksipqP65aZrzE2cpmC97EcMtNTcdbqrpWPFgV8jYolpPGZ1YgjKkcTzfs9wcLIeySzrDWV3nXP3dx3331c/szLKKyjj8i4nGi/KKSTzA7akYqjvUbcj+tVCdaQJxK4M1CXMWj6lCE0ZUWv6DFulEpERxZ+4W2/xje+dSNHLW/nx1/zWp7/jEvpgWQaWF85rMOFBQk24rhcmG7SFp0fESdCU3rWJKj0Mh5qat7/gQ/y5c99kWo04qLzz+N1r30Vu5e3S0g2fKFqKHK3IYfZ7MgdsQeKWsPDoVTTKxhR0yfDhkDPGHE6nZhW0ZeZEALLxghNjFuSTJW7w6k7DAQ1cT+PvTJuvHprcS4yAuqgbDciBVD5hl56Fs0W179xn4QkOGk0EARGCGosq02lucu54867eNd73sU999/LwrYFdh23m3/0yh/g7F3fRRZUFlrDS6PM1ZY6dVFSA3XylVjBqyOquXz685/j2zfcxGh1De+VU089lSue9xx279iOjWWn5ECoKvIksx5QgpEO7CuP5nv7RDD7vzCYhSYaYEhULhAi2bvxIZaKxlCHaHU/Go8Y9HpI8JEZ4PIu0NWTCicGmzZBEwLiDA0w0kYRg8NIIKgFCXWj0ZsxNn7bINktMk00pLoi7xUoQqOKE6FuPD1nmVQlNo+MyqbxFC5DknnFehPwzmiZ+mADkMn6WHcO+0nczKMiNAkmOhvMZntCWtaYPMebSAgvU6M8j71BUWBUjjQzlmFWRKekSKJGfUCsmUk6oshku+akNR8xMBaokl3uZLLOzt5QtPZkiZRdWRRrqSYTtveGwiRmT2rn1TgI7bQxZoOT2mMyi0lKKC2gVdvc1wdy0wJTpEOltPFRNvAxde7PgbIcU/T6eBxr1UQHeU/GbeacPq8BXdExA4meCUOs9JtkypoLjYk+otE5aSrUHwS8gdI3WOtYbSrNXE6RLqsKql6ijpymwZVWjQ5zJ9W4JMtszJJb2p48Ecwe7+GsU1eYPRnrqiHLHY0PSRRRaZqGzMWmtiYjiojWzqKeVvKZD7XH5JbaB0JqIk98rYXNZFKV2s8LcYBvaqxrsUw6FWHVme6GRG/LkJyTvFeCj8EVoCxHFL0elXosFhtlVAmqSCZMJGaX1aRmsRclZqgm0M+p0+x2NpjNBjTRdugYUGeZEOErpCFHq6zjachaB8i6jtlA0EgenLMH3RzMRKJCrFpJU8foJZklcHLjFZtnVESKmMXgm4aBdclaNKDSOr3rdNKYMmwB6oS6bY3OGxKYVYQsBS0zkyjFylGxVjZlZhuDGQJVU4HLsGRx6CBCUwWcM0zwie5lqWk0x0kPMOOAKQyNUXxrHa3MBBqTpqFxuNA6pDdeKazQ1OCyeC+TEHRgTCyufSwRWlWMIDF39xrJ9XbWhUGfCGaPr5bZjKHIfJrPTL9kPvhtac+31WNMEjqzxqWd3nwniRMeo0uO2TwgUDqXc2+id2Q2RT1MP3C2J+PTC2wLCTEdYXrjPcuGPk5rpadMxf06ccnUddts0ry5r71xI80+n87mrTNLNlOIQ3qFZYNcfucAf+SHuNHc+YgDj0ddL2x4p5Dez3SdR9nwIWpa0G5IuiTpebeBNQVp2WIdbPncZlyWWghIziw0JXRvNH22Zub5Pj5/ue/sQDZdko+kvPkozvWPeMDNmM5PF6se+Xrkse2i+f1kpoE5JHLy3M3JzE3MSFfzKPe8+eSLA4k2mLXPTVNA2/IweJTTX9loeTYbyKZRwf5vKtbPXoKdCWn694BOMt2RGIPVLFawNRhwErZcR+FR1qBsjMAb3qDt6crs4fYIz10ex/v5fwKrDo8jgk6+OQAAAABJRU5ErkJggg=="
HDM_SIDEBAR_BG_B64 = "/9j/4AAQSkZJRgABAQAAAQABAAD/2wCEAAkGBwgHBgkIBwgKCgkLDRYPDQwMDRsUFRAWIB0iIiAdHx8kKDQsJCYxJx8fLT0tMTU3Ojo6Iys/RD84QzQ5OjcBCgoKDQwNGg8PGjclHyU3Nzc3Nzc3Nzc3Nzc3Nzc3Nzc3Nzc3Nzc3Nzc3Nzc3Nzc3Nzc3Nzc3Nzc3Nzc3Nzc3N//AABEIAJAAwAMBEQACEQEDEQH/xAAbAAABBQEBAAAAAAAAAAAAAAAFAQIDBAYAB//EAD8QAAEDAgUCAwUGBQMEAgMAAAECAwQFEQAGEiExE0EUIlEVYXGRoQcjMkKB8BZSscHSJDPRJVOi8ZLhNVVy/8QAGgEAAgMBAQAAAAAAAAAAAAAAAwQBAgUABv/EAD8RAAIBAwICBwYEBQQCAQUAAAECAwAEERIhBTETQVFhcZHRFCKBoeHwFTKxwSNSU6LxBhZCQ3KSVCQzRGKC/9oADAMBAAIRAxEAPwCvl6iSa3VqO5KospykuOdRxxxv7tSChVv0vbBIozBqORnHfWhxC9iu0UKDsc712YaHJodWrDkSiSW6S251UONNfdpQEDUbk8XviJYzPp33x+9dw+9jtFZWB3NTUWoiKkICCouKFt+L4wL22Mh1E8q9JtIoNF6pU0xQlJTrDie3bGfa2plOx5VVE66yzSJUyW7Hp0CTLdQkLWlhF9IJNr/I49LDbsyBs4pW64hHbv0ZBJx+tGZuTFxsnUJ+HQXTWEqQZgbbu5+E3uPjbD7uX1KTtivMW5WKRHxyNCYbzkWoLRJiPx3o6hradTZQuLjbGVdW5Cac869XaXiXQOgEYrVR6l1oS5GkAjhN+cedktdEuijGPBArM1qpIfKZBaVt5dA5N+LY3LK0K+4DVZ5Ut4i55CjOTMpGoVySrMlAeEVMUJb8U3pGrVvax9Ma6Zhj0g756q81fTpdTa1HVjeswmBU6RCjt1OlTYyCsNpccbsgqN7C98DnhLszgin+H36RxpAVOeVaWi1IfdREoAAG6zjzV5aneTNbTpnelrVSAD0Qt7W2UMdZ2vKXNci43rMqgVKrQpCKbSZspIUWytDYKdQIuMelt4Sjq5IrF4jfxtG8ABzy6q0+cspKp9djLy3QXlRlRiHfCN3GvULXufTBXzNHpY7566z7CdLWYuRkYxtQWiVJLBU+G1XNwUkWIIuDfGPe2hf3M16eCVbiMOORrTyKl0YSJGkHVyPTGHHba5THU6BkispLeXJqCG4sV9959RCGmkhSjYXOPRWtsSmnI2oN3eJa6S459lGIOS1ScnV6RNoDorBWsww6394Nk2t+t8aySFCi52FeUuWSWR3A5mgjqJcKWzHnwJMNxaS4gPp06kggG259RjPlgKoWyMV6e1v4526NQQcVqqXUhKBSEBCWxyTztjzNzamPc75pp066E1mpCSktqbsWlHf1xo2Vr0ZDZ51bAQE1Wy9Q5Vcq1GclUaS5SXFlxxx5v7pTZQqxvfi9sb8UZgDHVvXnOIXsd0qqoOx68dlOzBQ5VEq9YXEospuktuBbbjbf3aUaU35PF746WNp9J1DNRw+9itEYMDueqtDKdbNMj5fhurRNp8FCkKIAUo/mSRuNhY/P1xB3YTOuVz9/SlFOlWhR8Pjb9vQ+NJFcQKTIy/NeW5MqdPcK1psVJ1DygDYcEkeu2Oz7xmQYXP39alslVhkbLAZ9fI/pWTqkJVEVSXXJrvSnsuuJEhCWyjQU6eCedWIkhV0JRd89W9aFlfSCULM+Fx14FVH5qHyyyzNSXHn22wQsKI1KAJt35wK3tvf95dqfvr5VgJhcatusVsaUwzlCRWZE12RIJDcNXWaShK99V06T/Ko7/DB9OtRGgHWawZZ8t00zcuvH6UiEvUqoTZ8qSXYEVKXGEm1nSrdA2+fwBGK4SRVjRfe665mkiZ5JH9zHd8fp41WzLSxVJtXrUOXIaSIHjtmUKbIS3sL3ve6SP0xddLlVcd3fUJPLArGI4zvy2rNMVJgMIBmIB07jqjn54Se2bVnRXp4ryDQCXGaIUmhqq9Oj1BUuQtp6oGOhDLKFJBSQQSokGx3392HVVYse7vjvrz9zdSzFo9WVzWkzFNXmB12TRXVIcakiO4lH50nZKt/Xf6DviqaIjqmXII2/al21yx6Ld8FTv+/lz86bWVxanQn6Q069rpsttwuR0JWtZ4JsSARqVz8MdGpTGcYbPhUu4MmqM7r5931rK1Jk0KtSqfJnai2hpQU5pQQVC5Gx7YFPAGRSi1sWF8WLidx3Z27aWmNGu1mLT4s/T1G3FKKNLijpTcCxNt9xjooAqsWTwqL+/KlBBIOvON61NFci06hMUl119K6lJccS4+gNqQfwgkAmwKk8/TBZFLg6QMLislJAr6nJDNn6/LlTsuTFZfcYl1l5S3XpJYbSvbSkbKVt34HzGOfRK+qFcADf77qqBJEmi4fJY7cvh58/iKzdWoblJpj09MyShtqoCOpt9pCQoqJJIUDwPX34sVSUkkb4o9vdSwlYw2Fz8qHv1NkNOBM5BOk2HVH6YSjtm1DKV6CW9gCMRIM79daXLVL9mTqRWpkqS4kQfHf7KQ2kKbsRqvfYqt78OtpTUqDu768088twqmVthvy2H2N6sqQ7VKhDqESSWYElKnHgLfdlO6h6+/4WxQ6I0aORffqQZJXSWN/c6xt8Pr4GnVWOxnCRRX4jsiMpIciAsNpWE7g+bUbjZIPzxbBQFJByxURzFW6aFued8ffX+1Y2PORHLzL01AcafcbKisIJ0qKQbX92AXFsdfuLtW7Y3yNCDLINW/M1apcM1k1d1E17owGWnFCOhLil6yq/J7afrgscKogLLvnrpG+vpDKyxPlcdWDmtbKWj2Uxl+G8tEylQG9KzYKNh5geRwAfdviT+YSuuVzWep2aKNsPjPp994pYrjaKVIy/NdWuZPgrU4obqSTulIHF7XP6jHcmMyjC5qrbqsLtlsf5++6s219oHUfW62whbmlStfsltJ1223Diu/ODTQ6ACMc+s/SgxyO53Zj4CmvfaCUr6i2R4j8SVLorQJI4N+tftzi4t0b8u/xqrSyrszEeNSzM/KkLJkxy821sypdHbcCU2HdTqbfLsOcUhUaQGIB7mrpFbJK5x4VXbz1H1BTUNGpPmuigsXHvuHtsGaML+Zsf/wBfSqIWY+6c/Cppf2iuPFvxqErVpurXSm3CFE+inRbYDucBjgDFicY6iD1VcyMgUaiD4fpXOZ/cMdSVxgW2wjppcpCEhPm81h1CNgdtxucV6NekUKRjx3q+qUKxLNnwOPjT4/2ir0hhpKjHVZDrYpDaUlB5GzyvftbFpIlRTjAbqyfpVFZpDuSR4VU/jmOnyrhR0r40ihsm3zeBPyGDJHldWr+6htlTjr8KtI+0N2NGS0lkssFILjbdJb0azsrbqJHp64AkayEhmHdvnaiMGQArkdu3XTY2f1toWuKylKrpAW3SG0m17q2DhG1kkeYY54gsgBwBjrNWV5HQkMx+FJH+0BDMpvptBperSs+yGm1AW7EOqN+NrYl4B0ZZMHs32rhIxfQ7HfntvTX8+l8h6fFSt4ga1rojS9/cpTwJHy+GLoiHZWH/ALUJlce8cj4UjGfSwevT4qUvpB0KRRGkb+mpLxIHrbEsiDZm/uqQGIyCT8KfJ+0Ft6U5rZS5vpQTR2nFEWHN3Um977WwOOAdGGfAz37UQyOJNKOfLelkfaCtzQ7KYSV3UNTlJbJtymwLgHdV/McQkQaQgYxjbff/ABXGR0T8zDfs2px+0N6RGLCmeswEXQhyko06xxsHFAjn0+OOkRUICkd+/VUKGcEksezbrqp/HEdXlagx1OHhBobO9u1w8SPlg7R4XVq+ZAoS5LYyfKrkn7RV6THdSrwyTpabNIbUnQOBu6n5WwGKESAZxq68HOKKzNGdiQPCmN5+c8KlDccBtevWG6S2oHfy+XqAG453O+IEa62DEY2x72/fV9UpUFS2fA/Kli/aMtkr8EkIOglOiltoJVcdkuG+xPcYmS3C6SMYzvk9Xd31RXZ8gsxPh+tQO54j6yp6GjUok3XQWLk/q/gqxq35T86GwZd2yPhU8LP5jKvGjllpy3WUiitt6kfFLqr/ACwKZBpIUgnsLUSNXyMk47hUbX2glTgdDA8RyVN0VpRBPO/WB7+nfFzbxqPex51CyytsrHypzv2gaZCXXI6EuaUq1iktqVrtvuXEnnjnFIYdYOcczyPV1fGrSSumBqYbdYrG5cekx6uy5FWUkE9TzWCkdwfji/EVieHS4z2eNX4THJJcYQ4HXS5kekyKs6uWtS/+3f8AKjsB7h/947hwiWHEYxuc13GI5Y7jDnIxtRJ6XUjlVDCn1EfnufMWuyPgNz672OwFkEjthd69PX8+2tWW2uRw/JbfH38qpZVefYqgWyohops+kEgOIvukkcX9RuMOcSSJ4wHGTWdwaGWSVtJwMb/tVWtuSnqk+5OcU86r854Ke1rcYLYiJYAqDAHOgcVjkS6Or4eFGKrLqTmXIzLz61abdYg3K0/kB/8A5+t9+MZ1rHbLdalXHPFa99b3AsclsnbIqtlCRJizluMqPQ0/fIKiAv03G+x3/S3BwzxRInVQ4yaU4JDLIzaTgfvQyoqkqqDy5ji3ZJWSpzuT2I9P0w1bdCLcBBhQOVIXscq3RVuedv2o3mGVUX6PDbkyCsIP3/m/3FW8pI9w222vvzjNsI7dJyVXHZ51scUt51tQS2eWfKo8oSZkVySphxYYKN0kkJK/yq29OffwdsF4okLlQy5NC4JBNIGYHC/vQV7xSpy1OuOLlFy6l6vMVX5vh9OiFv7o93HKsp45xdaP+eaOZqlzn48NMl4qQhP3ljsty34iO3w4FvfjO4bHAkp0rgmtjjEEywqS2R10uVJU6NGliO8UNqT93cmyHLbKt7vTg/piOJRwSSDUuTip4NBO8TMGwOqgTKpQmoU2twS+oLK1ebVfm/8AfGk/RG3O3u45VjJHP7UFH580azfKmSXGC+6SyE2Si+2v8yvXfn3cDGfwuOBCwVcGtXjUEqBSWyv71Jl2VUmKPLajPKQFf7BKt0Lt5iPiNt+Dvzvil/FbvOCVyeujcLt53tWYNjsoJTlSUVCO5DcU3J6gKHO4Pcn1740rkRG3IYZUj/FY9nHM12FH5s7/AL0TzdJlSZ6FPKPRCfuUBRIR6877nf6cYT4WkKBgowfv9Ke41DLGVLNkVZpEupNZckMsPqAJPQ8xuhO+sD4/S23OF7qK3a61Fc8s+NOWNtcPYkq2++KDURyS1Uo7kB0tOg7LT2He+NG/EbwFXGQf1rI4XHI90Am3b4VbzU/IkVMrfUenpswjUSG0X/Dc7m3qdz3wLhqxLGQgwc5NG4zFLFKNRyMbfv8AGr0eZUhlVTCX1hNvJY2IZ7pHuOx9RYAYTeO29s16ev51oxW1yeH6g2+NvD/GaGZcdksVVpcVZbI/HYkakdx+vr25w/xFI2hw4z2eNZfCYpZLj3DgY3rsxuyH6s85KcKirdFzslH5Uj4DHcOWJYcRjHbUcYjkjuMOcjAxV2LlzMURRWxFSlRFiSts/wBScBbidg/5m+R9KheH3yflU0eZykibCjvV6U+xIcWWEIaQhQvc2vpH/OEzxI62FsAVHbTicOeRMzMdQ3O/V8ao1WhZgblSIsVsORUkoQpXSupP0wWHiNkUBlIDdexoMtleFj0eSvVvVCLl3MUQrLEVKSvZRK21f1ODtxOwfm3yNBTh99H+VSPjR5rJzM5iI5V5UhmdJBbS0hKCm4udyBsOP3vhP8Sc6+gUFRy7/Snk4cXUNKxyOe/KhtTomYlvSWGmQuMVKSkktXKb7bi2Dw8RsCisxwfA0vJZXxYgZI8aqxaBmOJq6EVKdW5uptX9TgjcTsG/M3yNCXh98gwoIHjWgZybHkmC5UZkhqoSGw500pQUpKbE7gWHB/drpHiUulzEqlAdu8dfhTy8O1gO7HIxnfroLNomZHi40WEqZ1HT5mgSL7cW92G4+I2AUNnfwNKSWV+2RuR41DGoWZIaCliKACbkXaUePfizcSsGOWb5GqCwvlGAp860ScmxFSWlPTpKaitvxHTCU6drX3027H92uh+Jy9EWCrpz8u36U/8Ahuf4jMc8ue+eztoBJoeZZLZbcjoKL3ACmh/S2Hl4jYLght/A0k1jftkEHHjSRqHmSK2G2YyEovfctHf4nEtxHh7blvkahLC/UAAHHjWhXkyKiS6pqbJVUENl/plCdG4Nt9Nu4/8AfCB4nL0QZlXTn5dv0p78Mx74Y55c989nbWdkULMsppKH4qSkG4AU2mx/Q4fXiXD0OVb5GkWsb9xhgfOpIdEzGwGmfDJS1qF92iQL3O++KycRsDls5PgastlfDAIOPGjT2To0YznKdMku1GO2XOktKQklQJG5Fj2/fCY4lJoQyqoQ8+6m24bpy6Mc9W43NAJVAzJLCQ/GBCCSmym02+Rw8vE7BN1b5GkWsL9hgg+dWabRMxNux2FsJRG1hKiC0SlN97HfA5eI2GklTk+BokdlfagGBx40UdyczCYmOUiXIenRh0y24hIQVHSeSLHg/vfC/wCIv7ntCqFPPuphuGlVLRE5PLfnQGVl7McspL8UHSPLZbabfI4cXidgm6t8j6Um3D79huCfjV6l0OvrlRospoNxSQhRSWrpTa3O+Az8SshGTHgt1bHnRorK91jpMhfGrr+UUQYcl6gSX35TS+ipDqEpFwQTyNx8v+RfiJDqLkAKfHn99dFfhrImYScncb8x8KBysu5hlKSqRFSopFhZbY/ocOJxOwT8rfI+lJvw++f8wJ+NaifUZEVuOpuOlwuJuoWOx22xl2nDYJjIHfTpO3LcVt3vFLmBYikWrUMnntV+m1WaqOdMV1gW2B0+b4X/AL4zri3jichJA2ezNPwOblNckRU9+PWhcysy0yCDT1oJO5cuSfhbb6nGpacJtpUyZh8PrWXd8XurdgiW5x3/AEqSo1GREdbS1HS6FAE7Hn02wOy4bb3EbM8mnBx1Ue/4rcW0qJHFqBGeuiUKrTPCqPhXmgBcNnTdXwudv1tjOmgSNyqSZzzxmnoiZ06SSIqew49aEPViWJQbNPU3qO/UBvz7tvqcatvwm2eMsZgfD61lXHGLqKQRi3IHafpUlRqUiJK6TUYOpIHF7m/bjFLPhsE9v0ryaTk0e+4rc2910UcOsYHb+vKikaqzTDUfCut2GzPkurbtvt+tsZjwIjmMSAg8zvin0PSJ0rREN2HGf1oP7XlKmpZVBLSVKAJWDcXPPFv64104VbGAydMCQM9X771lPxi7W4WL2cgEgZP0p8ypyY81TDUTqp2ta9zsD6Ypb8NgltlmaXST6miXPFrmG8aFISy9oz2eVFU1WYIRUIjotwx5Ln6/u2MowIH6ISDT274rRHvJ0piOrs2z+tCI9Ukuzksrh9NBUAdV9Q9+NZuF262zSiUMQOrH+azI+L3TXSwGAqpOMn7xSv1SS1OUy3DLyQbDTcE/THJwu3a2WUyhSRXScXuku2hWEsoPMZ/xRZyqzPBX8I6T/wBgadVvn+/TGSIEL9H0gx274rUO0fSiI6uzbP60Jh1SRImhhyJ0kkm5UDcbd7gY1bjhlvHbGZZQxHUMd1Zltxe6muxC8OlT1nP+KQ1aUiatpMFTiUqI8lwTY88WxZuFW3QCTpgCRnfH7VC8XujctD0BIBIyPrtRaTVZnggoxXnLjdkabp+O+/6XxkpAjyBDJgDr3xWm/wDDTpViJbs2z+tDafUpEuV0lxemjc3sb/DfGld8Nt4bbpUk1H4Vn2XFbi4uuieHQN9znq+VRNViV4oo8ApwJJtoBvzz6YLccJtki1CYA9/0oUHGLqWVo2tyR2j60Xn1SYIoJjPPXA8gCfJ8bH+l8ZMMCSOFaQADlnOK1JWMC9IkRJPUMUOp1RflvLQ7GDQAJHPyONG84bBbxK6SasnHVSNhxW4uZmSSLQAM75qOHWZapFhT1rtx09iPjfb64vecKt4kysw36j9KDacXubhtD25x2j60TqVUlhlN4jr9xwnT5fj3/UA4zbe3SZ8PKBjtrTndrZdccRY91D6dUXpSHlOsBvQjUNjufTfGhd8Ngg6PQ+rUQOqkbHilxc9J0kWnSMjnvRNdLzJ4VrpUV3xAI6mopt77b97YXFg+o55eNMnikGkYO/hT3aXmASmenR3xHseqTpKt+Lb9v74qOHvjfnVjxWDOx2pWqXmDxboco8jw23T06dQPe+/f+3vxJ4fJpGOdcOK25JydvCo26XmXwrvUo73iLq6dikpt2vv24xJ4e+oY5eNVHFYMc96V2mZk8M10qO74i46tynTbvbfv9McOHuGOeVT+Kw42O9PdpeYPFNdKjSDHseoCU6r9rb2xA4fJpOefjXHisBI328K5ul5g8U51KM+I5FmtJTq999+/0tjm4fJpGP1rhxW31ZJ+VMbpeZTEd10Z4SPN0tOnTbtfftiTw9tQxyrhxaDHPeudpeZPCtdOjO9e46oUUhNu9vN/6vjhw99RzyrjxWDA338Ke5S8weKZDVHkeH0nqatGq54tviBYPpOefjXHitvq2Pyrm6XmDxL3Uoz/AIfSOnbQFX73837tiDw+TSMc64cVgycnbwqNFMzJ4Z7q0d3xF1dIJKdNu1/N2xY8PbUMcvGo/FYMbnfwrnKXmTwzXTor3iLjqBWnTbvbzd/pjl4e+o55VP4rBjY/KpHKZmDxTXTo7/h7HqlRQVX7W83bFRw+TSc86n8VgJGDtXN0vMHiXA5R5IjWHTtp1A9774k2EmgY5+NQOKwatz8qY1S8y+Gd6lGe69zo0qTpA7X37fW2JPD31DHKoHFINJ338KRyl5l8Kx06M94i6epcJ02723744cPfUc8qk8Vgxz3qR2mZg8Y106PI6Fj1NRTqv2tv2xUcPk0nPOuPFbfUN/lXM0vMCpToco8nw9h09JTquOb797/THNYSaRjn41w4rADz28KY1S8y+Ed6tHd8Tc9O2nT7r79sWPD31DHKoHFYMHfeucpmZPDNdKjOmQVJ6mop0kDm3m7/AN8cLBtRzy8ak8VgwMHfwp7lMzAJLRbo8nw9j1dWnV7rb/u/uxA4fJpOedSeKW+djt4VzVLzAJTuujP+HsOlp0ar97798ceHyaRjnXDisGrntTEUvMvhXepR3hJurp6SnTY8X83bHGwfWMcvGqjisGDk7+FeiKafkSFLD78dJjloJFrXP5xt+IY120ivMRuxGCtcGX47yFl999KY4aKDaxI/OduTjhg7ZrpGYDZaRDbz3hAlx1pLTutRSBZwbjSbjjEvjehwO2AMVxiPsdE+NkO6H1ukeXzBV/Kdvwi/9MV2zRnYj/jmmvB5+PJYbU6ytw/7jdtTfwuOcXOBzpaJ2BPumny4T7iJSkz30l1CEgICfu7clO3JucUBBptsgZA+/Knl5zxGsNqCtFuPri+BppTW3S7KaZGhPpRGKp75LSFghYTdwq4KtuR2xQ43FNBiVzimMdZhmOy4XnlN7dRy2pz42HOCAA5paR2LA6ad4SQ+l4+OkMlx9LgFk+UADyjb8Jtv+vrihxTCMW3IxSLQ6wJaVrddS65rSpdrNDbYWHG31xK4zQJnOkjTTlMPyH1LEh+OPDloJTawJ/MNvxY44FFiZmGNOK5LD0d9C+u9IR0A0pG1iR+fj8WIGK6RyNtNNbbdeEQIcdaS25rUUWssbixv23/pizAChQOdIGml8I+z0VibIdKH1OkeXzA38h2/CL7YptR3YjkM/fhTXg6/HksoU8ytw26jZ8zfwvghApaNyGPu5p8mE+4iUU1CQkvIQkBIT5NPJTtye+KAimixAyB9+VO6znX1dJdyi2n++LYGKV1sJM6TTY0KQ2iMpc99fSQtJ1hJ16u6tu3bFCRypoMSOX35UxlLsdmMy6XnlNm3Vctqc+NhzggA3pWR2LD3cU4RJDoePjpDWuQl0CyfKE28g2/CbD6+uBnA2ppGLDdcffhSKQ6wJWpbriXHNaSq1mhsNIsONsXXFLzu2MYp5ZfffWtMh9hJYLQSLWBP5xtzipA6qLEzEYK0iWXo7yVl999IYDWg2sSN9XHOOGDXSOwGy1XrqRGpb0uLTEynw2SbaAU+UnUdXNrDbnAvZoXb31FGWeVUADGky+BLpUeZJpiYkhTaT+Q67oB1DSTYG525x3s0Mbe6o8qhp5WXDNVlTDSoPWVDTKdShWhvygq34BNgP1xaWGORssM0O2ldI8KcUKyq4uqMSFT6OlgNyHkIWemQdLhGmyTyALE8HFXs7dCNKjyo63Mx2LGjLbDL630vNocAXcBQvi8kaOoDDNLQSOjsQaBU156RmSZCkUENR2m2yFFTR6dwvc2Nzq0iwHGKtZW4XVpGfCmBczZ/MaPFhoyOj009PpW022xbQvR6cbUuXbp9Wd6AVR96PmSDBj0FLsd1DyrhTQ62kI8wubgJ1G4Nr3GKiytipOkdXVTHtU2oe8aPusNMFlLLaUpLm4SLXxaKNEBCil55GdlLHNBc1uLpkZlcCjofC5DKFujp2SFOJSU2URuQbX7E4olnbtnKjypg3MwGzGipZaTALohpiuKSNTflJTvwSNji0UMcb+4AKDcyyPFhjmoK6lMWmPy41MTLfS2VG2gEeU+Y6ubemK+ywO+WUeVFWaRUADUlBQJFMZlSaYmI+WwbHQSbjdQ08X9Md7LDG+UUeVc08rIQWqx0GlwA6YaZTqUnS3sCrfgE7D9cTLDHI+XAoVtLIkYCnFCsrOKqkZ9c+kJZ0SHkJdJbIOlxSQmyTe4AAv3xD2dup91R5UcXMxG7GjTTDL630vNpWnqX8wvbF5YkkADDNLQSOjMVON6AUuQ/JzJOgyKCGYzTbKgSpo9LUHNzY3UFaRa24tijWVsEB0jypgXU5ONRo/0WfE9Hpp6XStottzi/Rp0WjG1L9I/tGrO+KBVJ52NmGHCj0EOxnG3TqCmh1LBG4ubjTc889sUFlbFSdI8qZ9qnz+Y0ddYZYW0hltKEld7AWxeKNEBCDFLXDs7KWOaC5pcVS2IyoFGS/wBSSyhbiS2ANTiQU2UeSDa/a+KR2Vux3UeVMm6n/mNFeg0iD1hDTFeWka29iU78EjY/piYYIo2yigUC5lkePDHNV6+kRaXImR6amVIDSjcaAU2STqOoi4FhtzivssDvllHlRVnlVAA1dQUCTS2JcimpiyC0FD8JKrpB1DSbAG/HOJ9lgR8qoyO6uaeV1wzVNWJsNugvvvSEtxXGCkOWKvxiwsBuecWSRSA45VZbeVpOiA3rqNMhroMZ5qSHY7LCUKd0lO6RY7HccY5pFALnlXPbSLJ0LDepKZUIkmmCYy6PDpSSpSj+G25vikcwmXUK6S0eB+iPw781QyjMgyYspEGWHz4p11f3akFIcWpQFlf192Clw+4q01tJb46Qc6vUyoQ5kiW1Ec1LacssE/0wCK5SUkL1fearLZvbgM3/ACofTahTns0VJLExLkh1ttvpdNQsW9WqyjsfxdvTDGsEBas9pMkfSsNjRD2hD9siF1f9T076b7fD44B7SnSdF1/e3jUexv0ftHVQ6qVCnNZppSX5oakNIdbDRaUQS7oCbqGw/D39cH6QAaalLOaRDKq7CiNTqEOE/Eblu6FvOWRvYfE+7C8tykBAbr+XfVY7J5wWUflqjm+ZBjwY6J0roEyWXU2aU4VdNxKzsnf8tr4YDhDvVobWW4JEY5VfqlQiR6WqW+6DHUE6VIP4gbWtgMswhXUfhVY7Vp36L72plXmRG6E++8+ERnWSlLmkq/ELDYbnngYusqkBxyNctvI0nRAb8qSjzYblCjvsSEuRmmQkuEEfhFjcHccY5pFClydhXNbSJJ0TDf1qSl1CHJpiZbDv+nSFFSln8Nib3xWOcSrqrpLV7d+ixVDJ8yDIhPtQJYkFMl1xV2lIIDi1LGyt/wA3PuwTWH3FWmtpbfHSDGavUyoQ5j8tqI5rcac0rF+/qPdzgMdyspKjq/TtqJLN7cBm/wCW9DqZPpzuaqomPNS5JdbaR0umoBJa1hVlHY7q7emD6wRpHVUvaTRoJWGxoj7Rh+2fBdT/AFQb1aSdrenxwD2lOk6Lr+9vGo9jfo/aOrl9fCh9TqNNZzRTUPzEokNNuNhotqNy5p03Vaw/D39cGDgDTUx2k0kZlUbCiFSqEKHIiNSndLjrlkC/u5OAy3KwkKev5VEVm9wC6/8AGqObZkGPEjNzZQjnxTTqLIUsq6a0rOyfhz78HDhNzXQ2stwSIxnFX6nUYcelmW+7aOQFJUk/iBta2AyTCFdf2ahLV536Lkf0qOsS4aaBJefkJbjOsFAcKSrdQsNhuecXSRSoccjXLbyNJ0KjeupE2G7QmH2ZIcjtshCndJTukWOx3HBxzSKAXPKue3kSTomG9ed/Z9Vmq9S15drocFyC0q+lVxxY/v64HPGtvNpH5G+R7PA07DNK8YuAMOmxyOY7aX7QKszQqUnL1D6h3PVWLqN773I/fyxEEa3E2kn3FO+/M9Xl+tTNLMkZnYZkbljqHb6UE+zXMSoTztJqiVmDJSW1BwGwB2P9cHv41if2iPlyYd3b8KBZGWaPomBDLup38vSthmR+nZJo8hFK6jkuaLXKtStO9rfPCpXpXEMZ58z2D69VNdNJKOnmGQnIAcz98685yfmSdRK4JTwdU28r726TvfvjQu7VDGpgxqTl3jspG1mlkdkuAdL9eDseo/tXqshuhUJl3M0UKU68lSmklV0havxWH6n6+uMwzKVDR/mbYDv6/Kn83Ep9nl/KvM9w+9q8fOYap/EQrFn+prvax/DfjGp7FD7L0Orfnn/9u2s/2yc3OvQdHLGDjT9717BFRQ67HZzLK1B2OhKnUhVklSdxcfv34zBMAhaT8y8/vvrQzPH/APTxbq3LwPP615VnLMc6tV0yWkupaZX90Eg22xpWdsgRmmI1Pz7h1CkLqeaN1jtwdKdeOZ6zXo+W36fnijx26r1G5cQG5CtKtPfn9MZygROYZDy5HtH0p7ppIh08AwG5jsP3yrH/AGlZiVNeRSqUlYgxkhI0A22/Yw1YRiVvaH5clH7/ABpW9MsUfQqCWbdiP0oz9n1XZrtLVl2uBwC46KzdJuPQn9/PAp0W3m0g+4527j6Hq76NDNK8YuAMSIMHPWO31rs/1hqg0pGXaH1DuS8oXUbnc3OOgjW4m0n8i8+89nwrpppkQ3BGZG2GByHb6UG+zPMKoLy6TVELMGUkpIWDbfbBr9Fif2hDsdmH7/Cg2Rlmj6FgQ67qd/L0rX5ifp+R6K+3S+o5MlgWKlalBPbjtzhUr0jiFDz5nsH1prppJR08w2XkMcz98685ybmOdRa8JTwdU26r70FJ7/v+mH7y1QRBoMBk5d/aPjSNpPNJIyXAOl+vB2PUa9VkpodCjvZljai48lS2kFdwFK5sMZ3TKVBj/M2w/fPh10+TPLmCX8q8z3CvIP4hqn8RmsWf6mu+mxFk3xp+ww+y9Fnfnnrz2+dZ5vJ/adeg6OWMf8a9fjIoVdYazLKKkuspSp1AVYFSfw3H6DGWJgFLS/mXn999aGbiP/6eLdW5eB515VnLMc6tVwymg6lplX3Vkm3/AKxpWVqgjZpsam5+n310hdTyxuqWwOlOvB3PbXo2XH6fnijsN1XqImQxbyq0q0973+GEAoicwueXI9o+nXT3TSxDp4Rs/MHqNY/7Scwma61SKWlYgxUhtOgG21gN/wBP6YZ4fGsz+0Py5L6/Glb1pYU6JASzbsf2o19n9XZr1LOXa4FixBaVcpNxxb9/1wKeNbebSPyN8j2fGjwzSvGJwMOuxyOY7a77QKuzQaWjLtCS4RqKnVbqNzzc/vtiIY1uJdJ/IvPvP0qJppY09oIzI2wx1Dt9KjNJQBIKK/T0LUbx1e1yej8fL5vjth/Xw/I/hr5D7FZxs+KkHd/7qf7LZMlC1V2Aprp2Wj2ydSl/zatO3wsfjjtfD8fkXyH38a4WfFR/PjxaohSE9FtJzBTuqHNTi/axstH8oTp8p99z8Mdr4cSf4a48B9/Cu9k4tjm/m338ankUyO4+ytuuQUtp2cQutlZWANgDoGn5HEB7DG6L5Cp9j4oCNOsDxaoF0gGO8lOYacl5S7tuGqkhtPoRp8x99x8MTr4fn/7a48BUex8Wxzf+6plUxovkiuwQwEWQ0a0SUr/m1aPpb9cdrsMY0L/6iu9j4pqz7/8AdUaaQ2ER9Vfp5UhR6xFWP3o9ANPl+O+O18P3/hr5Cu9j4tgbv/dXGkpDUlKMxU9JcXdr/qxs2L8KGnzH37Y4Pw/b+GvfsN/SuNnxXfd/7qeKVH66Ve3Kf0enYte2Tcq/m1aOPdb9cRr4fjGhf/UfpU+ycVzzfzamN0oIbZvmKB1krPWWmrlIdT6AafLtte5xYycPySEXu2FV9i4qQMl/7q5dIbIk6K/Tklduh/1cnpet/L5/piNfDxj3F8hUm04tvu/m1ORS2UyGnDXIGhDdloTWSCpf8wOjyj3WPxxGuwx+Rc/+IrvY+Knnrx4tSrpTKpDqvbsDpKRZCDWSSlf8xVo3+Fh8cTrsMbxr5Cp9j4r1F/7qYmkt6Y2uv09RQfvz7XI6vpby+T644vw/J/hr5CuFpxbAGX/urnaWHGngMx0/qlQLK11YqDaf5SNO/wAdvhjtfD8jMan4D7FV9j4sAcF/7qeaVH8Rq9twAz09Ib9sm+v+bVo491v1xGvh+PyLnwFW9k4rnm/m1MFKHSjIXmGnktqu8fa5s4PRI0+T474sZOH/ANNe7YbetV9j4sMHL56/zb1yqS2pEgCv09Klm7J9rEhoehGnzfTEa+H5H8NfIVPsnFt93/up6aYyHh/12D4costoVkjUv+a+nj3W/XHa+H6caF59g/Su9j4pqz7/AJt+tRJpCQ0yFZgpxeSu7qxVbBafQDT5T77n4Y4vw7J/hr5Cu9k4tgbv/dU7FNjtyH1uVuEptX+0hFbKFIFtwVaDqv8AAYgvYdSLn/xFT7HxQn3tZHi1QGjjoupTmGnB1Tl2l+1TZtH8pGnzH33HwxPScOGP4a48BUG04uQd3/uqT2UyJClJrsANaLIQKybpc/n1aNx7rD447XYYHuLn/wAR9/GpNnxXJ3fHi1NFKQfDlzMFPWtKiX1e1j978PL5frjtfD8k9GvkPs1HsfFgAMv5tQuRmxmKkKkU+U2CbXJRz88eYT/TsrnCyr8/SvRTcU6AAyRMM+HrXRs2MyklUeny1pGxIKP8sc/+nZY/zSr8/Suh4p04JjiJx4etXY9XkyQSxRpywObKb/ywvJwlYzhp1Hn6UU3rjnC3y9akM6eASaHOAAvfW3/lgY4bEdvaF+fpUe3P/Sb5etUzmMi96XN253R/lhkcBY/9y/P0qxvJB/0t8vWqjWdYTrobbiSVLUbAXTv9cGb/AEzOoyZF+fpScfG4pGCqhyfCrzdfW6sIbpUwqPA1I/ywu3BCoyZl+fpTntcn9Fvl61a8bP8A/wBFO/8Ak3/lgH4bF/8AIX+70qvtz/0m+XrVd+tuxlBL9ImoUexU3/lgqcHDjKzqfP0qwvJDyhb5etUpOcY0VQTIhSkKIuASnj54ZT/TczjKyr8/SlpuMLC2mSMg/CpmszoebDjVNmKQrg3Rv/5YG3AHU4aVc/H0osfEGlUOkTEHw9auNVKY82HGqJOUk8ELb/ywu3C0Q4a4XPx9Kub1xzhb5etc9U5jCCt2iTkpAuTqb/yxy8LRzhbhc/H0rheueULfL1qk7mhtltTjtOlpSkXJJRt/5YZX/T7u2lZlz8fSqScQaJdTxMB8PWoo+cYspZRHhSlqAuQNPHzxZ/8ATkyDLSr8/ShQ8YWZtMcbE/D1q7Hrb0lZSxSJq1DewU3/AJYXk4OIxl51Hn6UybyQf9LfL1qx42oW/wDwU7/5t/5YF+Gxf/IX5+lV9uf+k3y9aqOZgW2soXSpiVDkXR/lg68ELDImX5+lW9rk/ot8vWqTmdobThbdhykrSbFJ07fXDC/6ZnYZEi/P0pOTjcUbFWjOR4VcTmMnilzfmj/LATwJh/3L8/SnBeSH/pb5etWxPnEAihziD31t/wCWFjw2If8A5C/3elV9uf8Aot8vWopFXkxReRRpyAeCVN7/APlgkfCVk/LOp8/SpF455Qt8vWqcjNjMVIVIgS0AmwJKN/rhhP8ATskhwsqnz9KFNxToADJEwz4etdHzYzJSVMU+W4kGxIKP+cc/+npYzh5VHn6V0PFOnBMcTHHh60SytQ6gmoUKvVB+mMQCfEJQXldQgoI4KLX3HfHo8RQagCcnasie7mvUA07A52rs1USoGfX69BfpkiClPiVJD6uolKWwCLabflPfHYim0qSciutrqWzQ4UYPjVGmVB6OlDbdkJWQVE7fHGNdWyvknqr1C4dQT10UqtVKA2I7qVgiygcIWtmGzrFQqBRyoJT6fUazNmMU4wm/DMJfdXKcUgAKJHZJ9Meiihj6MM5PP9KzL7iEkM3RIoOR31qahRHJOTaFRYlQpip8BaNYW8oIWQCLAhN+SO2GOmRi2TsaxIklidX07rWW6c2k1aTEmpY8TFKQSyoqQbpvyQD9MJXUC6QAdjXobC9a5DFgBjxo9FqpXCUt11KXyPKMYMtpiXSo2p0oMjas9Vag++gOkILl0oF7gXJtjZs7VFbQORoN3MYIS6itlk2iyst1+VJr0umI1xktJQy8pSh5tVyCkdsPlo1XQhOc15u4kmunErLjbFY+Tl6q0KmRpEh2myIi3wwFR3lqVc37FAGJkSOXLqTmm7S9lg6O3ZR86J0mpuNuIZcKUNJG5OPPXdqCC45mvQMoNLVqmtbi2ELSthQ/XHWlqAA5G4qVXHVvQuNl+rV6lyX47tNYiIfLBXJeUlVxY8BJHf1x6CNYoSsjHevP3l5LMXgC7VsM50WVmSvxJFAlUxZbjKZU2+8pJuVargJSe2J1I6aGJpS3eW1fpVXO2POsbSZ77KS6EIDl1JVa9tjbbGfeWylih5CvR2kxuIVdhjNaKTVCiC2tl1JfFtQxix2mqUhhtRhGM7jagHTnVWrRokJLAlSlKALqilHlTfkA+mN+1t00EE4ApLiF49qF0DOfGtRT6K5GyZXaLLn0xM+ctZRoeUUIJAFiSkEbj0w6JkVlwdhXnpUlmcsV3besvUKdUqNOix6j4NYksqeaXEcUsFKSkb6kj+bC8sMfR6kJ59dbdjfSzS9E64wO+jdKqxV1BKcQgAWSB/8AePO3VoF/IK02TVuKF1SoPSAtDhSpKFHSfXGha2qoQRzNS38NSQOVXsqUWoCdl+vz5FMYgq/1SUF5XVKVNkAW02/MO+NkCKHUoJya8xPdTXiDKgAGlzTRKguoV2uwJFNfhAh9SeurqBIQBawTb8p74gCKbCknIqILuayQ4UEHei8uO9MZ/h4wXG4seKlDDiW1aA6nkAnt255A9McHKkTZ3zy++2gGNWVoWX3SOf32bGkix3ocf+HRBcciSYi25K1IOlTqxwSO1ttjtc44uzsZifezy++yuWIIFhAyAOf327/Ks9XstyKWqjPRaXI1vsOqlNxg46EKujQDcmx5xZ/4qEdefCi2VwLeUM5OMdud6HLp86W5Hj+yqkEOSGg4THWkaNY1b9tr4HDCUfUxHXT17xGKaEpGTmtmKc9lNypuUqlguPOIZYLXUWXGwdRJ1KNrXI7cnF9XSYUtjn9/pWVllBfGezfJqE0pdKkzazBiPrW4hJiNBBJSpXPb8tvmBjulaRUhc4A6/vsFWMKRO86DLEcv18/WurGX/biapVnqbZ1dO6qFpLgdL4QQE6b22IvxvcYmOTBCk9fZtiqOpXJXIz37/GsmzGqCWm0qpVTuAL/6Vf8AxhdrZixII863o+LW4QA56uqjNBykqVSGZ0mmK8QqetLviS4lTbQsQoJBH9MMs/RjGerq7aw3bpZGIyQT1nqotVIsvNpdU/DdYfZk/dFxBGpkm3yF/rfFUla3bUhzkeX+KiSBLhNDjAB8/wDO/wAqWY3IrECTSBSlOMR3W1w0vpWgLAGlRukix3KsRHhNLZ58+0VaQM5OFwRjHUD/AIrO1ygv0avy48KmznIvTaKFNNLcTqt5rHfv2xEiGVV5Z6+qnLC9jty+snBxjr7aSiUJ+r12JHm02a3G0Olanmltp1afLci3ftjoozEjcs7d9RxC8juNAjJwOfVWkhtP0mFFpBpSkMyXVrmpYStxLYOyTckknYKxLnWCwbly6s/4pVAVIyuSefWPs7fOmUqLKyn0fDw3n3nZB6qkIKtLINvrb6X74mSVrhiznGB5/wCaoluluuhBnJ8v8cvOhddykY1HdnRaarxInoDQjF1RW0bkrKSo8/S2LLJr2J6usddShMLg5xg9R6qCOxp6mlpFJqYJSQP9Kr/jC62zg5JHnW5LxW3KMBnr6q1dFy/7DTTKs3TLuopyXVKUXC51ykAp06rbk349cHkfJKg9fyrBjUtgtvgdvX3ffOl9kLqsqHWZ0V5C0NnxbWggqWni23fb9ScQZmjVoVOQevu+u1XEKSuk7DBHV99m/wAqn9nP5tdpjlWpYDjTi2ny7rbLbZOoFNlAEG1v0GOLdGWVWz8Pvqqu7KHIxzzvg/fLzNYxFOnxFyGPZVSKUPuBsiOtV0azp377WxWeEu+pSMbfpWpY8RihhCSE586IULLb9VNYflUqRrYYaMVuSHGg4q69YFiN7acXT+EihiOe+N6TvLgXE2qPPLA3xvWilR3prH8PeBWiJFiIbjOJbVpDqBwCb/AXO9hviodkImByc8vvt9KAY1dWhI2I5/fZ60sNh+Gz/DyYDjkN+KpD7im1BKnVb2JFtu1/eTjmcuTKTvnl99nKuEYULCq+7jn99u58qwkLNVbqL6mYQklzpL8gqbqe1tXmX25we76OFAztgZ7KVto+lchFzgdtdU81VuA54eUJDLym9QCam8ogHg7OEYtA8VyC0ZyPCqTxGAgSL86srruY0xXZ62XvClIdSr2s5ZKCLC1nByQe3uwtFc2wIhL5blypp7OTBk0ADxqGn5qrdQccREaecU2jWsGqvpsn13cG2Gbh4bcAyNjPdQIIWuDhF5d9RTM2VmmyPDTfEB1pAC0GpvHe97nSv0I2PuxFtomUyI2QeW3LuFVuEETBGXB8asy65mSLEclvIkMx9LatYqrhA81/+4SdQ2wulzbyzKkb777aef8AimGs3jjZmQdW+aZT8zV2oqJhNSF9EpU4ParqTa/Hmc4PGC3ckMK6ZH0k8jjNCtoGmOUTOO+oH831eHIVFlh/rNnSu9Te2PvIXYYNF0ckQkU5Hbjn30ORCkmhhv41bmVzMkBgPTW5CW0pS2VGquG6rXv5XL3IPHoMK21xbyuUV8k7/lxgUxPaPGmtkAA7+dMgZizDVEuqhIkrSlaAQKo4nT7t1332xa5lht5F6V8bcsc++qQW7TodCZ37agVnGrNzhHWHg4h3SpHtN4gnjTfWRzgzKskBdTgY5438aoFCzdGy5Oe2rEyv5jpcZC6gzITuEajVXSVqt6Jc578YHBcW07aImyfCrzW0sK65FwPGuh1/MdVjOKp7T6hct6xVnQUqI22U5z3x09xb27aZGwf/ABqYbWSddSLt41AnONWdmqjo65ccd0oR7UeAB4sFBYB3B39+CKqxwB2bI7cb0IqHmKKu/jtUs/MWYaYhpU5ElAWpYAVVHVX422Xfy2+uA28kNxIejfO3LHz+NFnt2gjHSJjftqSHXcx1BhT8JmQptSVISsVVzZXN/Mu+wBxW4nt4nCO+kju5j61aC1eVCypkeNVGc4Via+IkTr+IcOlFqm/ufcS5bDMvRxxdI5wPDlS8aa5NCrv41PUczV2nLCprL6OspSm0+1HVWF+PKvgcYFaPDONMb6iOZxRbmBod3XAPLepIldzHLity2G33Y9nFazVnACNVzf7wEaRtv64FJc28MrK777bY+nXREs5JY1ZUHjn75VVh5srU+R4aGJCnnE2Qj2k8N9jcal82B+Z9MMXISFRI7YAO+3VQII+lYoq5PjUk7NVbp60IltvoUtGtAFVfVqT2Ozh2OJt3huATEc47q6eF4MdIvzqdFezGqK1PSw94UJLxWas5YoA3uC5tzvt7sLy3Nvkw68Ny/L1/WjpZyYEhQFefOq1MzVW56xHiJkPvhGo6qo+nUByd3AMMXDRWyhpTgeFLwQmdiqLk+NdOzVW6fISzN8SlzpI8oqjx7fi8i+/NsVtOjmQsjahk8x8vhU3KdCwDpg47fnWeoLbi6k0631LsnqakIKiLccA/XBL1j0ekDOaNwuKNp9UjAAd+M11eacRUXHHNf35Ll1pI555GJs2zHpYYxU8VjjWfMTAg9+aIPRnxlxEf7+zZL+gtqsCQL9vQDvba/fCaSMLnpNPyrSkt7cWPR9INYHb8cc6qZaQ6KimS0XEmPZYUlBVvfa+x79jhq+ZtAUDOaz+ERRGQvKwAAxzxzqrVWVMTnEkLsTqGsEKIPx3wSzc9CARjFB4lEguT0bAg9ec0WqUZ9NDjslT5EU6rKaIG9tXb3Dni22Ebd29o1FcA5rWvILb2LQkg1DHXzx1c6r5WDyJxltF1KmU6QUNldyfXYj54Pfs2Aiik+DxQks8zADlzxQ6bGUxOWxZSiVWTcG5v69++GIJP4AJHIUndwoLoqjAgnnnt9KMVuO8KVERqeKIidIDjRAsTvuR6+vb3YRtHYSksvPurV4lDb+zDo3BK46+e2O2o8rodQ69Ja6wBSWrobKtjz2I449O2L35LEIBnFC4PFAFZ5mA7s45UJciLTLMSxK9Wj8JufTbnDiy5h1kVmvbqLrogwxnnnqoxmNp4xoq1KfWlhHRHUbKbD1vbn478YSsWZXIYc+6tXi0UBiVonBI5jP1rstsuhiStCnkiQgs/dtk7c8gGx355G/riL5mZgAM4ruExQLCWlcDPfjGKDoirXLEWxCteggJO3vtbDzS/wdYG5rKjt1N0ItQxnnnbHbmi2Z0OrWxIdL5AQGvvGym1htyB279+ecJ2BYEqwrS4xFAQHhYHHVn60+iR3jSZSQX0oleU6G1EEA7bgbb+nbbg4pdsxmBA5d1G4bDbi2PSOAT34xQiFGW/NbYssLC/PpSbgDm1hfth6aX+CSBzrJtIFN0EZgADz8KI5oQ65NTKdLpLw03cbKbEegsBaxHHvOF7AkAoRjrpzjEUAKyQsD8c1YpsZ9dCfZT1wJZC9KW1afLe29tu97c33wCeRjcawvLHVTdpb2/sWiSQBjnr5UJpLKpE5pI6mx1koBJAHww9duei2Gc1lcNiRrka2AAq3mVLpqS5LxWVSfOStGnzd9rD5cYHYsdBRhR+LRRLIGiYEY7eWPWrTEZ85cWwOvZwh/T01abi9he3oT7rm+FWkb2npNPI9lPx21v7DoLjWR2/HFD6E2t2otrb6l2fvLoSSbjj1w3esRFpUc6zOFxRvcZlbAHz7q6vNrTU3nHNZLx6mpaCkknnnHWTkx4Ixiu4rFGs2qNgQe/OMbUeTkh1J8tWbHwZWP74zPx5P6Z8xTv+35f5q0VJypAbp7CZkNFSf6mhboJToQb7m53tf97YUPEpJndlcoMDb9TTKcNjiUJIAT+vdQiqZNLlQkGNUWmWCohDQbXZI9OcFg45pjAdCx7cjfvoL8Dd21IQAeqqickPJ4qzY+DSx/fBvx5P6Z8xVP8Ab8v8wrTU/KlMEKG3Mgomuquh6UFKToHm3Ivc8n5fDCZ4jJKXdZCo7Oz77KZXh0UShZFBI+dAJ2TXXZT6kVNptpS1FLQbXZIvxzhiLjoEahoyTjtFAfgLsxKsAKgTkl5F9NXbHwaWP74J+PJ/SPmKr/t+X+atTDytSwxCbkU9EpRa++malJCVACxKb33sPmffhFuIyuskgkK93Z3d9NDh8KAK6jI+dZqRkx5x1y1VaCFKJCOmuwF+OcOJx5QoBjOfEUueASE5DACoxkp9IOmrtj4NLH98X/H0/pHzFV/2/L/MK1jeV6TpbQumtrJZClzLmwcFtim997D5+l8Z34hM0TSCXG/L9vHv9Kb9gh/KVGeX1rKqyU+q4VV2yL3t0l/840Px9P6R8xSv4BL/ADCkTkp5PFXbt6Bpdv6448eT+mfMVw/0/KDnUK1buV6UUOoRTW2z0dSJlyQXDfbTe+1z8vhjP/EJhEr9KTvy+PL601+Hwt7mkZ+96yhyS+oDVV2z8Wln++ND8fT+kfMUp+AS/wAwpzGTHWnG71VsoSsEo6S7EX374q/HlKnEZz4irLwCQHOoVppuVqYY81DFObirDX3M0kqClKBvZN9rX/e2EhxGWNY5DITuMjt5ZFMnh0LgqqjJz8KyyskvqtqqzRt6tLNvrh/8fT+kfMUr/t+X+cVPCya4zKYU5U23G0rSS30lWUL7jnvgc3HVaNgsZBxzyKtHwF1YMxBHXR+o5TpngZjUSCiC4myGJSlKVrG1yBe44+vxuuOIyxFHaQsOsdv+KMeGxSqVRQCflWZVkl5VtVVaVbi7Szb64c/H0/pHzFLf7fl/mFW6Vk4tVCMZNRaeYCxraLa/MPTnArjjoaMhUIPbkbVePgTxtrYggcxRiq5VgLp76IcNFNf6gQ26SV60C3mFjte373wEcSkhkRmcuOz4c/CjNw2OVSsYGf07qzhyQ8ojVVWjbi7Sz/fDf4+n9I+Ypb/b8v8AMK//2Q=="

import base64 as _b64, io as _io
try:
    from PIL import Image as _PILImage
    _PAGE_ICON = _PILImage.open(_io.BytesIO(_b64.b64decode(HDM_LOGO_B64)))
except Exception:
    _PAGE_ICON = None

st.set_page_config(page_title="Hospital Dashboard", layout="wide",
                   page_icon=_PAGE_ICON if _PAGE_ICON is not None else None)

# ── Theme (light default; runtime toggle lives in the sidebar) ──
DARK = bool(st.session_state.get("ui_dark", False))
if DARK:
    APP_BG, SIDE_BG, SIDE_BORDER = "#0E1A1A", "#0B1515", "#1E3636"
    BODY_FG, MUTED_FG, SECT_BORDER = "#E6F2F2", "#9FBABA", "#1B3433"
    CARD_BG, METRIC_BG, METRIC_LABEL = "#13201F", "#13201F", "#9FBABA"
    SIDE_OVERLAY = "rgba(11,21,21,0.66)"   # dark tint over the sidebar pattern
else:
    # sidebar uses a light red tint in light mode; charts/cards keep LIGHT_BG
    APP_BG, SIDE_BG, SIDE_BORDER = "#FFFFFF", "#FBE2E2", GRID
    BODY_FG, MUTED_FG, SECT_BORDER = "#0E2A2A", "#5E7373", LIGHT_BG
    CARD_BG, METRIC_BG, METRIC_LABEL = "#FFFFFF", LIGHT_BG, "#5E7373"
    SIDE_OVERLAY = "rgba(251,226,226,0.58)"   # light tint over the sidebar pattern

_DARK_OVERRIDES = (f"""
    .stApp, .stApp p, .stApp li, .stApp label,
    [data-testid="stWidgetLabel"], [data-testid="stWidgetLabel"] p,
    [data-testid="stMarkdownContainer"], [data-testid="stMarkdownContainer"] p,
    [data-testid="stCaptionContainer"], .stRadio label, .stCheckbox label,
    [data-testid="stExpander"] summary, [data-testid="stExpander"] p {{
        color:{BODY_FG}; }}
    [data-testid="stSidebar"] * {{ color:{BODY_FG}; }}
    /* keep card internals readable on their light/dark surfaces */
    .chart-card-title {{ color:{PRIMARY} !important; }}
    [data-testid="stMetricValue"] {{ color:{TEAL2} !important; }}
    [data-testid="stMetricLabel"], [data-testid="stMetricLabel"] * {{ color:{METRIC_LABEL} !important; }}

    /* dark-mode hover fixes: never flash a light background under light text */
    .stButton>button:hover:not(:disabled) {{ background:{METRIC_BG} !important;
        color:{TEAL2} !important; border-color:{TEAL2} !important; }}
    [data-testid="stFileUploaderDropzone"] {{ background:{CARD_BG} !important; }}
    /* selectbox / multiselect drop-down menu popover */
    ul[data-baseweb="menu"], [data-baseweb="popover"] ul,
    [data-baseweb="popover"] [role="listbox"] {{ background:{CARD_BG} !important; }}
    ul[data-baseweb="menu"] li, [data-baseweb="popover"] [role="option"] {{
        background:{CARD_BG} !important; color:{BODY_FG} !important; }}
    ul[data-baseweb="menu"] li:hover, [data-baseweb="popover"] [role="option"]:hover,
    [data-baseweb="popover"] [role="option"][aria-selected="true"] {{
        background:{PRIMARY} !important; color:#FFFFFF !important; }}
    /* date-picker calendar popover */
    [data-baseweb="calendar"] {{ background:{CARD_BG} !important; color:{BODY_FG} !important; }}
    /* inputs (incl. date pickers) and sidebar buttons: dark surface, light text */
    .stTextInput input, .stNumberInput input, .stDateInput input,
    [data-baseweb="input"], [data-baseweb="input"] input,
    [data-baseweb="base-input"], [data-baseweb="base-input"] input,
    [data-baseweb="select"] > div {{
        background:{CARD_BG} !important; color:{BODY_FG} !important; }}
    [data-testid="stSidebar"] .stButton>button {{
        background:{METRIC_BG} !important; color:{BODY_FG} !important;
        border:1px solid {SIDE_BORDER} !important; }}
""" if DARK else "")

_CS = "dark" if DARK else "light"
st.markdown(f"""
<style>
    :root, html, body, .stApp {{ color-scheme:{_CS} !important; }}
    .stApp {{ background:{APP_BG}; color:{BODY_FG}; }}
    /* trim the large default whitespace above the first element */
    .block-container {{ padding-top:1rem !important; }}
    header[data-testid="stHeader"] {{ background:transparent; height:0; }}
    /* Standalone-software chrome removal: no Deploy button, menu, footer or
       status widget — this runs as a clean desktop app, not a Streamlit page. */
    [data-testid="stToolbar"], [data-testid="stDecoration"],
    [data-testid="stStatusWidget"], [data-testid="stAppDeployButton"],
    .stDeployButton, #MainMenu, footer {{
        display:none !important; visibility:hidden !important; }}
    [data-testid="stSidebar"] {{
        background-image: linear-gradient({SIDE_OVERLAY}, {SIDE_OVERLAY}),
            url("data:image/jpeg;base64,{HDM_SIDEBAR_BG_B64}");
        background-size: auto, cover;
        background-position: center;
        background-repeat: repeat, no-repeat;
        background-attachment: fixed;
        border-right:1px solid {SIDE_BORDER}; }}
    /* sidebar "View" selector: no radio circles — clean items that highlight on hover */
    [data-testid="stSidebar"] div[role="radiogroup"] {{ gap:6px; }}
    [data-testid="stSidebar"] div[role="radiogroup"] > label {{
        display:flex; align-items:center; width:100%; margin:0;
        background:{CARD_BG}; border:1px solid {SIDE_BORDER}; border-radius:10px;
        padding:9px 14px; cursor:pointer; font-weight:600; color:{BODY_FG};
        transition:background .15s ease, color .15s ease, border-color .15s ease; }}
    /* hide the circular radio glyph entirely */
    [data-testid="stSidebar"] div[role="radiogroup"] > label > div:first-child {{
        display:none !important; }}
    /* highlight on hover */
    [data-testid="stSidebar"] div[role="radiogroup"] > label:hover {{
        background:{PRIMARY}; border-color:{PRIMARY}; color:#FFFFFF; }}
    [data-testid="stSidebar"] div[role="radiogroup"] > label:hover * {{ color:#FFFFFF; }}
    /* keep the chosen option filled so it's clear which view is active */
    [data-testid="stSidebar"] div[role="radiogroup"] > label:has(input:checked) {{
        background:{PRIMARY}; border-color:{PRIMARY}; color:#FFFFFF; font-weight:700; }}
    [data-testid="stSidebar"] div[role="radiogroup"] > label:has(input:checked) * {{
        color:#FFFFFF; }}
    .big-title {{ font-size:clamp(1.6rem,4.5vw,2.4rem); font-weight:800;
        background:linear-gradient(90deg,{PRIMARY},{TEAL2});
        -webkit-background-clip:text; -webkit-text-fill-color:transparent;
        background-clip:text; margin-bottom:.1rem; }}
    .sub {{ color:{MUTED_FG}; font-size:.95rem; margin-bottom:1rem; }}
    .section {{ font-size:1.15rem; font-weight:700; color:{PRIMARY};
        border-left:4px solid {TEAL2}; padding:.15rem 0 .15rem .6rem;
        margin:1.4rem 0 .8rem 0; }}
    [data-testid="stMetric"] {{ background:{METRIC_BG}; border:1px solid {SIDE_BORDER};
        border-radius:12px; padding:14px 16px; }}
    [data-testid="stMetricValue"] {{ color:{PRIMARY}; }}
    .stButton>button[kind="primary"], .stDownloadButton>button {{
        background:{PRIMARY}; border:1px solid {PRIMARY}; color:#FFF; font-weight:600; border-radius:8px; }}
    .pill {{ display:inline-block; padding:4px 12px; border-radius:999px;
        color:#FFF; font-weight:700; font-size:.85rem; margin:2px; }}
    [data-testid="stPlotlyChart"], [data-testid="stDataFrame"] {{ max-width:100% !important; }}
    /* unified chart cards: title sits in a header inside the same bordered card as
       the chart, so a title can never look detached or stacked over the content */
    [data-testid="stVerticalBlockBorderWrapper"] {{ border-radius:16px !important;
        border-color:{SIDE_BORDER} !important; box-shadow:0 2px 10px rgba(6,52,58,.05);
        background:{CARD_BG}; padding:6px 14px 4px 14px !important; margin-bottom:.7rem; }}
    .chart-card-title {{ font-size:1.06rem; font-weight:700; color:{PRIMARY};
        margin:.1rem 0 .15rem 0; line-height:1.2; }}
    .chart-card-sub {{ font-size:.82rem; color:{MUTED_FG}; margin:-.05rem 0 .2rem 0; }}

    /* ── Expanders styled like the sidebar view selector: rounded pill headers
       that fill with the brand teal on hover and when open ── */
    [data-testid="stExpander"] {{ max-width:46rem; border:none !important;
        background:transparent !important; box-shadow:none !important;
        overflow:visible; margin:.28rem 0 .5rem 0; }}
    [data-testid="stExpander"] details {{ border:none !important;
        background:transparent !important; }}
    [data-testid="stExpander"] summary {{ background:{CARD_BG};
        border:1px solid {DANGER}; border-radius:10px;
        padding:9px 14px !important; font-weight:600; font-size:.9rem;
        color:{BODY_FG}; min-height:0;
        transition:background .15s ease, color .15s ease, border-color .15s ease; }}
    [data-testid="stExpander"] summary:hover {{ background:{PRIMARY} !important;
        border-color:{DANGER} !important; color:#FFFFFF !important; }}
    [data-testid="stExpander"] summary:hover * {{ color:#FFFFFF !important; }}
    [data-testid="stExpander"] details[open] > summary {{ background:{PRIMARY};
        border-color:{DANGER}; color:#FFFFFF; font-weight:700;
        border-radius:10px 10px 0 0; }}
    [data-testid="stExpander"] details[open] > summary * {{ color:#FFFFFF; }}
    [data-testid="stExpander"] summary:hover svg,
    [data-testid="stExpander"] details[open] > summary svg {{
        fill:#FFFFFF !important; stroke:#FFFFFF; }}
    [data-testid="stExpander"] summary svg {{ opacity:.75; }}
    [data-testid="stExpander"] [data-testid="stExpanderDetails"] {{
        border:1px solid {DANGER}; border-top:none;
        border-radius:0 0 10px 10px; background:{CARD_BG};
        padding:8px 14px 12px 14px; }}

    /* Selectboxes, number/text inputs, date inputs: softer, rounder, teal focus */
    [data-baseweb="select"] > div, .stTextInput input, .stNumberInput input,
    .stDateInput input {{ border-radius:10px !important;
        border-color:{SIDE_BORDER} !important; transition:border-color .15s ease,
        box-shadow .15s ease; }}
    [data-baseweb="select"] > div:hover, .stTextInput input:hover,
    .stNumberInput input:hover, .stDateInput input:hover {{
        border-color:{TEAL2} !important; }}
    [data-baseweb="select"] > div:focus-within, .stTextInput input:focus,
    .stNumberInput input:focus, .stDateInput input:focus {{
        border-color:{PRIMARY} !important;
        box-shadow:0 0 0 2px rgba(0,104,104,.15) !important; }}
    [data-testid="stWidgetLabel"] p {{ font-weight:600; color:{BODY_FG};
        font-size:.86rem; }}

    /* Multiselect chips in the brand teal */
    [data-baseweb="tag"] {{ background:{PRIMARY} !important; border-radius:8px !important;
        color:#fff !important; }}

    /* Horizontal radios (roll-up, import mode) as pill toggles */
    div[role="radiogroup"][aria-orientation="horizontal"] {{ gap:8px; }}
    div[role="radiogroup"][aria-orientation="horizontal"] > label {{
        border:1px solid {SIDE_BORDER}; border-radius:999px; padding:5px 14px;
        background:{CARD_BG}; transition:all .15s ease; }}
    div[role="radiogroup"][aria-orientation="horizontal"] > label:hover {{
        border-color:{TEAL2}; }}
    div[role="radiogroup"][aria-orientation="horizontal"] > label:has(input:checked) {{
        background:{PRIMARY}; border-color:{PRIMARY}; }}
    div[role="radiogroup"][aria-orientation="horizontal"] > label:has(input:checked) * {{
        color:#fff !important; }}

    /* File uploader: dashed brand-tinted drop zone */
    [data-testid="stFileUploaderDropzone"] {{ border:1.5px dashed {TEAL2} !important;
        border-radius:12px !important; background:{LIGHT_BG} !important; }}

    /* Secondary buttons (move ▲▼⤒, reset, add table): crisp outline that fills on hover */
    .stButton>button {{ border-radius:9px; border:1px solid {SIDE_BORDER};
        font-weight:600; transition:all .13s ease; }}
    .stButton>button:hover:not(:disabled) {{ border-color:{PRIMARY};
        color:{PRIMARY}; background:{LIGHT_BG}; }}
    .stButton>button:disabled {{ opacity:.4; }}
    /* Primary + download buttons: subtle lift on hover */
    .stButton>button[kind="primary"]:hover, .stDownloadButton>button:hover {{
        filter:brightness(1.07); box-shadow:0 3px 10px rgba(0,104,104,.25); }}
    .stDownloadButton>button {{ border-radius:9px; }}

    {_DARK_OVERRIDES}
    /* phones & tablets: wrap columns two-up so cards, inputs and buttons stay
       readable. Charts render full-width (outside columns), so they're unaffected. */
    @media (max-width: 980px) {{
        [data-testid="stHorizontalBlock"] {{ flex-wrap:wrap !important; gap:0.55rem !important; }}
        [data-testid="stColumn"], [data-testid="column"] {{
            min-width:47% !important; flex:1 1 47% !important; }}
    }}
    /* small phones: a touch more compact */
    @media (max-width: 460px) {{
        .block-container {{ padding-left:0.6rem !important; padding-right:0.6rem !important; }}
        [data-testid="stMetric"] {{ padding:10px 12px; }}
        [data-testid="stMetricValue"] {{ font-size:1.15rem !important; }}
    }}
</style>
""", unsafe_allow_html=True)


# ──────────────────────────────────────────────
# DATABASE
# ──────────────────────────────────────────────
@st.cache_resource
def get_conn():
    conn = sqlite3.connect(DB_PATH, check_same_thread=False)
    conn.execute("PRAGMA journal_mode=WAL;")
    return conn


def init_db():
    c = get_conn()
    cols = ", ".join(f"{k} INTEGER DEFAULT 0" for k in FIELD_KEYS)
    c.executescript(f"""
    CREATE TABLE IF NOT EXISTS daily (
        entry_date TEXT PRIMARY KEY, {cols}, notes TEXT, updated_at TEXT);
    CREATE TABLE IF NOT EXISTS departments (
        entry_date TEXT, name TEXT, status TEXT, PRIMARY KEY (entry_date, name));
    CREATE TABLE IF NOT EXISTS medications (
        entry_date TEXT, name TEXT, stock INTEGER, unit TEXT, status TEXT,
        PRIMARY KEY (entry_date, name));
    CREATE TABLE IF NOT EXISTS tests (
        entry_date TEXT, name TEXT, available INTEGER, PRIMARY KEY (entry_date, name));
    CREATE TABLE IF NOT EXISTS blood_bank (
        entry_date TEXT, blood_type TEXT, units INTEGER, PRIMARY KEY (entry_date, blood_type));
    CREATE TABLE IF NOT EXISTS absent_specialists (
        entry_date TEXT, name TEXT, specialty TEXT, expected_return TEXT,
        PRIMARY KEY (entry_date, name));
    CREATE TABLE IF NOT EXISTS mortality (
        entry_date TEXT, ward TEXT, age INTEGER, los_days INTEGER,
        condition TEXT, cause TEXT, time_of_death TEXT);
    CREATE TABLE IF NOT EXISTS custom_table_rows (
        entry_date TEXT, table_id TEXT, payload TEXT,
        PRIMARY KEY (entry_date, table_id));
    CREATE TABLE IF NOT EXISTS users (
        email TEXT PRIMARY KEY, password_hash TEXT NOT NULL,
        created_at TEXT, last_login TEXT);
    CREATE TABLE IF NOT EXISTS observations (
        scope TEXT, ref TEXT, content TEXT, updated_at TEXT, updated_by TEXT,
        fmt TEXT, PRIMARY KEY (scope, ref));
    CREATE TABLE IF NOT EXISTS daily_notes (
        entry_date TEXT, note_id TEXT, title TEXT, content TEXT, updated_at TEXT,
        PRIMARY KEY (entry_date, note_id));
    CREATE TABLE IF NOT EXISTS settings (key TEXT PRIMARY KEY, value TEXT);
    """)
    c.commit()
    migrate()


def migrate():
    c = get_conn()
    existing = {r[1] for r in c.execute("PRAGMA table_info(daily)").fetchall()}
    for k in FIELD_KEYS:
        if k not in existing:
            c.execute(f"ALTER TABLE daily ADD COLUMN {k} INTEGER DEFAULT 0")
    med_cols = {r[1] for r in c.execute("PRAGMA table_info(medications)").fetchall()}
    if "status" not in med_cols:
        c.execute("ALTER TABLE medications ADD COLUMN status TEXT")
    mort_cols = {r[1] for r in c.execute("PRAGMA table_info(mortality)").fetchall()}
    if "condition" not in mort_cols:
        c.execute("ALTER TABLE mortality ADD COLUMN condition TEXT")
    if "cause" not in mort_cols:
        c.execute("ALTER TABLE mortality ADD COLUMN cause TEXT")
    if "time_of_death" not in mort_cols:
        c.execute("ALTER TABLE mortality ADD COLUMN time_of_death TEXT")
    obs_cols = {r[1] for r in c.execute("PRAGMA table_info(observations)").fetchall()}
    if "fmt" not in obs_cols:
        c.execute("ALTER TABLE observations ADD COLUMN fmt TEXT")
    c.commit()


def get_setting(key, default=""):
    row = get_conn().execute("SELECT value FROM settings WHERE key=?", (key,)).fetchone()
    return row[0] if row else default


def set_setting(key, value):
    c = get_conn()
    c.execute("INSERT INTO settings(key,value) VALUES(?,?) "
              "ON CONFLICT(key) DO UPDATE SET value=excluded.value", (key, str(value)))
    c.commit()


def get_dashboard_title():
    """The full heading shown on the dashboard. Editable in Settings; defaults to
    '<hospital name> — Dashboard' until customised."""
    t = (get_setting("dashboard_title", "") or "").strip()
    return t or f"{HOSPITAL_NAME} — Dashboard"


# ── Dashboard notes & observations (per day / week / month) ──
# Note formatting is block-level: a whole section note carries one alignment,
# font, size, colour, style and (optionally) one image. The same settings drive
# the dashboard, the TV view and every report (PDF and PPTX, daily and period).
NOTE_FONTS_CSS = {"Sans": "'Segoe UI', system-ui, sans-serif",
                  "Serif": "Georgia, 'Times New Roman', serif",
                  "Mono": "'Courier New', ui-monospace, monospace"}
NOTE_FONTS_PDF = {"Sans": "Helvetica", "Serif": "Times-Roman", "Mono": "Courier"}
NOTE_FONTS_PPTX = {"Sans": "Calibri", "Serif": "Georgia", "Mono": "Consolas"}
NOTE_ALIGN_PDF = {"left": 0, "center": 1, "right": 2}


def default_note_fmt():
    return {"align": "left", "bold": False, "italic": False, "underline": False,
            "color": "", "font": "Sans", "size": 14, "image": ""}


def parse_note_fmt(raw):
    """Turn a stored fmt value (JSON string or dict) into a safe, complete dict."""
    d = default_note_fmt()
    src = raw
    if isinstance(raw, str) and raw:
        try:
            src = json.loads(raw)
        except Exception:
            src = None
    if isinstance(src, dict):
        for k in d:
            if k in src and src[k] is not None:
                d[k] = src[k]
    if d["align"] not in ("left", "center", "right"):
        d["align"] = "left"
    if d["font"] not in NOTE_FONTS_CSS:
        d["font"] = "Sans"
    try:
        d["size"] = max(8, min(48, int(d["size"])))
    except Exception:
        d["size"] = 14
    for b in ("bold", "italic", "underline"):
        d[b] = bool(d[b])
    d["color"] = d["color"] if isinstance(d["color"], str) else ""
    d["image"] = d["image"] if isinstance(d["image"], str) else ""
    return d


def note_block_html(text, fmt, default_color):
    """Render a note (text + optional image) as themed HTML for the dashboard."""
    import html as _hh
    css = (f"text-align:{fmt['align']};"
           f"font-family:{NOTE_FONTS_CSS.get(fmt['font'], NOTE_FONTS_CSS['Sans'])};"
           f"font-size:{fmt['size']}px;color:{fmt['color'] or default_color};")
    if fmt["bold"]:
        css += "font-weight:700;"
    if fmt["italic"]:
        css += "font-style:italic;"
    if fmt["underline"]:
        css += "text-decoration:underline;"
    body = _hh.escape(text).replace("\n", "<br/>") if text else ""
    out = f'<div style="white-space:pre-wrap;{css}">{body}</div>'
    if fmt.get("image"):
        just = {"left": "flex-start", "center": "center",
                "right": "flex-end"}.get(fmt["align"], "flex-start")
        out += (f'<div style="display:flex;justify-content:{just};margin-top:8px;">'
                f'<img src="{fmt["image"]}" style="max-width:100%;max-height:340px;'
                f'border-radius:8px;"/></div>')
    return out


def _note_pdf_image(data_uri):
    """A reportlab Image flowable from a note's stored data-URI, or None."""
    try:
        import io as _iio
        from reportlab.platypus import Image as _RLImage
        from reportlab.lib.units import mm as _mmu
        from reportlab.lib.utils import ImageReader as _IR
        b = data_uri.split(",", 1)[1] if "," in data_uri else data_uri
        raw = _iio.BytesIO(base64.b64decode(b))
        iw, ih = _IR(raw).getSize()
        w = 120 * _mmu
        h = w * ih / iw
        if h > 80 * _mmu:
            h = 80 * _mmu
            w = h * iw / ih
        raw.seek(0)
        return _RLImage(raw, width=w, height=h)
    except Exception:
        return None


def get_observation(scope, ref):
    """The saved observation for one dashboard: scope is 'day' / 'week' /
    'month', ref identifies the period (ISO date, week-Monday, or YYYY-MM)."""
    row = get_conn().execute(
        "SELECT content, updated_at, updated_by, fmt FROM observations "
        "WHERE scope=? AND ref=?", (scope, str(ref))).fetchone()
    if not row:
        return {"content": "", "updated_at": "", "updated_by": "", "fmt": ""}
    return {"content": row[0] or "", "updated_at": row[1] or "",
            "updated_by": row[2] or "", "fmt": row[3] or ""}


def set_observation(scope, ref, content, fmt=None):
    c = get_conn()
    c.execute(
        "INSERT INTO observations(scope, ref, content, updated_at, updated_by, fmt) "
        "VALUES(?,?,?,?,?,?) ON CONFLICT(scope, ref) DO UPDATE SET "
        "content=excluded.content, updated_at=excluded.updated_at, "
        "updated_by=excluded.updated_by, fmt=excluded.fmt",
        (scope, str(ref), clean_text(content, multiline=True),
         now_local().isoformat(timespec="seconds"),
         str(st.session_state.get("user_email") or "analyst"),
         json.dumps(fmt) if fmt else ""))
    c.commit()


def get_daily_notes(day):
    """Custom titled notes recorded for one day (entered in Data Entry, like
    custom tables). Returns a list of {id, title, content}."""
    rows = get_conn().execute(
        "SELECT note_id, title, content FROM daily_notes WHERE entry_date=? "
        "ORDER BY rowid", (day.isoformat(),)).fetchall()
    return [{"id": r[0], "title": r[1] or "", "content": r[2] or ""} for r in rows]


def save_daily_notes(day, notes):
    """Replace all custom notes for a day. `notes` is a list of dicts with
    'id', 'title', 'content'; blank ones (no title and no content) are dropped."""
    ds = day.isoformat()
    c = get_conn()
    with c:
        c.execute("DELETE FROM daily_notes WHERE entry_date=?", (ds,))
        for n in (notes or []):
            title = clean_text(n.get("title", ""))
            content = clean_text(n.get("content", ""), multiline=True)
            if not title and not content:
                continue
            c.execute("INSERT INTO daily_notes(entry_date, note_id, title, "
                      "content, updated_at) VALUES(?,?,?,?,?)",
                      (ds, n.get("id") or pysecrets.token_hex(4), title, content,
                       now_local().isoformat(timespec="seconds")))


def get_daily_notes_range(start, end):
    """All custom notes across a date range, each tagged with its date — for the
    weekly / monthly reports."""
    rows = get_conn().execute(
        "SELECT entry_date, title, content FROM daily_notes "
        "WHERE entry_date BETWEEN ? AND ? ORDER BY entry_date, rowid",
        (start.isoformat(), end.isoformat())).fetchall()
    return [{"date": r[0], "title": r[1] or "", "content": r[2] or ""}
            for r in rows]


def _note_lookup(scope, period_ref, section_title):
    """Find the observation for a report section. A report heading often bundles
    one or more dashboard sections under a different name (e.g. the PDF's
    'Patients' is the dashboard's 'Patient Activity'), so match the heading to
    its dashboard title(s) via REPORT_SECTION_ALIASES and return the first note
    found — this is why a note written on the dashboard shows up in the report."""
    seen, cands = set(), []
    for cand in [section_title] + list(REPORT_SECTION_ALIASES.get(section_title, [])):
        if cand and cand not in seen:
            seen.add(cand)
            cands.append(cand)
    for cand in cands:
        obs = get_observation(scope, f"{period_ref}::{cand}")
        if obs["content"].strip() or parse_note_fmt(obs.get("fmt")).get("image"):
            return obs
    return get_observation(scope, f"{period_ref}::{section_title}")


def section_note(scope, period_ref, section_title):
    """The note attached to one specific chart/table section (or '')."""
    if not section_title:
        return ""
    return _note_lookup(scope, period_ref, section_title)["content"].strip()


def section_note_full(scope, period_ref, section_title):
    """(text, fmt_dict) for a section note — used by the report builders so the
    note's formatting and image travel into every PDF and PPTX."""
    if not section_title:
        return "", default_note_fmt()
    obs = _note_lookup(scope, period_ref, section_title)
    return obs["content"].strip(), parse_note_fmt(obs.get("fmt"))


def gather_general_notes(scope, period_ref, start=None, end=None):
    """Notes that are NOT tied to a specific chart/table: the period overview
    note and the custom titled notes. Per-section notes are placed under their
    own chart/table by the report builders, so they're excluded here."""
    out = []
    base = get_observation(scope, period_ref)["content"].strip()
    if base:
        out.append(("Overview", base))
    if scope == "day":
        for n in get_daily_notes(_date_from_ref(period_ref)):
            if n["content"].strip() or n["title"].strip():
                out.append((n["title"] or "Note", n["content"].strip()))
    elif start is not None and end is not None:
        for n in get_daily_notes_range(start, end):
            if n["content"].strip() or n["title"].strip():
                out.append(((n["title"] or "Note") + f" — {n['date']}",
                            n["content"].strip()))
    return out


def gather_report_notes(scope, period_ref, start=None, end=None):
    """Every note for a report as (label, text): the period overview note, each
    per-section chart/table note, and the custom titled notes. Retained for
    back-compat / CSV; the PDF & PPTX now place per-section notes under their
    own chart/table via section_note()."""
    out = []
    base = get_observation(scope, period_ref)["content"].strip()
    if base:
        out.append(("Overview", base))
    prefix = f"{period_ref}::"
    rows = get_conn().execute(
        "SELECT ref, content FROM observations WHERE scope=? AND ref LIKE ? "
        "ORDER BY ref", (scope, prefix + "%")).fetchall()
    for ref, content in rows:
        if content and content.strip():
            out.append((ref.split("::", 1)[1], content.strip()))
    if scope == "day":
        for n in get_daily_notes(_date_from_ref(period_ref)):
            if n["content"].strip() or n["title"].strip():
                out.append((n["title"] or "Note", n["content"].strip()))
    elif start is not None and end is not None:
        for n in get_daily_notes_range(start, end):
            if n["content"].strip() or n["title"].strip():
                out.append(((n["title"] or "Note") + f" — {n['date']}",
                            n["content"].strip()))
    return out


def _date_from_ref(ref):
    """Parse a 'YYYY-MM-DD' day ref back into a date (for custom-note lookup)."""
    from datetime import date as _d
    try:
        y, m, dd = (int(x) for x in ref.split("-")[:3])
        return _d(y, m, dd)
    except Exception:
        return now_local().date()


def render_observations(scope, ref):
    """A simple '📝 Notes' panel under each dashboard. Collapsed by default —
    click the word to open it. Everyone sees the saved note; a signed-in
    analyst edits it in place. The saved text is public and goes into the PDF."""
    import html as _h
    obs = get_observation(scope, ref)
    has = bool(obs["content"].strip())
    stamp = ""
    if obs["updated_at"]:
        stamp = "Last updated " + obs["updated_at"].replace("T", " at ")
        if obs["updated_by"]:
            stamp += f" by {obs['updated_by']}"
    label = "📝 Notes"
    if has:
        preview = " ".join(obs["content"].split())
        label += " — " + (preview[:48] + "…" if len(preview) > 48 else preview)
    with st.expander(label, expanded=False):
        if st.session_state.get("authed"):
            txt = st.text_area(
                "Write notes for this dashboard (shown publicly and included in "
                "the PDF report):",
                value=obs["content"], height=150, key=f"obs_{scope}_{ref}",
                placeholder="Anything worth remembering about this period — "
                            "incidents, context behind the numbers, follow-ups…")
            bcols = st.columns([1, 2])
            if bcols[0].button("💾 Save notes", key=f"obs_save_{scope}_{ref}",
                               use_container_width=True):
                set_observation(scope, ref, txt)
                st.session_state.auth_time = time.time()
                st.success("Notes saved.")
                st.rerun()
            if stamp:
                bcols[1].caption(stamp)
        else:
            if has:
                st.markdown(
                    f'<div style="white-space:pre-wrap;color:{BODY_FG};">'
                    f'{_h.escape(obs["content"])}</div>', unsafe_allow_html=True)
                if stamp:
                    st.caption(stamp)
            else:
                st.caption("No notes recorded for this period yet.")
            st.caption("🔒 Sign in under **Data Entry (Analyst)** to add notes.")


# ── Full backup & restore (protection against Streamlit Cloud disk resets) ──
def export_backup_bytes():
    """One JSON file containing every row of every table — a complete snapshot
    of the dashboard (daily entries, custom tables + their chart configs,
    observations, wards, users, settings). Safe to download and re-import."""
    c = get_conn()
    names = [r[0] for r in c.execute(
        "SELECT name FROM sqlite_master WHERE type='table' "
        "AND name NOT LIKE 'sqlite_%'").fetchall()]
    dump = {"format": "hdm-backup", "version": 1,
            "created": now_local().isoformat(timespec="seconds"), "tables": {}}
    for n in names:
        cur = c.execute(f'SELECT * FROM "{n}"')
        cols = [d[0] for d in cur.description]
        dump["tables"][n] = {"columns": cols,
                             "rows": [list(r) for r in cur.fetchall()]}
    return json.dumps(dump, ensure_ascii=False).encode("utf-8")


def restore_backup_bytes(data):
    """Load a backup produced by export_backup_bytes, replacing all current
    data. Unknown tables/columns in the file are skipped, so older backups
    keep working after upgrades. Returns (tables_restored, rows_restored);
    raises ValueError for files that are not Hospital Dashboard backups."""
    try:
        dump = json.loads(bytes(data).decode("utf-8"))
    except Exception:
        raise ValueError("That file could not be read as a backup (not valid JSON).")
    if not isinstance(dump, dict) or dump.get("format") != "hdm-backup":
        raise ValueError("That file is not a Hospital Dashboard backup.")
    c = get_conn()
    existing = {r[0] for r in c.execute(
        "SELECT name FROM sqlite_master WHERE type='table' "
        "AND name NOT LIKE 'sqlite_%'").fetchall()}
    n_tables = n_rows = 0
    try:
        for name, t in (dump.get("tables") or {}).items():
            if name not in existing or not isinstance(t, dict):
                continue
            cols = list(t.get("columns") or [])
            valid = {r[1] for r in c.execute(f'PRAGMA table_info("{name}")')}
            keep = [i for i, col in enumerate(cols) if col in valid]
            if not keep:
                continue
            c.execute(f'DELETE FROM "{name}"')
            kcols = [cols[i] for i in keep]
            rows = [[r[i] for i in keep] for r in (t.get("rows") or [])]
            if rows:
                c.executemany(
                    'INSERT INTO "{}" ({}) VALUES ({})'.format(
                        name, ",".join(f'"{k}"' for k in kcols),
                        ",".join("?" * len(kcols))), rows)
            n_tables += 1
            n_rows += len(rows)
        c.commit()
    except Exception:
        c.rollback()
        raise
    return n_tables, n_rows


# ── PowerPoint report decks ─────────────────────────────────────
PPTX_MIME = ("application/vnd.openxmlformats-officedocument."
             "presentationml.presentation")
try:
    import pptx as _pptx_probe                    # noqa: F401
    PPTX_OK = True
except Exception:
    PPTX_OK = False


def _pptx_toolkit():
    """Late imports plus shared slide helpers for the PPTX report builders
    (brand title slide, content headings, centred pictures, styled tables)."""
    import io as _io
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    BRAND = RGBColor.from_string(PRIMARY.lstrip("#"))
    INK_C = RGBColor.from_string(INK.lstrip("#"))
    WHITE = RGBColor.from_string("FFFFFF")
    SOFT = RGBColor.from_string("BFE3E3")

    def new_deck():
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
        return prs

    def blank(prs):
        return prs.slides.add_slide(prs.slide_layouts[6])

    def textbox(slide, l, t, w, h, text, size, color, bold=False,
                align=PP_ALIGN.LEFT, italic=False, underline=False,
                font="Calibri"):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        first = True
        for line in str(text).split("\n"):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            r = p.add_run()
            r.text = line
            r.font.name = font
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.underline = underline
            r.font.color.rgb = color
        return tb

    def note_box(slide, left, top, width, note_res):
        """Render a section note (text + optional image) with its saved
        formatting: alignment, font, size, colour, bold/italic/underline."""
        if not note_res:
            return
        if isinstance(note_res, tuple):
            txt, nf = note_res
        else:
            txt, nf = note_res, None
        nf = nf or {"align": "left", "bold": False, "italic": False,
                    "underline": False, "color": "", "font": "Sans",
                    "size": 14, "image": ""}
        _al = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
               "right": PP_ALIGN.RIGHT}.get(nf.get("align"), PP_ALIGN.LEFT)
        _col = RGBColor.from_string((nf.get("color") or INK).lstrip("#"))
        if txt:
            textbox(slide, left, top, width, Inches(0.7), "📝 " + txt,
                    max(9, min(28, int(nf.get("size", 14) * 0.8))), _col,
                    bold=nf.get("bold", False), align=_al,
                    italic=nf.get("italic", False),
                    underline=nf.get("underline", False),
                    font=NOTE_FONTS_PPTX.get(nf.get("font"), "Calibri"))
        if nf.get("image"):
            try:
                _b = (nf["image"].split(",", 1)[1] if "," in nf["image"]
                      else nf["image"])
                _bio = _io.BytesIO(base64.b64decode(_b))
                slide.shapes.add_picture(_bio, left, top + Inches(0.5),
                                         height=Inches(0.85))
            except Exception:
                pass

    def title_slide(prs, heading, sub, stamp):
        s = blank(prs)
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BRAND
        bg.line.fill.background()
        bg.shadow.inherit = False
        try:
            logo = _io.BytesIO(base64.b64decode(HDM_LOGO_B64))
            s.shapes.add_picture(logo, Inches(5.82), Inches(0.75),
                                 height=Inches(1.9))
        except Exception:
            pass
        textbox(s, Inches(0.8), Inches(3.05), Inches(11.73), Inches(1.0),
                heading, 40, WHITE, bold=True, align=PP_ALIGN.CENTER)
        textbox(s, Inches(0.8), Inches(4.2), Inches(11.73), Inches(0.7),
                sub, 20, SOFT, align=PP_ALIGN.CENTER)
        textbox(s, Inches(0.8), Inches(6.75), Inches(11.73), Inches(0.4),
                stamp, 12, SOFT, align=PP_ALIGN.CENTER)
        return s

    def content_slide(prs, heading):
        s = blank(prs)
        textbox(s, Inches(0.6), Inches(0.32), Inches(12.13), Inches(0.75),
                heading, 30, BRAND, bold=True)
        return s

    def picture(slide, bio_aspect, top=Inches(1.35), max_w=Inches(12.1),
                max_h=Inches(5.7)):
        bio, aspect = bio_aspect                  # aspect = height / width
        w_in = max_w.inches
        h_in = w_in * float(aspect)
        if h_in > max_h.inches:
            h_in = max_h.inches
            w_in = h_in / float(aspect)
        bio.seek(0)
        slide.shapes.add_picture(bio, Inches((13.333 - w_in) / 2), top,
                                 width=Inches(w_in), height=Inches(h_in))

    def add_table(slide, headers, rows, left, top, width, header_fills=None,
                  font=11, max_rows=12):
        rows = [[(str(v).strip() or "—") for v in row] for row in rows]
        clipped = max(0, len(rows) - max_rows)
        rows = rows[:max_rows]
        shp = slide.shapes.add_table(len(rows) + 1, len(headers), left, top,
                                     width, Inches(0.36 * (len(rows) + 1)))
        tbl = shp.table
        for j, htxt in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.text = str(htxt)
            cell.fill.solid()
            fill = None
            if header_fills and j < len(header_fills) and header_fills[j]:
                try:
                    fill = RGBColor.from_string(str(header_fills[j]).lstrip("#"))
                except Exception:
                    fill = None
            cell.fill.fore_color.rgb = fill or BRAND
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = "Calibri"
                    r.font.size = Pt(font + 1)
                    r.font.bold = True
                    r.font.color.rgb = WHITE
        for i, row in enumerate(rows, start=1):
            for j, val in enumerate(row):
                cell = tbl.cell(i, j)
                cell.text = val
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.name = "Calibri"
                        r.font.size = Pt(font)
                        r.font.color.rgb = INK_C
        if clipped:
            textbox(slide, left, top + Inches(0.36 * (len(rows) + 1) + 0.06),
                    width, Inches(0.3), f"… and {clipped} more rows in the CSV "
                    "export", 11, RGBColor.from_string("5E7373"))
        return tbl

    def df_rows(df, cols):
        out = []
        for _, r in df.iterrows():
            vals = [str(r.get(c, "") or "").strip() for c in cols]
            if any(vals):
                out.append(vals)
        return out

    def customs_slides(prs, customs, note_fn=None):
        for item in (customs or []):
            ctitle, cdf = item[0], item[1]
            tdef = item[2] if len(item) > 2 else None
            if cdf is None or getattr(cdf, "empty", True):
                continue
            sl = content_slide(prs, str(ctitle))
            top = Inches(1.35)
            dcfg = table_display(tdef) if tdef else {}
            if tdef and dcfg.get("mode") == "kpi":
                kitems = custom_table_kpis(tdef, cdf)
                if kitems:
                    add_table(sl, ["Metric", "Value"],
                              [(k, str(v)) for k, v in kitems],
                              Inches(2.17), top, Inches(9), font=13)
                    note_box(sl, Inches(0.6), Inches(6.6), Inches(12.13),
                             note_fn(str(ctitle)) if note_fn else None)
                    continue
            ch = custom_pdf_chart(tdef, cdf) if tdef else None
            if ch:
                picture(sl, ch, top=top, max_h=Inches(3.35))
                top = Inches(4.9)
            ccols = [c for c in cdf.columns if c != "Date"]
            ccolors = (dcfg.get("colors") or {}) if tdef else {}
            rows = df_rows(cdf, ccols)
            if rows and ccols:
                add_table(sl, ccols, rows, Inches(0.6), top, Inches(12.13),
                          header_fills=[ccolors.get(c) for c in ccols],
                          font=10, max_rows=(5 if ch else 11))
            note_box(sl, Inches(0.6), Inches(6.7), Inches(12.13),
                     note_fn(str(ctitle)) if note_fn else None)

    def finish(prs):
        bio = _io.BytesIO()
        prs.save(bio)
        return bio.getvalue()

    return locals()


@_maybe_cache(show_spinner=False, max_entries=48, ttl=900)
def build_day_pptx(day, scalars, depts, meds, tests, blood, hospital_name,
                   mortality=None, customs=None, observations=None,
                   note_scope=None, note_ref=None, absent=None):
    """The daily report as a branded 16:9 PowerPoint deck: title, At-a-Glance
    tiles, resources, departments & medications, tests & blood bank, mortality,
    the custom tables (with their configured charts and colours), and the
    day's notes & observations."""
    K = _pptx_toolkit()
    from pptx.util import Inches
    s0 = dict(scalars or {})

    def _i(key):
        try:
            return int(float(s0.get(key, 0) or 0))
        except (TypeError, ValueError):
            return 0

    prs = K["new_deck"]()
    K["title_slide"](prs, hospital_name, f"Daily Report — {day:%A, %d %B %Y}",
                     f"Generated {now_local():%d %b %Y %H:%M}")

    beds_t, beds_a = _i("beds_total"), _i("beds_available")
    occ_txt = f"{round((beds_t - beds_a) / beds_t * 100)}%" if beds_t else "—"
    amb = f'{_i("ambulances_available")}/{_i("ambulances_total")}'
    glance_items = [
        ("patient", "Patients in hospital", _i("current_inpatients")),
        ("admit", "New admissions", _i("admitted")),
        ("discharge", "Discharged", _i("discharged")),
        ("er", "ER visits", _i("er_visits")),
        ("surgery", "Surgeries", _i("surgeries")),
        ("birth", "Births", _i("births")),
        ("stillbirth", "Stillbirths", _i("stillbirths")),
        ("death", "Mortality", _i("deaths")),
        ("bed", "Bed occupancy", occ_txt),
        ("doctor", "Doctors on duty", _i("doctors")),
        ("nurse", "Nurses on duty", _i("nurses")),
        ("ambulance", "Ambulances", amb),
        ("oxygen", "Oxygen supply", f'{_i("oxygen_pct")}%'),
    ]
    # Hospital Performance slide first (banner + KPI cards + gauge + breakdown)
    try:
        _perf = health_summary(*perf_inputs_single(s0, depts, meds, tests, blood))
        _pimg = performance_report_image(_perf)
        if _pimg:
            slp = K["content_slide"](prs,
                                     f"Hospital Performance — {_perf['condition']}")
            K["picture"](slp, _pimg)
    except Exception:
        log.exception("PPTX: Performance slide failed")

    try:
        sl = K["content_slide"](prs, "At a Glance")
        K["picture"](sl, build_glance_image(glance_items))
    except Exception:
        log.exception("PPTX: At-a-Glance slide failed")

    res_rows = [
        ("Beds (available / total)", f"{beds_a} / {beds_t}"),
        ("ICU beds available", _i("icu_beds_available")),
        ("ICU patients", _i("icu_patients")),
        ("Doctors on duty", _i("doctors")),
        ("Nurses on duty", _i("nurses")),
        ("Support staff", _i("support_staff")),
        ("Specialists on call", _i("specialists_on_call")),
        ("Ambulances (available / total)", amb),
        ("Ambulance calls", _i("ambulance_calls")),
        ("Average ER wait (minutes)", _i("avg_er_wait_min")),
        ("Oxygen supply", f'{_i("oxygen_pct")}%'),
        ("Referrals out / back", f'{_i("referrals_out")} / {_i("referrals_back")}'),
    ]
    # ── Reorderable band: each content slide is built by a deferred callable so
    # the whole set can be emitted in the analyst's saved dashboard order. Custom
    # tables get one callable each, so they interleave with the built-in slides
    # exactly as they sit on the dashboard. ──
    band = _ReportBand()
    _note_fn = ((lambda tt: section_note_full(note_scope, note_ref, tt))
                if (note_scope and note_ref) else None)

    def _slide_resources():
        sl = K["content_slide"](prs, "Resources & Capacity")
        K["add_table"](sl, ["Metric", "Value"],
                       [(a, str(b)) for a, b in res_rows],
                       Inches(2.17), Inches(1.35), Inches(9), font=13, max_rows=13)
    band.add("Resources & Capacity", _slide_resources)

    def _slide_supplies():
        sl = K["content_slide"](prs, "Critical Supplies")
        K["picture"](sl, oxygen_gauge_image(_i("oxygen_pct")),
                     top=Inches(1.35), max_h=Inches(5.4))
    band.add("Critical Supplies", _slide_supplies)

    have_d = depts is not None and not getattr(depts, "empty", True)
    have_m = meds is not None and not getattr(meds, "empty", True)
    if have_d or have_m:
        def _slide_dept_meds():
            sl = K["content_slide"](prs, "Departments & Medications")
            if have_d:
                K["add_table"](sl, ["Department", "Status"],
                               K["df_rows"](depts, ["Department", "Status"]),
                               Inches(0.6), Inches(1.35), Inches(5.9))
            if have_m:
                K["add_table"](sl, ["Medication", "Status"],
                               K["df_rows"](meds, ["Medication", "Status"]),
                               Inches(6.83), Inches(1.35), Inches(5.9))
        band.add("Departments & Medications", _slide_dept_meds)

    have_t = tests is not None and not getattr(tests, "empty", True)
    have_b = blood is not None and not getattr(blood, "empty", True)
    if have_t or have_b:
        def _slide_tests_blood():
            sl = K["content_slide"](prs, "Medical Tests & Blood Bank")
            if have_t:
                trows = [[str(r.get("Test", "")),
                          "Yes" if bool(r.get("Available")) else "No"]
                         for _, r in tests.iterrows()
                         if str(r.get("Test", "") or "").strip()]
                K["add_table"](sl, ["Test", "Available"], trows,
                               Inches(0.6), Inches(1.35), Inches(5.9))
            if have_b:
                K["add_table"](sl, ["Blood Type", "Units"],
                               K["df_rows"](blood, ["Blood Type", "Units"]),
                               Inches(6.83), Inches(1.35), Inches(5.9))
        band.add("Medical Tests & Blood Bank", _slide_tests_blood)

    if absent is not None and not getattr(absent, "empty", True):
        _arows = [[str(r.get("Specialist", "") or "").strip(),
                   str(r.get("Specialty / Area", "") or "").strip(),
                   str(r.get("Expected return", "") or "").strip() or "—"]
                  for _, r in absent.iterrows()
                  if str(r.get("Specialist", "") or "").strip()]
        if _arows:
            def _slide_absent(_arows=_arows):
                sl = K["content_slide"](prs, "Absent Specialists")
                K["add_table"](sl, ["Specialist", "Specialty / Area",
                                     "Expected return"], _arows,
                               Inches(0.6), Inches(1.35), Inches(12.13))
            band.add("Absent Specialists", _slide_absent)

    if mortality is not None and not getattr(mortality, "empty", True):
        mc = mortality.columns
        cols = [("Ward" if "Ward" in mc else "ward"),
                ("Age" if "Age" in mc else "age"),
                ("Length of stay (days)" if "Length of stay (days)" in mc
                 else "los_days"),
                ("Time of death" if "Time of death" in mc else "time_of_death"),
                ("Condition" if "Condition" in mc else "condition"),
                ("Cause of death" if "Cause of death" in mc else "cause")]
        cols = [c for c in cols if c in mc]
        rows = K["df_rows"](mortality, cols)
        if rows:
            def _slide_mortality(cols=cols, rows=rows):
                sl = K["content_slide"](prs, "Mortality Register")
                K["add_table"](sl, ["Ward", "Age", "LOS (days)", "Time",
                                    "Condition", "Cause"][:len(cols)], rows,
                               Inches(0.6), Inches(1.35), Inches(12.13), font=10)
            band.add("Mortality Register", _slide_mortality)

    # Custom tables: one deferred slide-builder each (reusing customs_slides on a
    # single-item list) so they slot in by their saved dashboard position.
    for _item in (customs or []):
        band.add(_item[0],
                 (lambda it=_item: K["customs_slides"](prs, [it], note_fn=None)))

    # Emit the whole band in saved dashboard order.
    for _bt, _fn in band.items():
        _fn()
        if _note_fn and prs.slides:
            K["note_box"](prs.slides[-1], Inches(0.6), Inches(6.75),
                          Inches(12.13), _note_fn(_bt))

    note = str(s0.get("notes") or "").strip()
    _notes = _norm_notes(observations)
    if note or _notes:
        sl = K["content_slide"](prs, "Notes & Observations")
        _parts = []
        if note:
            _parts.append(note)
        for _lbl, _txt in _notes:
            _parts.append(f"{_lbl}: {_txt}" if _txt else str(_lbl))
        K["textbox"](sl, Inches(0.7), Inches(1.4), Inches(11.93), Inches(5.6),
                     "\n\n".join(_parts), 16, K["INK_C"])
    return K["finish"](prs)


@_maybe_cache(show_spinner=False, max_entries=48, ttl=900)
def build_period_pptx(start, end, daily, mortality, customs, hospital_name,
                      period_label, observations=None, note_scope=None,
                      note_ref=None, perf_summary=None, depts=None, meds=None,
                      tests=None, blood=None, absent=None):
    """The weekly / monthly report as a PowerPoint deck: title, At-a-Glance
    period summary, mortality by ward, custom tables (with their configured
    charts and colours), and the period's observations."""
    K = _pptx_toolkit()
    from pptx.util import Inches
    prs = K["new_deck"]()
    K["title_slide"](prs, hospital_name,
                     f"{period_label} • {start:%d %b %Y} – {end:%d %b %Y}",
                     f"Generated {now_local():%d %b %Y %H:%M}")

    d = daily if daily is not None else pd.DataFrame()
    srows = []
    for k, lbl, kind in PERIOD_SUMMARY:
        if d.empty or k not in d.columns:
            srows.append((lbl, "—"))
            continue
        ser = pd.to_numeric(d[k], errors="coerce").dropna()
        if ser.empty:
            srows.append((lbl, "—"))
        elif kind == "sum":
            srows.append((lbl, f"{int(ser.sum()):,}"))
        else:
            srows.append((lbl, f"{ser.mean():.1f}"))
    days_rep = int(d["entry_date"].nunique()) if "entry_date" in d.columns else 0
    srows.append(("Days reported", str(days_rep)))
    if perf_summary:
        try:
            _pimg = performance_report_image(perf_summary)
            if _pimg:
                slp = K["content_slide"](
                    prs, f"Hospital Performance — {perf_summary['condition']}")
                K["picture"](slp, _pimg)
        except Exception:
            log.exception("period PPTX: Performance slide failed")

    sl = K["content_slide"](prs, "At a Glance — Period Summary")
    K["add_table"](sl, ["Metric", "Value"], srows, Inches(2.17), Inches(1.3),
                   Inches(9), font=12, max_rows=14)

    # ── Reorderable band: Mortality-by-Ward and the custom-table slides, emitted
    # in the analyst's saved dashboard order. ──
    band = _ReportBand()
    _note_fn = ((lambda tt: section_note_full(note_scope, note_ref, tt))
                if (note_scope and note_ref) else None)

    def _trend_png(specs, title):
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        from PIL import Image as _PILImage
        import io as _io2
        dd = d.copy()
        if "entry_date" in dd.columns:
            dd["entry_date"] = pd.to_datetime(dd["entry_date"])
            dd = dd.sort_values("entry_date")
            xv = dd["entry_date"]
        else:
            xv = range(len(dd))
        fig, ax = plt.subplots(figsize=(9.0, 4.2), dpi=150)
        for col, lbl, color in specs:
            if col in dd.columns:
                ax.plot(xv, pd.to_numeric(dd[col], errors="coerce"), marker="o",
                        label=lbl, color=color, linewidth=2)
        ax.set_title(title, color=PRIMARY, fontsize=13, loc="left")
        ax.legend(fontsize=8, frameon=False)
        ax.grid(True, alpha=.25)
        for sp in ax.spines.values():
            sp.set_visible(False)
        try:
            fig.autofmt_xdate(rotation=30)
        except Exception:
            pass
        b = _io2.BytesIO()
        fig.savefig(b, format="png", bbox_inches="tight")
        plt.close(fig)
        b.seek(0)
        _im = _PILImage.open(b)
        ar = _im.height / _im.width
        b.seek(0)
        return (b, ar)

    def _mk_trend_slide(title, chart_title, specs):
        def _fn():
            sl = K["content_slide"](prs, title)
            K["picture"](sl, _trend_png(specs, chart_title), top=Inches(1.35))
        return _fn

    if not d.empty:
        band.add("Patient Flow", _mk_trend_slide(
            "Patient Flow", "Patients, admissions & discharges",
            [("current_inpatients", "Patients in hospital", PRIMARY),
             ("admitted", "New admissions", TEAL2),
             ("discharged", "Discharged", WARN)]))
        band.add("ER / Surgeries / ICU", _mk_trend_slide(
            "ER / Surgeries / ICU", "ER visits, surgeries & ICU",
            [("er_visits", "ER visits", PRIMARY),
             ("surgeries", "Surgeries", TEAL2),
             ("icu_patients", "ICU patients", DANGER)]))
        band.add("Beds & Occupancy", _mk_trend_slide(
            "Beds & Occupancy", "Bed availability",
            [("beds_available", "Beds available", PRIMARY),
             ("beds_total", "Total beds", TEAL2),
             ("icu_beds_available", "ICU beds available", WARN)]))
        band.add("Staffing", _mk_trend_slide(
            "Staffing", "Staff on duty",
            [("doctors", "Doctors", PRIMARY), ("nurses", "Nurses", TEAL2),
             ("support_staff", "Support", WARN),
             ("specialists_on_call", "Specialists", DANGER)]))
        band.add("Ambulances", _mk_trend_slide(
            "Ambulances", "Ambulance availability & calls",
            [("ambulances_available", "Available", PRIMARY),
             ("ambulance_calls", "Calls responded", TEAL2)]))

        _avg_ox = pd.to_numeric(d.get("oxygen_pct"), errors="coerce").mean() \
            if "oxygen_pct" in d.columns else 0

        def _slide_supplies(_ox=_avg_ox):
            sl = K["content_slide"](prs, "Critical Supplies")
            K["picture"](sl, oxygen_gauge_image(
                0 if pd.isna(_ox) else _ox,
                title="Oxygen supply level (period avg)"),
                top=Inches(1.35), max_h=Inches(5.4))
        band.add("Critical Supplies", _slide_supplies)

        band.add("Daily Trend", _mk_trend_slide(
            "Daily Trend", "Daily trend",
            [("current_inpatients", "Patients in hospital", PRIMARY),
             ("admitted", "New admissions", TEAL2),
             ("deaths", "Mortality", DANGER)]))

    # ── latest resource snapshots ──
    _blood_rows = snapshot_blood(blood)
    if _blood_rows:
        _rows = [(bt, str(u)) for bt, u in _blood_rows]

        def _slide_blood(_rows=_rows):
            sl = K["content_slide"](prs, "Blood Bank")
            K["add_table"](sl, ["Blood Type", "Units"], _rows,
                           Inches(2.17), Inches(1.35), Inches(9), font=13)
        band.add("Blood Bank", _slide_blood)

    _dept_rows = snapshot_status(depts, ("name", "Department"), ("status", "Status"))
    if _dept_rows:
        def _slide_dept(_rows=_dept_rows):
            sl = K["content_slide"](prs, "Department Status")
            K["add_table"](sl, ["Department", "Status"], _rows,
                           Inches(2.17), Inches(1.35), Inches(9))
        band.add("Department Status", _slide_dept)

    _med_rows = snapshot_status(meds, ("name", "Medication"), ("status", "Status"))
    if _med_rows:
        def _slide_meds(_rows=_med_rows):
            sl = K["content_slide"](prs, "Medication Availability")
            K["add_table"](sl, ["Medication", "Availability"], _rows,
                           Inches(2.17), Inches(1.35), Inches(9))
        band.add("Medication Availability", _slide_meds)

    _test_rows = snapshot_tests(tests)
    if _test_rows:
        _rows = [(n, "Yes" if av else "No") for n, av in _test_rows]

        def _slide_tests(_rows=_rows):
            sl = K["content_slide"](prs, "Medical Tests")
            K["add_table"](sl, ["Test", "Available"], _rows,
                           Inches(2.17), Inches(1.35), Inches(9))
        band.add("Medical Tests", _slide_tests)

    _absent_rows = snapshot_absent(absent)
    if _absent_rows:
        def _slide_absent(_rows=_absent_rows):
            sl = K["content_slide"](prs, "Absent Specialists")
            K["add_table"](sl, ["Specialist", "Specialty / Area", "Expected return"],
                           _rows, Inches(0.6), Inches(1.35), Inches(12.13))
        band.add("Absent Specialists", _slide_absent)

    if mortality is not None and not getattr(mortality, "empty", True):
        mc = mortality.columns
        wcol = "Ward" if "Ward" in mc else ("ward" if "ward" in mc else None)
        if wcol:
            counts = (mortality[mortality[wcol].astype(str).str.strip() != ""]
                      .groupby(wcol).size().sort_values(ascending=False))
            if len(counts):
                def _slide_mort_ward(counts=counts):
                    sl = K["content_slide"](prs, "Mortality by Ward")
                    K["add_table"](sl, ["Ward", "Deaths"],
                                   [(str(w), str(int(n))) for w, n in counts.items()],
                                   Inches(2.17), Inches(1.35), Inches(9), font=13)
                band.add("Mortality by Ward", _slide_mort_ward)

    for _item in (customs or []):
        band.add(_item[0],
                 (lambda it=_item: K["customs_slides"](prs, [it], note_fn=None)))

    for _bt, _fn in band.items():
        _fn()
        if _note_fn and prs.slides:
            K["note_box"](prs.slides[-1], Inches(0.6), Inches(6.75),
                          Inches(12.13), _note_fn(_bt))

    _notes = _norm_notes(observations)
    if _notes:
        sl = K["content_slide"](prs, "Notes & Observations")
        _parts = [f"{_lbl}: {_txt}" if _txt else str(_lbl)
                  for _lbl, _txt in _notes]
        K["textbox"](sl, Inches(0.7), Inches(1.4), Inches(11.93), Inches(5.6),
                     "\n\n".join(_parts), 16, K["INK_C"])
    return K["finish"](prs)


# ── Analyst logins (email + salted password hash) ──
def users_count():
    return get_conn().execute("SELECT COUNT(*) FROM users").fetchone()[0]


def get_user(email):
    row = get_conn().execute(
        "SELECT email, password_hash, created_at, last_login FROM users WHERE email=?",
        ((email or "").strip().lower(),)).fetchone()
    if not row:
        return None
    return {"email": row[0], "password_hash": row[1],
            "created_at": row[2], "last_login": row[3]}


def list_users():
    return [{"email": r[0], "created_at": r[1], "last_login": r[2]} for r in
            get_conn().execute(
                "SELECT email, created_at, last_login FROM users ORDER BY email").fetchall()]


def upsert_user(email, pw):
    """Create or update an analyst login, storing only a salted PBKDF2 hash."""
    email = (email or "").strip().lower()
    c = get_conn()
    c.execute("INSERT INTO users(email, password_hash, created_at) VALUES(?,?,?) "
              "ON CONFLICT(email) DO UPDATE SET password_hash=excluded.password_hash",
              (email, make_password_hash(pw),
               now_local().isoformat(timespec="seconds")))
    c.commit()


def delete_user(email):
    c = get_conn()
    c.execute("DELETE FROM users WHERE email=?", ((email or "").strip().lower(),))
    c.commit()


def set_user_last_login(email):
    c = get_conn()
    c.execute("UPDATE users SET last_login=? WHERE email=?",
              (now_local().isoformat(timespec="seconds"), (email or "").strip().lower()))
    c.commit()


def verify_login(email, pw):
    """Return (ok, who) for an email+password attempt.
    Order: a stored analyst account, else the master/bootstrap admin password
    (which always works when no accounts exist yet, or matches HOSPITAL_ADMIN_EMAIL
    when that is configured)."""
    email = (email or "").strip().lower()
    u = get_user(email)
    if u and _verify_pbkdf2(pw, u["password_hash"]):
        return True, email
    if verify_password(pw) and (users_count() == 0 or not ADMIN_EMAIL
                                or email == ADMIN_EMAIL):
        return True, (email or "admin")
    return False, None


def dates_with_data(start, end):
    """Set of YYYY-MM-DD strings in [start,end] that already have any data."""
    c = get_conn()
    p = (start.isoformat(), end.isoformat())
    s = set()
    for tbl in ("daily", "departments", "medications", "tests", "blood_bank",
                "absent_specialists", "mortality", "custom_table_rows"):
        for r in c.execute(f"SELECT DISTINCT entry_date FROM {tbl} "
                           f"WHERE entry_date BETWEEN ? AND ?", p).fetchall():
            s.add(r[0])
    return s


def _to_iso_date(v):
    """Coerce a date / Timestamp / string / NaT cell into an ISO date string ('' if blank)."""
    if v is None:
        return ""
    try:
        if pd.isna(v):
            return ""
    except (TypeError, ValueError):
        pass
    if isinstance(v, str):
        v = v.strip()
        if not v:
            return ""
        try:
            return date.fromisoformat(v[:10]).isoformat()
        except ValueError:
            try:
                return pd.to_datetime(v).date().isoformat()
            except Exception:
                return ""
    if hasattr(v, "date"):          # datetime / pandas Timestamp
        try:
            return v.date().isoformat()
        except Exception:
            return ""
    if isinstance(v, date):
        return v.isoformat()
    return ""


def _to_time_str(v):
    """Coerce a datetime.time / Timestamp / string / NaT cell into an 'HH:MM' string
    ('' if blank)."""
    if v is None:
        return ""
    try:
        if pd.isna(v):
            return ""
    except (TypeError, ValueError):
        pass
    if hasattr(v, "strftime"):          # datetime.time / datetime / Timestamp
        try:
            return v.strftime("%H:%M")
        except Exception:
            return ""
    s = str(v).strip()
    if not s:
        return ""
    try:                                # parse strings like '8:30', '08:30:00'
        return pd.to_datetime(s).strftime("%H:%M")
    except Exception:
        return s[:5]


def save_day(d, numeric, notes, dept_df, med_df, test_df, blood_df, absent_df=None,
             mortality_df=None):
    notes = clean_text(notes, multiline=True)
    dept_df, med_df, test_df = clean_frame(dept_df), clean_frame(med_df), clean_frame(test_df)
    blood_df, absent_df, mortality_df = (clean_frame(blood_df), clean_frame(absent_df),
                                         clean_frame(mortality_df))
    ds = d.isoformat()
    c = get_conn()
    cols = FIELD_KEYS + ["notes", "updated_at"]
    placeholders = ", ".join("?" for _ in cols) + ", ?"
    updates = ", ".join(f"{k}=excluded.{k}" for k in cols)
    vals = [int(numeric.get(k, 0)) for k in FIELD_KEYS] + \
           [notes, now_local().isoformat(timespec="seconds"), ds]

    def _int_or_none(v):
        try:
            if v is None or (isinstance(v, float) and pd.isna(v)):
                return None
            return int(v)
        except (TypeError, ValueError):
            return None

    # one atomic transaction: the whole day saves, or nothing changes
    with c:
        c.execute(f"INSERT INTO daily({', '.join(cols)}, entry_date) VALUES({placeholders}) "
                  f"ON CONFLICT(entry_date) DO UPDATE SET {updates}", vals)
        for tbl in ("departments", "medications", "tests", "blood_bank",
                    "absent_specialists", "mortality"):
            c.execute(f"DELETE FROM {tbl} WHERE entry_date=?", (ds,))
        for _, r in dept_df.iterrows():
            n = str(r.get("Department", "")).strip()
            if n:
                c.execute("INSERT OR REPLACE INTO departments VALUES(?,?,?)",
                          (ds, n, str(r.get("Status", "Operational"))))
        for _, r in med_df.iterrows():
            n = str(r.get("Medication", "")).strip()
            if n:
                status = str(r.get("Status", "Available")).strip() or "Available"
                c.execute("INSERT OR REPLACE INTO medications(entry_date, name, status) "
                          "VALUES(?,?,?)", (ds, n, status))
        for _, r in test_df.iterrows():
            n = str(r.get("Test", "")).strip()
            if n:
                c.execute("INSERT OR REPLACE INTO tests VALUES(?,?,?)",
                          (ds, n, 1 if r.get("Available", False) else 0))
        for _, r in blood_df.iterrows():
            bt = str(r.get("Blood Type", "")).strip()
            if bt:
                c.execute("INSERT OR REPLACE INTO blood_bank VALUES(?,?,?)",
                          (ds, bt, int(r.get("Units", 0) or 0)))
        if absent_df is not None:
            for _, r in absent_df.iterrows():
                nm = str(r.get("Specialist", "")).strip()
                if nm:
                    c.execute("INSERT OR REPLACE INTO absent_specialists VALUES(?,?,?,?)",
                              (ds, nm, str(r.get("Specialty / Area", "") or "").strip(),
                               _to_iso_date(r.get("Expected return"))))
        if mortality_df is not None:
            for _, r in mortality_df.iterrows():
                ward = str(r.get("Ward", "") or "").strip()
                if not ward:
                    continue
                c.execute("INSERT INTO mortality(entry_date, ward, age, los_days, condition, "
                          "cause, time_of_death) VALUES(?,?,?,?,?,?,?)",
                          (ds, ward, _int_or_none(r.get("Age")),
                           _int_or_none(r.get("Length of stay (days)")),
                           str(r.get("Condition", "") or "").strip(),
                           str(r.get("Cause of death", "") or "").strip(),
                           _to_time_str(r.get("Time of death"))))


def load_day(d):
    ds = d.isoformat()
    c = get_conn()
    cur = c.execute("SELECT * FROM daily WHERE entry_date=?", (ds,))
    row = cur.fetchone()
    scalars = dict(zip([x[0] for x in cur.description], row)) if row else None
    depts = pd.read_sql_query("SELECT name AS Department, status AS Status FROM departments "
                              "WHERE entry_date=?", c, params=(ds,))
    meds = pd.read_sql_query(
        "SELECT name AS Medication, COALESCE(status, "
        "CASE WHEN stock IS NULL THEN 'Available' WHEN stock<=0 THEN 'Not available' "
        "WHEN stock<=10 THEN 'Limited availability' ELSE 'Available' END) AS Status "
        "FROM medications WHERE entry_date=?", c, params=(ds,))
    tests = pd.read_sql_query("SELECT name AS Test, available FROM tests WHERE entry_date=?",
                              c, params=(ds,))
    if not tests.empty:
        tests["Available"] = tests["available"].astype(bool)
        tests = tests[["Test", "Available"]]
    blood = pd.read_sql_query("SELECT blood_type AS 'Blood Type', units AS Units FROM blood_bank "
                              "WHERE entry_date=?", c, params=(ds,))
    absent = pd.read_sql_query(
        "SELECT name AS 'Specialist', specialty AS 'Specialty / Area', "
        "expected_return AS 'Expected return' FROM absent_specialists WHERE entry_date=?",
        c, params=(ds,))
    if not absent.empty:
        absent["Expected return"] = pd.to_datetime(
            absent["Expected return"], errors="coerce").dt.date
    return scalars, depts, meds, tests, blood, absent


# tables that carry a day's figures and layouts forward when duplicating a day.
# mortality and per-day observation notes are deliberately excluded: those
# record what actually happened on a specific day and must not be copied.
_DUPLICATE_TABLES = ("departments", "medications", "tests", "blood_bank",
                     "absent_specialists", "custom_table_rows")


def day_has_data(d):
    """True if any figures, resource tables or custom tables are saved for day d."""
    ds = d.isoformat()
    c = get_conn()
    if c.execute("SELECT 1 FROM daily WHERE entry_date=? LIMIT 1", (ds,)).fetchone():
        return True
    for tbl in _DUPLICATE_TABLES:
        if c.execute(f"SELECT 1 FROM {tbl} WHERE entry_date=? LIMIT 1", (ds,)).fetchone():
            return True
    return False


def duplicate_day(src, dst):
    """Copy the figures, resource tables and custom tables saved for day src into
    day dst, overwriting whatever dst held. Mortality records and per-day notes
    are left untouched, being specific to the day they occurred."""
    s, t = src.isoformat(), dst.isoformat()
    c = get_conn()
    with c:
        cur = c.execute("SELECT * FROM daily WHERE entry_date=?", (s,))
        row = cur.fetchone()
        if row is not None:
            names = [dd[0] for dd in cur.description]
            data = dict(zip(names, row))
            data["entry_date"] = t
            data["updated_at"] = now_local().isoformat(timespec="seconds")
            ph = ", ".join("?" for _ in names)
            upd = ", ".join(f"{k}=excluded.{k}" for k in names if k != "entry_date")
            c.execute(f"INSERT INTO daily({', '.join(names)}) VALUES({ph}) "
                      f"ON CONFLICT(entry_date) DO UPDATE SET {upd}",
                      [data[k] for k in names])
        for tbl in _DUPLICATE_TABLES:
            c.execute(f"DELETE FROM {tbl} WHERE entry_date=?", (t,))
            cur2 = c.execute(f"SELECT * FROM {tbl} WHERE entry_date=?", (s,))
            names2 = [dd[0] for dd in cur2.description]
            idx = names2.index("entry_date")
            ph2 = ", ".join("?" for _ in names2)
            for r in cur2.fetchall():
                r = list(r)
                r[idx] = t
                c.execute(f"INSERT INTO {tbl}({', '.join(names2)}) VALUES({ph2})", r)


def load_range(start, end):
    c = get_conn()
    p = (start.isoformat(), end.isoformat())
    return (pd.read_sql_query("SELECT * FROM daily WHERE entry_date BETWEEN ? AND ? ORDER BY entry_date", c, params=p),
            pd.read_sql_query("SELECT * FROM departments WHERE entry_date BETWEEN ? AND ?", c, params=p),
            pd.read_sql_query(
                "SELECT entry_date, name, COALESCE(status, "
                "CASE WHEN stock IS NULL THEN 'Available' WHEN stock<=0 THEN 'Not available' "
                "WHEN stock<=10 THEN 'Limited availability' ELSE 'Available' END) AS status "
                "FROM medications WHERE entry_date BETWEEN ? AND ?", c, params=p),
            pd.read_sql_query("SELECT * FROM tests WHERE entry_date BETWEEN ? AND ?", c, params=p),
            pd.read_sql_query("SELECT * FROM blood_bank WHERE entry_date BETWEEN ? AND ?", c, params=p),
            pd.read_sql_query("SELECT * FROM absent_specialists WHERE entry_date BETWEEN ? AND ?", c, params=p))


# ── Wards (configurable list) + mortality register ──
def get_wards():
    """The configurable ward list (stored one per line in settings)."""
    raw = get_setting("wards", "")
    wards = [w.strip() for w in raw.replace(",", "\n").splitlines() if w.strip()]
    # de-duplicate, preserving order
    seen, out = set(), []
    for w in (wards or DEFAULT_WARDS):
        if w.lower() not in seen:
            seen.add(w.lower())
            out.append(w)
    return out


def set_wards(text):
    wards = [w.strip() for w in (text or "").replace(",", "\n").splitlines() if w.strip()]
    seen, out = set(), []
    for w in wards:
        if w.lower() not in seen:
            seen.add(w.lower())
            out.append(w)
    set_setting("wards", "\n".join(out))


_MORT_COLS = ["Ward", "Age", "Length of stay (days)", "Time of death",
              "Condition", "Cause of death"]


def _empty_mortality():
    return pd.DataFrame({"Ward": pd.Series(dtype="str"),
                         "Age": pd.Series(dtype="object"),
                         "Length of stay (days)": pd.Series(dtype="object"),
                         "Time of death": pd.Series(dtype="object"),
                         "Condition": pd.Series(dtype="str"),
                         "Cause of death": pd.Series(dtype="str")})


def load_mortality(d):
    """Mortality register rows for one day -> DataFrame[Ward, Age, Length of stay (days),
    Time of death, Condition, Cause of death]."""
    df = pd.read_sql_query(
        "SELECT ward AS 'Ward', age AS 'Age', los_days AS 'Length of stay (days)', "
        "COALESCE(time_of_death,'') AS 'Time of death', "
        "COALESCE(condition,'') AS 'Condition', COALESCE(cause,'') AS 'Cause of death' "
        "FROM mortality WHERE entry_date=?", get_conn(), params=(d.isoformat(),))
    return df if not df.empty else _empty_mortality()


def load_mortality_range(start, end):
    """Mortality register rows across a date range, including entry_date."""
    return pd.read_sql_query(
        "SELECT entry_date, ward, age, los_days, "
        "COALESCE(time_of_death,'') AS time_of_death, "
        "COALESCE(condition,'') AS condition, COALESCE(cause,'') AS cause FROM mortality "
        "WHERE entry_date BETWEEN ? AND ?", get_conn(),
        params=(start.isoformat(), end.isoformat()))


def mortality_stats(df):
    """Summarise a mortality DataFrame: total, average age, average length of stay,
    and per-ward counts. Accepts editor labels or raw column names."""
    if df is None or getattr(df, "empty", True):
        return {"total": 0, "avg_age": None, "avg_los": None, "by_ward": []}
    cols = list(df.columns)
    ward_c = "Ward" if "Ward" in cols else "ward"
    age_c = "Age" if "Age" in cols else "age"
    los_c = "Length of stay (days)" if "Length of stay (days)" in cols else "los_days"
    wcol = df[ward_c].astype(str).str.strip()
    valid = wcol != ""
    total = int(valid.sum())
    ages = pd.to_numeric(df.loc[valid, age_c], errors="coerce").dropna()
    los = pd.to_numeric(df.loc[valid, los_c], errors="coerce").dropna()
    by_ward = (wcol[valid].value_counts().sort_values(ascending=False)
               .items() if total else [])
    return {"total": total,
            "avg_age": round(float(ages.mean()), 1) if len(ages) else None,
            "avg_los": round(float(los.mean()), 1) if len(los) else None,
            "by_ward": list(by_ward)}


def mortality_by_ward_fig(df, title="Mortality by ward"):
    """Horizontal bar of recorded deaths per ward. Returns None when empty."""
    st_ = mortality_stats(df)
    if not st_["by_ward"]:
        return None
    wards = [w for w, _ in st_["by_ward"]]
    counts = [int(c) for _, c in st_["by_ward"]]
    fig = go.Figure(go.Bar(x=counts, y=wards, orientation="h", marker_color=DANGER))
    fig.update_layout(title=title, yaxis=dict(autorange="reversed"))
    return style_fig(fig, h=max(300, 42 * len(wards) + 120))



def mortality_ward_breakdown(df):
    """Per-ward summary -> DataFrame[Ward, Deaths, Ages recorded, Average age].
    'Deaths' is the count of records for that ward; 'Ages recorded' lists each
    recorded age."""
    empty = pd.DataFrame(columns=["Ward", "Deaths", "Ages recorded", "Average age"])
    if df is None or getattr(df, "empty", True):
        return empty
    cols = list(df.columns)
    ward_c = "Ward" if "Ward" in cols else "ward"
    age_c = "Age" if "Age" in cols else "age"
    sub = df[df[ward_c].astype(str).str.strip() != ""].copy()
    if sub.empty:
        return empty
    sub[ward_c] = sub[ward_c].astype(str).str.strip()
    rows = []
    for ward, g in sub.groupby(ward_c, sort=False):
        ages = pd.to_numeric(g[age_c], errors="coerce").dropna()
        ages_list = sorted(int(a) for a in ages)
        rows.append({"Ward": ward, "Deaths": int(len(g)),
                     "Ages recorded": ", ".join(str(a) for a in ages_list) if ages_list else "—",
                     "Average age": round(float(ages.mean()), 1) if len(ages) else "—"})
    return (pd.DataFrame(rows, columns=["Ward", "Deaths", "Ages recorded", "Average age"])
            .sort_values("Deaths", ascending=False, kind="stable").reset_index(drop=True))


def mortality_ward_table_fig(df, title="Mortality by ward (count & ages)"):
    """Table fig of the per-ward breakdown: ward, death count, ages recorded, avg age."""
    bd = mortality_ward_breakdown(df)
    if bd.empty:
        return None
    fig = go.Figure(go.Table(
        header=dict(values=[f"<b>{c}</b>" for c in bd.columns],
                    fill_color=PRIMARY, font=dict(color="#FFFFFF", size=14),
                    align="left", height=32),
        cells=dict(values=[bd[c].tolist() for c in bd.columns], fill_color="#FFFFFF",
                   font=dict(color=INK, size=13), align="left", height=28)))
    fig.update_layout(title=title, paper_bgcolor="#FFFFFF",
                      margin=dict(l=8, r=8, t=44, b=8), title_font=dict(color=PRIMARY))
    return fig


# ── User-defined editable tables (multiple; each has a title, columns and per-day rows) ──
DEFAULT_CUSTOM_COLS = ["Column 1", "Column 2"]


# ── Paste hygiene ───────────────────────────────────────────────
# Text pasted from Word / Excel / WhatsApp carries invisible characters
# (non-breaking & zero-width spaces, BOMs, soft hyphens), smart quotes and
# stray tabs. Left alone they make "ICU" and "ICU\u00a0" two different chart
# categories with two different colours. Everything typed or pasted into a
# table goes through clean_text() on its way to the database.
_PASTE_MAP = {0x00A0: " ", 0x2007: " ", 0x202F: " ", 0x2000: " ", 0x2001: " ",
              0x2002: " ", 0x2003: " ", 0x2004: " ", 0x2005: " ", 0x2006: " ",
              0x2008: " ", 0x2009: " ", 0x200A: " ", 0x3000: " ",
              0x2018: "'", 0x2019: "'", 0x201A: "'", 0x2032: "'",
              0x201C: '"', 0x201D: '"', 0x201E: '"',
              0x2026: "...",
              0x200B: None, 0x200C: None, 0x200D: None, 0x2060: None,
              0xFEFF: None, 0x00AD: None, 0x200E: None, 0x200F: None}


def clean_text(v, multiline=False):
    """Normalise pasted text: odd unicode spaces become normal spaces, smart
    quotes become straight quotes, zero-width/control characters are removed,
    and whitespace is collapsed and trimmed. Multiline mode (notes,
    observations) keeps the line breaks."""
    if v is None:
        return ""
    s = str(v).translate(_PASTE_MAP)
    s = "".join(ch for ch in s if ch in "\n\t" or ord(ch) >= 32)
    if multiline:
        return "\n".join(" ".join(ln.split()) for ln in s.split("\n")).strip()
    return " ".join(s.split())


def clean_frame(df):
    """Paste hygiene for every text cell of an edited table; other value
    types (numbers, dates, checkboxes) pass through untouched."""
    if df is None or getattr(df, "empty", True):
        return df
    out = df.copy()
    for c in out.columns:
        if str(out[c].dtype) in ("object", "string", "str"):
            out[c] = out[c].map(lambda v: clean_text(v) if isinstance(v, str) else v)
    return out


def _dedupe(cols):
    seen, out = set(), []
    for c in cols:
        c = clean_text(c)
        if c and c.lower() not in seen:
            seen.add(c.lower())
            out.append(c)
    return out


def get_custom_tables():
    """List of custom tables: [{id, title, columns}], from settings."""
    raw = get_setting("custom_tables", "")
    if not raw:
        return []
    try:
        data = json.loads(raw)
    except Exception:
        return []
    out = []
    for t in data:
        tid = str(t.get("id") or "").strip()
        if not tid:
            continue
        cols = _dedupe(t.get("columns", [])) or list(DEFAULT_CUSTOM_COLS)
        disp = t.get("display") if isinstance(t.get("display"), dict) else {}
        out.append({"id": tid, "title": (str(t.get("title") or "Untitled table").strip()
                                         or "Untitled table"), "columns": cols,
                    "display": disp})
    return out


def set_custom_tables(tables):
    set_setting("custom_tables", json.dumps(tables))


def get_custom_table(tid):
    return next((t for t in get_custom_tables() if t["id"] == tid), None)


def add_custom_table(title=None, columns=None, display=None):
    tables = get_custom_tables()
    tid = pysecrets.token_hex(4)
    title = clean_text(title) or f"Table {len(tables) + 1}"
    tables.append({"id": tid, "title": title,
                   "columns": _dedupe(columns or []) or list(DEFAULT_CUSTOM_COLS),
                   "display": dict(display or {})})
    set_custom_tables(tables)
    return tid


def update_custom_table(tid, title=None, columns=None, display=None):
    tables = get_custom_tables()
    for t in tables:
        if t["id"] == tid:
            if title is not None and clean_text(title):
                t["title"] = clean_text(title)
            if columns is not None:
                if isinstance(columns, str):
                    columns = columns.replace(",", "\n").splitlines()
                t["columns"] = _dedupe(columns) or list(DEFAULT_CUSTOM_COLS)
            if display is not None and isinstance(display, dict):
                t["display"] = dict(display)
    set_custom_tables(tables)


def replicate_custom_table(tid):
    """Duplicate a table's title + columns into a brand-new table (rows not copied)."""
    src = get_custom_table(tid)
    if not src:
        return None
    return add_custom_table("Copy of " + src["title"], list(src["columns"]),
                            dict(src.get("display") or {}))


def delete_custom_table(tid):
    set_custom_tables([t for t in get_custom_tables() if t["id"] != tid])
    c = get_conn()
    c.execute("DELETE FROM custom_table_rows WHERE table_id=?", (tid,))
    c.commit()


def _empty_rows(cols):
    return pd.DataFrame({c: pd.Series(dtype="str") for c in cols})


def load_custom_rows(tid, day, cols=None):
    if cols is None:
        t = get_custom_table(tid)
        cols = t["columns"] if t else list(DEFAULT_CUSTOM_COLS)
    row = get_conn().execute(
        "SELECT payload FROM custom_table_rows WHERE entry_date=? AND table_id=?",
        (day.isoformat(), tid)).fetchone()
    data = []
    if row and row[0]:
        try:
            data = json.loads(row[0])
        except Exception:
            data = []
    df = pd.DataFrame(data)
    for c in cols:
        if c not in df.columns:
            df[c] = ""
    df = df.reindex(columns=cols)
    return df if not df.empty else _empty_rows(cols)


def save_custom_rows(tid, day, df, cols=None):
    df = clean_frame(df)
    if cols is None:
        t = get_custom_table(tid)
        cols = t["columns"] if t else []
    records = []
    if df is not None:
        for _, r in df.iterrows():
            rec = {}
            for c in cols:
                v = r.get(c, "")
                rec[c] = "" if (v is None or (isinstance(v, float) and pd.isna(v))) else str(v)
            if any(str(v).strip() for v in rec.values()):
                records.append(rec)
    c = get_conn()
    if records:
        c.execute("INSERT INTO custom_table_rows(entry_date, table_id, payload) VALUES(?,?,?) "
                  "ON CONFLICT(entry_date, table_id) DO UPDATE SET payload=excluded.payload",
                  (day.isoformat(), tid, json.dumps(records)))
    else:
        c.execute("DELETE FROM custom_table_rows WHERE entry_date=? AND table_id=?",
                  (day.isoformat(), tid))
    c.commit()


# ── Import from CSV / Excel ─────────────────────────────────────
def parse_import_file(name, data):
    """Read an uploaded CSV or Excel file into a clean DataFrame (all columns
    as trimmed text so mapping and storage are predictable). Returns
    (df, error_message): on success error is '', on failure df is None and the
    message explains what to fix."""
    import io as _io
    nm = str(name or "").lower()
    try:
        if nm.endswith((".xlsx", ".xls", ".xlsm")):
            try:
                df = pd.read_excel(_io.BytesIO(data))
            except ImportError:
                return None, ("Reading Excel needs the openpyxl package. Add "
                              "`openpyxl` to requirements.txt, or save the sheet "
                              "as CSV and import that.")
        else:
            raw = data.decode("utf-8-sig", errors="replace")
            try:
                df = pd.read_csv(_io.StringIO(raw))
                if df.shape[1] <= 1:          # probably not comma-delimited
                    alt = pd.read_csv(_io.StringIO(raw), sep=None, engine="python")
                    if alt.shape[1] > df.shape[1]:
                        df = alt
            except Exception:
                df = pd.read_csv(_io.StringIO(raw), sep=None, engine="python")
    except Exception as e:
        log.exception("import parse failed")
        return None, f"Could not read that file: {e}"
    if df is None or df.empty or not len(df.columns):
        return None, "That file has no rows or no columns to import."
    df = df.dropna(axis=1, how="all")                 # drop empty columns
    df.columns = [clean_text(c) or f"Column {i + 1}"
                  for i, c in enumerate(df.columns)]
    for c in df.columns:                              # everything as clean text
        df[c] = df[c].map(lambda v: "" if (v is None or (isinstance(v, float)
                          and pd.isna(v))) else clean_text(v))
    df = df.dropna(axis=0, how="all")
    df = df[~(df.apply(lambda r: all(str(x).strip() == "" for x in r), axis=1))]
    return df.reset_index(drop=True), ""


def _norm_header(s):
    return "".join(ch for ch in str(s).lower() if ch.isalnum())


def suggest_column_map(target_cols, source_cols):
    """Best-guess mapping {target: source} by matching normalised header names
    (exact first, then contains), so obvious columns line up automatically and
    the analyst only fixes the rest. Each source column is used at most once."""
    smap = {_norm_header(s): s for s in source_cols}
    used, out = set(), {}
    for t in target_cols:
        nt = _norm_header(t)
        pick = None
        if nt in smap and smap[nt] not in used:
            pick = smap[nt]
        else:
            for s in source_cols:
                if s in used:
                    continue
                ns = _norm_header(s)
                if ns and (ns in nt or nt in ns):
                    pick = s
                    break
        if pick:
            out[t] = pick
            used.add(pick)
    return out


def load_custom_rows_range(tid, start, end, cols=None):
    if cols is None:
        t = get_custom_table(tid)
        cols = t["columns"] if t else list(DEFAULT_CUSTOM_COLS)
    rows = get_conn().execute(
        "SELECT entry_date, payload FROM custom_table_rows WHERE table_id=? "
        "AND entry_date BETWEEN ? AND ? ORDER BY entry_date",
        (tid, start.isoformat(), end.isoformat())).fetchall()
    out = []
    for ds, payload in rows:
        try:
            for rec in (json.loads(payload) if payload else []):
                item = {"Date": ds}
                for c in cols:
                    item[c] = rec.get(c, "")
                out.append(item)
        except Exception:
            continue
    return pd.DataFrame(out, columns=["Date"] + cols) if out else pd.DataFrame(columns=["Date"] + cols)


# ── Custom table intelligence: type detection, charts & colors ──
CUSTOM_MODES = [("table", "Table"),
                ("column", "Column chart (vertical bars)"),
                ("bar", "Bar chart (horizontal bars)"),
                ("stacked_column", "Stacked column chart"),
                ("stacked_bar", "Stacked bar chart"),
                ("grouped", "Grouped columns — two-level axis"),
                ("line", "Line graph"),
                ("scatter", "Scatter plot"),
                ("funnel", "Funnel"),
                ("pie", "Pie chart"), ("donut", "Donut chart"),
                ("kpi", "KPI cards"), ("table+chart", "Table + chart")]
CUSTOM_AGGS = [("sum", "Sum"), ("avg", "Average"), ("count", "Count"),
               ("latest", "Latest entry")]
CHART_PALETTE = [PRIMARY, TEAL2, WARN, DANGER, OK_GREEN, ACCENT,
                 "#8E44AD", "#D35400", "#2C7A7B", "#4A5568"]
TYPE_LABEL = {"numeric": "Numbers", "date": "Dates", "text": "Text"}
TYPE_ICON = {"numeric": "🔢", "date": "📅", "text": "🔤"}
TYPE_COLOR = {"numeric": TEAL2, "date": WARN, "text": "#7A8B8B"}


def _safe(default=None):
    """Never let a custom-table rendering error take down a dashboard: on any
    unexpected exception the wrapped function logs it and returns `default`,
    so callers fall back to their plain-table / empty behaviour."""
    def deco(fn):
        def wrapped(*args, **kwargs):
            try:
                return fn(*args, **kwargs)
            except Exception:
                log.exception("custom table function %s failed; falling back",
                              fn.__name__)
                return default
        wrapped.__name__ = fn.__name__
        wrapped.__doc__ = fn.__doc__
        return wrapped
    return deco


def table_display(t):
    """A custom table's saved display config merged over safe defaults, so
    every table keeps working (plain table) until the analyst customises it."""
    cfg = dict(t.get("display") or {}) if isinstance(t, dict) else {}
    cfg.setdefault("mode", "table")
    cfg.setdefault("category", "(auto)")
    cfg.setdefault("values", [])
    cfg.setdefault("pct_cols", [])
    cfg.setdefault("subcategory", "(auto)")
    cfg.setdefault("agg", "sum")
    cfg.setdefault("colors", {})
    cfg.setdefault("cat_colors", {})
    cfg.setdefault("color_by_cat", False)
    cfg.setdefault("row_stripe", "")
    cfg.setdefault("col_tint", False)
    cfg.setdefault("formulas", {})
    cfg.setdefault("summary", "")
    return cfg


def _cell_text(series):
    """The non-empty cells of a column, as trimmed strings."""
    s = series.astype(str).str.strip()
    return s[(s != "") & (s.str.lower() != "nan") & (s.str.lower() != "none")]


# ── Safe spreadsheet-style formulas for custom tables ──────────────
# Formulas never use eval(). An AST walker allows only numbers, the four
# arithmetic operators, parentheses, a handful of functions and references to
# the table's own columns. Anything else is rejected, so a formula can't run
# arbitrary code.
import ast as _ast

_FORMULA_FUNCS = {
    "sum": lambda xs: float(sum(xs)),
    "avg": lambda xs: (float(sum(xs)) / len(xs)) if xs else 0.0,
    "average": lambda xs: (float(sum(xs)) / len(xs)) if xs else 0.0,
    "min": lambda xs: float(min(xs)) if xs else 0.0,
    "max": lambda xs: float(max(xs)) if xs else 0.0,
    "count": lambda xs: float(len(xs)),
    "abs": lambda xs: abs(float(xs[0])) if xs else 0.0,
    "round": lambda xs: round(float(xs[0]), int(xs[1]) if len(xs) > 1 else 0),
}


def _to_number(v):
    """Coerce a single cell to a float (commas / % tolerated); blank -> 0.0."""
    try:
        s = str(v).strip().replace(",", "").replace("%", "")
        return float(s) if s not in ("", "nan", "none", "None") else 0.0
    except (TypeError, ValueError):
        return 0.0


def _colvals_number(df, col):
    """All numeric values of a column (for SUM/AVG-style aggregates)."""
    if df is None or col not in getattr(df, "columns", []):
        return []
    return [_to_number(v) for v in _cell_text(df[col])]


class _FormulaError(Exception):
    pass


def eval_formula(expr, row=None, df=None, colmap=None):
    """Evaluate a spreadsheet formula safely. `row` supplies per-row column
    values (bare column names); functions like SUM(Col) aggregate the whole
    column via `df`. `colmap` maps normalised names -> real column names so
    references are case/space-insensitive. Raises _FormulaError on anything
    unsupported."""
    expr = str(expr or "").strip()
    if expr.startswith("="):
        expr = expr[1:]
    if not expr:
        return 0.0
    colmap = colmap or {}

    def _resolve(name):
        return colmap.get(_norm_header(name), name)

    def _node(n):
        if isinstance(n, _ast.Expression):
            return _node(n.body)
        if isinstance(n, _ast.Constant):
            if isinstance(n.value, (int, float)):
                return float(n.value)
            raise _FormulaError("only numbers are allowed")
        if isinstance(n, _ast.BinOp):
            a, b = _node(n.left), _node(n.right)
            if isinstance(n.op, _ast.Add):
                return a + b
            if isinstance(n.op, _ast.Sub):
                return a - b
            if isinstance(n.op, _ast.Mult):
                return a * b
            if isinstance(n.op, _ast.Div):
                return a / b if b != 0 else 0.0
            if isinstance(n.op, _ast.Mod):
                return a % b if b != 0 else 0.0
            if isinstance(n.op, _ast.Pow):
                return a ** b
            raise _FormulaError("unsupported operator")
        if isinstance(n, _ast.UnaryOp):
            v = _node(n.operand)
            if isinstance(n.op, _ast.USub):
                return -v
            if isinstance(n.op, _ast.UAdd):
                return v
            raise _FormulaError("unsupported unary operator")
        if isinstance(n, _ast.Name):
            real = _resolve(n.id)
            if row is not None and real in row:
                return _to_number(row[real])
            raise _FormulaError(f"unknown column “{n.id}”")
        if isinstance(n, _ast.Call):
            if not isinstance(n.func, _ast.Name):
                raise _FormulaError("unsupported function call")
            fn = n.func.id.lower()
            if fn not in _FORMULA_FUNCS:
                raise _FormulaError(f"unknown function “{n.func.id}”")
            if fn in ("sum", "avg", "average", "min", "max", "count") \
                    and len(n.args) == 1 and isinstance(n.args[0], _ast.Name) \
                    and df is not None:
                return _FORMULA_FUNCS[fn](_colvals_number(df, _resolve(n.args[0].id)))
            return _FORMULA_FUNCS[fn]([_node(a) for a in n.args])
        raise _FormulaError("unsupported expression")

    try:
        tree = _ast.parse(expr, mode="eval")
        return float(_node(tree))
    except _FormulaError:
        raise
    except Exception:
        raise _FormulaError("could not parse the formula")


def apply_formulas(t, df):
    """Return a copy of `df` with each formula column filled in per row.
    Formula columns live in the table's display config under 'formulas'
    ({column: expression}). A bad formula leaves the cell blank, never crashes."""
    cfg = table_display(t)
    formulas = cfg.get("formulas") or {}
    if df is None or getattr(df, "empty", True) or not formulas:
        return df
    out = df.copy()
    colmap = {_norm_header(c): c for c in out.columns}
    for col, expr in formulas.items():
        if not str(expr or "").strip():
            continue
        if col not in out.columns:
            out[col] = ""
        vals = []
        for _, r in out.iterrows():
            try:
                v = eval_formula(expr, row=r.to_dict(), df=out, colmap=colmap)
                vals.append(f"{int(v):,}" if float(v).is_integer()
                            else f"{v:,.2f}")
            except _FormulaError:
                vals.append("")
        out[col] = vals
    return out


def summary_row(t, df):
    """Compute the configured summary aggregate for each numeric column (e.g. a
    totals row). The aggregate is chosen under 'summary' ('', 'sum', 'avg',
    'min', 'max', 'count'); '' means no row. Returns {column: value} or None."""
    cfg = table_display(t)
    agg = cfg.get("summary") or ""
    if not agg or df is None or getattr(df, "empty", True):
        return None
    fn = {"sum": lambda xs: sum(xs),
          "avg": lambda xs: (sum(xs) / len(xs)) if xs else 0,
          "average": lambda xs: (sum(xs) / len(xs)) if xs else 0,
          "min": lambda xs: min(xs) if xs else 0,
          "max": lambda xs: max(xs) if xs else 0,
          "count": lambda xs: len(xs)}.get(agg)
    if not fn:
        return None
    types = infer_column_types(df, [c for c in df.columns if c != "Date"])
    label = {"sum": "Total", "avg": "Average", "average": "Average",
             "min": "Minimum", "max": "Maximum", "count": "Count"}[agg]
    out = {}
    labelled = False
    for c in df.columns:
        if c == "Date":
            out[c] = ""
        elif types.get(c) == "numeric":
            v = fn(_colvals_number(df, c))
            out[c] = f"{int(v):,}" if float(v).is_integer() else f"{v:,.2f}"
        else:
            out[c] = "" if labelled else label
            labelled = True
    return out


@_safe(default={})
def infer_column_types(df, cols=None):
    """Understand what kind of data each column holds: 'numeric', 'date' or
    'text'. A column counts as numeric/date when at least 80% of its non-empty
    cells parse as that type (commas and % signs are tolerated in numbers)."""
    import warnings
    types = {}
    if cols is None:
        cols = [] if df is None else [c for c in df.columns if c != "Date"]
    for c in cols:
        if df is None or c not in getattr(df, "columns", []):
            types[c] = "text"
            continue
        vals = _cell_text(df[c])
        if vals.empty:
            types[c] = "text"
            continue
        num = pd.to_numeric(vals.str.replace(",", "", regex=False)
                            .str.replace("%", "", regex=False), errors="coerce")
        if num.notna().mean() >= 0.8:
            types[c] = "numeric"
            continue
        with warnings.catch_warnings():
            warnings.simplefilter("ignore")
            dt = pd.to_datetime(vals, errors="coerce", dayfirst=True)
        types[c] = "date" if dt.notna().mean() >= 0.8 else "text"
    return types


def _looks_pct(df, col):
    """True when a column's cells mostly end in % (used to pre-suggest the
    'percentage columns' setting)."""
    try:
        vals = _cell_text(df[col])
        return (not vals.empty) and float(vals.str.endswith("%").mean()) >= 0.6
    except Exception:
        return False


def type_badges_html(types):
    """Small pills showing the detected data type of each column."""
    if not types:
        return ""
    pills = "".join(
        f'<span class="pill" style="background:{TYPE_COLOR[k]};font-size:.74rem;">'
        f'{TYPE_ICON[k]} {c} · {TYPE_LABEL[k]}</span>'
        for c, k in types.items())
    return ('<div style="margin:.15rem 0 .35rem 0;">'
            '<span style="font-size:.78rem;color:' + MUTED_FG + ';font-weight:600;">'
            'Detected data types:&nbsp;</span>' + pills + '</div>')


def _series_color(cfg, name, i):
    return (cfg.get("colors") or {}).get(name) or CHART_PALETTE[i % len(CHART_PALETTE)]


def _cat_color(cfg, label, i):
    """Colour for one category value (a pie slice, funnel stage, or an
    individual bar/point): the analyst's saved category colour first, then a
    matching column colour, then a palette default."""
    cc = cfg.get("cat_colors") or {}
    return cc.get(str(label)) or (cfg.get("colors") or {}).get(str(label)) \
        or CHART_PALETTE[i % len(CHART_PALETTE)]


def _tint_hex(color, f=0.85):
    """Blend a hex colour toward white by factor f (0..1) for soft cell fills."""
    try:
        c = str(color).lstrip("#")
        r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
        r = int(r + (255 - r) * f)
        g = int(g + (255 - g) * f)
        b = int(b + (255 - b) * f)
        return f"#{r:02X}{g:02X}{b:02X}"
    except Exception:
        return "#FFFFFF"


@_safe(default=None)
def custom_table_agg(t, df):
    """Aggregate a custom table for charting: pick the category column, the
    numeric value columns, and apply the configured roll-up (Sum / Average /
    Count / Latest entry). Works for a single day's rows or a Week/Month range
    (which includes a Date column). Returns {cat_label, cats, series, is_date}
    or None when there is nothing chartable."""
    if df is None or getattr(df, "empty", True):
        return None
    cfg = table_display(t)
    cols = [c for c in t.get("columns", []) if c in df.columns]
    if not cols:
        return None
    types = infer_column_types(df, cols)
    has_date = "Date" in df.columns
    work = df.copy()

    agg = cfg.get("agg", "sum")
    if agg == "latest" and has_date:      # only the most recent entry date's rows
        work = work[work["Date"] == work["Date"].max()]

    # category column: configured, or auto (first text col, else the Date)
    cat = cfg.get("category") or "(auto)"
    if cat == "(Date)" and has_date:
        cat_col, is_date = "Date", True
    elif cat in cols:
        cat_col, is_date = cat, types.get(cat) == "date"
    else:
        text_cols = [c for c in cols if types[c] == "text" and not _cell_text(work[c]).empty]
        date_cols = [c for c in cols if types[c] == "date"]
        if text_cols:
            cat_col, is_date = text_cols[0], False
        elif has_date:
            cat_col, is_date = "Date", True
        elif date_cols:
            cat_col, is_date = date_cols[0], True
        else:
            cat_col, is_date = None, False

    # numeric value columns: configured ones that still exist, else all numeric
    vals = [v for v in (cfg.get("values") or [])
            if v in cols and types.get(v) == "numeric"]
    if not vals:
        vals = [c for c in cols if types[c] == "numeric" and c != cat_col]
    if cat_col is None and not vals:
        return None

    for v in vals:
        work["__" + v] = pd.to_numeric(
            work[v].astype(str).str.strip().str.replace(",", "", regex=False)
            .str.replace("%", "", regex=False), errors="coerce")

    if cat_col is not None:
        key = work[cat_col].astype(str).str.strip()
        if is_date:
            import warnings
            with warnings.catch_warnings():
                warnings.simplefilter("ignore")
                parsed = pd.to_datetime(key, errors="coerce",
                                        dayfirst=(cat_col != "Date"))
            work["__cat"] = parsed.dt.strftime("%d %b %y")
            work["__ord"] = parsed
        else:
            work["__cat"] = key.where(key != "", "—")
        work = work[work["__cat"].notna()]
        if work.empty:
            return None
        g = work.groupby("__cat", sort=False)
        pctset = {v for v in vals if v in (cfg.get("pct_cols") or [])}
        series = {}
        if vals and agg != "count":
            for v in vals:
                fn = "mean" if (agg == "avg" or v in pctset) else "sum"
                series[v] = g["__" + v].agg(fn)
        if agg == "count" or not series:
            series = {"Count": g.size()}
        if is_date:                        # dates chronologically
            order = g["__ord"].min().sort_values().index.tolist()
        else:                              # categories by the first series, desc
            first = next(iter(series.values()))
            order = first.sort_values(ascending=False).index.tolist()
        cats = [c for c in order if str(c).strip()]
        return {"cat_label": ("Date" if cat_col == "Date" else cat_col),
                "cats": cats, "is_date": is_date,
                "pct": {n: (n in pctset) for n in series},
                "series": {n: [float(round(float(s.get(c, 0) or 0), 2)) for c in cats]
                           for n, s in series.items()}}

    # no category at all: one aggregate value per numeric column
    pctset = {v for v in vals if v in (cfg.get("pct_cols") or [])}
    out = {}
    for v in vals:
        s = work["__" + v].dropna()
        if agg == "avg" or v in pctset:
            out[v] = [round(float(s.mean()), 2)] if len(s) else [0.0]
        elif agg == "count":
            out[v] = [float(len(s))]
        else:
            out[v] = [float(s.sum())]
    return {"cat_label": "", "cats": ["All"], "is_date": False,
            "pct": {n: (n in pctset) for n in out}, "series": out}


@_safe(default=None)
def custom_table_agg2(t, df):
    """Two-level aggregation for the grouped (clustered) chart: rows group by
    the primary Category (outer axis level) and a Sub-category (inner level),
    values aggregate per pair with the same Sum/Average/Count/Latest and
    percentage rules as the single-level charts. Returns None when a second
    categorical column can't be resolved (callers fall back to plain columns)."""
    if df is None or getattr(df, "empty", True):
        return None
    cfg = table_display(t)
    cols = [c for c in t.get("columns", []) if c in df.columns]
    if len(cols) < 2:
        return None
    types = infer_column_types(df, cols)
    has_date = "Date" in df.columns
    work = df.copy()
    agg = cfg.get("agg", "sum")
    if agg == "latest" and has_date:
        work = work[work["Date"] == work["Date"].max()]

    text_cols = [c for c in cols if types[c] == "text"
                 and not _cell_text(work[c]).empty]
    cat = cfg.get("category") or "(auto)"
    if cat == "(Date)" and has_date:
        cat_col, is_date = "Date", True
    elif cat in cols:
        cat_col, is_date = cat, types.get(cat) == "date"
    elif text_cols:
        cat_col, is_date = text_cols[0], False
    elif has_date:
        cat_col, is_date = "Date", True
    else:
        return None

    sub = cfg.get("subcategory") or "(auto)"
    if sub in cols and sub != cat_col:
        sub_col = sub
    else:
        rest = [c for c in text_cols if c != cat_col]
        sub_col = rest[0] if rest else None
    if sub_col is None:
        return None

    vals = [v for v in (cfg.get("values") or []) if v in cols
            and types.get(v) == "numeric" and v not in (cat_col, sub_col)]
    if not vals:
        vals = [c for c in cols if types[c] == "numeric"
                and c not in (cat_col, sub_col)]
    pctset = {v for v in vals if v in (cfg.get("pct_cols") or [])}
    for v in vals:
        work["__" + v] = pd.to_numeric(
            work[v].astype(str).str.strip().str.replace(",", "", regex=False)
            .str.replace("%", "", regex=False), errors="coerce")

    key1 = work[cat_col].astype(str).str.strip()
    if is_date:
        import warnings
        with warnings.catch_warnings():
            warnings.simplefilter("ignore")
            parsed = pd.to_datetime(key1, errors="coerce",
                                    dayfirst=(cat_col != "Date"))
        work["__g"] = parsed.dt.strftime("%d %b %y")
        work["__gord"] = parsed
    else:
        work["__g"] = key1.where(key1 != "", "—")
    key2 = work[sub_col].astype(str).str.strip()
    work["__s"] = key2.where(key2 != "", "—")
    work = work[work["__g"].notna()]
    if work.empty:
        return None
    g = work.groupby(["__g", "__s"], sort=False)
    series = {}
    if vals and agg != "count":
        for v in vals:
            fn = "mean" if (agg == "avg" or v in pctset) else "sum"
            series[v] = g["__" + v].agg(fn)
    if agg == "count" or not series:
        series = {"Count": g.size()}
    first = next(iter(series.values()))
    if is_date:
        gorder = work.groupby("__g")["__gord"].min().sort_values().index.tolist()
    else:
        gorder = (first.groupby(level=0).sum()
                  .sort_values(ascending=False).index.tolist())
    sorder = (first.groupby(level=1).sum()
              .sort_values(ascending=False).index.tolist())
    pairs = [(gg, ss) for gg in gorder for ss in sorder if (gg, ss) in first.index]
    if not pairs:
        return None
    return {"cat_label": ("Date" if cat_col == "Date" else cat_col),
            "sub_label": sub_col, "groups": list(gorder),
            "pairs": pairs, "is_date": is_date,
            "pct": {n: (n in pctset) for n in series},
            "series": {n: [float(round(float(s.get(p, 0) or 0), 2)) for p in pairs]
                       for n, s in series.items()}}


@_safe(default=[])
def custom_table_kpis(t, df):
    """KPI-card values for a custom table: one aggregate per numeric column
    (following the table's roll-up setting) plus the number of entries."""
    if df is None or getattr(df, "empty", True):
        return []
    cfg = table_display(t)
    agg = cfg.get("agg", "sum")
    cols = [c for c in t.get("columns", []) if c in df.columns]
    if not cols:
        return []
    types = infer_column_types(df, cols)
    work = df.copy()
    if agg == "latest" and "Date" in work.columns:
        work = work[work["Date"] == work["Date"].max()]
    vals = [v for v in (cfg.get("values") or [])
            if v in cols and types.get(v) == "numeric"] \
        or [c for c in cols if types[c] == "numeric"]
    lbl = {"sum": "total", "avg": "average", "count": "count", "latest": "latest"}[agg]
    pctset = {v for v in vals if v in (cfg.get("pct_cols") or [])}
    items = []
    for v in vals:
        s = pd.to_numeric(work[v].astype(str).str.strip()
                          .str.replace(",", "", regex=False)
                          .str.replace("%", "", regex=False), errors="coerce").dropna()
        if not len(s):
            continue
        use_avg = agg == "avg" or v in pctset
        val = round(float(s.mean()), 2) if use_avg else \
            float(len(s)) if agg == "count" else float(s.sum())
        val = int(val) if float(val).is_integer() else val
        lbl_i = "average" if (v in pctset and agg in ("sum", "avg")) else lbl
        sfx = "%" if v in pctset else ""
        items.append((f"{v} ({lbl_i})", f"{val:,}{sfx}"))
    mask = pd.Series(False, index=work.index)
    for c in cols:
        mask = mask | _cell_text(work[c]).reindex(work.index).notna()
    items.append(("Entries", int(mask.sum())))
    return items


@_safe(default=None)
def custom_kpi_fig(items, title):
    """Big-screen KPI card grid for a custom table. Labels are annotations (not
    indicator titles) so presentation mode never strips them."""
    if not items:
        return None
    import math
    vals = []
    for lbl, v in items:
        try:
            vals.append((lbl, float(str(v).replace(",", ""))))
        except (TypeError, ValueError):
            continue
    if not vals:
        return None
    n = len(vals)
    ncol = min(3, n)
    nrow = math.ceil(n / ncol)
    fig = go.Figure()
    ann = []
    for i, (lbl, v) in enumerate(vals):
        r, c = divmod(i, ncol)
        x0, x1 = c / ncol + 0.03, (c + 1) / ncol - 0.03
        y1 = 1 - r / nrow - 0.06
        y0 = 1 - (r + 1) / nrow + 0.10
        fig.add_trace(go.Indicator(
            mode="number", value=v,
            number={"font": {"color": PRIMARY, "size": 42}},
            domain={"x": [x0, x1], "y": [max(0.0, y0), max(0.05, y1)]}))
        ann.append(dict(x=(x0 + x1) / 2, y=min(1.0, y1 + 0.05), xref="paper",
                        yref="paper", text=f"<b>{lbl}</b>", showarrow=False,
                        font=dict(size=13, color="#5E7373")))
    fig.update_layout(title=title, annotations=ann, paper_bgcolor="#FFFFFF",
                      margin=dict(l=10, r=10, t=48, b=10),
                      title_font=dict(color=PRIMARY), height=max(240, nrow * 150))
    return fig


@_safe(default=None)
def custom_table_chart_fig(t, df, title):
    """Chart a custom table according to its saved display mode. Returns None
    when the mode is Table / KPI cards, or when there is nothing chartable
    (callers then fall back to the plain table)."""
    cfg = table_display(t)
    mode = cfg.get("mode", "table")
    if mode in ("table", "kpi"):
        return None
    a = custom_table_agg(t, df)
    if not a or not a["cats"] or not a["series"]:
        return None
    cats = [str(c) for c in a["cats"]]
    names = list(a["series"].keys())
    if mode in ("pie", "donut"):
        pvals = a["series"][names[0]]
        pairs = [(c, v) for c, v in zip(cats, pvals) if v > 0]
        if not pairs:
            return None
        pcolors = [_cat_color(cfg, c, i) for i, (c, _) in enumerate(pairs)]
        ptitle = f"{title} — {names[0]}" if title else names[0]
        fig = go.Figure(go.Pie(
            labels=[p[0] for p in pairs], values=[p[1] for p in pairs],
            hole=0.45 if mode == "donut" else 0.0,
            marker=dict(colors=pcolors, line=dict(color="#FFFFFF", width=2)),
            textinfo="label+value", textfont=dict(size=13),
            texttemplate=("%{label}<br>%{value:,.1f}%"
                          if (a.get("pct") or {}).get(names[0]) else None)))
        fig.update_layout(
            title=ptitle, paper_bgcolor="#FFFFFF",
            font=dict(color=INK, size=13), title_font=dict(color=PRIMARY, size=16),
            margin=dict(l=14, r=14, t=64, b=24), height=380,
            legend=dict(orientation="h", yanchor="bottom", y=-0.14, x=0.5,
                        xanchor="center", font=dict(size=11)))
        return fig
    chart = mode
    if mode == "table+chart":
        chart = "line" if a["is_date"] else "column"
    if chart == "funnel":
        fig = go.Figure()
        for i, nm in enumerate(names):
            stage_cols = ([_cat_color(cfg, c, j) for j, c in enumerate(cats)]
                          if len(names) == 1 else _series_color(cfg, nm, i))
            fig.add_trace(go.Funnel(
                y=cats, x=a["series"][nm], name=nm,
                marker=dict(color=stage_cols),
                textinfo="value+percent initial", textfont=dict(size=13)))
        fig.update_layout(title=title, paper_bgcolor="#FFFFFF",
                          plot_bgcolor="#FFFFFF", font=dict(color=INK, size=13),
                          height=380, margin=dict(l=14, r=18, t=70, b=40),
                          title_font=dict(color=PRIMARY, size=16),
                          showlegend=len(names) > 1,
                          legend=dict(orientation="h", yanchor="bottom", y=1.02,
                                      x=1, xanchor="right", font=dict(size=12)))
        return fig
    if chart == "grouped":
        a2 = custom_table_agg2(t, df)
        if a2 and a2["pairs"] and a2["series"]:
            gnames = list(a2["series"].keys())
            gx = [str(g) for g, _ in a2["pairs"]]
            sx = [str(s) for _, s in a2["pairs"]]
            subs = list(dict.fromkeys(sx))
            by_sub = len(gnames) == 1 and bool(cfg.get("color_by_cat"))
            submap = {s: _cat_color(cfg, s, i) for i, s in enumerate(subs)}
            fig = go.Figure()
            for i, nm in enumerate(gnames):
                col = ([submap[s] for s in sx] if by_sub
                       else _series_color(cfg, nm, i))
                fig.add_bar(x=[gx, sx], y=a2["series"][nm], name=nm,
                            marker_color=col)
            fig.update_layout(title=title, barmode="group",
                              showlegend=len(gnames) > 1)
            fig = style_fig(fig)
            gpct = a2.get("pct") or {}
            for tr in fig.data:
                if tr.type == "bar":
                    tr.textposition = "inside"
                    tr.insidetextanchor = "middle"
                    tr.insidetextfont = dict(color="#FFFFFF", size=13)
                    if gpct.get(getattr(tr, "name", "")):
                        tr.texttemplate = "%{y:,.1f}%"
            if gnames and all(gpct.get(n) for n in gnames):
                fig.update_yaxes(ticksuffix="%")
            fig.update_xaxes(tickfont=dict(size=11))
            fig.update_layout(uniformtext=dict(minsize=8, mode="hide"))
            return fig
        chart = "column"          # no usable sub-category: plain columns
    horizontal = chart in ("bar", "stacked_bar")
    stacked = chart in ("stacked_bar", "stacked_column")
    by_cat = (len(names) == 1 and bool(cfg.get("color_by_cat"))
              and chart != "line")
    fig = go.Figure()
    for i, nm in enumerate(names):
        col = ([_cat_color(cfg, c, j) for j, c in enumerate(cats)]
               if by_cat else _series_color(cfg, nm, i))
        if chart == "line":
            fig.add_scatter(x=cats, y=a["series"][nm], name=nm, mode="lines+markers",
                            line=dict(color=col, width=3), marker=dict(size=7))
        elif chart == "scatter":
            fig.add_scatter(x=cats, y=a["series"][nm], name=nm, mode="markers",
                            marker=dict(size=12, color=col,
                                        line=dict(color="#FFFFFF", width=1)))
        elif horizontal:
            fig.add_bar(y=cats, x=a["series"][nm], name=nm, orientation="h",
                        marker_color=col)
        else:
            fig.add_bar(x=cats, y=a["series"][nm], name=nm, marker_color=col)
    fig.update_layout(title=title, barmode="stack" if stacked else "group",
                      showlegend=len(names) > 1)
    if horizontal:
        fig.update_yaxes(autorange="reversed")
    fig = style_fig(fig)
    pct = a.get("pct") or {}
    for tr in fig.data:                     # numbers always INSIDE the bars
        if tr.type == "bar":
            tr.textposition = "inside"
            tr.insidetextanchor = "middle"
            tr.insidetextfont = dict(color="#FFFFFF", size=14)
            if pct.get(getattr(tr, "name", "")):
                tr.texttemplate = ("%{x:,.1f}%"
                                   if getattr(tr, "orientation", None) == "h"
                                   else "%{y:,.1f}%")
        elif tr.type == "scatter" and pct.get(getattr(tr, "name", "")) \
                and "text" in (getattr(tr, "mode", "") or ""):
            tr.texttemplate = "%{y:,.1f}%"
    if names and all(pct.get(n) for n in names):
        if horizontal:
            fig.update_xaxes(ticksuffix="%")
        else:
            fig.update_yaxes(ticksuffix="%")
    fig.update_layout(uniformtext=dict(minsize=8, mode="hide"))
    return fig


def custom_df_for_display(t, df):
    """The frame to show for a custom table: formula columns filled in and the
    summary/totals row appended when configured. Used by the dashboards and the
    'show underlying data' panels so what's on screen matches the reports."""
    if df is None or getattr(df, "empty", True):
        return df
    out = apply_formulas(t, df)
    sr = summary_row(t, out)
    if sr:
        import pandas as _pd
        out = _pd.concat([out, _pd.DataFrame([sr])], ignore_index=True)
    return out


def style_custom_df(t, df):
    """A pandas Styler for st.dataframe applying the table's saved colours:
    row striping and per-column cell tints. Falls back to the plain frame."""
    cfg = table_display(t)
    tcolors = cfg.get("colors") or {}
    stripe = cfg.get("row_stripe") or ""
    tint = bool(cfg.get("col_tint"))
    if not stripe and not (tint and tcolors):
        return df
    try:
        css = pd.DataFrame("", index=df.index, columns=df.columns)
        if stripe:
            for i, ix in enumerate(df.index):
                if i % 2 == 1:
                    css.loc[ix, :] = f"background-color:{stripe}"
        if tint:
            for c, col in tcolors.items():
                if c in css.columns and col:
                    css[c] = f"background-color:{_tint_hex(col, 0.80)}"
        return df.style.apply(lambda _: css, axis=None)
    except Exception:
        return df


@_safe(default=None)
def custom_rows_fig(df, title, t=None):
    """Table figure for a custom table (drops a leading Date column if present).
    When the table definition `t` is given, its saved colours are applied to
    the column headers, cell tints and row stripes."""
    if df is None or getattr(df, "empty", True):
        return None
    show = df.drop(columns=[c for c in ["Date"] if c in df.columns]) if "Date" in df.columns else df
    cols = list(show.columns)
    if not cols:
        return None
    cfg = table_display(t) if t else {}
    tcolors = (cfg.get("colors") or {}) if cfg else {}
    stripe = (cfg.get("row_stripe") or "") if cfg else ""
    tint = bool(cfg.get("col_tint")) if cfg else False
    header_fill = [tcolors.get(c) or PRIMARY for c in cols]
    nrows = len(show)
    fill_cols = []
    for c in cols:
        if tint and tcolors.get(c):
            base = _tint_hex(tcolors[c], 0.80)
            alt = _tint_hex(tcolors[c], 0.62) if stripe else base
        else:
            base = "#FFFFFF"
            alt = stripe or "#FFFFFF"
        fill_cols.append([base if i % 2 == 0 else alt for i in range(nrows)])
    fig = go.Figure(go.Table(
        header=dict(values=[f"<b>{c}</b>" for c in cols],
                    fill_color=header_fill, font=dict(color="#FFFFFF", size=13),
                    align="left", height=32),
        cells=dict(values=[show[c].astype(str).tolist() for c in cols],
                   fill_color=fill_cols,
                   font=dict(color=INK, size=12), align="left", height=28)))
    fig.update_layout(title=title, paper_bgcolor="#FFFFFF",
                      margin=dict(l=8, r=8, t=44, b=8), title_font=dict(color=PRIMARY))
    return fig


@_safe(default=[])
def custom_table_tv_blocks(t, df, title):
    """Big-screen panels for one custom table according to its display mode:
    a chart, KPI cards, a coloured table, or chart + table."""
    if df is None or getattr(df, "empty", True):
        return []
    cfg = table_display(t)
    mode = cfg.get("mode", "table")
    blocks = []
    if mode == "kpi":
        kf = custom_kpi_fig(custom_table_kpis(t, df), title)
        if kf is not None:
            blocks.append((title, kf))
    elif mode != "table":
        cf = custom_table_chart_fig(t, df, title)
        if cf is not None:
            blocks.append((title, cf))
    if mode in ("table", "table+chart") or not blocks:
        tf = custom_rows_fig(df, title, t)
        if tf is not None:
            blocks.append((title if not blocks else f"{title} — data", tf))
    return blocks


def custom_table_block(t, df, suffix=""):
    """A (section-title, render-callable) block for one custom table, so it
    flows through the same ordered pipeline as the built-in chart sections
    (show_charts) — meaning the layout editor moves it on the dashboards, the
    TV slideshow and the reports alike. The callable renders the table in
    normal mode; its .tv_blocks attribute yields the big-screen panels."""
    if df is None or getattr(df, "empty", True):
        return None

    def _render():
        render_custom_table_normal(t, df, suffix)
    _render.tv_blocks = lambda: custom_table_tv_blocks(
        t, df, t["title"] + (suffix or ""))
    return (t["title"], _render)


def render_custom_table_normal(t, df, suffix=""):
    """Normal-mode rendering of one custom table: section header, then the
    saved display mode (chart / KPI cards / coloured table / chart + table),
    always with the underlying data reachable."""
    if df is None or getattr(df, "empty", True):
        return
    st.markdown(f'<div class="section">{t["title"]}{suffix}</div>',
                unsafe_allow_html=True)
    cfg = table_display(t)
    mode = cfg.get("mode", "table")
    shown = False
    if mode == "kpi":
        items = custom_table_kpis(t, df)
        if items:
            render_kpis([(f"📊 {l}", v) for l, v in items], False)
            shown = True
    elif mode != "table":
        fig = custom_table_chart_fig(t, df, "")
        if fig is not None:
            st.plotly_chart(fig, use_container_width=True, theme=None,
                            config={"displayModeBar": False})
            shown = True
            if mode == "table+chart":
                _disp = custom_df_for_display(t, df)
                st.dataframe(style_custom_df(t, _disp), use_container_width=True,
                             hide_index=True)
                return
    if not shown or mode == "table":
        _disp = custom_df_for_display(t, df)
        st.dataframe(style_custom_df(t, _disp), use_container_width=True,
                     hide_index=True)
    else:
        with st.expander("📄 Show the underlying data"):
            st.dataframe(style_custom_df(t, custom_df_for_display(t, df)),
                         use_container_width=True, hide_index=True)


def custom_pdf_chart(t, df):
    """Matplotlib chart image for a custom table's PDF report section,
    following its saved display mode. Returns (BytesIO, aspect) or None."""
    cfg = table_display(t)
    mode = cfg.get("mode", "table")
    if mode in ("table", "kpi"):
        return None
    a = custom_table_agg(t, df)
    if not a or not a["cats"] or not a["series"]:
        return None
    try:
        import io as _io2
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import numpy as np
        cats = [str(c) for c in a["cats"]]
        names = list(a["series"].keys())
        chart = mode
        if mode == "table+chart":
            chart = "line" if a["is_date"] else "column"
        _pcts = a.get("pct") or {}
        a2 = None
        if chart == "grouped":
            a2 = custom_table_agg2(t, df)
            if a2 and a2["pairs"] and a2["series"]:
                _pcts = a2.get("pct") or {}
            else:
                chart = "column"
        by_cat = len(names) == 1 and bool(cfg.get("color_by_cat"))

        def _mcolor(nm, i):
            if by_cat:
                return [_cat_color(cfg, c, j) for j, c in enumerate(cats)]
            return _series_color(cfg, nm, i)

        fw, fh = 7.0, 3.0
        fig, ax = plt.subplots(figsize=(fw, fh), dpi=150)
        if chart in ("pie", "donut"):
            pvals = a["series"][names[0]]
            pairs = [(c, v) for c, v in zip(cats, pvals) if v > 0]
            if not pairs:
                plt.close(fig)
                return None
            pcols = [_cat_color(cfg, c, i) for i, (c, _) in enumerate(pairs)]
            wprops = dict(width=0.45) if chart == "donut" else None
            ax.pie([p[1] for p in pairs], labels=[p[0] for p in pairs], colors=pcols,
                   autopct="%1.0f%%", textprops=dict(fontsize=8), wedgeprops=wprops)
            ax.set_title(names[0], fontsize=10, color=PRIMARY)
        elif chart == "funnel":
            fvals = a["series"][names[0]]
            order = sorted(range(len(cats)), key=lambda i: -fvals[i])
            fc = [cats[i] for i in order]
            fv = [fvals[i] for i in order]
            top = max(fv) if fv else 1
            lefts = [(top - v) / 2 for v in fv]
            fcols = [_cat_color(cfg, c, i) for i, c in enumerate(fc)]
            ypos = np.arange(len(fc))[::-1]
            ax.barh(ypos, fv, left=lefts, color=fcols, height=0.8)
            ax.set_yticks(ypos)
            ax.set_yticklabels(fc, fontsize=8)
            for yi, v, l in zip(ypos, fv, lefts):
                ax.text(l + v / 2, yi, f"{v:,.0f}", ha="center", va="center",
                        fontsize=8, color="#FFFFFF", fontweight="bold")
            ax.set_xticks([])
            ax.set_title(names[0], fontsize=10, color=PRIMARY)
        elif chart == "line":
            for i, nm in enumerate(names):
                ax.plot(cats, a["series"][nm], marker="o", label=nm,
                        color=_series_color(cfg, nm, i), linewidth=2)
            if len(names) > 1:
                ax.legend(fontsize=7, frameon=False)
            ax.grid(True, alpha=.25)
            fig.autofmt_xdate(rotation=30)
        elif chart == "scatter":
            for i, nm in enumerate(names):
                ax.scatter(cats, a["series"][nm], label=nm, s=45,
                           color=_mcolor(nm, i))
            if len(names) > 1:
                ax.legend(fontsize=7, frameon=False)
            ax.grid(True, alpha=.25)
            ax.tick_params(axis="x", rotation=30)
        elif chart == "grouped":
            def _fmt_v(v):
                try:
                    f = float(v)
                except (TypeError, ValueError):
                    return ""
                if f == 0:
                    return ""
                return f"{int(f):,}" if f.is_integer() else f"{f:,.1f}"
            pairs = a2["pairs"]
            gser = a2["series"]
            gnames = list(gser.keys())
            subs = list(dict.fromkeys(str(s) for _, s in pairs))
            submap = {s: _cat_color(cfg, s, i) for i, s in enumerate(subs)}
            plt.close(fig)
            fw = min(11.5, max(7.0, 1.2 + 0.6 * len(pairs)))
            fh = 3.4
            fig, ax = plt.subplots(figsize=(fw, fh), dpi=150)
            # slot positions with a visible gap between groups
            xs, prev_g, pos = [], None, -1.0
            gspan = {}
            for gg, ss in pairs:
                pos += 1.0 if gg == prev_g else (1.9 if prev_g is not None else 1.0)
                xs.append(pos)
                gspan.setdefault(gg, [pos, pos])[1] = pos
                prev_g = gg
            xs = np.array(xs)
            n = max(1, len(gnames))
            w = 0.8 / n
            vmax = max((max(v) for v in gser.values() if v), default=0) or 1
            single_sub = n == 1 and bool(cfg.get("color_by_cat"))
            for i, nm in enumerate(gnames):
                vals = gser[nm]
                colr = ([submap[str(s)] for _, s in pairs] if single_sub
                        else _series_color(cfg, nm, i))
                bx = xs + i * w
                ax.bar(bx, vals, width=w, label=nm, color=colr)
                for xi, vv in zip(bx, vals):
                    txt = _fmt_v(vv)
                    if txt and _pcts.get(nm):
                        txt += "%"
                    if not txt:
                        continue
                    if vv >= 0.15 * vmax:
                        ax.text(xi, vv / 2, txt, ha="center", va="center",
                                color="#FFFFFF", fontsize=7, fontweight="bold")
                    else:
                        ax.text(xi, vv + vmax * 0.015, txt, ha="center",
                                va="bottom", color=INK, fontsize=7,
                                fontweight="bold")
            mid = w * (n - 1) / 2
            ax.set_xticks(xs + mid)
            ax.set_xticklabels([str(s) for _, s in pairs], fontsize=7,
                               rotation=30, ha="right")
            xtrans = ax.get_xaxis_transform()
            glist = list(gspan.items())
            for gi, (gg, (x0, x1)) in enumerate(glist):
                ax.text((x0 + x1) / 2 + mid, -0.34, str(gg), transform=xtrans,
                        ha="center", va="top", fontsize=8.5, fontweight="bold",
                        color=PRIMARY)
                if gi:
                    ax.axvline((x0 - 0.95) + mid, color="#D8D8D8",
                               linewidth=0.8, linestyle=(0, (3, 3)))
            if len(gnames) > 1:
                ax.legend(fontsize=7, frameon=False)
        elif chart in ("bar", "stacked_bar"):
            def _fmt_v(v):
                try:
                    f = float(v)
                except (TypeError, ValueError):
                    return ""
                if f == 0:
                    return ""
                return f"{int(f):,}" if f.is_integer() else f"{f:,.1f}"
            ypos = np.arange(len(cats))
            if chart == "stacked_bar":
                totals = np.zeros(len(cats))
                for nm in names:
                    totals += np.array(a["series"][nm], dtype=float)
                smax = float(totals.max()) if len(totals) else 1.0
                left = np.zeros(len(cats))
                for i, nm in enumerate(names):
                    v = np.array(a["series"][nm], dtype=float)
                    ax.barh(ypos, v, left=left, label=nm, color=_mcolor(nm, i))
                    for yi, vv, ll in zip(ypos, v, left):
                        txt = _fmt_v(vv)
                        if txt and _pcts.get(nm):
                            txt += "%"
                        if txt and vv >= 0.06 * (smax or 1):
                            ax.text(ll + vv / 2, yi, txt, ha="center", va="center",
                                    color="#FFFFFF", fontsize=7, fontweight="bold")
                    left += v
                ax.set_yticks(ypos)
            else:
                xmax = max((max(a["series"][nm]) for nm in names
                            if a["series"][nm]), default=0) or 1
                h = 0.8 / max(1, len(names))
                for i, nm in enumerate(names):
                    vals = a["series"][nm]
                    ax.barh(ypos + i * h, vals, height=h, label=nm,
                            color=_mcolor(nm, i))
                    for yi, vv in zip(ypos + i * h, vals):
                        txt = _fmt_v(vv)
                        if txt and _pcts.get(nm):
                            txt += "%"
                        if not txt:
                            continue
                        if vv >= 0.15 * xmax:
                            ax.text(vv / 2, yi, txt, ha="center", va="center",
                                    color="#FFFFFF", fontsize=7, fontweight="bold")
                        else:
                            ax.text(vv + xmax * 0.012, yi, txt, ha="left",
                                    va="center", color=INK, fontsize=7,
                                    fontweight="bold")
                ax.set_yticks(ypos + h * (len(names) - 1) / 2)
            ax.set_yticklabels(cats, fontsize=8)
            ax.invert_yaxis()
            if len(names) > 1:
                ax.legend(fontsize=7, frameon=False)
        else:                              # column / stacked_column
            def _fmt_v(v):
                try:
                    f = float(v)
                except (TypeError, ValueError):
                    return ""
                if f == 0:
                    return ""
                return f"{int(f):,}" if f.is_integer() else f"{f:,.1f}"
            x = np.arange(len(cats))
            if chart == "stacked_column":
                totals = np.zeros(len(cats))
                for nm in names:
                    totals += np.array(a["series"][nm], dtype=float)
                smax = float(totals.max()) if len(totals) else 1.0
                bottom = np.zeros(len(cats))
                for i, nm in enumerate(names):
                    v = np.array(a["series"][nm], dtype=float)
                    ax.bar(x, v, bottom=bottom, label=nm, color=_mcolor(nm, i))
                    for xi, vv, bb in zip(x, v, bottom):
                        txt = _fmt_v(vv)
                        if txt and _pcts.get(nm):
                            txt += "%"
                        if txt and vv >= 0.06 * (smax or 1):
                            ax.text(xi, bb + vv / 2, txt, ha="center", va="center",
                                    color="#FFFFFF", fontsize=7, fontweight="bold")
                    bottom += v
                ax.set_xticks(x)
            else:
                ymax = max((max(a["series"][nm]) for nm in names
                            if a["series"][nm]), default=0) or 1
                w = 0.8 / max(1, len(names))
                for i, nm in enumerate(names):
                    vals = a["series"][nm]
                    ax.bar(x + i * w, vals, width=w, label=nm,
                           color=_mcolor(nm, i))
                    for xi, vv in zip(x + i * w, vals):
                        txt = _fmt_v(vv)
                        if txt and _pcts.get(nm):
                            txt += "%"
                        if not txt:
                            continue
                        if vv >= 0.15 * ymax:
                            ax.text(xi, vv / 2, txt, ha="center", va="center",
                                    color="#FFFFFF", fontsize=7, fontweight="bold")
                        else:
                            ax.text(xi, vv + ymax * 0.015, txt, ha="center",
                                    va="bottom", color=INK, fontsize=7,
                                    fontweight="bold")
                ax.set_xticks(x + w * (len(names) - 1) / 2)
            ax.set_xticklabels(cats, fontsize=7, rotation=30, ha="right")
            if len(names) > 1:
                ax.legend(fontsize=7, frameon=False)
        for sp in ax.spines.values():
            sp.set_visible(False)
        ax.tick_params(labelsize=7)
        fig.tight_layout()
        if chart == "grouped":
            fig.subplots_adjust(bottom=0.34)
        b = _io2.BytesIO()
        fig.savefig(b, format="png", bbox_inches="tight")
        plt.close(fig)
        b.seek(0)
        return (b, fh / fw)
    except Exception:
        log.exception("custom table PDF chart failed")
        return None


# Interactive chart preview with right-click recolouring (data entry only).
INTERACTIVE_CHART_HTML = """<!DOCTYPE html><html><head><meta charset="utf-8">
<style>
  html,body { margin:0; padding:0; background:transparent;
    font-family:'Source Sans Pro',system-ui,-apple-system,Segoe UI,Roboto,sans-serif; }
  #chart { width:100%; height:400px; }
  #cpick { display:none; position:fixed; z-index:99; background:#FFFFFF;
    border:1px solid #CFE6E6; border-radius:12px; padding:12px; width:180px;
    box-shadow:0 6px 18px rgba(6,52,58,.28); }
  #cplabel { font-weight:700; color:#06343A; margin-bottom:8px; font-size:13px;
    max-width:170px; overflow:hidden; text-overflow:ellipsis; white-space:nowrap; }
  #cpinput { width:100%; height:36px; border:1px solid #CFE6E6; border-radius:8px;
    cursor:pointer; background:#FFFFFF; padding:2px; }
  .cpbtn { flex:1; border:none; border-radius:8px; padding:7px 0; cursor:pointer;
    font-weight:700; font-size:12px; }
  #cpapply { background:#006868; color:#FFFFFF; }
  #cpcancel { background:#F0FAFA; color:#06343A; }
</style></head><body>
<div id="chart"></div>
<div id="cpick">
  <div id="cplabel"></div>
  <input type="color" id="cpinput" value="#006868">
  <div style="display:flex;gap:6px;margin-top:9px;">
    <button class="cpbtn" id="cpapply">Apply</button>
    <button class="cpbtn" id="cpcancel">Cancel</button>
  </div>
</div>
__PLOTLY_SCRIPT__
<script>
(function(){
  var D = __DATA__;
  var gd = document.getElementById('chart');
  var box = document.getElementById('cpick');
  Plotly.newPlot(gd, D.fig.data, D.fig.layout,
                 {responsive:true, displayModeBar:false});
  var hovered = null;
  gd.on('plotly_hover', function(ev){
    hovered = (ev.points && ev.points[0]) ? ev.points[0] : null;
  });
  gd.on('plotly_unhover', function(){
    setTimeout(function(){ hovered = null; }, 350);
  });
  gd.addEventListener('contextmenu', function(e){
    if (!hovered) return;
    e.preventDefault();
    var pt = hovered;
    var tr = (gd.data && gd.data[pt.curveNumber]) || {};
    var pieLike = (tr.type === 'pie' || tr.type === 'funnel');
    var horizontal = (tr.orientation === 'h');
    var label = (pt.label !== undefined && pt.label !== null)
        ? pt.label : (horizontal ? pt.y : pt.x);
    var series = tr.name || '';
    var target, disp;
    if (pieLike || D.ntraces === 1) {
      target = 'cat:' + label; disp = String(label);
    } else {
      target = 'col:' + series; disp = series + ' (whole series)';
    }
    document.getElementById('cplabel').textContent = '\\uD83C\\uDFA8 ' + disp;
    box.style.left = Math.min(e.clientX, window.innerWidth - 200) + 'px';
    box.style.top = Math.min(e.clientY, window.innerHeight - 150) + 'px';
    box.style.display = 'block';
    box.dataset.target = target;
  });
  document.getElementById('cpcancel').onclick = function(){
    box.style.display = 'none';
  };
  document.getElementById('cpapply').onclick = function(){
    try {
      var u = new URL(window.parent.location.href);
      u.searchParams.set('cchart', D.tid);
      u.searchParams.set('ctarget', box.dataset.target || '');
      u.searchParams.set('ccolor', document.getElementById('cpinput').value);
      window.parent.location.replace(u.toString());
    } catch (err) {}
  };
  document.addEventListener('click', function(e){
    if (!box.contains(e.target)) box.style.display = 'none';
  });
})();
</script></body></html>"""


@_safe(default=None)
def apply_chart_color_pick(tid, target, color):
    """Persist a right-click colour pick from the interactive chart preview.
    `target` is 'cat:<category label>' (one slice / stage / bar / point) or
    'col:<column name>' (a whole series). Saving a category colour on a bar,
    column or scatter chart also switches on colour-by-category so the pick is
    visible. Returns the saved key, or None when the pick is invalid."""
    t = get_custom_table(tid)
    color = str(color or "")
    target = str(target or "")
    if not t or not color.startswith("#") or ":" not in target:
        return None
    kind, key = target.split(":", 1)
    if not key:
        return None
    d = table_display(t)
    if kind == "cat":
        d["cat_colors"] = {**d.get("cat_colors", {}), key: color}
        if d.get("mode") not in ("pie", "donut", "funnel"):
            d["color_by_cat"] = True
    elif kind == "col":
        d["colors"] = {**d.get("colors", {}), key: color}
    else:
        return None
    update_custom_table(tid, display=d)
    return key


def render_interactive_chart_preview(t, df):
    """The data-entry chart preview: right-click any bar, slice, stage or
    point to recolour it on the spot (the pick saves through the page URL and
    applies everywhere the table is displayed). Falls back to a static chart
    when the interactive component isn't available. Returns False when there
    is nothing chartable."""
    fig = custom_table_chart_fig(t, df, "")
    if fig is None:
        return False
    try:
        from streamlit.components.v1 import html as _cihtml
        import json as _cijson
        ntraces = len([tr for tr in fig.data
                       if getattr(tr, "type", "") != "table"])
        payload = _cijson.dumps({
            "tid": t["id"], "ntraces": ntraces,
            "fig": _plain_arrays(_cijson.loads(fig.to_json()))})
        st.caption("💡 Right-click any bar, slice, stage or point on the "
                   "preview to change its color — the pick saves instantly "
                   "and applies on every dashboard and report.")
        _cihtml(INTERACTIVE_CHART_HTML.replace("__DATA__", payload).replace("__PLOTLY_SCRIPT__", plotly_script_tag()), height=440)
    except Exception:
        st.plotly_chart(fig, use_container_width=True, theme=None,
                        config={"displayModeBar": False})
    return True


# ── At-a-Glance trend lines over a period (week / month) ──
GLANCE_TREND = [
    ("current_inpatients", "Patients in hospital"),
    ("admitted", "New admissions"),
    ("discharged", "Discharged"),
    ("er_visits", "ER visits"),
    ("surgeries", "Surgeries"),
    ("births", "Births"),
    ("stillbirths", "Stillbirths"),
    ("deaths", "Mortality"),
    ("doctors", "Doctors on duty"),
    ("nurses", "Nurses on duty"),
    ("oxygen_pct", "Oxygen supply (%)"),
]

# (key, label, "sum"|"avg") for the period report summary
PERIOD_SUMMARY = [
    ("current_inpatients", "Patients in hospital (avg/day)", "avg"),
    ("admitted", "New admissions (total)", "sum"),
    ("discharged", "Discharged (total)", "sum"),
    ("er_visits", "ER visits (total)", "sum"),
    ("surgeries", "Surgeries (total)", "sum"),
    ("births", "Births (total)", "sum"),
    ("stillbirths", "Stillbirths (total)", "sum"),
    ("deaths", "Mortality (total)", "sum"),
    ("doctors", "Doctors on duty (avg/day)", "avg"),
    ("nurses", "Nurses on duty (avg/day)", "avg"),
    ("oxygen_pct", "Oxygen supply (avg %)", "avg"),
]


def trend_fig(daily, items, title="Daily trend"):
    """Line chart of one or more At-a-Glance metrics across the loaded period.
    `items` is a list of (column_key, label). Returns None when nothing to plot."""
    if daily is None or getattr(daily, "empty", True):
        return None
    d = daily.copy()
    d["entry_date"] = pd.to_datetime(d["entry_date"])
    d = d.sort_values("entry_date")
    palette = [PRIMARY, DANGER, WARN, TEAL2, OK_GREEN, ACCENT, INK]
    fig = go.Figure()
    plotted = 0
    for k, lbl in items:
        if k not in d.columns:
            continue
        fig.add_scatter(x=d["entry_date"], y=pd.to_numeric(d[k], errors="coerce"),
                        name=lbl, mode="lines+markers",
                        line=dict(color=palette[plotted % len(palette)], width=3),
                        marker=dict(size=7))
        plotted += 1
    if not plotted:
        return None
    fig.update_layout(title=title)
    return style_fig(fig)


init_db()
HOSPITAL_NAME = get_setting("hospital_name", "General Hospital")


# ──────────────────────────────────────────────
# HELPERS
# ──────────────────────────────────────────────
def week_bounds(any_day):
    monday = any_day - timedelta(days=any_day.weekday())
    return monday, monday + timedelta(days=6)


def month_bounds(any_day):
    first = any_day.replace(day=1)
    nxt = first.replace(year=first.year + 1, month=1) if first.month == 12 \
        else first.replace(month=first.month + 1)
    return first, nxt - timedelta(days=1)


def day_label(dstr):
    return datetime.fromisoformat(dstr).strftime("%a %d")


def style_fig(fig, h=380, hide_axis_titles=True):
    """Clean styling: black plot border, no gridlines, and bar values shown inside
    the bar when they fit, otherwise on top of the bar (never hidden)."""
    for tr in fig.data:
        if tr.type == "bar":
            horizontal = getattr(tr, "orientation", None) == "h"
            tr.text = tr.x if horizontal else tr.y
            tr.texttemplate = "%{x:,.0f}" if horizontal else "%{y:,.0f}"
            tr.textposition = "auto"            # inside if it fits, else on top of the bar
            tr.insidetextanchor = "middle"
            tr.textangle = 0
            tr.cliponaxis = False               # don't clip labels drawn above bars
            tr.insidetextfont = dict(color="#FFFFFF", size=16)
            tr.outsidetextfont = dict(color=INK, size=16)
        elif tr.type == "scatter":
            # also print the value at each point on line charts (skip when too dense)
            horizontal = getattr(tr, "orientation", None) == "h"
            seq = getattr(tr, "x", None) if horizontal else getattr(tr, "y", None)
            try:
                n = len(seq) if seq is not None else 0
            except TypeError:
                n = 0
            if 0 < n <= 16:
                mode = getattr(tr, "mode", None) or "lines+markers"
                if "text" not in mode:
                    tr.mode = mode + "+text"
                tr.texttemplate = "%{x:,.0f}" if horizontal else "%{y:,.0f}"
                tr.textposition = "top center"
                tr.cliponaxis = False
                tr.textfont = dict(size=14, color=INK)
    fig.update_layout(
        paper_bgcolor="#FFFFFF", plot_bgcolor="#FFFFFF",
        font=dict(color=INK, size=13), height=h,
        margin=dict(l=14, r=18, t=70, b=60),
        title=dict(x=0.01, xanchor="left", y=0.98, yanchor="top",
                   font=dict(size=16, color=PRIMARY)),
        legend=dict(orientation="h", yanchor="bottom", y=1.02, x=1, xanchor="right",
                    font=dict(size=12)),
        bargap=0.30, bargroupgap=0.12, autosize=True,
    )
    fig.update_xaxes(showline=True, linecolor="#000000", linewidth=1.4, mirror=True,
                     showgrid=False, zeroline=False, automargin=True,
                     tickfont=dict(size=12), title_text=("" if hide_axis_titles else None))
    fig.update_yaxes(showline=True, linecolor="#000000", linewidth=1.4, mirror=True,
                     showgrid=False, zeroline=False, automargin=True,
                     tickfont=dict(size=12), title_text=("" if hide_axis_titles else None))
    return fig


# CSS for big-screen presentation mode: hides app chrome and compacts spacing so
# the dashboard fills a TV / large display. It never hides overflow, so content can
# never be clipped or overlap — at worst it grows slightly, it never covers a title.
TV_CSS = """
<style>
    [data-testid="stSidebar"], [data-testid="collapsedControl"] { display:none !important; }
    header[data-testid="stHeader"], [data-testid="stToolbar"] { display:none !important; }
    .block-container { max-width:100% !important; min-height:100vh !important;
        padding:0.5rem 1.2rem 0.8rem 1.2rem !important; animation: tvfade .5s ease both; }
    @keyframes tvfade { from { opacity:0; } to { opacity:1; } }
    .big-title { font-size:clamp(1.4rem, 2.1vw, 3rem) !important; margin:0 !important; }
    .sub { display:none !important; }
    /* professional presentation header bar (full-width, branded, underlined) */
    .tvhead { display:flex; align-items:center; gap:clamp(10px,1.4vw,20px); width:100%;
        padding:2px 2px 12px 2px; border-bottom:2px solid #CFE6E6; margin:.1rem 0 .55rem 0; }
    .tvlogo { height:clamp(36px,4.6vw,70px); width:auto; border-radius:10px; flex:0 0 auto;
        box-shadow:0 1px 6px rgba(6,52,58,.16); }
    .tvname { font-weight:800; color:#006868; line-height:1.05;
        font-size:clamp(1.35rem,2.8vw,2.8rem);
        background:linear-gradient(90deg,#006868,#02A6A6);
        -webkit-background-clip:text; -webkit-text-fill-color:transparent; background-clip:text;
        white-space:nowrap; overflow:hidden; text-overflow:ellipsis; }
    .tvhead-meta { margin-left:auto; text-align:right; flex:0 0 auto; line-height:1.25;
        font-weight:700; color:#06343A; font-size:clamp(0.82rem,1.15vw,1.4rem); }
    .tvhead-upd { display:block; font-weight:600; color:#5E7373;
        font-size:clamp(0.7rem,0.9vw,1.05rem); }
    .tvinfo { font-size:clamp(0.82rem, 1.05vw, 1.3rem); font-weight:600; color:#5E7373;
        margin-top:3px; }
    /* compact, understated Exit button, right-aligned above the header */
    [data-testid="stColumn"]:last-child .stButton>button {
        background:#FFFFFF !important; color:#06343A !important;
        border:1px solid #CFE6E6 !important; border-radius:10px !important;
        font-weight:700 !important; padding:2px 10px !important; }
    [data-testid="stColumn"]:last-child .stButton>button:hover {
        border-color:#02A6A6 !important; color:#006868 !important; }
    .section { font-size:clamp(0.9rem, 1.2vw, 1.5rem) !important; font-weight:700 !important;
        margin:0.2rem 0 0.15rem 0 !important; padding-bottom:0.1rem !important; }
    [data-testid="stMetric"] { padding:8px 12px !important; text-align:center; }
    [data-testid="stMetricValue"] { font-size:clamp(1.15rem, 1.8vw, 2.8rem) !important; }
    [data-testid="stMetricLabel"] p { font-size:clamp(0.68rem, 0.9vw, 1.1rem) !important; }
    .pill { font-size:clamp(0.72rem, 0.95vw, 1.15rem) !important;
        padding:3px 12px !important; margin:1px !important; }
    [data-testid="stHorizontalBlock"] { gap:0.6rem !important; }
    @media (min-width: 1000px) {
        [data-testid="stColumn"], [data-testid="column"] {
            min-width:0 !important; flex:1 1 0 !important; }
    }
    /* each panel is a self-contained card: title sits above its own chart, with a
       border separating panels so a title can never look like it overlaps a chart */
    .tvcard { border:1px solid #CFE6E6; border-radius:12px; padding:6px 8px 2px 8px;
        background:#FFFFFF; margin-bottom:0.4rem; }
    [data-testid="stPlotlyChart"] { border:1px solid #E3F1F1; border-radius:10px;
        padding:4px 4px 0 4px; background:#FFFFFF; }
    .tvchart-title { font-size:clamp(0.85rem, 1.2vw, 1.5rem); font-weight:700;
        color:#006868; text-align:center; line-height:1.15; margin:0.1rem 0 0.15rem 0; }
    /* slideshow footer (normal flow, so it can never cover a chart) */
    .tvfoot { display:flex; align-items:center; gap:16px; padding:6px 4px 2px 4px;
        margin-top:0.3rem; border-top:1px solid #CFE6E6; }
    .tvprogress { flex:1; height:7px; background:#E3F1F1; border-radius:999px; overflow:hidden; }
    .tvbar { height:100%; width:0; border-radius:999px;
        background:linear-gradient(90deg,#006868,#02A6A6);
        animation-name: tvgrow; animation-timing-function:linear; animation-fill-mode:forwards; }
    @keyframes tvgrow { from { width:0; } to { width:100%; } }
    .tvdots { display:flex; gap:7px; }
    .tvdot { width:10px; height:10px; border-radius:50%; background:#CFE6E6; }
    .tvdot.on { background:#006868; transform:scale(1.18); }
    .tvmeta { font-size:clamp(0.7rem, 0.85vw, 1rem); color:#5E7373; white-space:nowrap; }
</style>
"""


def _tv_kpi_per_row(vw):
    return 6 if vw >= 1500 else 4 if vw >= 1000 else 3


def render_kpis(items, tv):
    """KPI cards. On the big screen the cards-per-row adapt to the viewport width;
    in normal mode they wrap responsively via CSS."""
    if tv:
        vw, _vh = _tv_viewport()
        per_row = _tv_kpi_per_row(vw)
    else:
        per_row = 4
    for i in range(0, len(items), per_row):
        cols = st.columns(per_row)
        for col, (label, value) in zip(cols, items[i:i + per_row]):
            col.metric(label, value)


def dept_status_fig(pairs):
    """One full-width bar per department, coloured by status with the status named
    inside it. Bar length is constant so a Closed or Limited department is shown
    just as clearly as an Operational one — the colour and label carry the meaning."""
    if not pairs:
        return None
    names = [p[0] for p in pairs]
    colors = [STATUS_COLOR.get(p[1], "#777") for p in pairs]
    fig = go.Figure(go.Bar(
        x=[1] * len(names), y=names, orientation="h", marker_color=colors,
        text=[p[1] for p in pairs], textposition="inside", insidetextanchor="middle",
        textfont=dict(color="#FFFFFF"), cliponaxis=False,
        hovertext=[f"{n}: {s}" for n, s in pairs], hoverinfo="text"))
    fig.update_layout(title="Department status", paper_bgcolor="#FFFFFF",
                      plot_bgcolor="#FFFFFF", font=dict(color=INK), showlegend=False,
                      margin=dict(l=12, r=14, t=48, b=24), title_font=dict(color=PRIMARY),
                      bargap=0.24)
    fig.update_xaxes(range=[0, 1], showticklabels=False, showgrid=False, zeroline=False,
                     showline=True, linecolor="#000000", linewidth=1.4, mirror=True)
    fig.update_yaxes(autorange="reversed", automargin=True, showgrid=False,
                     showline=True, linecolor="#000000", linewidth=1.4, mirror=True)
    return fig


def med_status_fig(pairs):
    """Compact availability grid: one small coloured circle per medication
    (green = Available, amber = Limited availability, red = Not available) laid
    out side by side to save space, with the name under each circle and a colour
    key. Scales to many medications by wrapping into rows."""
    if not pairs:
        return None
    import math
    n = len(pairs)
    cols = max(1, min(n, 6))
    rows = math.ceil(n / cols)
    xs, ys, names, cols_color, hov = [], [], [], [], []
    for i, (name, status) in enumerate(pairs):
        s = status if status in MED_STATUS_COLOR else "Not available"
        xs.append(i % cols)
        ys.append(-(i // cols))
        names.append(str(name))
        cols_color.append(MED_STATUS_COLOR[s])
        hov.append(f"{name}: {s}")
    fig = go.Figure()
    # the medication circles
    fig.add_scatter(
        x=xs, y=ys, mode="markers+text", text=names,
        textposition="bottom center", textfont=dict(size=11, color=INK),
        marker=dict(size=30, color=cols_color, line=dict(color="#FFFFFF", width=2)),
        hovertext=hov, hoverinfo="text", showlegend=False, cliponaxis=False)
    # colour key (always shows all three) via legend-only points
    for s in MED_STATUSES:
        fig.add_scatter(x=[None], y=[None], mode="markers", name=s,
                        marker=dict(size=13, color=MED_STATUS_COLOR[s]),
                        showlegend=True, hoverinfo="skip")
    fig.update_layout(
        title="Medication availability", paper_bgcolor="#FFFFFF", plot_bgcolor="#FFFFFF",
        font=dict(color=INK), title_font=dict(color=PRIMARY),
        margin=dict(l=8, r=8, t=54, b=48),
        legend=dict(orientation="h", yanchor="top", y=-0.02, x=0.5, xanchor="center",
                    font=dict(size=12)),
        height=max(180, rows * 104 + 96))
    pad = 0.6
    fig.update_xaxes(visible=False, showgrid=False, zeroline=False,
                     range=[-pad, (cols - 1) + pad])
    fig.update_yaxes(visible=False, showgrid=False, zeroline=False,
                     range=[-(rows - 1) - 0.78, 0.78])
    return fig


def tests_fig(avail, unavail):
    """One coloured bar per test — green available, red not available."""
    names = list(avail) + list(unavail)
    if not names:
        return None
    colors = [OK_GREEN] * len(avail) + [DANGER] * len(unavail)
    marks = ["✓ " + n for n in avail] + ["✕ " + n for n in unavail]
    fig = go.Figure(go.Bar(x=[1] * len(names), y=names, orientation="h", marker_color=colors,
                           text=marks, textposition="inside", insidetextanchor="start",
                           textfont=dict(color="#FFFFFF")))
    fig.update_layout(title="Tests available (green) / not available (red)",
                      paper_bgcolor="#FFFFFF", plot_bgcolor="#FFFFFF", font=dict(color=INK),
                      showlegend=False, margin=dict(l=12, r=14, t=48, b=24),
                      title_font=dict(color=PRIMARY))
    fig.update_xaxes(range=[0, 1], showticklabels=False, showgrid=False, zeroline=False,
                     showline=True, linecolor="#000000", linewidth=1.4, mirror=True)
    fig.update_yaxes(showticklabels=False, autorange="reversed", showgrid=False,
                     showline=True, linecolor="#000000", linewidth=1.4, mirror=True)
    return fig


# ── PERFORMANCE / HOSPITAL CONDITION ────────────────────────────────
# A field counts as RED (out / closed), YELLOW (limited / low) or GREEN
# (available / operational). The weighted average of these decides whether the
# hospital is operating at a Critical, Medium or Stable condition.
LOW_BLOOD = 5        # blood units: 1..below this = limited (yellow)


def health_summary(statuses, tests_avail, tests_unavail, med_statuses, blood_units, oxygen):
    """Classify every status field red/yellow/green and derive the overall
    operating condition. Returns counts, a 0–100 health score, the condition
    label/colour, and a per-category breakdown."""
    def tally(items, fn):
        r = y = g = 0
        for it in items:
            c = fn(it)
            r += c == "r"; y += c == "y"; g += c == "g"
        return r, y, g

    dr, dy, dg = tally(statuses, lambda s: "g" if s == "Operational"
                       else "y" if s == "Limited" else "r")
    tr, ty, tg = int(tests_unavail), 0, int(tests_avail)       # tests: avail=green, not=red
    mr, my, mg = tally(med_statuses, lambda s: "g" if s == "Available"
                       else "y" if s == "Limited availability" else "r")
    br, by, bg = tally(blood_units, lambda u: "r" if u <= 0
                       else "y" if u < LOW_BLOOD else "g")
    o_r, o_y, o_g = (1, 0, 0) if oxygen < 25 else (0, 1, 0) if oxygen < 50 else (0, 0, 1)

    cats = [("Departments", dr, dy, dg), ("Tests", tr, ty, tg),
            ("Medications", mr, my, mg), ("Blood bank", br, by, bg),
            ("Oxygen", o_r, o_y, o_g)]
    red = dr + tr + mr + br + o_r
    yellow = dy + ty + my + by + o_y
    green = dg + tg + mg + bg + o_g
    total = red + yellow + green
    score = round((green * 2 + yellow) / (2 * total) * 100) if total else 0
    if total == 0:
        condition, color = "No data", "#777777"
    elif score >= 75:
        condition, color = "Stable", OK_GREEN
    elif score >= 50:
        condition, color = "Medium", WARN
    else:
        condition, color = "Critical", DANGER
    return dict(red=red, yellow=yellow, green=green, total=total, score=score,
                condition=condition, color=color, cats=cats)


def performance_report_image(summary):
    """A matplotlib PNG that reproduces the on-screen Hospital Performance
    dashboard for the PDF and PPTX reports: the coloured condition banner
    ("Hospital condition: X · operational health N%"), the four KPI cards
    (Critical / Limited / Stable counts + Operational health %), the health
    gauge with its percentage, and the status-breakdown bars. Returns
    (BytesIO, aspect) or None."""
    if not summary or not summary.get("total"):
        return None
    import matplotlib.pyplot as plt
    from matplotlib.patches import FancyBboxPatch
    import numpy as _np
    import io as _io

    score = summary["score"]
    color = summary["color"]
    cond = summary["condition"]
    cats = summary["cats"]

    fig = plt.figure(figsize=(9.2, 5.4), dpi=150)
    # rows: banner (thin) | KPI cards | charts (tall)
    gs = fig.add_gridspec(3, 4, height_ratios=[0.5, 0.7, 2.4],
                          hspace=0.45, wspace=0.22,
                          left=0.04, right=0.97, top=0.97, bottom=0.06)

    # ── condition banner (full width) ──
    axban = fig.add_subplot(gs[0, :])
    axban.axis("off")
    axban.add_patch(FancyBboxPatch(
        (0, 0), 1, 1, boxstyle="round,pad=0.02,rounding_size=0.12",
        transform=axban.transAxes, facecolor=color, edgecolor="none"))
    axban.text(0.5, 0.52,
               f"Hospital condition: {cond}      ·      operational health {score}%",
               transform=axban.transAxes, ha="center", va="center",
               color="#FFFFFF", fontsize=14, fontweight="bold")

    # ── KPI cards ──
    kpis = [("Critical (out / closed)", summary["red"], DANGER),
            ("Limited / low", summary["yellow"], WARN),
            ("Stable (available)", summary["green"], OK_GREEN),
            ("Operational health", f"{score}%", PRIMARY)]
    for i, (label, val, col) in enumerate(kpis):
        axk = fig.add_subplot(gs[1, i])
        axk.axis("off")
        axk.add_patch(FancyBboxPatch(
            (0.02, 0.05), 0.96, 0.9,
            boxstyle="round,pad=0.02,rounding_size=0.10",
            transform=axk.transAxes, facecolor=LIGHT_BG,
            edgecolor=col, linewidth=1.4))
        axk.text(0.5, 0.62, str(val), transform=axk.transAxes, ha="center",
                 va="center", fontsize=19, fontweight="bold", color=col)
        axk.text(0.5, 0.22, label, transform=axk.transAxes, ha="center",
                 va="center", fontsize=7.5, color=INK, wrap=True)

    # ── gauge (semicircular) ──
    axg = fig.add_subplot(gs[2, 0:2], projection="polar")
    axg.set_theta_zero_location("W")
    axg.set_theta_direction(-1)
    axg.set_thetamin(0)
    axg.set_thetamax(180)
    for lo, hi, bc in [(0, 50, "#fde2e2"), (50, 75, "#fdf0db"), (75, 100, "#e3f5ea")]:
        axg.barh(1, width=_np.radians((hi - lo) * 1.8),
                 left=_np.radians(lo * 1.8), height=0.42, color=bc,
                 edgecolor="white", linewidth=1)
    ang = _np.radians(score * 1.8)
    axg.plot([ang, ang], [0, 1.15], color=INK, linewidth=2.6,
             solid_capstyle="round", zorder=5)
    axg.scatter([ang], [0], s=60, color=color, zorder=6)
    axg.set_ylim(0, 1.4)
    axg.set_yticks([])
    axg.set_xticks([])
    axg.spines["polar"].set_visible(False)
    axg.set_title("Operational health", color=PRIMARY, fontsize=11,
                  fontweight="bold", pad=6)
    axg.text(_np.radians(90), -0.32, f"{score}%", ha="center", va="center",
             fontsize=22, fontweight="bold", color=color)
    axg.text(_np.radians(90), -0.68, cond, ha="center", va="center",
             fontsize=11, color=INK)

    # ── status breakdown (stacked bars) ──
    axb = fig.add_subplot(gs[2, 2:4])
    names = [c[0] for c in cats]
    y = _np.arange(len(names))
    greens = _np.array([c[3] for c in cats], dtype=float)
    yellows = _np.array([c[2] for c in cats], dtype=float)
    reds = _np.array([c[1] for c in cats], dtype=float)
    axb.barh(y, greens, color=OK_GREEN, label="OK")
    axb.barh(y, yellows, left=greens, color=WARN, label="Limited")
    axb.barh(y, reds, left=greens + yellows, color=DANGER, label="Critical")
    for i, (g, ye, r) in enumerate(zip(greens, yellows, reds)):
        x = 0
        for v, tc in ((g, "#FFFFFF"), (ye, "#5A3A00"), (r, "#FFFFFF")):
            if v > 0:
                axb.text(x + v / 2, i, str(int(v)), ha="center", va="center",
                         color=tc, fontsize=8, fontweight="bold")
            x += v
    axb.set_yticks(y)
    axb.set_yticklabels(names, fontsize=9)
    axb.invert_yaxis()
    axb.set_title("Field status by category", color=PRIMARY, fontsize=11,
                  fontweight="bold")
    axb.legend(fontsize=7, frameon=False, ncol=3, loc="lower right")
    for sp in ("top", "right", "left"):
        axb.spines[sp].set_visible(False)
    axb.tick_params(labelsize=8)
    axb.set_xticks([])

    buf = _io.BytesIO()
    fig.savefig(buf, format="png", facecolor="white", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf, 5.4 / 9.2


def oxygen_gauge_image(pct, title="Oxygen supply level"):
    """A matplotlib semicircular gauge (0–100%) reproducing the on-dashboard
    oxygen-supply gauge for the reports: coloured zones (red < 25, amber 25–50,
    green ≥ 50), a needle, and the big percentage read-out. Returns
    (BytesIO, aspect)."""
    import matplotlib.pyplot as plt
    import numpy as _np
    import io as _io

    try:
        pct = int(round(float(pct)))
    except (TypeError, ValueError):
        pct = 0
    pct = max(0, min(100, pct))
    col = OK_GREEN if pct >= 50 else WARN if pct >= 25 else DANGER

    fig = plt.figure(figsize=(6.4, 3.9), dpi=150)
    axg = fig.add_subplot(111, projection="polar")
    axg.set_theta_zero_location("W")
    axg.set_theta_direction(-1)
    axg.set_thetamin(0)
    axg.set_thetamax(180)
    # coloured zones (deg = value * 1.8 to span the half-circle)
    for lo, hi, bc in [(0, 25, "#fde2e2"), (25, 50, "#fdf0db"), (50, 100, "#e3f5ea")]:
        axg.barh(1, width=_np.radians((hi - lo) * 1.8),
                 left=_np.radians(lo * 1.8), height=0.42, color=bc,
                 edgecolor="white", linewidth=1)
    ang = _np.radians(pct * 1.8)
    axg.plot([ang, ang], [0, 1.15], color=INK, linewidth=2.6,
             solid_capstyle="round", zorder=5)
    axg.scatter([ang], [0], s=60, color=col, zorder=6)
    axg.set_ylim(0, 1.4)
    axg.set_yticks([])
    axg.set_xticks([])
    axg.spines["polar"].set_visible(False)
    axg.set_title(title, color=PRIMARY, fontsize=12, fontweight="bold", pad=8)
    axg.text(_np.radians(90), -0.34, f"{pct}%", ha="center", va="center",
             fontsize=26, fontweight="bold", color=col)
    axg.text(_np.radians(90), -0.72,
             "Stable" if pct >= 50 else "Low" if pct >= 25 else "Critical",
             ha="center", va="center", fontsize=11, color=INK)

    buf = _io.BytesIO()
    fig.savefig(buf, format="png", facecolor="white", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf, 3.9 / 6.4


def _period_latest(df, datecol="entry_date"):
    """Rows of a period resource frame for its most recent day — the 'as of'
    snapshot the week / month dashboard shows. Returns the frame unchanged if it
    has no date column, or None/empty passed straight through."""
    if df is None or getattr(df, "empty", True) or datecol not in df.columns:
        return df
    return df[df[datecol] == df[datecol].max()]


def _first_col(df, *names):
    """First of `names` present in df.columns (period frames use raw DB column
    names like 'blood_type'; day frames use display names like 'Blood Type')."""
    cols = getattr(df, "columns", [])
    for n in names:
        if n in cols:
            return n
    return names[-1]


def snapshot_blood(df):
    """Latest [(blood_type, units:int)] rows, or []."""
    df = _period_latest(df)
    if df is None or getattr(df, "empty", True):
        return []
    bt, bu = _first_col(df, "blood_type", "Blood Type"), _first_col(df, "units", "Units")
    out = []
    for _, r in df.iterrows():
        nm = str(r.get(bt, "") or "").strip()
        if not nm:
            continue
        try:
            u = int(r.get(bu, 0) or 0)
        except (TypeError, ValueError):
            u = 0
        out.append((nm, u))
    return out


def snapshot_status(df, name_names, status_names):
    """Latest [(name, status)] rows for a name/status table (departments,
    medications), or []."""
    df = _period_latest(df)
    if df is None or getattr(df, "empty", True):
        return []
    nc, sc = _first_col(df, *name_names), _first_col(df, *status_names)
    out = []
    for _, r in df.iterrows():
        nm = str(r.get(nc, "") or "").strip()
        if nm:
            out.append((nm, str(r.get(sc, "") or "").strip()))
    return out


def snapshot_tests(df):
    """Latest [(test, available:bool)] rows, or []."""
    df = _period_latest(df)
    if df is None or getattr(df, "empty", True):
        return []
    nc, ac = _first_col(df, "name", "Test"), _first_col(df, "available", "Available")
    out = []
    for _, r in df.iterrows():
        nm = str(r.get(nc, "") or "").strip()
        if nm:
            out.append((nm, bool(r.get(ac))))
    return out


def snapshot_absent(df):
    """Latest [(specialist, specialty, expected_return)] rows, or []."""
    df = _period_latest(df)
    if df is None or getattr(df, "empty", True):
        return []
    nc = _first_col(df, "name", "Specialist")
    sp = _first_col(df, "specialty", "Specialty / Area")
    er = _first_col(df, "expected_return", "Expected return")
    out = []
    for _, r in df.iterrows():
        nm = str(r.get(nc, "") or "").strip()
        if nm:
            out.append((nm, str(r.get(sp, "") or "").strip(),
                        str(r.get(er, "") or "").strip() or "—"))
    return out


def performance_figs(summary):
    """The performance dashboard panels: an operational-health gauge and a
    red/yellow/green breakdown by category."""
    color, score = summary["color"], summary["score"]
    gz = go.Figure(go.Indicator(
        mode="gauge+number", value=score, number={"suffix": "%"},
        title={"text": "Operational health"},
        gauge={"axis": {"range": [0, 100]}, "bar": {"color": color},
               "steps": [{"range": [0, 50], "color": "#fde2e2"},
                         {"range": [50, 75], "color": "#fdf0db"},
                         {"range": [75, 100], "color": "#e3f5ea"}],
               "threshold": {"line": {"color": INK, "width": 3},
                             "thickness": 0.75, "value": score}}))
    gz.update_layout(paper_bgcolor="#FFFFFF", plot_bgcolor="#FFFFFF",
                     font=dict(color=INK), margin=dict(l=22, r=22, t=58, b=12),
                     title=dict(text="Operational health", x=0.5, xanchor="center",
                                font=dict(color=PRIMARY)))

    cats = summary["cats"]
    names = [c[0] for c in cats]
    fb = go.Figure()
    fb.add_bar(y=names, x=[c[3] for c in cats], name="Available / OK",
               orientation="h", marker_color=OK_GREEN)
    fb.add_bar(y=names, x=[c[2] for c in cats], name="Limited / low",
               orientation="h", marker_color=WARN)
    fb.add_bar(y=names, x=[c[1] for c in cats], name="Out / closed",
               orientation="h", marker_color=DANGER)
    fb.update_layout(barmode="stack",
                     title="Field status by category (green = ok, amber = limited, red = critical)")
    fb.update_yaxes(autorange="reversed")
    return [(f"Hospital Status — {summary['condition']}", gz),
            ("Status Breakdown", style_fig(fb, h=360))]


def perf_inputs_single(s, depts, meds, tests, blood):
    statuses = list(depts["Status"]) if not depts.empty else []
    ta = int(tests["Available"].sum()) if not tests.empty else 0
    tu = int((~tests["Available"]).sum()) if not tests.empty else 0
    med_statuses = list(meds["Status"]) if not meds.empty else []
    blood_units = [int(x) for x in blood["Units"]] if not blood.empty else []
    return statuses, ta, tu, med_statuses, blood_units, int(s.get("oxygen_pct", 0))


def perf_inputs_range(latest, depts, meds, tests, blood):
    statuses = []
    if not depts.empty:
        last_d = sorted(depts["entry_date"].unique())[-1]
        statuses = list(depts[depts["entry_date"] == last_d]["status"])
    ta = tu = 0
    if not tests.empty:
        last_t = tests["entry_date"].max()
        tt = tests[tests["entry_date"] == last_t]
        ta = int((tt["available"] == 1).sum()); tu = int((tt["available"] == 0).sum())
    med_statuses = []
    if not meds.empty:
        last_m = meds["entry_date"].max()
        med_statuses = list(meds[meds["entry_date"] == last_m]["status"])
    blood_units = []
    if not blood.empty:
        last_b = blood["entry_date"].max()
        blood_units = [int(x) for x in blood[blood["entry_date"] == last_b]["units"]]
    oxygen = int(latest["oxygen_pct"]) if latest is not None else 0
    return statuses, ta, tu, med_statuses, blood_units, oxygen


def render_performance_normal(summary, figs):
    """Top-of-page performance dashboard for the scrolling (non-presentation) view."""
    st.markdown('<div class="section">Hospital Performance</div>', unsafe_allow_html=True)
    st.markdown(
        f'<div style="background:{summary["color"]};color:#fff;border-radius:14px;'
        f'padding:14px 18px;text-align:center;margin-bottom:.6rem;font-weight:800;'
        f'font-size:clamp(1.1rem,2.4vw,1.9rem);box-shadow:0 2px 10px rgba(6,52,58,.12);">'
        f'Hospital condition: {summary["condition"]}'
        f'<span style="font-weight:600;font-size:.62em;">'
        f' &nbsp;·&nbsp; operational health {summary["score"]}%</span></div>',
        unsafe_allow_html=True)
    render_kpis([("🟥 Critical (out / closed)", summary["red"]),
                 ("🟨 Limited / low", summary["yellow"]),
                 ("🟩 Stable (available)", summary["green"]),
                 ("❤️ Operational health", f'{summary["score"]}%')], False)
    for title, fig in figs:
        _render_chart_card(title, fig)


TV_SLIDESHOW_HTML = """<!DOCTYPE html><html><head><meta charset="utf-8">
<style>
  * { box-sizing:border-box; }
  html,body { margin:0; padding:0; background:transparent;
    font-family:'Source Sans Pro',system-ui,-apple-system,Segoe UI,Roboto,sans-serif; }
  #wrap { position:relative; width:100%; height:100vh; display:flex; flex-direction:column; }
  #stage { position:relative; flex:1 1 auto; min-height:0; }
  .slide { position:absolute; inset:0; opacity:0; transition:opacity .65s ease;
    display:grid; gap:12px; padding:2px; pointer-events:none; }
  .slide.active { opacity:1; pointer-events:auto; }
  .cell { display:flex; flex-direction:column; min-height:0; min-width:0; background:#FFFFFF;
    border:1px solid #CFE6E6; border-radius:14px; padding:8px 8px 2px 8px;
    box-shadow:0 2px 10px rgba(6,52,58,.05); }
  .ctitle { font-weight:700; color:#006868; text-align:center; line-height:1.12; margin:0 0 3px 0; }
  .plot { flex:1 1 auto; min-height:0; width:100%; }
  #foot { display:flex; align-items:center; gap:16px; padding:8px 8px 4px 8px;
    border-top:1px solid #CFE6E6; flex:0 0 auto; }
  #bar { flex:1; height:8px; background:#E3F1F1; border-radius:999px; overflow:hidden; }
  #barfill { height:100%; width:0; border-radius:999px;
    background:linear-gradient(90deg,#006868,#02A6A6); }
  #dots { display:flex; gap:8px; }
  .dot { width:11px; height:11px; border-radius:50%; background:#CFE6E6; transition:transform .3s,background .3s; }
  .dot.on { background:#006868; transform:scale(1.25); }
  #meta { color:#5E7373; white-space:nowrap; }
  /* floating fullscreen control (auto-hides like PowerPoint) */
  #fsbtn { position:fixed; top:10px; right:12px; z-index:50; border:none; cursor:pointer;
    background:rgba(0,104,104,.92); color:#fff; font-size:18px; line-height:1;
    width:42px; height:42px; border-radius:10px; box-shadow:0 2px 8px rgba(6,52,58,.25);
    opacity:0; transition:opacity .3s; }
  #fsbtn.show { opacity:1; }
  /* start-presentation splash */
  #splash { position:fixed; inset:0; z-index:60; display:flex; flex-direction:column;
    align-items:center; justify-content:center; gap:14px; text-align:center;
    background:linear-gradient(135deg,#F0FAFA,#FFFFFF); }
  #splash h2 { margin:0; color:#006868; font-size:clamp(1.2rem,2.4vw,2rem); font-weight:800; }
  #splash p { margin:0; color:#5E7373; font-size:clamp(.85rem,1.3vw,1.1rem); }
  #startbtn { margin-top:6px; border:none; cursor:pointer; color:#fff; font-weight:700;
    font-size:clamp(1rem,1.6vw,1.3rem); padding:14px 26px; border-radius:12px;
    background:linear-gradient(90deg,#006868,#02A6A6); box-shadow:0 6px 18px rgba(0,104,104,.3); }
</style></head><body>
<button id="fsbtn" title="Full screen (Esc to exit)">&#x26F6;</button>
<div id="splash">
  <h2>Ready to present</h2>
  <p>Show only the dashboards on the whole screen.</p>
  <button id="startbtn">&#9654;&nbsp; Start full-screen presentation</button>
  <p style="opacity:.8;">Press <b>Esc</b> any time to exit.</p>
</div>
<div id="wrap">
  <div id="stage"></div>
  <div id="foot"><div id="bar"><div id="barfill"></div></div>
    <div id="dots"></div><div id="meta"></div></div>
</div>
__PLOTLY_SCRIPT__
<script>
(function(){
  var D = __DATA__;
  var seconds = D.seconds || 30, panels = D.panels || [];
  var stage = document.getElementById('stage'),
      dotsEl = document.getElementById('dots'),
      metaEl = document.getElementById('meta'),
      barfill = document.getElementById('barfill'),
      fsbtn = document.getElementById('fsbtn'),
      splash = document.getElementById('splash'),
      startbtn = document.getElementById('startbtn');
  var st = {cols:0, per:0, pages:1, idx:0, slides:[], timer:null};

  // ---- make the iframe fill from its top to the bottom of the viewport ----
  function fit(){
    try{
      var fe = window.frameElement;
      if(fe){
        var top = fe.getBoundingClientRect().top;
        var avail = Math.max(260, window.parent.innerHeight - top - 6);
        fe.style.height = avail + 'px'; fe.height = avail;
      }
    }catch(e){}
  }
  function dims(){
    return { w: document.documentElement.clientWidth || window.innerWidth,
             h: document.documentElement.clientHeight || window.innerHeight };
  }
  function calc(){
    var d = dims();
    var cols = d.w < 760 ? 1 : 2;
    var stageH = Math.max(200, d.h - 54);
    var rowsFit = Math.max(1, Math.floor(stageH / 300));
    var per = Math.max(1, Math.min(6, cols * rowsFit));
    var fs = Math.min(2.2, Math.max(1.0, d.h / 1000));
    return {cols:cols, per:per, fs:fs};
  }
  function startBar(){
    barfill.style.transition='none'; barfill.style.width='0%';
    void barfill.offsetWidth;
    barfill.style.transition='width '+seconds+'s linear'; barfill.style.width='100%';
  }
  function show(i){
    for(var s=0;s<st.slides.length;s++) st.slides[s].classList.toggle('active', s===i);
    var dl=dotsEl.children;
    for(var d=0;d<dl.length;d++) dl[d].classList.toggle('on', d===i);
    var s0=i*st.per+1, s1=Math.min(panels.length,(i+1)*st.per);
    metaEl.textContent='Showing '+s0+'\\u2013'+s1+' of '+panels.length+
      ' \\u00b7 slide '+(i+1)+'/'+st.pages+' \\u00b7 every '+seconds+'s';
    startBar();
    var pl=st.slides[i].querySelectorAll('.plot');
    for(var q=0;q<pl.length;q++){ try{Plotly.Plots.resize(pl[q]);}catch(e){} }
  }
  function timer(){ clearInterval(st.timer); if(st.pages>1) st.timer=setInterval(function(){ st.idx=(st.idx+1)%st.pages; show(st.idx); }, seconds*1000); }
  function build(){
    var c=calc(); st.cols=c.cols; st.per=c.per;
    var titlePx=Math.round(15*c.fs);
    stage.innerHTML=''; dotsEl.innerHTML=''; st.slides=[];
    st.pages=Math.max(1, Math.ceil(panels.length/c.per));
    if(st.idx>=st.pages) st.idx=0;
    for(var p=0;p<st.pages;p++){
      var slide=document.createElement('div'); slide.className='slide';
      var count=Math.min(c.per, panels.length-p*c.per);
      var rows=Math.max(1, Math.ceil(count/c.cols));
      slide.style.gridTemplateColumns='repeat('+c.cols+', 1fr)';
      slide.style.gridTemplateRows='repeat('+rows+', 1fr)';
      for(var k=p*c.per;k<Math.min(panels.length,(p+1)*c.per);k++){
        var cell=document.createElement('div'); cell.className='cell';
        var t=document.createElement('div'); t.className='ctitle';
        t.style.fontSize=titlePx+'px'; t.textContent=panels[k].title;
        var pd=document.createElement('div'); pd.className='plot';
        cell.appendChild(t); cell.appendChild(pd); slide.appendChild(cell);
        var lay=Object.assign({}, panels[k].fig.layout);
        lay.font=Object.assign({}, lay.font||{}, {size:Math.round(13*c.fs)});
        var data=(panels[k].fig.data||[]).map(function(tr){
          var t=Object.assign({}, tr);
          if(t.type==='bar'){
            t.insidetextfont=Object.assign({}, t.insidetextfont||{}, {size:Math.round(18*c.fs)});
            t.outsidetextfont=Object.assign({}, t.outsidetextfont||{}, {size:Math.round(18*c.fs)});
          } else if(t.type==='scatter' && (t.text || (t.texttemplate && String(t.mode||'').indexOf('text')>=0))){
            t.textfont=Object.assign({}, t.textfont||{}, {size:Math.round(15*c.fs)});
          }
          return t;
        });
        try{ Plotly.newPlot(pd, data, lay, {responsive:true, displayModeBar:false}); }catch(e){}
      }
      stage.appendChild(slide); st.slides.push(slide);
      var dot=document.createElement('div'); dot.className='dot'; dotsEl.appendChild(dot);
    }
    show(st.idx); timer();
  }
  function onResize(){
    fit(); var c=calc();
    if(c.cols!==st.cols || c.per!==st.per){ build(); }
    else { var all=document.querySelectorAll('.plot');
      for(var r=0;r<all.length;r++){ try{Plotly.Plots.resize(all[r]);}catch(e){} } }
  }

  // ---- fullscreen (true PowerPoint-style presentation) ----
  function isFs(){ try{ return !!(window.parent.document.fullscreenElement||window.parent.document.webkitFullscreenElement); }catch(e){ return false; } }
  function enterFs(){
    try{ var el=window.parent.document.documentElement;
      (el.requestFullscreen||el.webkitRequestFullscreen||function(){}).call(el); }catch(e){}
  }
  function exitFsApi(){
    try{ var doc=window.parent.document;
      (doc.exitFullscreen||doc.webkitExitFullscreen||function(){}).call(doc); }catch(e){}
  }
  function exitPresentation(){
    try{ var u=new URL(window.parent.location.href);
      ['tv','v','d','t0','vw','vh'].forEach(function(k){ u.searchParams.delete(k); });
      window.parent.location.replace(u.toString());
    }catch(e){}
  }
  var wasFs=false;
  function onFsChange(){
    var f=isFs();
    fsbtn.innerHTML = f ? '&#x2715;' : '&#x26F6;';
    fsbtn.title = f ? 'Exit full screen (Esc)' : 'Full screen';
    if(f){ wasFs=true; splash.style.display='none'; }
    if(!f && wasFs){ exitPresentation(); return; }   // Esc / exit -> leave presentation
    setTimeout(onResize, 60);
  }
  startbtn.addEventListener('click', enterFs);
  fsbtn.addEventListener('click', function(){ isFs() ? exitFsApi() : enterFs(); });
  if(isFs()){ splash.style.display='none'; wasFs=true; }

  // auto-hide the floating control, reveal on mouse move
  var hideT;
  function poke(){ fsbtn.classList.add('show'); clearTimeout(hideT);
    hideT=setTimeout(function(){ fsbtn.classList.remove('show'); }, 2500); }
  ['mousemove','click','keydown'].forEach(function(ev){ document.addEventListener(ev, poke); });
  try{ ['mousemove','click','keydown'].forEach(function(ev){ window.parent.document.addEventListener(ev, poke); }); }catch(e){}
  poke();

  fit(); build();
  var rzT;
  function deb(){ clearTimeout(rzT); rzT=setTimeout(onResize, 150); }
  window.addEventListener('resize', deb);
  try{ window.parent.addEventListener('resize', deb); }catch(e){}
  document.addEventListener('fullscreenchange', onFsChange);
  try{
    window.parent.document.addEventListener('fullscreenchange', onFsChange);
    window.parent.document.addEventListener('webkitfullscreenchange', onFsChange);
  }catch(e){}
})();
</script></body></html>"""


def _plain_arrays(obj):
    """Plotly serialises numeric arrays as base64 'typed arrays'
    ({'bdata': ..., 'dtype': ..., 'shape': ...}). The plotly.js loaded for the
    presentation slideshow can mis-read those, which made the big-screen charts
    show wrong/empty values. Decode them back into plain number lists so the
    presented data is exactly what was entered, on any plotly.js version."""
    if isinstance(obj, dict):
        if "bdata" in obj and "dtype" in obj:
            try:
                import base64
                import numpy as np
                raw = base64.b64decode(obj["bdata"])
                arr = np.frombuffer(raw, dtype=np.dtype(obj["dtype"]))
                shape = obj.get("shape")
                if shape:
                    if isinstance(shape, str):
                        shape = tuple(int(s) for s in shape.split(",") if s.strip())
                    arr = arr.reshape(shape)
                return arr.tolist()
            except Exception:
                return obj
        return {k: _plain_arrays(v) for k, v in obj.items()}
    if isinstance(obj, list):
        return [_plain_arrays(v) for v in obj]
    return obj


def _clear_fig_titles(fig, top=48):
    """Remove a figure's own title (and any gauge/Indicator title) so the single
    styled header above the chart is the only title — never two titles stacked.
    Also tightens the top margin so no empty band is left where the title was."""
    try:
        fig.update_layout(title_text="", margin=dict(t=top))
    except Exception:
        pass
    try:
        for tr in getattr(fig, "data", []):
            if getattr(tr, "type", "") == "indicator":
                tr.title = {"text": ""}
    except Exception:
        pass
    return fig


def _render_chart_card(title, fig, key=None):
    """One chart as a self-contained bordered card: the title is the card header and
    the chart sits directly beneath it inside the same border — never detached."""
    with st.container(border=True):
        if title:
            st.markdown(f'<div class="chart-card-title">{title}</div>', unsafe_allow_html=True)
        st.plotly_chart(_clear_fig_titles(fig, top=28), use_container_width=True,
                        theme=None, config={"displayModeBar": False}, key=key)


def render_tv_slideshow(blocks, seconds=30):
    """Render the whole rotating dashboard as one self-contained component.

    Every chart is drawn once with Plotly.js inside a single iframe; the slides
    then cross-fade in the browser on a timer. Because nothing on the Streamlit
    page re-runs to advance a slide, the screen never blanks or reloads."""
    blocks = [b for b in blocks if b is not None and b[1] is not None]
    if not blocks:
        return

    vw, vh = _tv_viewport()
    cols_n = 1 if vw < 760 else 2
    kpi_pr = _tv_kpi_per_row(vw)
    kpi_rows = (12 + kpi_pr - 1) // kpi_pr
    reserve_outside = 120 + kpi_rows * 64               # header + info line + KPI band
    total_iframe = max(320, vh - reserve_outside)
    stage_h = total_iframe - 48                         # leave room for the footer
    rows_fit = max(1, int(stage_h // 300))              # ~300px-tall charts
    per_page = max(1, min(6, cols_n * rows_fit))        # panels per slide — device adaptive
    fs = min(2.0, max(1.0, vh / 1000.0))                # scale text with screen

    import json
    panels = []
    for title, fig in blocks:
        try:
            _clear_fig_titles(fig, top=24)        # also clears gauge/Indicator titles
            fig.update_layout(
                autosize=True,
                margin=dict(l=16, r=14, t=24, b=44),
                font=dict(size=int(13 * fs)),
                legend=dict(orientation="h", yanchor="bottom", y=1.0, x=1,
                            xanchor="right", font=dict(size=int(10 * fs))))
            fig.update_xaxes(tickfont=dict(size=int(12 * fs)), automargin=True)
            fig.update_yaxes(tickfont=dict(size=int(12 * fs)), automargin=True)
            fig.update_traces(insidetextfont=dict(color="#FFFFFF", size=int(17 * fs)),
                              outsidetextfont=dict(color=INK, size=int(17 * fs)),
                              selector=dict(type="bar"))
            fig.update_traces(textfont=dict(color=INK, size=int(15 * fs)),
                              selector=dict(type="scatter"))
            try:
                fig.layout.height = None
            except Exception:
                pass
            panels.append({"title": title,
                           "fig": _plain_arrays(json.loads(fig.to_json()))})
        except Exception:
            pass

    if not panels:
        # Graceful fallback (e.g. Plotly serialization unavailable): show the first
        # slide's charts statically — still no reload.
        for title, fig in blocks[:per_page]:
            if title:
                st.markdown(f'<div class="tvchart-title">{title}</div>',
                            unsafe_allow_html=True)
            st.plotly_chart(fig, use_container_width=True, theme=None,
                            config={"displayModeBar": False})
        return

    payload = json.dumps({"panels": panels, "seconds": seconds})
    html = TV_SLIDESHOW_HTML.replace("__DATA__", payload).replace("__PLOTLY_SCRIPT__", plotly_script_tag())
    try:
        from streamlit.components.v1 import html as _html
        # initial height; the component then self-fits to fill the screen (incl. fullscreen)
        _html(html, height=total_iframe, scrolling=False)
    except Exception:
        for title, fig in blocks[:per_page]:
            st.plotly_chart(fig, use_container_width=True, theme=None,
                            config={"displayModeBar": False})


# Chart sections the analyst can reorder (everything that flows through show_charts,
# i.e. the cards after the Performance dashboard and the At-a-Glance KPIs).
REORDERABLE_SECTIONS = [
    "Patient Flow", "Patient Activity", "ER / Surgeries / ICU", "Beds & Occupancy",
    "Staffing", "Staff on Duty", "Ambulances", "Critical Supplies", "Blood Bank",
    "Medication Availability", "Department Status", "Tests",
    "Specialist Availability", "Mortality by Ward", "Mortality — ward detail",
    "Daily Trend",
]


def reorderable_sections():
    """All orderable section names: the built-in chart sections plus one entry
    per custom table (by title), so new tables always join the layout editor."""
    names = list(REORDERABLE_SECTIONS)
    for t in get_custom_tables():
        if t["title"] not in names:
            names.append(t["title"])
    return names


def get_section_order():
    """The saved chart-section order, padded with any sections not yet placed
    (so new built-in sections and newly added custom tables always appear,
    at the end)."""
    raw = get_setting("section_order", "")
    try:
        saved = json.loads(raw) if raw else []
    except Exception:
        saved = []
    allowed = reorderable_sections()
    known = [s for s in saved if s in allowed]
    rest = [s for s in allowed if s not in known]
    return known + rest


def set_section_order(order):
    allowed = set(reorderable_sections())
    set_setting("section_order",
                json.dumps([s for s in order if s in allowed]))


def _section_key(title):
    """Map a rendered block title back to its layout-section name. Custom
    table panels gain ' (this week)' / ' (this month)' / ' — data' suffixes
    on the roll-up dashboards and in Table + chart mode."""
    t = str(title)
    for suf in (" (this week)", " (this month)"):
        if t.endswith(suf):
            t = t[: -len(suf)]
            break
    if t.endswith(" — data"):
        t = t[: -len(" — data")]
    return t


def _apply_section_order(blocks):
    """Globally reorder the managed chart sections to match the layout editor,
    while unmanaged blocks (e.g. the Performance gauge) keep their anchor
    positions. Works across the whole list — a custom table can move above or
    below any built-in chart. Blocks from the same section (a custom table's
    chart + its data panel) keep their relative order (stable sort)."""
    order = get_section_order()
    rank = {t: i for i, t in enumerate(order)}
    managed_q = [b for b in blocks if b and _section_key(b[0]) in rank]
    managed_q.sort(key=lambda b: rank[_section_key(b[0])])   # stable
    it = iter(managed_q)
    out = []
    for b in blocks:
        if b and _section_key(b[0]) in rank:
            out.append(next(it))
        else:
            out.append(b)
    return out


def order_custom_tables(tables):
    """Custom tables sorted by the saved dashboard-layout order, so the
    scrolling dashboards and reports follow the layout editor too."""
    rank = {s: i for i, s in enumerate(get_section_order())}
    return sorted(tables, key=lambda t: rank.get(t["title"], len(rank)))


# Each report (PDF + PPTX) emits sections under its own headings, which don't
# always match the dashboard section names the analyst reorders. This maps a
# report heading to the dashboard layout-section(s) it represents so the reports
# can be re-sequenced to follow get_section_order() — the same order the
# on-dashboard ⬆️/⬇️ arrows produce. A heading that isn't listed falls back to
# its own text (this is how custom tables, keyed by their title, slot in).
REPORT_SECTION_ALIASES = {
    # ── daily PDF headings ──
    "Patients": ["Patient Flow", "Patient Activity", "ER / Surgeries / ICU"],
    "Capacity": ["Beds & Occupancy"],
    "Staff on Duty": ["Staffing", "Staff on Duty"],
    "Ambulances & Emergency": ["Ambulances"],
    "Critical Supplies": ["Critical Supplies"],
    "Blood Bank": ["Blood Bank"],
    "Department Status": ["Department Status"],
    "Medication Availability": ["Medication Availability"],
    "Medical Tests": ["Tests"],
    "Absent Specialists": ["Specialist Availability"],
    "Mortality Register": ["Mortality by Ward", "Mortality — ward detail"],
    # ── daily PPTX headings (a slide bundles several dashboard sections) ──
    "Resources & Capacity": ["Beds & Occupancy", "Staffing", "Staff on Duty",
                             "Ambulances"],
    "Departments & Medications": ["Department Status", "Medication Availability"],
    "Medical Tests & Blood Bank": ["Tests", "Blood Bank"],
    # ── period (week / month) headings ──
    "Trends": ["Daily Trend"],
    "Mortality by ward": ["Mortality by Ward", "Mortality — ward detail"],
    "Mortality by Ward": ["Mortality by Ward", "Mortality — ward detail"],
}


def report_section_rank(title):
    """Sort key for a report section so PDF/PPTX reports follow the analyst's
    saved dashboard order. Built-in headings resolve through the alias map (a
    slide that bundles several dashboard sections ranks by the earliest of
    them); a custom table ranks by its own title's saved position. Anything
    unplaced keeps a stable spot at the end of the reorderable band."""
    order = get_section_order()
    rank = {s: i for i, s in enumerate(order)}
    keys = REPORT_SECTION_ALIASES.get(title, [title])
    hits = [rank[k] for k in keys if k in rank]
    return min(hits) if hits else len(order) + 1


class _ReportBand:
    """Collects the reorderable middle sections of a report and hands them back
    in the analyst's saved dashboard order (stable within equal ranks, so the
    default layout is preserved exactly). Payloads are opaque: the PDF builders
    store lists of flowables, the PPTX builders store slide-building callables."""

    def __init__(self):
        self._items = []

    def add(self, title, payload):
        # Skip empties so absent sections don't leave gaps or stray headings.
        if payload:
            self._items.append(
                (report_section_rank(title), len(self._items), title, payload))

    def payloads(self):
        return [p for _r, _s, _t, p in sorted(self._items,
                                              key=lambda x: (x[0], x[1]))]

    def items(self):
        """(title, payload) pairs in saved dashboard order — lets the PPTX
        builders attach each section's note to the slide it just created."""
        return [(t, p) for _r, _s, t, p in sorted(self._items,
                                                  key=lambda x: (x[0], x[1]))]


def import_targets():
    """Every place imported data can land: the six built-in table sections, the
    scalar field groups (Patients, Capacity, …), and each custom table. Each
    entry is (label, kind, key, columns)."""
    out = [
        ("Departments", "table", "departments", ["Department", "Status"]),
        ("Medications", "table", "medications", ["Medication", "Status"]),
        ("Medical Tests", "table", "tests", ["Test", "Available"]),
        ("Blood Bank", "table", "blood_bank", ["Blood Type", "Units"]),
        ("Absent Specialists", "table", "absent_specialists",
         ["Specialist", "Specialty / Area", "Expected return"]),
        ("Mortality Register", "table", "mortality", list(_MORT_COLS)),
    ]
    for g in FIELD_GROUPS:
        labels = [lbl for k, lbl, grp in DAILY_FIELDS if grp == g]
        out.append((f"Daily figures — {g}", "scalars", g, labels))
    for t in get_custom_tables():
        out.append((f"Custom table — {t['title']}", "custom", t["id"],
                    list(t["columns"])))
    return out


def import_into_section(entry_date, kind, key, mapped, mode):
    """Write mapped rows into any data-entry section for one day, re-saving the
    whole day atomically so nothing else is disturbed. `kind` is 'table',
    'scalars' or 'custom'; `mode` is 'replace' or 'append' (scalars always
    replace the mapped fields). Returns a short summary string."""
    if kind == "custom":
        cols = get_custom_table(key)["columns"] if get_custom_table(key) else \
            list(mapped.columns)
        if mode == "append":
            ex = load_custom_rows(key, entry_date, cols)
            ex = ex[~(ex.apply(lambda r: all(str(x).strip() == "" for x in r),
                               axis=1))]
            mapped = pd.concat([ex, mapped], ignore_index=True)
        save_custom_rows(key, entry_date, mapped, cols)
        return f"{len(mapped):,} rows"

    # built-in sections: load the current day, swap the target, re-save
    scalars, depts, meds, tests, blood, absent = load_day(entry_date)
    mort = load_mortality(entry_date)
    numeric = {k: int((scalars or {}).get(k, 0) or 0) for k in FIELD_KEYS}
    notes = (scalars or {}).get("notes", "") or ""
    section_df = {"departments": depts, "medications": meds, "tests": tests,
                  "blood_bank": blood, "absent_specialists": absent,
                  "mortality": mort}

    if kind == "scalars":
        # map imported headers (already renamed to field labels) to field keys
        lbl_to_key = {lbl: k for k, lbl, grp in DAILY_FIELDS if grp == key}
        n = 0
        for lbl, kf in lbl_to_key.items():
            if lbl in mapped.columns:
                col = pd.to_numeric(mapped[lbl].astype(str).str.replace(
                    ",", "", regex=False).str.replace("%", "", regex=False),
                    errors="coerce").dropna()
                if len(col):
                    numeric[kf] = int(round(float(col.iloc[0])))
                    n += 1
        save_day(entry_date, numeric, notes, depts, meds, tests, blood, absent,
                 mort)
        return f"{n} field(s)"

    if kind == "table":
        if key == "tests" and "Available" in mapped.columns:
            mapped["Available"] = mapped["Available"].astype(str).str.strip()\
                .str.lower().isin(["yes", "true", "1", "available", "y"])
        if mode == "append":
            base = section_df[key]
            base = base[~(base.apply(
                lambda r: all(str(x).strip() == "" for x in r), axis=1))]
            mapped = pd.concat([base, mapped], ignore_index=True)
        section_df[key] = mapped
        save_day(entry_date, numeric, notes, section_df["departments"],
                 section_df["medications"], section_df["tests"],
                 section_df["blood_bank"], section_df["absent_specialists"],
                 section_df["mortality"])
        return f"{len(mapped):,} rows"
    return "nothing"


def render_day_downloads(entry_date, key_prefix="dl"):
    """CSV / PDF / PPTX download buttons for one day's saved data, loaded fresh
    from the database. Shared by the Data Entry form and the sidebar Reporting
    tab. Returns True if there was data to offer, False otherwise."""
    _s, _dp, _md, _ts, _bl, _ab = load_day(entry_date)
    _mort = load_mortality(entry_date)
    _customs = [(t["title"], load_custom_rows(t["id"], entry_date, t["columns"]), t)
                for t in order_custom_tables(get_custom_tables())]
    _has_c = any(not item[1].empty for item in _customs)
    if not (_s or not _dp.empty or not _md.empty or not _ts.empty or not _bl.empty
            or not _mort.empty or _has_c):
        st.caption(f"No saved data for {entry_date:%d %b %Y} yet — save the day "
                   "first, then download.")
        return False
    _obs = gather_general_notes("day", entry_date.isoformat())
    _nref = entry_date.isoformat()
    dl = st.columns(3)
    dl_button(dl[0], 
        "⬇️ CSV (data)",
        day_csv_bytes(entry_date, _s, _dp, _md, _ts, _bl, _mort, _customs),
        file_name=f"hospital_{entry_date:%Y%m%d}.csv", mime="text/csv",
        use_container_width=True, key=f"{key_prefix}_csv")
    dl_button(dl[1], 
        "📄 PDF (report)",
        build_day_pdf(entry_date, _s, _dp, _md, _ts, _bl, HOSPITAL_NAME, _mort,
                      _customs, _obs, note_scope="day", note_ref=_nref,
                      absent=_ab),
        file_name=f"hospital_report_{entry_date:%Y%m%d}.pdf",
        mime="application/pdf", use_container_width=True, key=f"{key_prefix}_pdf")
    if PPTX_OK:
        dl_button(dl[2], 
            "📊 PPTX (slides)",
            build_day_pptx(entry_date, _s, _dp, _md, _ts, _bl, HOSPITAL_NAME,
                           _mort, _customs, _obs, note_scope="day",
                           note_ref=_nref, absent=_ab),
            file_name=f"hospital_slides_{entry_date:%Y%m%d}.pptx",
            mime=PPTX_MIME, use_container_width=True, key=f"{key_prefix}_pptx")
    else:
        dl[2].caption("Add `python-pptx` for PPTX.")
    return True


def render_import_tool(entry_date, in_sidebar=False):
    """📥 Import data from a CSV or Excel file into any data-entry section. The
    upload's columns are shown with a preview; the analyst maps each target
    column to an imported column (auto-matched by name), previews, then imports
    the rows for the selected day. Analyst-only; never crashes the page. When
    `in_sidebar` is True it renders inline (the caller supplies the expander)."""
    import contextlib
    _wrap = (contextlib.nullcontext() if in_sidebar
             else st.expander("📥 Import data from CSV or Excel"))
    _kp = "sideimp" if in_sidebar else "import"
    with _wrap:
        st.caption("Upload a spreadsheet, match its columns to any section of "
                   "the system — the built-in tables, the daily figures, or your "
                   "custom tables — and import the rows for the selected day.")
        up = st.file_uploader("Choose a CSV or Excel file",
                              type=["csv", "xlsx", "xls", "xlsm"],
                              key=f"{_kp}_upload")
        if up is None:
            return
        df, err = parse_import_file(up.name, up.getvalue())
        if err:
            st.error(err)
            return
        st.success(f"Read **{len(df):,} rows** and **{len(df.columns)} columns** "
                   f"from {up.name}.")
        st.caption("Preview of the imported data (first 5 rows):")
        st.dataframe(df.head(5), use_container_width=True, hide_index=True)

        targets = import_targets()
        labels = [t[0] for t in targets]
        pick = st.selectbox("Import into which section?", labels,
                            key=f"{_kp}_target")
        tlabel, kind, key, tcols = targets[labels.index(pick)]
        src_cols = list(df.columns)
        auto = suggest_column_map(tcols, src_cols)
        if kind == "scalars":
            st.markdown("**Match each figure to an imported column** — the value "
                        "from the first data row is used:")
        else:
            st.markdown("**Map this section's columns to the imported columns:**")
        st.caption("“(leave blank)” skips a column. Matching names are selected "
                   "for you.")
        options = ["(leave blank)"] + src_cols
        colmap = {}
        _per = 1 if in_sidebar else 2
        for i in range(0, len(tcols), _per):
            row = st.columns(_per)
            for cc, tcol in zip(row, tcols[i:i + _per]):
                with cc:
                    default = auto.get(tcol, "(leave blank)")
                    idx = options.index(default) if default in options else 0
                    sel = st.selectbox(f"➡️ {tcol}", options, index=idx,
                                       key=f"{_kp}_map_{kind}_{key}_{tcol}")
                    if sel != "(leave blank)":
                        colmap[tcol] = sel
        if not colmap:
            st.warning("Map at least one column to import.")
            return

        mapped = pd.DataFrame({tcol: df[src].astype(str)
                               for tcol, src in colmap.items()})
        for c in tcols:
            if c not in mapped.columns:
                mapped[c] = ""
        mapped = mapped.reindex(columns=tcols)
        if kind == "scalars":
            st.caption("Value that will be saved (from the first data row):")
            _one = {c: (mapped[c].iloc[0] if len(mapped) else "") for c in colmap}
            st.dataframe(pd.DataFrame([_one]), use_container_width=True,
                         hide_index=True)
        else:
            st.caption(f"Preview of what will be imported into “{tlabel}” "
                       f"(first 5 of {len(mapped):,} rows):")
            st.dataframe(mapped.head(5), use_container_width=True,
                         hide_index=True)

        if kind == "scalars":
            mode = "replace"
            st.caption("Daily figures are single values, so this replaces the "
                       "mapped fields for the day.")
        else:
            mode_lbl = st.radio(
                f"How should this apply to {entry_date:%d %b %Y}?",
                ["Replace that section's rows", "Add to existing rows"],
                key=f"{_kp}_mode", horizontal=not in_sidebar)
            mode = "append" if mode_lbl.startswith("Add") else "replace"
        if st.button(f"📥 Import into “{tlabel}”", key=f"{_kp}_go",
                     use_container_width=True):
            try:
                summary = import_into_section(entry_date, kind, key, mapped, mode)
                st.session_state.auth_time = time.time()
                st.success(f"Imported {summary} into “{tlabel}” for "
                           f"{entry_date:%d %b %Y}. Scroll down to review, and "
                           "download a backup afterwards.")
                st.rerun()
            except Exception as e:
                log.exception("import failed")
                st.error(f"Could not import: {e}")


def render_section_note(scope, period_ref, section_title, uid=None):
    """A compact '📝 Notes' expander for one specific chart/table section.
    Stored in the same observations table, keyed per period AND per section so
    every chart and table can carry its own note. Analyst edits inline; public
    read-only. Kept lightweight so it sits cleanly beneath each card.

    `uid` (optional) only makes the Streamlit widget keys unique when the same
    section title appears more than once on a page; the stored note stays keyed
    by (scope, period::title) so it still loads and saves correctly."""
    if not section_title:
        return
    ref = f"{period_ref}::{section_title}"
    wkey = f"{scope}_{ref}" if uid is None else f"{scope}_{ref}_{uid}"
    obs = get_observation(scope, ref)
    _fmt0 = parse_note_fmt(obs.get("fmt"))
    has = bool(obs["content"].strip()) or bool(_fmt0.get("image"))
    label = "📝 Notes"
    if obs["content"].strip():
        preview = " ".join(obs["content"].split())
        label += " — " + (preview[:40] + "…" if len(preview) > 40 else preview)
    with st.expander(label, expanded=False):
        if st.session_state.get("authed"):
            # ── Formatting toolbar (like Word's Home tab), above the text ──
            st.caption("Formatting")
            _aligns = ["left", "center", "right"]
            r1 = st.columns([3, 2, 1.3])
            _align = r1[0].radio(
                "Align", _aligns, index=_aligns.index(_fmt0["align"]),
                horizontal=True, key=f"nf_align_{wkey}",
                format_func=lambda v: {"left": "⯇ Left", "center": "≡ Center",
                                       "right": "Right ⯈"}[v])
            _font = r1[1].selectbox("Font", list(NOTE_FONTS_CSS.keys()),
                                    index=list(NOTE_FONTS_CSS).index(_fmt0["font"]),
                                    key=f"nf_font_{wkey}")
            _size = r1[2].number_input("Size", 8, 48, int(_fmt0["size"]),
                                       key=f"nf_size_{wkey}")
            r2 = st.columns(4)
            _bold = r2[0].checkbox("Bold", value=_fmt0["bold"], key=f"nf_bold_{wkey}")
            _ital = r2[1].checkbox("Italic", value=_fmt0["italic"], key=f"nf_ital_{wkey}")
            _und = r2[2].checkbox("Underline", value=_fmt0["underline"],
                                  key=f"nf_und_{wkey}")
            _cc = r2[3].checkbox("Custom colour", value=bool(_fmt0["color"]),
                                 key=f"nf_cc_{wkey}")
            _color = ""
            if _cc:
                _color = st.color_picker("Text colour",
                                         value=_fmt0["color"] or "#0E2A2A",
                                         key=f"nf_color_{wkey}")
            _img = _fmt0.get("image", "")
            _up = st.file_uploader("Insert a picture (optional)",
                                   type=["png", "jpg", "jpeg"], key=f"nf_img_{wkey}")
            if _up is not None:
                _mime = ("image/png" if _up.name.lower().endswith("png")
                         else "image/jpeg")
                _img = (f"data:{_mime};base64,"
                        + base64.b64encode(_up.getvalue()).decode())
            if _img:
                ic = st.columns([3, 1])
                ic[0].image(_img, width=180)
                if ic[1].checkbox("Remove picture", key=f"nf_imgrm_{wkey}"):
                    _img = ""
            txt = st.text_area(
                f"Notes for “{section_title}” (shown publicly, included in every "
                "report):", value=obs["content"], height=110,
                key=f"secnote_{wkey}",
                placeholder="Context, caveats or follow-ups for this chart…")
            _pf = {"align": _align, "bold": _bold, "italic": _ital,
                   "underline": _und, "color": _color, "font": _font,
                   "size": int(_size), "image": _img}
            if txt.strip() or _img:
                st.caption("Preview")
                st.markdown(note_block_html(txt, parse_note_fmt(_pf), BODY_FG),
                            unsafe_allow_html=True)
            if st.button("💾 Save", key=f"secnote_save_{wkey}"):
                set_observation(scope, ref, txt, _pf)
                st.session_state.auth_time = time.time()
                st.success("Saved.")
                st.rerun()
        else:
            if has:
                st.markdown(note_block_html(obs["content"], _fmt0, BODY_FG),
                            unsafe_allow_html=True)
            else:
                st.caption("No notes yet.")


def render_move_controls(section_name, visible_sections, key_prefix="", uid=""):
    """A small ⬆️ / ⬇️ toolbar shown under a chart/table on the dashboard, only
    to a signed-in analyst (never in TV mode or to the public). Clicking moves
    the whole section up or down; the new order is saved and used everywhere.
    `uid` keeps the button keys unique if a section is drawn more than once."""
    if not st.session_state.get("authed"):
        return
    if section_name not in visible_sections:
        return
    pos = visible_sections.index(section_name)
    is_first = pos == 0
    is_last = pos == len(visible_sections) - 1
    c = st.columns([1, 1, 10])
    if c[0].button("⬆️", key=f"mv_up_{key_prefix}_{section_name}_{uid}",
                   help="Move this section up", disabled=is_first,
                   use_container_width=True):
        # find the visible neighbour above and swap past it in the saved order
        target = visible_sections[pos - 1]
        _swap_sections(section_name, target)
        st.session_state.auth_time = time.time()
        st.rerun()
    if c[1].button("⬇️", key=f"mv_down_{key_prefix}_{section_name}_{uid}",
                   help="Move this section down", disabled=is_last,
                   use_container_width=True):
        target = visible_sections[pos + 1]
        _swap_sections(section_name, target)
        st.session_state.auth_time = time.time()
        st.rerun()


def _swap_sections(a, b):
    """Swap the saved-order positions of two sections (used by the on-dashboard
    move buttons so a section jumps past its visible neighbour)."""
    order = get_section_order()
    if a in order and b in order:
        ia, ib = order.index(a), order.index(b)
        order[ia], order[ib] = order[ib], order[ia]
        set_section_order(order)


def show_charts(blocks, tv, seconds=30, per_page=4, note_scope=None,
                note_ref=None):
    """Normal mode: one chart per row under its section header, each with its
    own 📝 Notes expander when note_scope/note_ref are given. A signed-in
    analyst also sees ⬆️/⬇️ controls to reorder each section in place.
    Big-screen mode: an auto-rotating slideshow — `per_page` panels at a time,
    advancing every `seconds` and looping, sized to fill the screen cleanly."""
    blocks = [b for b in blocks if b is not None and b[1] is not None]
    blocks = _apply_section_order(blocks)
    if not tv:
        # the distinct, reorderable sections currently on screen, in view order
        visible = []
        for title, _ in blocks:
            sec = _section_key(title)
            if sec in reorderable_sections() and sec not in visible:
                visible.append(sec)
        # A section title (built-in or custom) can legitimately appear more than
        # once on a page — e.g. a custom table named like a built-in section.
        # Render each section's 📝 note and ⬆️/⬇️ controls only once, and give
        # every chart a unique key, so Streamlit never sees a duplicate key.
        seen_notes, seen_moves = set(), set()
        for i, (title, fig) in enumerate(blocks):
            if callable(fig):        # custom-table block renders itself in place
                fig()
            else:
                _render_chart_card(title, fig,
                                   key=f"chart_{note_scope or 'd'}_{i}")
            if note_scope and note_ref and title and title not in seen_notes:
                seen_notes.add(title)
                render_section_note(note_scope, note_ref, title, uid=i)
            sec = _section_key(title)
            if sec not in seen_moves:
                seen_moves.add(sec)
                render_move_controls(sec, visible,
                                     key_prefix=(note_scope or "day"), uid=i)
        return

    # Big-screen slideshow can't run per-widget callables; expand each custom
    # table into its own figure panels first, then hand plain figures on.
    expanded = []
    for title, fig in blocks:
        if callable(fig):
            expanded.extend(getattr(fig, "tv_blocks", lambda: [])())
        else:
            expanded.append((title, fig))
    render_tv_slideshow([b for b in expanded if b and b[1] is not None], seconds)


def render_tests_block(avail, unavail, suffix=""):
    """Available / not-available test pills, stacked (safe inside columns)."""
    st.markdown(f"**Available ({len(avail)}){suffix}**")
    if avail:
        st.markdown(" ".join(f'<span class="pill" style="background:{OK_GREEN}">✓ {n}</span>'
                             for n in avail), unsafe_allow_html=True)
    else:
        st.caption("None reported.")
    st.markdown(f"**Not available ({len(unavail)})**")
    if unavail:
        st.markdown(" ".join(f'<span class="pill" style="background:{DANGER}">✕ {n}</span>'
                             for n in unavail), unsafe_allow_html=True)
    else:
        st.caption("None.")


def _qp_date(key):
    """Read a YYYY-MM-DD value from the URL query params, defaulting to today."""
    try:
        return date.fromisoformat(st.query_params.get(key, ""))
    except Exception:
        return date.today()


def enter_tv():
    """Enter big-screen slideshow mode. State is stored in the URL so it survives
    the periodic auto-refresh that drives the slideshow."""
    view_sel = st.session_state.get("pub_view", "Day")
    if view_sel.startswith("Week"):
        code, key = "w", "wk_ref"
    elif view_sel.startswith("Month"):
        code, key = "m", "mo_ref"
    else:
        code, key = "s", "dash_day"
    sel = st.session_state.get(key, date.today())
    st.query_params["tv"] = "1"
    st.query_params["v"] = code
    try:
        st.query_params["d"] = sel.isoformat()
    except Exception:
        st.query_params["d"] = date.today().isoformat()
    st.query_params["t0"] = str(int(time.time()))   # slideshow start time
    st.rerun()


def exit_tv():
    for k in ("tv", "v", "d", "t0", "vw", "vh"):
        try:
            del st.query_params[k]
        except Exception:
            pass
    st.session_state.pop("tv_mode", None)
    st.rerun()


def _tv_viewport():
    """Actual browser viewport (px) reported by the client, with safe defaults."""
    def _int(key, default):
        try:
            return int(st.query_params.get(key, "0") or 0) or default
        except Exception:
            return default
    return _int("vw", 1280), _int("vh", 800)


def inject_tv_autosize():
    """Measure the real viewport on the client and store it in the URL, so the
    server can size the slideshow to whatever screen it's displayed on. Re-runs
    on window resize / entering fullscreen, then reloads to re-fit."""
    try:
        from streamlit.components.v1 import html as _html
        _html("""
        <script>
        (function(){
          function sync(){
            try{
              if (window.parent.document.fullscreenElement) return;  // never reload in fullscreen
              var w = window.parent.innerWidth, h = window.parent.innerHeight;
              var u = new URL(window.parent.location.href);
              var sw = parseInt(u.searchParams.get('vw')||'0');
              var sh = parseInt(u.searchParams.get('vh')||'0');
              if (Math.abs(sw-w) > 40 || Math.abs(sh-h) > 40) {
                u.searchParams.set('vw', w);
                u.searchParams.set('vh', h);
                window.parent.location.replace(u.toString());
              }
            } catch(e) {}
          }
          sync();
          if (!window.parent.__tvResize) {
            window.parent.__tvResize = 1;
            window.parent.addEventListener('resize', function(){
              clearTimeout(window.parent.__rz);
              window.parent.__rz = setTimeout(sync, 400);
            });
          }
        })();
        </script>
        """, height=0)
    except Exception:
        pass


def day_picker(state_key, title):
    """Calendar date picker + a Mon–Sun button strip (✓ marks days with data).
    Returns (selected_date, monday, sunday, set_of_dates_with_data)."""
    if state_key not in st.session_state:
        st.session_state[state_key] = date.today()
    picked = st.date_input(title, value=st.session_state[state_key])
    if picked != st.session_state[state_key]:
        st.session_state[state_key] = picked
        st.rerun()
    day = st.session_state[state_key]
    mon, sun = week_bounds(day)
    have = dates_with_data(mon, sun)
    st.caption(f"Week of {mon:%d %b} – {sun:%d %b %Y}  ·  ✓ = data saved for that day")
    cols = st.columns(7)
    for i in range(7):
        dy = mon + timedelta(days=i)
        mark = " ✓" if dy.isoformat() in have else ""
        selected = (dy == day)
        if cols[i].button(f"{dy:%a} {dy:%d}{mark}", key=f"{state_key}_b{i}",
                          type="primary" if selected else "secondary",
                          use_container_width=True):
            st.session_state[state_key] = dy
            st.rerun()
    return st.session_state[state_key], mon, sun, have


def dept_pills(pairs):
    st.markdown(" ".join(
        f'<span class="pill" style="background:{STATUS_COLOR.get(s,"#777")}">{n}: {s}</span>'
        for n, s in pairs), unsafe_allow_html=True)


def _absent_rows(absent_df):
    """Normalise an absent-specialists dataframe (from load_day or load_range) into a
    list of (name, area, return_date_or_None, overdue_bool), sorted by return date.
    Tolerates either the editor's column labels or the raw database column names."""
    out = []
    if absent_df is None or getattr(absent_df, "empty", True):
        return out
    cols = list(absent_df.columns)
    name_col = "Specialist" if "Specialist" in cols else "name"
    area_col = "Specialty / Area" if "Specialty / Area" in cols else "specialty"
    ret_col = "Expected return" if "Expected return" in cols else "expected_return"
    today = date.today()
    for _, r in absent_df.iterrows():
        nm = str(r.get(name_col, "") or "").strip()
        if not nm:
            continue
        area = str(r.get(area_col, "") or "").strip()
        rv = r.get(ret_col)
        rd = None
        try:
            missing = rv is None
            if not missing:
                try:
                    missing = bool(pd.isna(rv))
                except (TypeError, ValueError):
                    missing = False
            if not missing:
                if isinstance(rv, str):
                    rd = date.fromisoformat(rv[:10]) if rv.strip() else None
                elif isinstance(rv, datetime):
                    rd = rv.date()
                elif isinstance(rv, date):
                    rd = rv
                elif hasattr(rv, "date"):
                    rd = rv.date()
        except Exception:
            rd = None
        out.append((nm, area, rd, rd is not None and rd < today))
    out.sort(key=lambda x: (x[2] is None, x[2] or date.max))
    return out


def render_absent_block(absent_df):
    """Normal-mode 'Specialist Availability' panel: one card per absent specialist
    (name · specialty · expected-return date), highlighted red when overdue. Shows an
    all-clear message when no specialists are recorded as away."""
    rows = _absent_rows(absent_df)
    if not rows:
        st.markdown(
            f'<div style="background:{OK_GREEN};color:#fff;border-radius:10px;'
            f'padding:8px 14px;font-weight:700;display:inline-block;">'
            f'✓ All specialists available</div>', unsafe_allow_html=True)
        return
    cards = []
    for nm, area, rd, overdue in rows:
        bg = "#FDECEC" if overdue else CARD_BG
        bd = DANGER if overdue else GRID
        when = rd.strftime("%d %b %Y") if rd else "—"
        tag = (f'<span style="color:{DANGER};font-weight:700;"> · overdue</span>'
               if overdue else "")
        meta = area if area else "—"
        cards.append(
            f'<div style="border:1px solid {bd};background:{bg};border-radius:12px;'
            f'padding:8px 12px;min-width:180px;">'
            f'<div style="font-weight:700;color:{BODY_FG};">{nm}</div>'
            f'<div style="font-size:.82rem;color:{MUTED_FG};">{meta}</div>'
            f'<div style="font-size:.82rem;color:{BODY_FG};">Back: <b>{when}</b>{tag}</div>'
            f'</div>')
    st.markdown('<div style="display:flex;flex-wrap:wrap;gap:8px;">'
                + "".join(cards) + '</div>', unsafe_allow_html=True)


def absent_fig(absent_df):
    """Big-screen panel: absent specialists as a table (specialist, specialty, expected
    return, status) with overdue rows tinted red. Returns None when none are away."""
    rows = _absent_rows(absent_df)
    if not rows:
        return None
    names = [r[0] for r in rows]
    areas = [r[1] or "—" for r in rows]
    backs = [r[2].strftime("%d %b %Y") if r[2] else "—" for r in rows]
    status = ["Overdue" if r[3] else "Away" for r in rows]
    rowcols = ["#FDECEC" if r[3] else "#FFFFFF" for r in rows]
    fig = go.Figure(go.Table(
        header=dict(values=["<b>Specialist</b>", "<b>Specialty / Area</b>",
                            "<b>Expected return</b>", "<b>Status</b>"],
                    fill_color=PRIMARY, font=dict(color="#FFFFFF", size=14),
                    align="left", height=32),
        cells=dict(values=[names, areas, backs, status],
                   fill_color=[rowcols] * 4,
                   font=dict(color=INK, size=13), align="left", height=28)))
    fig.update_layout(title="Specialist availability", paper_bgcolor="#FFFFFF",
                      margin=dict(l=8, r=8, t=44, b=8), title_font=dict(color=PRIMARY))
    return fig


def day_dataframe(day, scalars, depts, meds, tests, blood, mortality=None, customs=None):
    """Flatten one day's full record into a tidy CSV-ready table:
    Date, Section, Item, Value, Detail."""
    ds = day.isoformat()
    s = scalars or {}
    label_map = {k: lbl for k, lbl, _ in DAILY_FIELDS}
    group_map = {k: grp for k, lbl, grp in DAILY_FIELDS}
    rows = []
    for k in FIELD_KEYS:
        rows.append({"Date": ds, "Section": group_map[k], "Item": label_map[k],
                     "Value": int(s.get(k, 0) or 0), "Detail": ""})
    if s.get("beds_total"):
        occ = (s["beds_total"] - s["beds_available"]) / s["beds_total"] * 100
        rows.append({"Date": ds, "Section": "Capacity", "Item": "Bed occupancy (%)",
                     "Value": round(occ, 1), "Detail": ""})
    for _, r in depts.iterrows():
        rows.append({"Date": ds, "Section": "Department", "Item": r["Department"],
                     "Value": r["Status"], "Detail": ""})
    for _, r in meds.iterrows():
        rows.append({"Date": ds, "Section": "Medication", "Item": r["Medication"],
                     "Value": str(r.get("Status", "") or ""), "Detail": ""})
    for _, r in tests.iterrows():
        rows.append({"Date": ds, "Section": "Test", "Item": r["Test"],
                     "Value": "Available" if r["Available"] else "Not available", "Detail": ""})
    for _, r in blood.iterrows():
        rows.append({"Date": ds, "Section": "Blood Bank", "Item": r["Blood Type"],
                     "Value": int(r["Units"] or 0), "Detail": "units"})
    if mortality is not None and not getattr(mortality, "empty", True):
        for _, r in mortality.iterrows():
            ward = str(r.get("Ward", r.get("ward", "")) or "").strip()
            if not ward:
                continue
            age = r.get("Age", r.get("age"))
            los = r.get("Length of stay (days)", r.get("los_days"))
            cond = str(r.get("Condition", r.get("condition", "")) or "").strip()
            cause = str(r.get("Cause of death", r.get("cause", "")) or "").strip()
            tod = _to_time_str(r.get("Time of death", r.get("time_of_death", "")))
            try:
                age_v = "" if age is None or (isinstance(age, float) and pd.isna(age)) else int(age)
            except (TypeError, ValueError):
                age_v = ""
            try:
                los_v = "" if los is None or (isinstance(los, float) and pd.isna(los)) else int(los)
            except (TypeError, ValueError):
                los_v = ""
            detail = []
            if tod:
                detail.append(f"time of death {tod}")
            if los_v != "":
                detail.append(f"length of stay {los_v} days")
            if cond:
                detail.append(f"condition: {cond}")
            if cause:
                detail.append(f"cause: {cause}")
            rows.append({"Date": ds, "Section": "Mortality", "Item": ward,
                         "Value": f"age {age_v}" if age_v != "" else "age —",
                         "Detail": "; ".join(detail)})
    if customs:
        for item in customs:
            ctitle, cdf = item[0], item[1]
            if cdf is None or getattr(cdf, "empty", True):
                continue
            ccols = [c for c in cdf.columns if c != "Date"]
            for _, r in cdf.iterrows():
                cells = [(c, str(r.get(c, "") or "").strip()) for c in ccols]
                if not any(v for _, v in cells):
                    continue
                first = cells[0][1] if cells else ""
                rows.append({"Date": ds, "Section": ctitle, "Item": first, "Value": "",
                             "Detail": "; ".join(f"{c}: {v}" for c, v in cells if v)})
    if s.get("notes"):
        rows.append({"Date": ds, "Section": "Notes", "Item": "Notes",
                     "Value": s.get("notes"), "Detail": ""})
    return pd.DataFrame(rows, columns=["Date", "Section", "Item", "Value", "Detail"])


def day_csv_bytes(day, scalars, depts, meds, tests, blood, mortality=None, customs=None):
    return day_dataframe(day, scalars, depts, meds, tests, blood,
                         mortality, customs).to_csv(index=False).encode("utf-8")


def _draw_icon(ax, key, cx, cy, s, color):
    """Draw a simple vector icon centered at (cx,cy), spanning ~s, for the PDF
    'At a Glance' tiles (font-independent so it always renders)."""
    import matplotlib.patches as mp
    WHITE = "#FFFFFF"

    def P(u, v):
        return (cx + u * s, cy + v * s)

    def circle(u, v, r, fill=True, col=None):
        col = col or color
        ax.add_patch(mp.Circle(P(u, v), r * s, facecolor=(col if fill else "none"),
                               edgecolor=col, linewidth=1.3, zorder=5))

    def rect(u, v, w, h, fill=True, rounded=0.0, col=None):
        col = col or color
        fc = col if fill else "none"
        if rounded > 0:
            ax.add_patch(mp.FancyBboxPatch((cx + u * s, cy + v * s), w * s, h * s,
                         boxstyle=f"round,pad=0,rounding_size={rounded*s}",
                         facecolor=fc, edgecolor=col, linewidth=1.3, zorder=5))
        else:
            ax.add_patch(mp.Rectangle((cx + u * s, cy + v * s), w * s, h * s,
                         facecolor=fc, edgecolor=col, linewidth=1.3, zorder=5))

    def line(u1, v1, u2, v2, w=1.8, col=None):
        col = col or color
        ax.plot([cx + u1 * s, cx + u2 * s], [cy + v1 * s, cy + v2 * s],
                color=col, linewidth=w, solid_capstyle="round", zorder=6)

    def poly(pts, fill=True, col=None):
        col = col or color
        ax.add_patch(mp.Polygon([P(u, v) for u, v in pts], closed=True,
                     facecolor=(col if fill else "none"), edgecolor=col, linewidth=1.3, zorder=5))

    def wedge(u, v, r, t1, t2, col=None):
        col = col or color
        ax.add_patch(mp.Wedge(P(u, v), r * s, t1, t2, facecolor=col, edgecolor=col, zorder=5))

    def plus(u, v, hw, hh, col):
        rect(u - hw, v - hh * 0.32, hw * 2, hh * 0.64, col=col)
        rect(u - hw * 0.32, v - hh, hw * 0.64, hh * 2, col=col)

    if key == "patient":
        circle(0, 0.28, 0.16)
        poly([(-0.28, -0.34), (0.28, -0.34), (0.18, 0.06), (-0.18, 0.06)])
    elif key == "bed":
        rect(-0.34, -0.16, 0.68, 0.16, rounded=0.04)
        rect(-0.40, -0.16, 0.06, 0.34)
        rect(-0.30, 0.00, 0.16, 0.09, rounded=0.03, col=WHITE)
        line(-0.36, -0.16, -0.36, -0.30); line(0.30, -0.16, 0.30, -0.30)
    elif key == "admit":
        line(0, 0.34, 0, -0.04, w=2.4)
        poly([(-0.12, -0.02), (0.12, -0.02), (0, -0.20)])
        line(-0.28, -0.28, 0.28, -0.28, w=2.4)
    elif key == "discharge":
        line(0, -0.30, 0, 0.12, w=2.4)
        poly([(-0.12, 0.12), (0.12, 0.12), (0, 0.30)])
        line(-0.28, -0.30, 0.28, -0.30, w=2.4)
    elif key == "er":
        plus(0, 0, 0.34, 0.34, color)
    elif key == "surgery":
        line(-0.30, -0.28, 0.10, 0.12, w=3.2)
        poly([(0.10, 0.12), (0.32, 0.32), (0.30, 0.12)])
    elif key == "birth":
        circle(0, 0.20, 0.15)
        poly([(-0.22, -0.34), (0.22, -0.34), (0.15, 0.04), (-0.15, 0.04)])
        circle(0, 0.20, 0.05, col=WHITE)
    elif key == "stillbirth":
        circle(0, 0.20, 0.15, fill=False)
        poly([(-0.22, -0.34), (0.22, -0.34), (0.15, 0.04), (-0.15, 0.04)], fill=False)
    elif key == "death":
        rect(-0.24, -0.34, 0.48, 0.42)
        wedge(0, 0.08, 0.24, 0, 180)
        plus(0, 0.04, 0.10, 0.14, WHITE)
    elif key == "doctor":
        circle(0, 0.28, 0.15)
        poly([(-0.28, -0.34), (0.28, -0.34), (0.18, 0.05), (-0.18, 0.05)])
        plus(0, -0.14, 0.09, 0.11, WHITE)
    elif key == "nurse":
        plus(0, 0.40, 0.06, 0.07, color)
        circle(0, 0.22, 0.15)
        poly([(-0.28, -0.34), (0.28, -0.34), (0.18, 0.00), (-0.18, 0.00)])
    elif key == "ambulance":
        rect(-0.40, -0.18, 0.58, 0.34, rounded=0.04)
        poly([(0.18, -0.18), (0.40, -0.18), (0.40, 0.04), (0.18, 0.12)])
        rect(0.20, -0.04, 0.14, 0.10, col=WHITE)
        plus(-0.18, 0.02, 0.10, 0.10, WHITE)
        circle(-0.20, -0.22, 0.08); circle(0.18, -0.22, 0.08)
    elif key == "oxygen":
        rect(-0.16, -0.34, 0.32, 0.58, rounded=0.10)
        rect(-0.06, 0.22, 0.12, 0.10)
        rect(-0.16, -0.12, 0.32, 0.10, col=WHITE)
    elif key == "calendar":
        rect(-0.32, -0.30, 0.64, 0.56, fill=False, rounded=0.04)
        rect(-0.32, 0.12, 0.64, 0.14)
        line(-0.18, 0.20, -0.18, 0.34, w=2.4); line(0.18, 0.20, 0.18, 0.34, w=2.4)
    else:
        circle(0, 0, 0.22, fill=False)


def _pdf_esc(v):
    """Escape user-entered text for ReportLab paragraphs: an '&' in a ward
    name like 'A&E', or a '<' in a note, must never be parsed as XML markup
    (unescaped, it crashes the whole report)."""
    import html as _h
    return _h.escape(str("" if v is None else v))


def build_glance_image(items, color=PRIMARY, ncol=4):
    """Render an 'At a Glance' grid of KPI tiles (icon + value + label) as a PNG.
    Returns (BytesIO, aspect_ratio)."""
    import io
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import matplotlib.patches as mp

    nrow = (len(items) + ncol - 1) // ncol
    fw = 7.4
    fh = 1.5 * nrow
    fig, ax = plt.subplots(figsize=(fw, fh))
    ax.set_xlim(0, ncol); ax.set_ylim(0, nrow); ax.axis("off")
    for i, (key, label, value) in enumerate(items):
        c = i % ncol
        r = i // ncol
        x = c
        y = nrow - 1 - r
        ax.add_patch(mp.FancyBboxPatch((x + 0.05, y + 0.05), 0.90, 0.90,
                     boxstyle="round,pad=0,rounding_size=0.10",
                     facecolor="#F0FAFA", edgecolor="#CFE6E6", linewidth=1.1))
        _draw_icon(ax, key, x + 0.5, y + 0.68, 0.28, color)
        ax.text(x + 0.5, y + 0.40, str(value), ha="center", va="center",
                fontsize=15, fontweight="bold", color=color)
        ax.text(x + 0.5, y + 0.16, label, ha="center", va="center",
                fontsize=7.3, color="#5E7373")
    fig.subplots_adjust(left=0.01, right=0.99, top=0.99, bottom=0.01)
    b = io.BytesIO()
    fig.savefig(b, format="png", dpi=150, bbox_inches="tight")
    plt.close(fig)
    b.seek(0)
    return (b, fh / fw)


def _pdf_header_table_style(extra=None, fontsize=9):
    """Shared report table style: a PRIMARY header row with white bold text and
    zebra-striped body rows. `extra` appends extra TableStyle commands."""
    from reportlab.lib import colors
    from reportlab.platypus import TableStyle
    cmds = [
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(PRIMARY)),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#BFE0E0")),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor(LIGHT_BG)]),
        ("FONTSIZE", (0, 0), (-1, -1), fontsize),
    ]
    if extra:
        cmds.extend(extra)
    return TableStyle(cmds)


@_maybe_cache(show_spinner=False, max_entries=48, ttl=900)
def _norm_notes(observations):
    """Normalise a builder's `observations` argument into a list of
    (label, text) pairs. Accepts the new list form (from gather_report_notes) or
    a plain string (back-compatible) which becomes a single 'Overview' block."""
    if not observations:
        return []
    if isinstance(observations, str):
        return [("Overview", observations)] if observations.strip() else []
    out = []
    for item in observations:
        if isinstance(item, (list, tuple)) and len(item) == 2:
            lbl, txt = item
            if str(txt).strip() or str(lbl).strip():
                out.append((str(lbl), str(txt)))
        elif str(item).strip():
            out.append(("Note", str(item)))
    return out


def build_day_pdf(day, scalars, depts, meds, tests, blood, hospital_name, mortality=None,
                  customs=None, observations=None, note_scope=None, note_ref=None,
                  absent=None):
    """A themed one-page-plus PDF report for a single day."""
    import io
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
                                    KeepTogether,
                                    TableStyle, HRFlowable, Image as RLImage)
    s = scalars or {}

    # ── optional charts (matplotlib → PNG) ──
    charts = {}
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        def _bar(pairs, title, horizontal=False, color=PRIMARY):
            labels = [p[0] for p in pairs]
            vals = [p[1] for p in pairs]
            fw, fh = (6.6, max(2.4, 0.45 * len(labels) + 1)) if horizontal else (6.6, 3.0)
            fig, ax = plt.subplots(figsize=(fw, fh))
            if horizontal:
                bars = ax.barh(labels, vals, color=color); ax.invert_yaxis()
            else:
                bars = ax.bar(labels, vals, color=color)
            ax.bar_label(bars, padding=3, fontsize=8)
            ax.set_title(title, fontsize=11, fontweight="bold", color=PRIMARY)
            ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
            ax.tick_params(labelsize=8)
            fig.tight_layout()
            b = io.BytesIO(); fig.savefig(b, format="png", dpi=150, bbox_inches="tight")
            plt.close(fig); b.seek(0)
            return (b, fh / fw)

        charts["patients"] = _bar(
            [("Admissions", int(s.get("admitted", 0))), ("Discharges", int(s.get("discharged", 0))),
             ("ER visits", int(s.get("er_visits", 0))), ("ICU", int(s.get("icu_patients", 0))),
             ("Surgeries", int(s.get("surgeries", 0))), ("Births", int(s.get("births", 0))),
             ("Stillbirths", int(s.get("stillbirths", 0))),
             ("Mortality", int(s.get("deaths", 0))), ("Referrals out", int(s.get("referrals_out", 0))),
             ("Referrals back", int(s.get("referrals_back", 0)))],
            "Patient activity", horizontal=True, color=PRIMARY)
        charts["staff"] = _bar(
            [("Doctors", int(s.get("doctors", 0))), ("Nurses", int(s.get("nurses", 0))),
             ("Support", int(s.get("support_staff", 0))), ("Specialists", int(s.get("specialists_on_call", 0)))],
            "Staff on duty", color=TEAL2)
        if not blood.empty and blood["Units"].sum() > 0:
            charts["blood"] = _bar([(r["Blood Type"], int(r["Units"])) for _, r in blood.iterrows()],
                                   "Blood bank units by type", color=DANGER)
    except Exception:
        charts = {}

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, leftMargin=18 * mm, rightMargin=18 * mm,
                            topMargin=16 * mm, bottomMargin=14 * mm,
                            title=f"Daily Report {day.isoformat()}")
    base = getSampleStyleSheet()
    title_style = ParagraphStyle("t", parent=base["Title"], fontSize=20,
                                 textColor=colors.HexColor(INK), spaceAfter=2, alignment=TA_CENTER)
    sub_style = ParagraphStyle("s", parent=base["Normal"], fontSize=10.5, textColor=colors.grey,
                               alignment=TA_CENTER, spaceAfter=10)
    h_style = ParagraphStyle("h", parent=base["Heading2"], fontSize=12.5,
                             textColor=colors.HexColor(PRIMARY), spaceBefore=12, spaceAfter=4)
    body = ParagraphStyle("b", parent=base["Normal"], fontSize=9.5, leading=13)
    small = ParagraphStyle("sm", parent=base["Normal"], fontSize=8.5, textColor=colors.grey)

    def hr():
        return HRFlowable(width="100%", thickness=0.6, color=colors.HexColor("#BFE0E0"),
                          spaceBefore=2, spaceAfter=8)

    def section(t):
        return [Paragraph(_pdf_esc(t), h_style), hr()]

    def _note_flow(title):
        """A section's note block — its text (with alignment, font, size, colour
        and style) plus any inserted image, or [] if none."""
        if not (note_scope and note_ref):
            return []
        txt, nf = section_note_full(note_scope, note_ref, title)
        if not txt and not nf.get("image"):
            return []
        import html as _h2
        note_sty = ParagraphStyle(
            "secnote", parent=body,
            fontSize=max(7.0, min(20.0, float(nf["size"]) * 0.62)),
            fontName=NOTE_FONTS_PDF.get(nf["font"], "Helvetica"),
            textColor=colors.HexColor(nf["color"] or INK),
            backColor=colors.HexColor(LIGHT_BG), borderPadding=6, leftIndent=2,
            spaceBefore=3, spaceAfter=2,
            alignment=NOTE_ALIGN_PDF.get(nf["align"], 0))
        _t = _h2.escape(txt).replace("\n", "<br/>")
        if nf["bold"]:
            _t = f"<b>{_t}</b>"
        if nf["italic"]:
            _t = f"<i>{_t}</i>"
        if nf["underline"]:
            _t = f"<u>{_t}</u>"
        flow = [Spacer(1, 3)]
        if txt:
            flow.append(Paragraph("📝 " + _t, note_sty))
        if nf.get("image"):
            _im = _note_pdf_image(nf["image"])
            if _im is not None:
                _im.hAlign = {"left": "LEFT", "center": "CENTER",
                              "right": "RIGHT"}.get(nf["align"], "LEFT")
                flow += [Spacer(1, 3), _im]
        return flow

    def keep_section(title, *flowables):
        """A heading, its divider, and the first following flowable held
        together so a section title is never orphaned at a page bottom while
        its chart or table slips to the next page. The section's own note (if
        any) is appended directly beneath it."""
        flat = []
        for f in flowables:
            flat.extend(f if isinstance(f, list) else [f])
        flat += _note_flow(title)
        head = [Paragraph(_pdf_esc(title), h_style), hr()]
        if not flat:
            return head
        return [KeepTogether(head + [flat[0]])] + flat[1:]

    def kv_table(rows, widths=(95 * mm, 65 * mm)):
        data = [[Paragraph(f"<b>{k}</b>", body), Paragraph(_pdf_esc(v), body)] for k, v in rows]
        t = Table(data, colWidths=widths)
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (0, -1), colors.HexColor("#D8F0F0")),
            ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#BFE0E0")),
            ("ROWBACKGROUNDS", (0, 0), (-1, -1), [colors.white, colors.HexColor(LIGHT_BG)]),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("LEFTPADDING", (0, 0), (-1, -1), 6), ("RIGHTPADDING", (0, 0), (-1, -1), 6),
            ("TOPPADDING", (0, 0), (-1, -1), 3), ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
        ]))
        return t

    def img(item, width=170 * mm):
        b, ar = item
        return RLImage(b, width=width, height=width * ar)

    def grp_rows(group):
        return [(lbl, int(s.get(k, 0))) for k, lbl, g in DAILY_FIELDS if g == group]

    story = []
    story.append(Paragraph(_pdf_esc(hospital_name), title_style))
    story.append(Paragraph(f"Daily Report &nbsp;•&nbsp; {day:%A, %d %B %Y}", sub_style))
    story.append(Paragraph(f"Generated {now_local():%d %b %Y %H:%M}", small))
    story.append(Spacer(1, 6))

    # ── At a Glance (KPI tiles with icons) ──
    occ_txt = "—"
    if s.get("beds_total"):
        occ_txt = f"{(s['beds_total'] - s['beds_available']) / s['beds_total'] * 100:.0f}%"
    ox = int(s.get("oxygen_pct", 0))
    avail = int(s.get("ambulances_available", 0))
    total = int(s.get("ambulances_total", 0))
    glance_items = [
        ("patient", "Patients in hospital", int(s.get("current_inpatients", 0))),
        ("admit", "New admissions", int(s.get("admitted", 0))),
        ("discharge", "Discharged", int(s.get("discharged", 0))),
        ("er", "ER visits", int(s.get("er_visits", 0))),
        ("surgery", "Surgeries", int(s.get("surgeries", 0))),
        ("birth", "Births", int(s.get("births", 0))),
        ("stillbirth", "Stillbirths", int(s.get("stillbirths", 0))),
        ("death", "Mortality", int(s.get("deaths", 0))),
        ("bed", "Bed occupancy", occ_txt),
        ("doctor", "Doctors on duty", int(s.get("doctors", 0))),
        ("nurse", "Nurses on duty", int(s.get("nurses", 0))),
        ("ambulance", "Ambulances", f"{avail}/{total}"),
        ("oxygen", "Oxygen supply", f"{ox}%"),
    ]
    # Hospital Performance (condition banner + KPI cards + gauge + breakdown),
    # shown first — before At a Glance — to mirror the on-screen dashboard.
    try:
        _perf = health_summary(*perf_inputs_single(s, depts, meds, tests, blood))
        _pimg = performance_report_image(_perf)
        if _pimg:
            story += keep_section(f"Hospital Performance — {_perf['condition']}",
                                  img(_pimg))
    except Exception:
        log.exception("PDF: failed to build the Performance section")

    try:
        story += keep_section("At a Glance", img(build_glance_image(glance_items)))
    except Exception:
        log.exception("PDF: failed to build the At-a-Glance image section")

    # ── Reorderable band: everything below At-a-Glance and above the Notes is
    # collected here and flushed in the analyst's saved dashboard order, so the
    # on-dashboard ⬆️/⬇️ arrows drive this report too. ──
    band = _ReportBand()

    # Patients + occupancy
    _sec = keep_section("Patients", kv_table(grp_rows("Patients")))
    if charts.get("patients"):
        _sec += [Spacer(1, 6), img(charts["patients"])]
    band.add("Patients", _sec)

    # Capacity (with occupancy)
    cap = grp_rows("Capacity")
    if s.get("beds_total"):
        occ = (s["beds_total"] - s["beds_available"]) / s["beds_total"] * 100
        cap = cap + [("Bed occupancy (%)", f"{occ:.0f}%")]
    band.add("Capacity", keep_section("Capacity", kv_table(cap)))

    # Staff
    _sec = keep_section("Staff on Duty", kv_table(grp_rows("Staff")))
    if charts.get("staff"):
        _sec += [Spacer(1, 6), img(charts["staff"])]
    band.add("Staff on Duty", _sec)

    # Ambulances + supplies
    band.add("Ambulances & Emergency",
             keep_section("Ambulances & Emergency", kv_table(grp_rows("Ambulances & Emergency"))))
    _sec = keep_section("Critical Supplies", kv_table(grp_rows("Critical Supplies")))
    _sec += [Spacer(1, 6), img(oxygen_gauge_image(s.get("oxygen_pct", 0)))]
    band.add("Critical Supplies", _sec)

    # Blood bank
    if not blood.empty:
        bt = [["Blood Type", "Units"]] + [[r["Blood Type"], str(int(r["Units"]))]
                                          for _, r in blood.iterrows()]
        t = Table(bt, colWidths=(80 * mm, 80 * mm))
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(PRIMARY)),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#BFE0E0")),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor(LIGHT_BG)]),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
        ]))
        _sec = keep_section("Blood Bank", t)
        if charts.get("blood"):
            _sec += [Spacer(1, 6), img(charts["blood"])]
        band.add("Blood Bank", _sec)

    # Departments (status colour-coded)
    if not depts.empty:
        rows = [["Department", "Status"]] + [[r["Department"], r["Status"]]
                                             for _, r in depts.iterrows()]
        t = Table(rows, colWidths=(110 * mm, 50 * mm))
        ts = [("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(PRIMARY)),
              ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
              ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
              ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#BFE0E0")),
              ("FONTSIZE", (0, 0), (-1, -1), 9)]
        for i, (_, r) in enumerate(depts.iterrows(), start=1):
            col = STATUS_COLOR.get(r["Status"], "#777777")
            ts.append(("BACKGROUND", (1, i), (1, i), colors.HexColor(col)))
            ts.append(("TEXTCOLOR", (1, i), (1, i), colors.white))
        t.setStyle(TableStyle(ts))
        band.add("Department Status", keep_section("Department Status", t))

    # Medications
    if not meds.empty:
        rows = [["Medication", "Availability"]] + [
            [r["Medication"], str(r.get("Status", "") or "")]
            for _, r in meds.iterrows()]
        t = Table(rows, colWidths=(110 * mm, 50 * mm))
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(PRIMARY)),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#BFE0E0")),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor(LIGHT_BG)]),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
        ]))
        band.add("Medication Availability", keep_section("Medication Availability", t))

    # Tests
    if not tests.empty:
        rows = [["Test", "Available"]] + [
            [r["Test"], "Yes" if r["Available"] else "No"] for _, r in tests.iterrows()]
        t = Table(rows, colWidths=(120 * mm, 40 * mm))
        ts = [("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(PRIMARY)),
              ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
              ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
              ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#BFE0E0")),
              ("FONTSIZE", (0, 0), (-1, -1), 9)]
        for i, (_, r) in enumerate(tests.iterrows(), start=1):
            col = OK_GREEN if r["Available"] else DANGER
            ts.append(("BACKGROUND", (1, i), (1, i), colors.HexColor(col)))
            ts.append(("TEXTCOLOR", (1, i), (1, i), colors.white))
        t.setStyle(TableStyle(ts))
        band.add("Medical Tests", keep_section("Medical Tests", t))

    # Absent Specialists (specialty / expected return)
    if absent is not None and not getattr(absent, "empty", True):
        _arows = [["Specialist", "Specialty / Area", "Expected return"]]
        for _, r in absent.iterrows():
            nm = str(r.get("Specialist", "") or "").strip()
            if not nm:
                continue
            _arows.append([nm, str(r.get("Specialty / Area", "") or "").strip(),
                           str(r.get("Expected return", "") or "").strip() or "—"])
        if len(_arows) > 1:
            at = Table(_arows, colWidths=(60 * mm, 65 * mm, 45 * mm))
            at.setStyle(_pdf_header_table_style())
            band.add("Absent Specialists", keep_section("Absent Specialists", at))

    # Mortality register (ward / age / length of stay)
    if mortality is not None and not getattr(mortality, "empty", True):
        _ms = mortality_stats(mortality)
        if _ms["total"]:
            summ = [("Deaths recorded", _ms["total"]),
                    ("Average age",
                     "—" if _ms["avg_age"] is None else _ms["avg_age"]),
                    ("Average length of stay (days)",
                     "—" if _ms["avg_los"] is None else _ms["avg_los"])]
            _sec = keep_section("Mortality Register", kv_table(summ))
            _sec.append(Spacer(1, 6))
            # per-ward breakdown: count of deaths and the ages recorded
            _bd = mortality_ward_breakdown(mortality)
            if not _bd.empty:
                brows = [["Ward", "Deaths", "Ages recorded", "Average age"]]
                for _, br in _bd.iterrows():
                    brows.append([str(br["Ward"]), str(br["Deaths"]),
                                  str(br["Ages recorded"]), str(br["Average age"])])
                bt = Table(brows, colWidths=(55 * mm, 20 * mm, 60 * mm, 25 * mm))
                bt.setStyle(_pdf_header_table_style())
                _sec += [bt, Spacer(1, 8)]
            _cols = list(mortality.columns)
            _wc = "Ward" if "Ward" in _cols else "ward"
            _ac = "Age" if "Age" in _cols else "age"
            _lc = "Length of stay (days)" if "Length of stay (days)" in _cols else "los_days"
            _cc = "Condition" if "Condition" in _cols else "condition"
            _xc = "Cause of death" if "Cause of death" in _cols else "cause"
            _tc = "Time of death" if "Time of death" in _cols else "time_of_death"
            cell = ParagraphStyle("mc", parent=body, fontSize=8.5, leading=11)
            mrows = [["Ward", "Age", "LOS", "Time", "Condition", "Cause of death"]]
            for _, r in mortality.iterrows():
                w = str(r.get(_wc, "") or "").strip()
                if not w:
                    continue
                a = r.get(_ac)
                l = r.get(_lc)
                try:
                    a = "—" if a is None or (isinstance(a, float) and pd.isna(a)) else int(a)
                except (TypeError, ValueError):
                    a = "—"
                try:
                    l = "—" if l is None or (isinstance(l, float) and pd.isna(l)) else int(l)
                except (TypeError, ValueError):
                    l = "—"
                tod = _to_time_str(r.get(_tc, "")) or "—"
                cond = str(r.get(_cc, "") or "").strip() or "—"
                cause = str(r.get(_xc, "") or "").strip() or "—"
                mrows.append([Paragraph(_pdf_esc(w), cell), str(a), str(l), tod,
                              Paragraph(_pdf_esc(cond), cell),
                              Paragraph(_pdf_esc(cause), cell)])
            mt = Table(mrows, colWidths=(30 * mm, 11 * mm, 14 * mm, 16 * mm, 50 * mm, 59 * mm))
            mt.setStyle(_pdf_header_table_style(
                fontsize=8.5, extra=[("VALIGN", (0, 0), (-1, -1), "TOP")]))
            _sec.append(mt)
            band.add("Mortality Register", _sec)

    # User-defined custom tables (with their configured chart & colors)
    if customs:
        chdr = ParagraphStyle("cch", parent=body, fontSize=8.5, leading=11,
                              textColor=colors.white)
        ccell = ParagraphStyle("cc2", parent=body, fontSize=8.5, leading=11)
        for item in customs:
            ctitle, cdf = item[0], item[1]
            tdef = item[2] if len(item) > 2 else None
            if cdf is None or getattr(cdf, "empty", True):
                continue
            ccols = [c for c in cdf.columns if c != "Date"]
            if not ccols:
                continue
            dcfg = table_display(tdef) if tdef else {}
            _content = []
            if tdef and dcfg.get("mode") == "kpi":
                kitems = custom_table_kpis(tdef, cdf)
                if kitems:
                    _content.append(kv_table(kitems))
                    _content.append(Spacer(1, 6))
            elif tdef:
                ch = custom_pdf_chart(tdef, cdf)
                if ch:
                    _content.append(img(ch))
                    _content.append(Spacer(1, 6))
            ccolors = (dcfg.get("colors") or {}) if tdef else {}
            _cdf_disp = custom_df_for_display(tdef, cdf) if tdef else cdf
            crows = [[Paragraph(f"<b>{_pdf_esc(c)}</b>", chdr) for c in ccols]]
            for _, r in _cdf_disp.iterrows():
                vals = [str(r.get(c, "") or "").strip() for c in ccols]
                if not any(vals):
                    continue
                crows.append([Paragraph(_pdf_esc(v) or "—", ccell) for v in vals])
            if len(crows) > 1:
                extra = [("VALIGN", (0, 0), (-1, -1), "TOP")]
                for ci, c in enumerate(ccols):
                    if ccolors.get(c):
                        try:
                            extra.append(("BACKGROUND", (ci, 0), (ci, 0),
                                          colors.HexColor(ccolors[c])))
                        except Exception:
                            pass
                cw = [180.0 / len(ccols) * mm] * len(ccols)
                ct = Table(crows, colWidths=cw)
                ct.setStyle(_pdf_header_table_style(extra=extra))
                _content.append(ct)
            band.add(ctitle,
                     keep_section(ctitle, *_content) if _content else section(ctitle))

    # Flush the reorderable band in saved dashboard order.
    for _fl in band.payloads():
        story += _fl

    # Notes
    if s.get("notes"):
        story += keep_section("Notes", Paragraph(_pdf_esc(s["notes"]).replace("\n", "<br/>"), body))

    # All dashboard notes: overview, per-chart/table notes, and custom notes
    _notes = _norm_notes(observations)
    if _notes:
        import html as _h
        _blocks = []
        for _lbl, _txt in _notes:
            _blocks.append(Paragraph(f'<b>{_h.escape(str(_lbl))}</b>', body))
            _blocks.append(Paragraph(
                _h.escape(str(_txt)).replace("\n", "<br/>"), body))
            _blocks.append(Spacer(1, 5))
        story += keep_section("Notes & Observations", *_blocks)

    story.append(Spacer(1, 10))
    story.append(hr())
    story.append(Paragraph(
        "Automatically generated by the Hospital Dashboard. Figures as entered by the analyst for this day.",
        small))

    doc.build(story)
    return buf.getvalue()


@_maybe_cache(show_spinner=False, max_entries=48, ttl=900)
def build_period_pdf(start, end, daily, mortality, customs, hospital_name, period_label,
                     observations=None, note_scope=None, note_ref=None,
                     perf_summary=None, depts=None, meds=None, tests=None,
                     blood=None, absent=None):
    """Weekly / monthly report: At-a-Glance period summary, per-metric trend
    charts, the oxygen gauge, the latest resource snapshots (blood bank,
    departments, medications, tests, absent specialists), mortality by ward, and
    any custom tables — mirroring the on-screen week / month dashboard."""
    import io as _io
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
                                    KeepTogether,
                                    TableStyle, Image)
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    buf = _io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, topMargin=16 * mm, bottomMargin=14 * mm,
                            leftMargin=14 * mm, rightMargin=14 * mm)
    base = getSampleStyleSheet()
    title_style = ParagraphStyle("pt", parent=base["Title"], fontSize=20,
                                 textColor=colors.HexColor(PRIMARY))
    sub_style = ParagraphStyle("ps", parent=base["Normal"], fontSize=10.5, textColor=colors.grey)
    h_style = ParagraphStyle("ph", parent=base["Heading2"], fontSize=12.5,
                             textColor=colors.HexColor(PRIMARY), spaceBefore=10)
    body = ParagraphStyle("pb", parent=base["Normal"], fontSize=9.5, leading=13)

    def section(t):
        return [Spacer(1, 6), Paragraph(_pdf_esc(t), h_style)]

    def _note_flow(title):
        if not (note_scope and note_ref):
            return []
        txt, nf = section_note_full(note_scope, note_ref, title)
        if not txt and not nf.get("image"):
            return []
        import html as _h2
        note_sty = ParagraphStyle(
            "secnote2", parent=body,
            fontSize=max(7.0, min(20.0, float(nf["size"]) * 0.62)),
            fontName=NOTE_FONTS_PDF.get(nf["font"], "Helvetica"),
            textColor=colors.HexColor(nf["color"] or INK),
            backColor=colors.HexColor(LIGHT_BG), borderPadding=6, leftIndent=2,
            spaceBefore=3, spaceAfter=2,
            alignment=NOTE_ALIGN_PDF.get(nf["align"], 0))
        _t = _h2.escape(txt).replace("\n", "<br/>")
        if nf["bold"]:
            _t = f"<b>{_t}</b>"
        if nf["italic"]:
            _t = f"<i>{_t}</i>"
        if nf["underline"]:
            _t = f"<u>{_t}</u>"
        flow = [Spacer(1, 3)]
        if txt:
            flow.append(Paragraph("📝 " + _t, note_sty))
        if nf.get("image"):
            _im = _note_pdf_image(nf["image"])
            if _im is not None:
                _im.hAlign = {"left": "LEFT", "center": "CENTER",
                              "right": "RIGHT"}.get(nf["align"], "LEFT")
                flow += [Spacer(1, 3), _im]
        return flow

    def keep_section(title, *flowables):
        """Keep a heading with its first following flowable (chart/table) so it
        is never split across a page break; append the section's own note."""
        flat = []
        for f in flowables:
            flat.extend(f if isinstance(f, list) else [f])
        flat += _note_flow(title)
        head = [Spacer(1, 6), Paragraph(_pdf_esc(title), h_style)]
        if not flat:
            return head
        return [KeepTogether(head + [flat[0]])] + flat[1:]

    story = [Paragraph(_pdf_esc(hospital_name), title_style),
             Paragraph(f"{period_label} &nbsp;•&nbsp; {start:%d %b %Y} – {end:%d %b %Y}",
                       sub_style),
             Paragraph(f"Generated {now_local():%d %b %Y %H:%M}", sub_style),
             Spacer(1, 4)]

    d = daily.copy() if daily is not None else pd.DataFrame()
    days_reported = int(d["entry_date"].nunique()) if not d.empty else 0

    def _num(col):
        if d.empty or col not in d.columns:
            return pd.Series([], dtype=float)
        return pd.to_numeric(d[col], errors="coerce").dropna()

    # Hospital Performance first (banner + KPI cards + gauge + breakdown)
    if perf_summary:
        try:
            _pimg = performance_report_image(perf_summary)
            if _pimg:
                _pb, _par = _pimg
                story += keep_section(
                    f"Hospital Performance — {perf_summary['condition']}",
                    Image(_pb, width=178 * mm, height=178 * mm * _par))
        except Exception:
            log.exception("period PDF: failed to build the Performance section")

    story += section("At a Glance — period summary")
    srows = [["Metric", "Value"]]
    for k, lbl, kind in PERIOD_SUMMARY:
        ser = _num(k)
        if ser.empty:
            val = "—"
        elif kind == "sum":
            val = f"{int(ser.sum()):,}"
        else:
            val = f"{ser.mean():.1f}"
        srows.append([lbl, val])
    srows.append(["Days reported", str(days_reported)])
    stab = Table(srows, colWidths=(120 * mm, 60 * mm))
    stab.setStyle(_pdf_header_table_style())
    story = story[:-1] + [KeepTogether(story[-1:] + [stab])]

    # ── Reorderable band: Trends, Mortality-by-ward and the custom tables are
    # collected and flushed in the analyst's saved dashboard order. ──
    band = _ReportBand()

    if not d.empty:
        dd = d.copy()
        dd["entry_date"] = pd.to_datetime(dd["entry_date"])
        dd = dd.sort_values("entry_date")
        x = dd["entry_date"]

        def _chart(specs, title):
            fig, ax = plt.subplots(figsize=(7.2, 2.7), dpi=150)
            drew = False
            for col, lbl, color in specs:
                if col in dd.columns:
                    ax.plot(x, pd.to_numeric(dd[col], errors="coerce"), marker="o",
                            label=lbl, color=color, linewidth=2)
                    drew = True
            ax.set_title(title, color=PRIMARY, fontsize=11, loc="left")
            if drew:
                ax.legend(fontsize=7, frameon=False)
            ax.grid(True, alpha=.25)
            for sp in ax.spines.values():
                sp.set_visible(False)
            fig.autofmt_xdate(rotation=30)
            b = _io.BytesIO()
            fig.savefig(b, format="png", bbox_inches="tight")
            plt.close(fig)
            b.seek(0)
            return Image(b, width=178 * mm, height=66 * mm)

        # Per-metric trend charts, one per dashboard section (so they mirror the
        # on-screen week / month dashboard and follow the saved section order).
        band.add("Patient Flow", keep_section(
            "Patient Flow",
            _chart([("current_inpatients", "Patients in hospital", PRIMARY),
                    ("admitted", "New admissions", TEAL2),
                    ("discharged", "Discharged", WARN)],
                   "Patients, admissions & discharges")))
        band.add("ER / Surgeries / ICU", keep_section(
            "ER / Surgeries / ICU",
            _chart([("er_visits", "ER visits", PRIMARY),
                    ("surgeries", "Surgeries", TEAL2),
                    ("icu_patients", "ICU patients", DANGER)],
                   "ER visits, surgeries & ICU")))
        band.add("Beds & Occupancy", keep_section(
            "Beds & Occupancy",
            _chart([("beds_available", "Beds available", PRIMARY),
                    ("beds_total", "Total beds", TEAL2),
                    ("icu_beds_available", "ICU beds available", WARN)],
                   "Bed availability")))
        band.add("Staffing", keep_section(
            "Staffing",
            _chart([("doctors", "Doctors", PRIMARY), ("nurses", "Nurses", TEAL2),
                    ("support_staff", "Support", WARN),
                    ("specialists_on_call", "Specialists", DANGER)],
                   "Staff on duty")))
        band.add("Ambulances", keep_section(
            "Ambulances",
            _chart([("ambulances_available", "Available", PRIMARY),
                    ("ambulance_calls", "Calls responded", TEAL2)],
                   "Ambulance availability & calls")))

        # Oxygen gauge (period average), mirroring the dashboard's gauge.
        _avg_ox = pd.to_numeric(dd.get("oxygen_pct"), errors="coerce").mean() \
            if "oxygen_pct" in dd.columns else 0
        _ob, _oar = oxygen_gauge_image(0 if pd.isna(_avg_ox) else _avg_ox,
                                       title="Oxygen supply level (period avg)")
        band.add("Critical Supplies", keep_section(
            "Critical Supplies", Image(_ob, width=120 * mm, height=120 * mm * _oar)))

        # Daily trend (inpatients / admissions / mortality) + births breakdown.
        _sec = keep_section("Daily Trend",
                            _chart([("current_inpatients", "Patients in hospital", PRIMARY),
                                    ("admitted", "New admissions", TEAL2),
                                    ("deaths", "Mortality", DANGER)],
                                   "Daily trend"))
        _sec.append(Spacer(1, 6))
        _sec.append(_chart([("births", "Births", PRIMARY),
                            ("stillbirths", "Stillbirths", WARN),
                            ("deaths", "Mortality", DANGER)],
                           "Births, stillbirths & mortality"))
        band.add("Daily Trend", _sec)

    # ── Latest resource snapshots (as of the last day with data), mirroring the
    # "as of <day>" resource blocks on the week / month dashboard. ──
    _blood_rows = snapshot_blood(blood)
    if _blood_rows:
        _t = Table([["Blood Type", "Units"]] + [[bt, str(u)] for bt, u in _blood_rows],
                   colWidths=(80 * mm, 80 * mm))
        _t.setStyle(_pdf_header_table_style())
        band.add("Blood Bank", keep_section("Blood Bank", _t))

    _dept_rows = snapshot_status(depts, ("name", "Department"), ("status", "Status"))
    if _dept_rows:
        _t = Table([["Department", "Status"]] + [[n, s] for n, s in _dept_rows],
                   colWidths=(110 * mm, 50 * mm))
        _ts = [("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(PRIMARY)),
               ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
               ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
               ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#BFE0E0")),
               ("FONTSIZE", (0, 0), (-1, -1), 9)]
        for _i2, (_n, _s) in enumerate(_dept_rows, start=1):
            col = STATUS_COLOR.get(_s, "#777777")
            _ts.append(("BACKGROUND", (1, _i2), (1, _i2), colors.HexColor(col)))
            _ts.append(("TEXTCOLOR", (1, _i2), (1, _i2), colors.white))
        _t.setStyle(TableStyle(_ts))
        band.add("Department Status", keep_section("Department Status", _t))

    _med_rows = snapshot_status(meds, ("name", "Medication"), ("status", "Status"))
    if _med_rows:
        _t = Table([["Medication", "Availability"]] + [[n, s] for n, s in _med_rows],
                   colWidths=(110 * mm, 50 * mm))
        _t.setStyle(_pdf_header_table_style())
        band.add("Medication Availability", keep_section("Medication Availability", _t))

    _test_rows = snapshot_tests(tests)
    if _test_rows:
        _t = Table([["Test", "Available"]]
                   + [[n, "Yes" if av else "No"] for n, av in _test_rows],
                   colWidths=(120 * mm, 40 * mm))
        _tst = [("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(PRIMARY)),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#BFE0E0")),
                ("FONTSIZE", (0, 0), (-1, -1), 9)]
        for _i2, (_n, av) in enumerate(_test_rows, start=1):
            col = OK_GREEN if av else DANGER
            _tst.append(("BACKGROUND", (1, _i2), (1, _i2), colors.HexColor(col)))
            _tst.append(("TEXTCOLOR", (1, _i2), (1, _i2), colors.white))
        _t.setStyle(TableStyle(_tst))
        band.add("Medical Tests", keep_section("Medical Tests", _t))

    _absent_rows = snapshot_absent(absent)
    if _absent_rows:
        _t = Table([["Specialist", "Specialty / Area", "Expected return"]]
                   + [list(r) for r in _absent_rows],
                   colWidths=(60 * mm, 65 * mm, 45 * mm))
        _t.setStyle(_pdf_header_table_style())
        band.add("Absent Specialists", keep_section("Absent Specialists", _t))

    if mortality is not None and not getattr(mortality, "empty", True):
        bd = mortality_ward_breakdown(mortality)
        if not bd.empty:
            mrows = [list(bd.columns)] + bd.astype(str).values.tolist()
            mt = Table(mrows, colWidths=(55 * mm, 20 * mm, 60 * mm, 25 * mm))
            mt.setStyle(_pdf_header_table_style())
            band.add("Mortality by ward", keep_section("Mortality by ward", mt))

    if customs:
        chdr = ParagraphStyle("pch", parent=body, fontSize=8, leading=10,
                              textColor=colors.white)
        ccell = ParagraphStyle("pcc", parent=body, fontSize=8, leading=10)
        for item in customs:
            ctitle, cdf = item[0], item[1]
            tdef = item[2] if len(item) > 2 else None
            if cdf is None or getattr(cdf, "empty", True):
                continue
            ccols = list(cdf.columns)
            dcfg = table_display(tdef) if tdef else {}
            _content = []
            if tdef and dcfg.get("mode") == "kpi":
                kitems = custom_table_kpis(tdef, cdf)
                if kitems:
                    krows = [["Metric", "Value"]] + [[k, str(v)] for k, v in kitems]
                    kt = Table(krows, colWidths=(120 * mm, 58 * mm))
                    kt.setStyle(_pdf_header_table_style())
                    _content.append(kt)
                    _content.append(Spacer(1, 6))
            elif tdef:
                ch = custom_pdf_chart(tdef, cdf)
                if ch:
                    cb, car = ch
                    _content.append(Image(cb, width=172 * mm, height=172 * mm * car))
                    _content.append(Spacer(1, 6))
            ccolors = (dcfg.get("colors") or {}) if tdef else {}
            _cdf_disp = custom_df_for_display(tdef, cdf) if tdef else cdf
            crows = [[Paragraph(f"<b>{_pdf_esc(c)}</b>", chdr) for c in ccols]]
            for _, r in _cdf_disp.iterrows():
                vals = [str(r.get(c, "") or "").strip() for c in ccols]
                if not any(vals):
                    continue
                crows.append([Paragraph(_pdf_esc(v) or "—", ccell) for v in vals])
            if len(crows) > 1:
                extra = [("VALIGN", (0, 0), (-1, -1), "TOP")]
                for ci, c in enumerate(ccols):
                    if ccolors.get(c):
                        try:
                            extra.append(("BACKGROUND", (ci, 0), (ci, 0),
                                          colors.HexColor(ccolors[c])))
                        except Exception:
                            pass
                cw = [178.0 / len(ccols) * mm] * len(ccols)
                ct = Table(crows, colWidths=cw)
                ct.setStyle(_pdf_header_table_style(extra=extra))
                _content.append(ct)
            band.add(ctitle,
                     keep_section(ctitle, *_content) if _content else section(ctitle))

    # Flush the reorderable band in saved dashboard order.
    for _fl in band.payloads():
        story += _fl

    _notes = _norm_notes(observations)
    if _notes:
        import html as _h
        _blocks = []
        for _lbl, _txt in _notes:
            _blocks.append(Paragraph(f'<b>{_h.escape(str(_lbl))}</b>', body))
            _blocks.append(Paragraph(
                _h.escape(str(_txt)).replace("\n", "<br/>"), body))
            _blocks.append(Spacer(1, 5))
        story += keep_section("Notes & Observations", *_blocks)

    doc.build(story)
    buf.seek(0)
    return buf.getvalue()


def render_period_downloads(daily, mort, custom_rng, start, end, period_word,
                            csv_name, pdf_name, period_label, perf_summary=None,
                            depts=None, meds=None, tests=None, blood=None,
                            absent=None):
    """Shared CSV + PDF download row for the Week and Month roll-up dashboards."""
    if daily.empty:
        return
    obs_ref = start.isoformat() if period_word == "week" else f"{start:%Y-%m}"
    obs_text = gather_general_notes(period_word, obs_ref, start=start, end=end)
    dl = st.columns(3)
    dl_button(dl[0], 
        f"⬇️ Download this {period_word}'s data (CSV)",
        daily.to_csv(index=False).encode("utf-8"),
        file_name=csv_name, mime="text/csv", use_container_width=True)
    dl_button(dl[1], 
        f"📄 Download this {period_word}'s report (PDF)",
        build_period_pdf(start, end, daily, mort,
                         [(t["title"], cr, t) for t, cr in custom_rng],
                         HOSPITAL_NAME, period_label, observations=obs_text,
                         note_scope=period_word, note_ref=obs_ref,
                         perf_summary=perf_summary, depts=depts, meds=meds,
                         tests=tests, blood=blood, absent=absent),
        file_name=pdf_name, mime="application/pdf", use_container_width=True)
    if PPTX_OK:
        dl_button(dl[2], 
            f"📊 Download this {period_word}'s slides (PPTX)",
            build_period_pptx(start, end, daily, mort,
                              [(t["title"], cr, t) for t, cr in custom_rng],
                              HOSPITAL_NAME, period_label,
                              observations=obs_text, note_scope=period_word,
                              note_ref=obs_ref, perf_summary=perf_summary,
                              depts=depts, meds=meds, tests=tests, blood=blood,
                              absent=absent),
            file_name=(pdf_name[:-4] + ".pptx" if pdf_name.endswith(".pdf")
                       else pdf_name + ".pptx"),
            mime=PPTX_MIME, use_container_width=True)
    else:
        dl[2].caption("Add `python-pptx` to requirements.txt to enable "
                      "PowerPoint downloads.")


def render_rollup_normal(daily, mort, custom_rng, tests, latest_dept_status,
                         latest_absent, period_word, trend_key):
    """Shared normal-mode body for the Week and Month roll-up dashboards: the
    At-a-Glance trend chart plus the department / specialist / mortality /
    custom-table / tests sections. `period_word` is 'week' or 'month'."""
    st.markdown(f'<div class="section">Trends over this {period_word}</div>',
                unsafe_allow_html=True)
    opts = [(k, l) for k, l in GLANCE_TREND if k in daily.columns]
    labels = [l for _, l in opts]
    pick = st.multiselect(
        f"Pick At-a-Glance metrics to chart across the {period_word}", labels,
        default=labels[:3], key=trend_key)
    tf = trend_fig(daily, [(k, l) for k, l in opts if l in pick],
                   f"Daily trend — this {period_word}")
    if tf is not None:
        st.plotly_chart(tf, use_container_width=True)
    else:
        st.caption("Pick one or more metrics above (data builds up as days are entered).")

    if latest_dept_status is not None:
        st.caption(f"Latest department status — {latest_dept_status[0]}:")
        dept_pills(latest_dept_status[1])
    st.markdown('<div class="section">Specialist Availability</div>', unsafe_allow_html=True)
    render_absent_block(latest_absent)

    ms = mortality_stats(mort)
    if ms["total"]:
        suffix = "wk" if period_word == "week" else "mo"
        st.markdown(f'<div class="section">Mortality register (this {period_word})</div>',
                    unsafe_allow_html=True)
        render_kpis([(f"⚰️ Deaths recorded ({suffix})", ms["total"]),
                     ("🎂 Average age", "—" if ms["avg_age"] is None else ms["avg_age"]),
                     ("🛏️ Avg length of stay (days)",
                      "—" if ms["avg_los"] is None else ms["avg_los"])], False)
        st.caption("Mortality by ward — count and ages recorded:")
        st.dataframe(mortality_ward_breakdown(mort), use_container_width=True, hide_index=True)

    for t, cr in custom_rng:
        render_custom_table_normal(t, cr, f" (this {period_word})")

    if not tests.empty:
        st.markdown('<div class="section">Medical Tests Available</div>', unsafe_allow_html=True)
        latest_t = tests["entry_date"].max()
        tt = tests[tests["entry_date"] == latest_t]
        render_tests_block(tt[tt["available"] == 1]["name"].tolist(),
                           tt[tt["available"] == 0]["name"].tolist(),
                           suffix=f" — {day_label(latest_t)}")
st.sidebar.markdown(
    f'<div style="text-align:center;background:transparent;'
    f'border-radius:12px;padding:6px 8px;margin-bottom:8px;">'
    f'<div style="font-size:.66rem;font-weight:700;letter-spacing:.05em;'
    f'text-transform:uppercase;color:{PRIMARY};margin-bottom:6px;">'
    f'Powered by Health Data Matrics</div>'
    f'<img src="data:image/png;base64,{HDM_LOGO_B64}" style="width:150px;max-width:92%;"/>'
    f'</div>', unsafe_allow_html=True)
st.sidebar.markdown(f"#### {HOSPITAL_NAME}")

# ── Big-screen slideshow control, shown above the View selector. No analyst
# sign-in required. ──
with st.sidebar.container(border=True):
    if st.query_params.get("tv") == "1":
        st.caption("Big-screen mode is on.")
        if st.button("✕ Exit big-screen mode", use_container_width=True,
                     key="tab_tv_exit"):
            exit_tv()
    else:
        st.markdown(
            f'<div style="font-size:.85rem;line-height:1.4;color:{BODY_FG};">'
            f'Full-screen, auto-rotating slideshow of the current dashboard '
            f'— ideal for a waiting-room or ward TV.</div>',
            unsafe_allow_html=True)
        if st.button("📺 Display on big screen", use_container_width=True,
                     key="tab_tv_start"):
            enter_tv()

mode = st.sidebar.radio("View", ["Public Dashboard", "Data Entry (Analyst)", "Settings"])
st.sidebar.toggle("🌙 Dark mode", key="ui_dark",
                  help="Switch between light and dark display. Light is the default.")
st.sidebar.divider()

# auto sign-out on inactivity + manual logout
if st.session_state.get("authed"):
    if time.time() - st.session_state.get("auth_time", 0) > SESSION_TIMEOUT:
        st.session_state.pop("authed", None)
        st.session_state.pop("auth_time", None)
    elif st.sidebar.button("🔒 Log out", use_container_width=True):
        for _k in ("authed", "auth_time", "fails", "lock_until", "user_email"):
            st.session_state.pop(_k, None)
        st.rerun()

st.sidebar.markdown(
    f'<div style="font-size:.85rem;font-weight:600;color:{BODY_FG};'
    f'margin:6px 0 2px;">Data entry done by M&amp;E team.</div>',
    unsafe_allow_html=True)

# Product disclaimer, pinned to the bottom of the sidebar
st.sidebar.markdown(
    f'<div style="margin-top:16px;padding-top:10px;border-top:1px solid {SIDE_BORDER};'
    f'font-size:.8rem;line-height:1.5;color:{BODY_FG};text-align:center;">'
    f'<span style="font-weight:700;color:{PRIMARY};">This is a product of '
    f'Health Data Matrics.</span><br>'
    f'<span style="color:{BODY_FG};">© {now_local():%Y} Health Data Matrics. '
    f'All rights reserved.</span></div>', unsafe_allow_html=True)


def require_analyst():
    """Show the analyst login and stop the script unless already signed in."""
    if st.session_state.get("authed"):
        return
    now = time.time()
    lock_until = st.session_state.get("lock_until", 0)
    if now < lock_until:
        st.error(f"Too many attempts. Try again in {int(lock_until - now)} seconds.")
        st.stop()
    email = st.text_input("Email", placeholder="you@hospital.org",
                          help="Sign in with your analyst email, or the master "
                               "admin password if you haven't created a login yet.")
    pw = st.text_input("Password", type="password")
    if USING_DEFAULT_PASSWORD and users_count() == 0:
        st.warning("⚠️ **This deployment is using the default password.** Anyone "
                   "who guesses it can enter data. Before real use, set a "
                   "`HOSPITAL_ADMIN_PASSWORD_HASH` secret (Settings shows how) and "
                   "reboot.")
    code = ""
    if two_factor_enabled():
        code = st.text_input("Authenticator code (6 digits)", max_chars=6,
                             help="From your authenticator app (Google Authenticator, "
                                  "Authy, 1Password, etc.).")
    if st.button("Unlock", type="primary"):
        pw_ok, who = verify_login(email, pw)
        otp_ok = (not two_factor_enabled()) or totp_verify(TOTP_SECRET, code)
        if pw_ok and otp_ok:
            st.session_state.authed = True
            st.session_state.auth_time = time.time()
            st.session_state.user_email = who
            st.session_state.fails = 0
            st.session_state.pop("lock_until", None)
            try:
                set_setting("last_login", now_local().isoformat(timespec="seconds"))
                if get_user(who):
                    set_user_last_login(who)
            except Exception:
                log.warning("Login succeeded but last-login bookkeeping failed",
                            exc_info=True)
            st.rerun()
        else:
            fails = st.session_state.get("fails", 0) + 1
            st.session_state.fails = fails
            time.sleep(min(3.0, THROTTLE_BASE * fails))
            if fails >= MAX_FAILS:
                st.session_state.lock_until = time.time() + LOCK_SECONDS
                st.session_state.fails = 0
                st.error(f"Too many failed attempts. Locked for {LOCK_SECONDS} seconds.")
            else:
                left = MAX_FAILS - fails
                msg = "Incorrect credentials." if two_factor_enabled() else "Incorrect password."
                st.error(f"{msg} {left} attempt(s) left before a temporary lock.")
    st.stop()


if mode == "Settings":
    st.markdown('<div class="big-title">⚙️ Settings</div>', unsafe_allow_html=True)
    st.markdown('<div class="sub">Manage analyst accounts, security and the hospital name.</div>',
                unsafe_allow_html=True)
    require_analyst()
    _me = st.session_state.get("user_email")
    if _me:
        st.caption(f"Signed in as **{_me}**.")
    tab_acct, tab_sec, tab_hosp = st.tabs(["Account", "Security", "Hospital"])

    # ── Account: create analyst logins + set/change passwords ──
    with tab_acct:
        st.markdown("**Create an analyst account**")
        st.caption("Analysts sign in with their email and a password. Passwords are stored "
                   "only as salted PBKDF2 hashes — never in plain text.")
        nu_email = st.text_input("Email (username)", key="nu_email",
                                 placeholder="analyst@hospital.org")
        nu_pw = st.text_input("Password", type="password", key="nu_pw")
        nu_pw2 = st.text_input("Confirm password", type="password", key="nu_pw2")
        if st.button("Create account", type="primary", key="create_acct"):
            if not _valid_email(nu_email):
                st.warning("Enter a valid email address.")
            elif len(nu_pw) < 8:
                st.warning("Password must be at least 8 characters (12+ recommended).")
            elif nu_pw != nu_pw2:
                st.warning("The two passwords don't match.")
            else:
                _existed = get_user(nu_email) is not None
                upsert_user(nu_email, nu_pw)
                st.success(("Updated" if _existed else "Created")
                           + f" login for {nu_email.strip().lower()}.")

        _users = list_users()
        st.markdown("**Create a new password**")
        if _users:
            st.caption("Set a new password for an existing analyst account.")
            cp_email = st.selectbox("Account", [u["email"] for u in _users], key="cp_email")
            cp_pw = st.text_input("New password", type="password", key="cp_pw")
            cp_pw2 = st.text_input("Confirm new password", type="password", key="cp_pw2")
            if st.button("Update password", key="update_pw"):
                if len(cp_pw) < 8:
                    st.warning("Password must be at least 8 characters (12+ recommended).")
                elif cp_pw != cp_pw2:
                    st.warning("The two passwords don't match.")
                else:
                    upsert_user(cp_email, cp_pw)
                    st.success(f"Password updated for {cp_email}.")
        else:
            st.caption("No analyst accounts yet — create one above first.")

        if _users:
            st.markdown("**Current logins**")
            st.caption(", ".join(u["email"] for u in _users))
            _del = st.selectbox("Remove a login",
                                ["—"] + [u["email"] for u in _users], key="del_user")
            if st.button("Remove selected login", key="rm_login") and _del != "—":
                delete_user(_del)
                st.success(f"Removed {_del}. Reload to refresh the list.")
        st.caption("⚠️ On Streamlit Community Cloud the database resets when the app "
                   "restarts, so accounts created here are cleared on restart. The master "
                   "admin password always signs you back in to recreate them. For permanent "
                   "accounts, a Postgres database is recommended.")

    # ── Security: master admin password hash + 2FA ──
    with tab_sec:
        _last = get_setting("last_login", "")
        st.caption(("🔐 Two-factor authentication is **on**." if two_factor_enabled()
                    else "🔓 Two-factor authentication is off (optional).")
                   + (f"  •  Last sign-in: {_last}." if _last else ""))
        if not ADMIN_HASH:
            st.caption("⚠️ Using a plaintext master password. Generate a salted hash below "
                       "and store it as `HOSPITAL_ADMIN_PASSWORD_HASH` for stronger protection.")
        st.markdown("**Master admin password (hashed)**")
        st.caption("Generate a salted PBKDF2 hash, then set it as the "
                   "`HOSPITAL_ADMIN_PASSWORD_HASH` secret (or env var) and restart. "
                   "The plaintext is never stored.")
        npw = st.text_input("New master password", type="password", key="newpw")
        if st.button("Generate hash", key="gen_hash"):
            if len(npw) < 8:
                st.warning("Use at least 8 characters (12+ recommended).")
            else:
                st.code(make_password_hash(npw), language="text")

        st.markdown("**Two-factor authentication (optional)**")
        st.caption("Generate a secret, store it as `HOSPITAL_ADMIN_TOTP_SECRET`, add it to "
                   "an authenticator app, and restart.")
        if st.button("Generate 2FA secret", key="gen_2fa"):
            _sec = _b32_secret()
            _label = (HOSPITAL_NAME or "Hospital").replace(" ", "%20")
            st.code(_sec, language="text")
            st.caption("otpauth link (add to your authenticator app):")
            st.code(f"otpauth://totp/{_label}:analyst?secret={_sec}&issuer={_label}",
                    language="text")

    # ── Hospital: public name + dashboard heading ──
    with tab_hosp:
        new_name = st.text_input("Hospital name (shown publicly)", value=HOSPITAL_NAME)
        if st.button("Save name", key="save_name"):
            set_setting("hospital_name", new_name.strip() or "General Hospital")
            st.success("Saved. Reload to update it everywhere.")

        st.divider()
        new_title = st.text_input(
            "Dashboard heading", value=get_dashboard_title(),
            help="The full heading shown on the dashboard, including the “— Dashboard” "
                 "text. Edit all of it however you like.")
        bt = st.columns(2)
        if bt[0].button("Save heading", key="save_title"):
            set_setting("dashboard_title", new_title.strip())
            st.success("Saved. Reload to update the dashboard heading.")
        if bt[1].button("Reset to default", key="reset_title"):
            set_setting("dashboard_title", "")
            st.success("Reset. The heading will follow the hospital name again.")

    st.markdown('<div class="section">💾 Backup &amp; Restore</div>',
                unsafe_allow_html=True)
    st.caption("⚠️ **Why data can disappear:** Streamlit Cloud's disk is temporary — "
               "the database resets whenever the app restarts, redeploys or wakes "
               "from sleep. Your entries save correctly; it's the disk that gets "
               "wiped. **Download a backup after entering data**, and restore it "
               "here after any reset. (Moving to a managed Postgres database makes "
               "persistence automatic — the recommended permanent fix.)")
    _bk = st.columns(2)
    dl_button(_bk[0], 
        "⬇️ Download full backup (everything)", export_backup_bytes(),
        file_name=f"hospital_backup_{now_local():%Y%m%d_%H%M}.json",
        mime="application/json", use_container_width=True,
        help="Daily entries, custom tables and their chart configs, observations, "
             "wards, analyst accounts and settings — all of it, in one file.")
    with _bk[1]:
        _bk_up = st.file_uploader("Restore from a backup file", type=["json"],
                                  key="bk_upload")
        if _bk_up is not None:
            if st.button("♻️ Restore this backup (replaces ALL current data)",
                         key="bk_restore", use_container_width=True):
                try:
                    _bt, _br = restore_backup_bytes(_bk_up.getvalue())
                    st.success(f"Restored {_br:,} rows across {_bt} tables. "
                               "All dashboards are back.")
                    st.rerun()
                except ValueError as _be:
                    st.error(str(_be))


# ══════════════════════════════════════════════
# DATA ENTRY MODE
# ══════════════════════════════════════════════
if mode == "Data Entry (Analyst)":
    st.markdown('<div class="big-title">Daily Data Entry</div>', unsafe_allow_html=True)
    st.markdown('<div class="sub">Pick a day from the calendar, fill in that day\'s figures, '
                'and save. Each saved day powers its own dashboard.</div>', unsafe_allow_html=True)

    require_analyst()

    # Right-click colour picks from the interactive chart previews arrive via
    # query params (set by the preview itself); save, reset the affected
    # pickers so they show the new colour, then clean the URL.
    _cc_tid = st.query_params.get("cchart")
    if _cc_tid:
        _cc_saved = apply_chart_color_pick(
            _cc_tid, st.query_params.get("ctarget", ""),
            st.query_params.get("ccolor", ""))
        if _cc_saved:
            st.session_state.pop(f"ccolcat_{_cc_tid}_{_cc_saved}", None)
            st.session_state.pop(f"ccol_{_cc_tid}_{_cc_saved}", None)
            st.session_state.pop(f"cbycat_{_cc_tid}", None)
            st.success(f"🎨 Color saved for “{_cc_saved}”.")
        for _ck in ("cchart", "ctarget", "ccolor"):
            try:
                del st.query_params[_ck]
            except Exception:
                pass

    st.markdown('<div class="section">Select the day</div>', unsafe_allow_html=True)
    entry_date, _, _, _ = day_picker("entry_day", "Pick a day (calendar)")
    st.info(f"Entering data for **{entry_date:%A, %d %b %Y}**")

    render_import_tool(entry_date)

    scalars, depts, meds, tests, blood, absent = load_day(entry_date)
    if scalars:
        st.caption(f"An entry already exists (updated {scalars.get('updated_at','—')}). "
                   "Saving overwrites it.")

    if st.session_state.pop("_dup_done", None):
        st.success("Previous day copied in. Review the values below and press "
                   "💾 Save this day to keep them.")

    # ── Duplicate the previous day into this day ──
    _prev_day = entry_date - timedelta(days=1)
    with st.expander("📋 Copy from the previous day", expanded=False):
        if not day_has_data(_prev_day):
            st.caption(f"No saved data was found for {_prev_day:%A, %d %b %Y}, "
                       "so there is nothing to copy.")
        else:
            st.caption(f"Pull the figures, resource tables and custom tables saved for "
                       f"**{_prev_day:%A, %d %b %Y}** into this day, ready to review and "
                       "save. Mortality records are not copied.")
            _dup_ok = True
            if scalars:
                _dup_ok = st.checkbox("This day already has saved data — overwrite it.",
                                      key="dup_overwrite")
            if st.button("📋 Duplicate previous day into this day",
                         disabled=not _dup_ok, key="dup_prev_btn"):
                duplicate_day(_prev_day, entry_date)
                for _k in list(st.session_state.keys()):
                    if _k.startswith(("in_", "ced_")) or _k in (
                            "blood_ed", "dept_ed", "med_ed", "test_ed",
                            "absent_ed", "notes_ta"):
                        del st.session_state[_k]
                st.session_state["_dup_done"] = True
                st.rerun()

    def num_val(key):
        return int(scalars[key]) if scalars and scalars.get(key) is not None else 0

    numeric = {}
    for group in FIELD_GROUPS:
        st.markdown(f'<div class="section">{group}</div>', unsafe_allow_html=True)
        fields = [f for f in DAILY_FIELDS if f[2] == group]
        for i in range(0, len(fields), 3):
            cols = st.columns(3)
            for col, (key, label, _) in zip(cols, fields[i:i+3]):
                maxv = 100 if key == "oxygen_pct" else 1_000_000
                numeric[key] = col.number_input(label, 0, maxv, num_val(key), key=f"in_{key}")

    if numeric.get("beds_total", 0) > 0:
        occ = numeric["beds_total"] - numeric["beds_available"]
        st.caption(f"➡️ Calculated bed occupancy: **{occ:,} / {numeric['beds_total']:,} "
                   f"({occ / numeric['beds_total'] * 100:.0f}%)**")

    # gentle, non-blocking sanity checks so obvious typos are caught before saving
    _warnings = []
    if numeric.get("beds_available", 0) > numeric.get("beds_total", 0):
        _warnings.append("Beds available is greater than total beds.")
    if numeric.get("icu_beds_available", 0) > numeric.get("beds_total", 0):
        _warnings.append("ICU beds available is greater than total beds.")
    if numeric.get("ambulances_available", 0) > numeric.get("ambulances_total", 0):
        _warnings.append("Ambulances available is greater than the fleet total.")
    if numeric.get("icu_patients", 0) > numeric.get("current_inpatients", 0):
        _warnings.append("ICU patients is greater than total inpatients.")
    if _warnings:
        st.warning("Please double-check: " + "  ".join("• " + w for w in _warnings)
                   + "  (You can still save — these are just reminders.)")

    st.markdown('<div class="section">Blood Bank (units by type)</div>', unsafe_allow_html=True)
    if blood.empty:
        blood = pd.DataFrame({"Blood Type": BLOOD_TYPES, "Units": [0] * len(BLOOD_TYPES)})
    blood_edit = st.data_editor(
        blood, num_rows="dynamic", use_container_width=True, key="blood_ed",
        column_config={"Blood Type": st.column_config.SelectboxColumn("Blood Type", options=BLOOD_TYPES, required=True),
                       "Units": st.column_config.NumberColumn("Units", min_value=0, step=1)})

    st.markdown('<div class="section">Operational Departments</div>', unsafe_allow_html=True)
    if depts.empty:
        depts = pd.DataFrame({"Department": DEFAULT_DEPARTMENTS,
                              "Status": ["Operational"] * len(DEFAULT_DEPARTMENTS)})
    dept_edit = st.data_editor(
        depts, num_rows="dynamic", use_container_width=True, key="dept_ed",
        column_config={"Department": st.column_config.TextColumn("Department", required=True),
                       "Status": st.column_config.SelectboxColumn("Status", options=DEPT_STATUSES, required=True)})

    st.markdown('<div class="section">Medication Availability</div>', unsafe_allow_html=True)
    st.caption("Set each medication's availability. Add or remove rows as needed.")
    if meds.empty:
        meds = pd.DataFrame(DEFAULT_MEDICATIONS, columns=["Medication", "Status"])
    med_edit = st.data_editor(
        meds, num_rows="dynamic", use_container_width=True, key="med_ed",
        column_config={"Medication": st.column_config.TextColumn("Medication", required=True),
                       "Status": st.column_config.SelectboxColumn("Status", options=MED_STATUSES,
                                                                  required=True)})

    st.markdown('<div class="section">Medical Tests Available Today</div>', unsafe_allow_html=True)
    if tests.empty:
        tests = pd.DataFrame({"Test": DEFAULT_TESTS, "Available": [True] * len(DEFAULT_TESTS)})
    test_edit = st.data_editor(
        tests, num_rows="dynamic", use_container_width=True, key="test_ed",
        column_config={"Test": st.column_config.TextColumn("Test", required=True),
                       "Available": st.column_config.CheckboxColumn("Available")})

    st.markdown('<div class="section">Absent / Away Specialists</div>', unsafe_allow_html=True)
    st.caption("Record specialists who are away: their name, specialty area, and the date "
               "they are expected back. Leave empty if everyone is available.")
    if absent.empty:
        absent = pd.DataFrame({"Specialist": pd.Series(dtype="str"),
                               "Specialty / Area": pd.Series(dtype="str"),
                               "Expected return": pd.Series(dtype="object")})
    absent_edit = st.data_editor(
        absent, num_rows="dynamic", use_container_width=True, key="absent_ed",
        column_config={
            "Specialist": st.column_config.TextColumn("Specialist", required=True),
            "Specialty / Area": st.column_config.SelectboxColumn(
                "Specialty / Area", options=SPECIALTY_AREAS),
            "Expected return": st.column_config.DateColumn(
                "Expected return", format="DD MMM YYYY")})

    notes = st.text_area("Notes (optional)", value=scalars["notes"] if scalars else "",
                         key="notes_ta")

    # Mortality register data-entry removed from this form; any previously saved
    # rows for the day are preserved untouched on save.
    mortality_edit = load_mortality(entry_date)

    # ── User-defined editable tables (create, rename, replicate, edit columns, add rows) ──
    st.markdown('<div class="section">Custom tables</div>', unsafe_allow_html=True)
    st.caption("Build your own tables: set a title, rename or add columns, add rows, and "
               "replicate a table to start a new one from the same layout. Each table "
               "detects its data types and can be shown as a chart (column, bar, "
               "stacked, line, scatter, funnel, pie, donut or KPI cards) with your "
               "own colors.")
    if st.button("➕ Add a new table"):
        add_custom_table()
        st.rerun()
    custom_edits = []   # (table_id, columns, edited_df)
    _ctables = get_custom_tables()
    for _t in _ctables:
        st.markdown(f'<div style="font-weight:700;margin:.4rem 0 .2rem;">{_t["title"]}</div>',
                    unsafe_allow_html=True)
        with st.expander(f"🧩 Customize this table — title, columns, replicate"):
            _nt = st.text_input("Table title", value=_t["title"], key=f"ttl_{_t['id']}")
            _nc = st.text_area("Columns — one per line (add a line to add a column; rename or "
                               "delete lines to rename/remove)", value="\n".join(_t["columns"]),
                               key=f"cols_{_t['id']}", height=150)
            bcol = st.columns(2)
            if bcol[0].button("💾 Save layout", key=f"savelo_{_t['id']}"):
                update_custom_table(_t["id"], _nt, _nc)
                st.rerun()
            if bcol[1].button("⧉ Replicate", key=f"rep_{_t['id']}"):
                replicate_custom_table(_t["id"])
                st.success("Replicated. A copy was added below.")
                st.rerun()
        # ── Delete this table (tick to confirm, then delete) ──
        _dcol = st.columns([3, 2])
        _delok = _dcol[0].checkbox("Confirm delete", key=f"delok_{_t['id']}")
        if _dcol[1].button("🗑️ Delete table", key=f"del_{_t['id']}",
                           disabled=not _delok, use_container_width=True):
            delete_custom_table(_t["id"])
            st.rerun()
        _rows = load_custom_rows(_t["id"], entry_date, _t["columns"])
        _ed = st.data_editor(_rows, num_rows="dynamic", use_container_width=True,
                             key=f"ced_{_t['id']}")
        _ed = clean_frame(_ed)                 # paste hygiene, applied live
        _ttypes = infer_column_types(_ed, _t["columns"])
        st.markdown(type_badges_html(_ttypes), unsafe_allow_html=True)

        # ── Excel-style formulas & totals ──
        with st.expander("🧮 Formulas & totals — calculate columns and add a "
                         "summary row"):
            _fcfg = table_display(_t)
            _fid = _t["id"]
            st.caption("**Formula columns** compute a value for each row from "
                       "your other columns. Write an expression using column "
                       "names, e.g. `Cases + Recovered`, `Deaths / Cases * 100`, "
                       "or `(A - B) / B * 100`. Leave blank for a normal column. "
                       "Supported: + − × ÷, brackets, and SUM / AVG / MIN / MAX / "
                       "COUNT / ROUND.")
            _formulas = dict(_fcfg.get("formulas") or {})
            _new_formulas = {}
            for _fc in _t["columns"]:
                _val = st.text_input(
                    f"ƒ  {_fc}", value=_formulas.get(_fc, ""),
                    key=f"formula_{_fid}_{_fc}",
                    placeholder="(normal column — no formula)")
                if str(_val).strip():
                    _new_formulas[_fc] = clean_text(_val)
            _sum_opts = [("", "None"), ("sum", "Sum (total)"),
                         ("avg", "Average"), ("min", "Minimum"),
                         ("max", "Maximum"), ("count", "Count")]
            _sum_keys = [k for k, _ in _sum_opts]
            _cur_sum = _fcfg.get("summary", "")
            _sum_sel = st.selectbox(
                "Summary row (aggregate shown under the table for every numeric "
                "column)", [l for _, l in _sum_opts],
                index=_sum_keys.index(_cur_sum) if _cur_sum in _sum_keys else 0,
                key=f"summary_{_fid}")
            _sum_key = _sum_keys[[l for _, l in _sum_opts].index(_sum_sel)]
            if st.button("💾 Save formulas & totals", key=f"save_formula_{_fid}"):
                _merged = dict(_fcfg)
                _merged["formulas"] = _new_formulas
                _merged["summary"] = _sum_key
                update_custom_table(_fid, display=_merged)
                st.session_state.auth_time = time.time()
                st.success("Saved. Formulas apply on the dashboards and in reports.")
                st.rerun()
            # live preview of formulas on the current edits
            _ed_calc = apply_formulas(_t, _ed)
            if (_new_formulas or _sum_key) and not _ed_calc.empty:
                st.caption("Preview with formulas applied:")
                _prev = _ed_calc.copy()
                _sr = summary_row({**_t, "display": {**_fcfg, "summary": _sum_key}},
                                  _ed_calc)
                if _sr:
                    _prev = pd.concat([_prev, pd.DataFrame([_sr])],
                                      ignore_index=True)
                st.dataframe(_prev, use_container_width=True, hide_index=True)

        # formulas are computed values — store the calculated frame for the day
        _ed = apply_formulas(_t, _ed)
        custom_edits.append((_t["id"], _t["columns"], _ed))

        with st.expander("🎨 Chart & colors — how this table shows on dashboards "
                         "and reports"):
            _dcfg = table_display(_t)
            _tid = _t["id"]
            _mode_keys = [k for k, _ in CUSTOM_MODES]
            _mode_lbls = [lbl for _, lbl in CUSTOM_MODES]
            _agg_keys = [k for k, _ in CUSTOM_AGGS]
            _agg_lbls = [lbl for _, lbl in CUSTOM_AGGS]
            _num_cols = [c for c in _t["columns"] if _ttypes.get(c) == "numeric"]
            r1 = st.columns(2)
            _mi = _mode_keys.index(_dcfg["mode"]) if _dcfg["mode"] in _mode_keys else 0
            _mode_sel = r1[0].selectbox(
                "Display as", _mode_lbls, index=_mi, key=f"cmode_{_tid}",
                help="How this table is represented on the Day, Week and Month "
                     "dashboards, the TV presentation and the PDF reports.")
            _ai = _agg_keys.index(_dcfg["agg"]) if _dcfg["agg"] in _agg_keys else 0
            _agg_sel = r1[1].selectbox(
                "Week / Month roll-up", _agg_lbls, index=_ai, key=f"cagg_{_tid}",
                help="How daily rows combine on the weekly and monthly dashboards "
                     "and reports.")
            _mode_key_sel = _mode_keys[_mode_lbls.index(_mode_sel)]
            _agg_key_sel = _agg_keys[_agg_lbls.index(_agg_sel)]
            r2 = st.columns(2)
            _cat_opts = ["(auto)", "(Date)"] + _t["columns"]
            _ci = _cat_opts.index(_dcfg["category"]) if _dcfg["category"] in _cat_opts else 0
            _cat_sel = r2[0].selectbox(
                "Category (bars / slices / X-axis)", _cat_opts, index=_ci,
                key=f"ccat_{_tid}",
                help="(auto) picks the first text column. (Date) charts a trend "
                     "across the days of the week or month.")
            _val_opts = _num_cols or list(_t["columns"])
            _val_def = [v for v in _dcfg["values"] if v in _val_opts] or _num_cols
            _vals_sel = r2[1].multiselect("Values (numeric columns)", _val_opts,
                                          default=_val_def, key=f"cvals_{_tid}")
            _pct_opts = _num_cols or list(_t["columns"])
            _auto_pct = [c for c in _pct_opts if _looks_pct(_ed, c)]
            _pct_def = [c for c in _dcfg.get("pct_cols", []) if c in _pct_opts] \
                or _auto_pct
            _pcts_sel = st.multiselect(
                "Percentage columns (%)", _pct_opts, default=_pct_def,
                key=f"cpct_{_tid}",
                help="Marked columns show a % sign on charts, axes and KPI cards "
                     "(dashboards and reports), and always roll up as AVERAGES on "
                     "Week/Month — percentages are never summed. Columns whose "
                     "cells end in % are suggested automatically.")
            if _mode_key_sel == "grouped":
                _sub_opts = ["(auto)"] + [c for c in _t["columns"]
                                          if c != _cat_sel]
                _si = (_sub_opts.index(_dcfg.get("subcategory"))
                       if _dcfg.get("subcategory") in _sub_opts else 0)
                _sub_sel = st.selectbox(
                    "Sub-category (inner cluster of the two-level axis)",
                    _sub_opts, index=_si, key=f"csub_{_tid}",
                    help="For the grouped chart: Category is the outer axis "
                         "level, this Sub-category clusters inside each group. "
                         "(auto) uses the next text column.")
            else:
                _sub_sel = _dcfg.get("subcategory", "(auto)")
            st.caption("Column colors — used for chart series, table headers and "
                       "cell tints:")
            _new_colors = {}
            for _cs in range(0, len(_t["columns"]), 4):
                _crow = st.columns(4)
                for _cc, _cname in zip(_crow, _t["columns"][_cs:_cs + 4]):
                    _cdefault = _dcfg["colors"].get(_cname) or \
                        CHART_PALETTE[_t["columns"].index(_cname) % len(CHART_PALETTE)]
                    _new_colors[_cname] = _cc.color_picker(
                        _cname, value=_cdefault, key=f"ccol_{_tid}_{_cname}")
            _tmp_t = {**_t, "display": {**_dcfg, "mode": _mode_key_sel,
                                        "agg": _agg_key_sel, "category": _cat_sel,
                                        "values": list(_vals_sel)}}
            _agg_now = custom_table_agg(_tmp_t, _ed)
            _cat_list = [str(c) for c in (_agg_now["cats"] if _agg_now else [])
                         if str(c).strip() and str(c) != "All"][:12]
            _new_cat_colors = dict(_dcfg.get("cat_colors") or {})
            _by_cat = bool(_dcfg.get("color_by_cat"))
            if _cat_list:
                _by_cat = st.checkbox(
                    "Color each bar / point individually by category",
                    value=_by_cat, key=f"cbycat_{_tid}",
                    help="With one value column, every bar or point takes its own "
                         "category colour below. Pie, donut and funnel always use "
                         "the category colours.")
                if _by_cat or _mode_key_sel in ("pie", "donut", "funnel"):
                    st.caption("Category colors — pick any colour for each slice, "
                               "stage, bar or point:")
                    for _gs in range(0, len(_cat_list), 4):
                        _grow = st.columns(4)
                        for _gc, _gname in zip(_grow, _cat_list[_gs:_gs + 4]):
                            _gdefault = _new_cat_colors.get(_gname) or \
                                CHART_PALETTE[_cat_list.index(_gname)
                                              % len(CHART_PALETTE)]
                            _new_cat_colors[_gname] = _gc.color_picker(
                                _gname, value=_gdefault,
                                key=f"ccolcat_{_tid}_{_gname}")
            r3 = st.columns(3)
            _use_stripe = r3[0].checkbox("Stripe alternate rows",
                                         value=bool(_dcfg["row_stripe"]),
                                         key=f"cstripe_on_{_tid}")
            _stripe = r3[1].color_picker("Stripe color",
                                         value=_dcfg["row_stripe"] or "#F0FAFA",
                                         key=f"cstripe_{_tid}")
            _tint = r3[2].checkbox("Tint cells with column colors",
                                   value=bool(_dcfg["col_tint"]),
                                   key=f"ctint_{_tid}")
            _pending = {"mode": _mode_key_sel, "agg": _agg_key_sel,
                        "category": _cat_sel, "values": list(_vals_sel),
                        "pct_cols": list(_pcts_sel),
                        "subcategory": _sub_sel,
                        "colors": _new_colors, "cat_colors": _new_cat_colors,
                        "color_by_cat": bool(_by_cat),
                        "row_stripe": (_stripe if _use_stripe else ""),
                        "col_tint": bool(_tint)}
            if st.button("💾 Save chart & colors", key=f"cdisp_save_{_tid}"):
                update_custom_table(_tid, display=_pending)
                st.success("Saved. Dashboards, TV mode and reports now use "
                           "this display.")
                st.rerun()
            _prev_t = {**_t, "display": _pending}
            if _pending["mode"] == "kpi":
                _pk = custom_table_kpis(_prev_t, _ed)
                if _pk:
                    st.caption("Preview:")
                    render_kpis([(f"📊 {l}", v) for l, v in _pk], False)
            elif _pending["mode"] != "table":
                if not render_interactive_chart_preview(_prev_t, _ed):
                    st.caption("The chart preview appears once the table has "
                               "chartable data (a category plus at least one "
                               "numeric column).")
            else:
                st.caption("Preview:")
                st.dataframe(style_custom_df(_prev_t, _ed),
                             use_container_width=True, hide_index=True)
    if not _ctables:
        st.caption("No custom tables yet — click “➕ Add a new table” to create one.")

    # ── Custom notes (free-form titled notes, saved with the day, in reports) ──
    st.markdown('<div class="section">📝 Custom notes</div>', unsafe_allow_html=True)
    st.caption("Add any number of titled notes for this day — context, incidents, "
               "decisions. They save with the day and appear in the PDF and "
               "PowerPoint reports.")
    _nkey = f"cnotes_{entry_date.isoformat()}"
    if _nkey not in st.session_state:
        _existing = get_daily_notes(entry_date)
        st.session_state[_nkey] = _existing or [
            {"id": pysecrets.token_hex(4), "title": "", "content": ""}]
    _notes_state = st.session_state[_nkey]
    _remove_idx = None
    for _i, _n in enumerate(_notes_state):
        with st.container(border=True):
            _hc = st.columns([5, 1])
            _n["title"] = _hc[0].text_input(
                "Note title", value=_n.get("title", ""),
                key=f"cnote_t_{_n['id']}", placeholder="e.g. Incident report")
            if _hc[1].button("🗑️", key=f"cnote_del_{_n['id']}",
                             help="Remove this note"):
                _remove_idx = _i
            _n["content"] = st.text_area(
                "Note", value=_n.get("content", ""), key=f"cnote_c_{_n['id']}",
                height=100, label_visibility="collapsed",
                placeholder="Write the note here…")
    if _remove_idx is not None:
        _notes_state.pop(_remove_idx)
        st.rerun()
    if st.button("➕ Add another note", key="add_cnote"):
        _notes_state.append({"id": pysecrets.token_hex(4), "title": "",
                             "content": ""})
        st.rerun()

    custom_notes_edit = _notes_state

    if st.button("💾 Save this day", type="primary", use_container_width=True):
        save_day(entry_date, numeric, notes, dept_edit, med_edit, test_edit, blood_edit,
                 absent_edit, mortality_edit)
        for _tid, _tcols, _ted in custom_edits:
            save_custom_rows(_tid, entry_date, _ted, _tcols)
        save_daily_notes(entry_date, custom_notes_edit)
        st.session_state.auth_time = time.time()   # keep the session alive on activity
        st.success(f"Saved entry for {entry_date:%A, %d %b %Y}.")
        st.balloons()
        st.caption("💡 Saved to the database. Because Streamlit Cloud's disk resets "
                   "on every app restart, download a backup when you finish "
                   "(**Settings → 💾 Backup & Restore**) so today's work can always "
                   "be brought back.")

    # download the saved data for the selected day (reloaded fresh from the database)
    st.markdown('<div class="section">Download this day</div>',
                unsafe_allow_html=True)
    render_day_downloads(entry_date, key_prefix="de")


# ══════════════════════════════════════════════
# PUBLIC DASHBOARD MODE
# ══════════════════════════════════════════════
elif mode == "Public Dashboard":
    tv = st.query_params.get("tv") == "1"
    if tv:
        st.markdown(TV_CSS, unsafe_allow_html=True)
        inject_tv_autosize()

    if tv:
        # ---- professional presentation header ----
        _v = st.query_params.get("v")
        view = {"w": "Week", "m": "Month"}.get(_v, "Day")
        _sel = _qp_date("d")
        if view.startswith("Week"):
            _ws, _we = week_bounds(_sel)
            _info = f"Week · {_ws:%d %b} – {_we:%d %b %Y}"
        elif view.startswith("Month"):
            _ms, _me = month_bounds(_sel)
            _info = f"Month · {_ms:%B %Y}"
        else:
            _info = f"Day · {_sel:%A, %d %b %Y}"
        # compact exit affordance, then a full-width branded header bar
        ec = st.columns([8.6, 1.4])
        with ec[1]:
            if st.button("✕ Exit", use_container_width=True):
                exit_tv()
        st.markdown(
            '<div class="tvhead">'
            f'<img class="tvlogo" src="data:image/png;base64,{HDM_LOGO_B64}"/>'
            f'<div class="tvname">{HOSPITAL_NAME}</div>'
            f'<div class="tvhead-meta">{_info}'
            f'<span class="tvhead-upd">updated {now_local():%H:%M}</span></div>'
            '</div>', unsafe_allow_html=True)
    else:
        st.markdown(
            '<div class="big-title">'
            f'<img src="data:image/png;base64,{HDM_LOGO_B64}" '
            'style="height:2.2em;width:auto;vertical-align:middle;'
            'margin-right:.5em;border-radius:6px;"/>'
            f'{get_dashboard_title()}</div>',
            unsafe_allow_html=True)
        view = st.radio("View", ["Day", "Week", "Month"],
                        horizontal=True, key="pub_view")

    # ──────────────────────────────────────────
    # SINGLE-DAY DASHBOARD
    # ──────────────────────────────────────────
    if view == "Day":
        if tv:
            day = _qp_date("d")
        else:
            day, _, _, _ = day_picker("dash_day", "Choose a day (calendar)")
            st.markdown(f'<div class="sub">Showing <b>{day:%A, %d %b %Y}</b> &nbsp;•&nbsp; '
                        f'updated {now_local():%H:%M}</div>', unsafe_allow_html=True)

        scalars, depts, meds, tests, blood, absent = load_day(day)
        mort = load_mortality(day)
        _custom_day = [(t, load_custom_rows(t["id"], day, t["columns"]))
                       for t in order_custom_tables(get_custom_tables())]
        _has_custom = any(not df.empty for _, df in _custom_day)
        if (not scalars and depts.empty and meds.empty and tests.empty
                and blood.empty and absent.empty and mort.empty and not _has_custom):
            st.info(f"No data recorded for {day:%A, %d %b %Y}. An analyst can add it in "
                    "the **Data Entry** view.")
            if not tv:
                render_observations("day", day.isoformat())
            st.stop()
        s = scalars or {k: 0 for k in FIELD_KEYS}

        # download everything entered for this specific day. These files can
        # include record-level detail (e.g. individual mortality cause/time), so
        # they require an analyst sign-in even on the public dashboard. Aggregate
        # dashboards stay public; downloads do not.
        if not tv and st.session_state.get("authed"):
            dl = st.columns(3)
            dl_button(dl[0], 
                "⬇️ Download this day's data (CSV)",
                day_csv_bytes(day, scalars, depts, meds, tests, blood, mort,
                              [(t["title"], df, t) for t, df in _custom_day]),
                file_name=f"hospital_{day:%Y%m%d}.csv", mime="text/csv",
                use_container_width=True)
            dl_button(dl[1], 
                "📄 Download this day's report (PDF)",
                build_day_pdf(day, scalars, depts, meds, tests, blood, HOSPITAL_NAME, mort,
                              [(t["title"], df, t) for t, df in _custom_day],
                              gather_general_notes("day", day.isoformat()),
                              note_scope="day", note_ref=day.isoformat(),
                              absent=absent),
                file_name=f"hospital_report_{day:%Y%m%d}.pdf", mime="application/pdf",
                use_container_width=True)
            if PPTX_OK:
                dl_button(dl[2], 
                    "📊 Download this day's slides (PPTX)",
                    build_day_pptx(day, scalars, depts, meds, tests, blood,
                                   HOSPITAL_NAME, mort,
                                   [(t["title"], df, t) for t, df in _custom_day],
                                   gather_general_notes("day", day.isoformat()),
                                   note_scope="day", note_ref=day.isoformat(),
                                   absent=absent),
                    file_name=f"hospital_slides_{day:%Y%m%d}.pptx",
                    mime=PPTX_MIME, use_container_width=True)
            else:
                dl[2].caption("Add `python-pptx` to requirements.txt to enable "
                              "PowerPoint downloads.")
        elif not tv and not st.session_state.get("authed"):
            st.caption("🔒 Sign in as an analyst (Data Entry) to download this "
                       "day's data and reports.")

        # ── Performance dashboard (always shown first) ──
        perf_summary = health_summary(*perf_inputs_single(s, depts, meds, tests, blood))
        perf_figs_list = performance_figs(perf_summary)
        if tv:
            perf_blocks = perf_figs_list
        else:
            render_performance_normal(perf_summary, perf_figs_list)
            perf_blocks = []

        # KPI cards
        st.markdown('<div class="section">At a Glance</div>', unsafe_allow_html=True)
        occ_txt = "—"
        if s.get("beds_total"):
            occ_txt = f"{(s['beds_total'] - s['beds_available']) / s['beds_total'] * 100:.0f}%"
        kpis = [
            ("🤒 Patients in hospital", int(s.get("current_inpatients", 0))),
            ("📥 New admissions", int(s.get("admitted", 0))),
            ("📤 Discharged", int(s.get("discharged", 0))),
            ("🚨 ER visits", int(s.get("er_visits", 0))),
            ("🔪 Surgeries", int(s.get("surgeries", 0))),
            ("👶 Births", int(s.get("births", 0))),
            ("🕊️ Stillbirths", int(s.get("stillbirths", 0))),
            ("⚰️ Mortality", int(s.get("deaths", 0))),
            ("🛏️ Bed occupancy", occ_txt),
            ("🧑‍⚕️ Doctors on duty", int(s.get("doctors", 0))),
            ("👩‍⚕️ Nurses on duty", int(s.get("nurses", 0))),
            ("🚑 Ambulances",
             f"{int(s.get('ambulances_available', 0))}/{int(s.get('ambulances_total', 0))}"),
            ("🫁 Oxygen supply", f"{int(s.get('oxygen_pct', 0))}%"),
        ]
        render_kpis(kpis, tv)

        day_blocks = []

        # Patient activity (horizontal bars — labels inside)
        cats = [("Admissions", "admitted"), ("Discharges", "discharged"),
                ("ER visits", "er_visits"), ("ICU patients", "icu_patients"),
                ("Surgeries", "surgeries"), ("Births", "births"),
                ("Stillbirths", "stillbirths"),
                ("Mortality", "deaths"), ("Referrals out", "referrals_out"),
                ("Referrals back", "referrals_back")]
        pdf = pd.DataFrame({"Metric": [c[0] for c in cats],
                            "Count": [int(s.get(c[1], 0)) for c in cats]})
        figp = px.bar(pdf, x="Count", y="Metric", orientation="h",
                      title="Patients today", color_discrete_sequence=[PRIMARY])
        figp.update_layout(showlegend=False, yaxis=dict(autorange="reversed"))
        day_blocks.append(("Patient Activity", style_fig(figp, h=420)))

        # Staff
        sdf = pd.DataFrame({"Role": ["Doctors", "Nurses", "Support", "Specialists"],
                            "Count": [int(s.get("doctors", 0)), int(s.get("nurses", 0)),
                                      int(s.get("support_staff", 0)),
                                      int(s.get("specialists_on_call", 0))]})
        figs = px.bar(sdf, x="Role", y="Count", title="Staff today",
                      color_discrete_sequence=[TEAL2])
        figs.update_layout(showlegend=False)
        day_blocks.append(("Staff on Duty", style_fig(figs, h=340)))

        # Oxygen gauge
        ox = int(s.get("oxygen_pct", 0))
        ox_color = OK_GREEN if ox >= 50 else WARN if ox >= 25 else DANGER
        figo = go.Figure(go.Indicator(
            mode="gauge+number", value=ox, number={"suffix": "%"},
            title={"text": "Oxygen supply level"},
            gauge={"axis": {"range": [0, 100]}, "bar": {"color": ox_color},
                   "steps": [{"range": [0, 25], "color": "#fde2e2"},
                             {"range": [25, 50], "color": "#fdf0db"},
                             {"range": [50, 100], "color": "#e3f5ea"}]}))
        figo.update_layout(paper_bgcolor="#FFFFFF", height=300,
                           margin=dict(l=20, r=20, t=60, b=10))
        day_blocks.append(("Critical Supplies", figo))

        # Blood bank
        if not blood.empty and blood["Units"].sum() > 0:
            b = blood.copy()
            b["low"] = b["Units"] < 5
            figb = px.bar(b, x="Blood Type", y="Units", color="low",
                          color_discrete_map={True: DANGER, False: PRIMARY},
                          title="Blood bank units by type (red = low, <5)")
            figb.update_layout(showlegend=False)
            day_blocks.append(("Blood Bank", style_fig(figb, h=340)))

        # Medications (availability status)
        if not meds.empty:
            day_blocks.append(("Medication Availability",
                               med_status_fig(list(zip(meds["Medication"], meds["Status"])))))

        if tv:
            if not depts.empty:
                day_blocks.append(("Department Status",
                                   dept_status_fig(list(zip(depts["Department"],
                                                            depts["Status"])))))
            if not tests.empty:
                day_blocks.append(("Tests",
                                   tests_fig(tests[tests["Available"]]["Test"].tolist(),
                                             tests[~tests["Available"]]["Test"].tolist())))
            _abs_fig = absent_fig(absent)
            if _abs_fig is not None:
                day_blocks.append(("Specialist Availability", _abs_fig))

        # Mortality register (by ward), shown on screen and in presentation mode
        _mort_fig = mortality_by_ward_fig(mort, "Mortality by ward")
        if _mort_fig is not None:
            day_blocks.append(("Mortality by Ward", _mort_fig))
        if tv:
            _mort_tbl = mortality_ward_table_fig(mort)
            if _mort_tbl is not None:
                day_blocks.append(("Mortality — ward detail", _mort_tbl))
        for _t, _cr in _custom_day:
            _blk = custom_table_block(_t, _cr)
            if _blk is not None:
                day_blocks.append(_blk)
        show_charts(perf_blocks + day_blocks, tv,
                    note_scope="day", note_ref=day.isoformat())

        if not tv:
            if not depts.empty:
                st.markdown('<div class="section">Department Status</div>',
                            unsafe_allow_html=True)
                dept_pills(list(zip(depts["Department"], depts["Status"])))
            st.markdown('<div class="section">Specialist Availability</div>',
                        unsafe_allow_html=True)
            render_absent_block(absent)
            _ms = mortality_stats(mort)
            if _ms["total"]:
                st.markdown('<div class="section">Mortality register</div>',
                            unsafe_allow_html=True)
                render_kpis([("⚰️ Deaths recorded", _ms["total"]),
                             ("🎂 Average age",
                              "—" if _ms["avg_age"] is None else _ms["avg_age"]),
                             ("🛏️ Avg length of stay (days)",
                              "—" if _ms["avg_los"] is None else _ms["avg_los"])], False)
                st.caption("Mortality by ward — count and ages recorded:")
                st.dataframe(mortality_ward_breakdown(mort),
                             use_container_width=True, hide_index=True)
                st.caption("Full register (with length of stay):")
                st.dataframe(mort, use_container_width=True, hide_index=True)
            if not tests.empty:
                st.markdown('<div class="section">Medical Tests Available</div>',
                            unsafe_allow_html=True)
                render_tests_block(tests[tests["Available"]]["Test"].tolist(),
                                   tests[~tests["Available"]]["Test"].tolist())
            render_observations("day", day.isoformat())

    # ──────────────────────────────────────────
    # WEEKLY ROLL-UP DASHBOARD
    # ──────────────────────────────────────────
    elif view.startswith("Week"):
        if tv:
            ref_day = _qp_date("d")
        else:
            ref_day = st.date_input("Show week containing", value=date.today(), key="wk_ref")
        start, end = week_bounds(ref_day)
        if not tv:
            st.markdown(f'<div class="sub">Week of <b>{start:%d %b}</b> – <b>{end:%d %b %Y}</b> '
                        f'&nbsp;•&nbsp; updated {now_local():%H:%M}</div>',
                        unsafe_allow_html=True)

        daily, depts, meds, tests, blood, absent = load_range(start, end)
        mort = load_mortality_range(start, end)
        if all(x.empty for x in (daily, depts, meds, tests, blood, absent)) and mort.empty:
            st.info("No data recorded for this week yet.")
            if not tv:
                render_observations("week", start.isoformat())
            st.stop()

        _custom_rng = [(t, load_custom_rows_range(t["id"], start, end, t["columns"]))
                       for t in order_custom_tables(get_custom_tables())]
        # download the rolled-up week (shown at the top, like the Day view)
        d = daily.copy()
        if not d.empty:
            d["Day"] = d["entry_date"].map(day_label)
        latest = d.iloc[-1] if not d.empty else None
        perf_summary = health_summary(*perf_inputs_range(latest, depts, meds, tests, blood))
        if not tv and st.session_state.get("authed"):
            render_period_downloads(daily, mort, _custom_rng, start, end, "week",
                                    f"weekly_summary_{start:%Y%m%d}.csv",
                                    f"weekly_report_{start:%Y%m%d}.pdf", "Weekly Report",
                                    perf_summary=perf_summary, depts=depts, meds=meds,
                                    tests=tests, blood=blood, absent=absent)
        elif not tv:
            st.caption("🔒 Sign in as an analyst (Data Entry) to download this "
                       "week's data and reports.")

        # ── Performance dashboard (always shown first) ──
        perf_figs_list = performance_figs(perf_summary)
        if tv:
            perf_blocks = perf_figs_list
        else:
            render_performance_normal(perf_summary, perf_figs_list)
            perf_blocks = []

        st.markdown('<div class="section">This Week at a Glance</div>', unsafe_allow_html=True)
        g = lambda col: int(d[col].sum()) if not d.empty else 0
        avg = lambda col: round(d[col].mean(), 1) if not d.empty else 0
        occ_txt = "—"
        if latest is not None and latest["beds_total"]:
            occ_txt = f"{(latest['beds_total']-latest['beds_available'])/latest['beds_total']*100:.0f}%"
        kpis = [
            ("🤒 Patients in hospital",
             int(latest["current_inpatients"]) if latest is not None else 0),
            ("📥 Admissions (wk)", f"{g('admitted'):,}"),
            ("📤 Discharged (wk)", f"{g('discharged'):,}"),
            ("🚨 ER visits (wk)", f"{g('er_visits'):,}"),
            ("🔪 Surgeries (wk)", f"{g('surgeries'):,}"),
            ("👶 Births (wk)", f"{g('births'):,}"),
            ("🕊️ Stillbirths (wk)", f"{g('stillbirths'):,}"),
            ("⚰️ Mortality (wk)", f"{g('deaths'):,}"),
            ("🛏️ Bed occupancy", occ_txt),
            ("🧑‍⚕️ Avg doctors/day", avg("doctors")),
            ("👩‍⚕️ Avg nurses/day", avg("nurses")),
            ("🚑 Avg ambulances", avg("ambulances_available")),
            ("📅 Days reported", f"{d['entry_date'].nunique() if not d.empty else 0}/7"),
        ]
        render_kpis(kpis, tv)

        wk_blocks = []
        if not d.empty:
            fig_p = go.Figure()
            fig_p.add_bar(x=d["Day"], y=d["admitted"], name="Admissions", marker_color=PRIMARY)
            fig_p.add_bar(x=d["Day"], y=d["discharged"], name="Discharges", marker_color=TEAL2)
            fig_p.update_layout(title="Admissions vs Discharges", barmode="group")
            wk_blocks.append(("Patient Flow", style_fig(fig_p)))

            fig_act = go.Figure()
            for col, name, color in [("er_visits", "ER visits", DANGER),
                                     ("surgeries", "Surgeries", PRIMARY),
                                     ("icu_patients", "ICU patients", WARN)]:
                fig_act.add_scatter(x=d["Day"], y=d[col], name=name, mode="lines+markers",
                                    line=dict(color=color))
            fig_act.update_layout(title="ER / Surgeries / ICU")
            wk_blocks.append(("ER / Surgeries / ICU", style_fig(fig_act)))

            fig_bed = go.Figure()
            fig_bed.add_bar(x=d["Day"], y=d["beds_available"], name="Beds available",
                            marker_color=TEAL2)
            fig_bed.add_scatter(x=d["Day"], y=d["beds_total"], name="Total beds",
                                mode="lines+markers", line=dict(color=INK, dash="dot"))
            fig_bed.update_layout(title="Beds Available vs Total")
            wk_blocks.append(("Beds & Occupancy", style_fig(fig_bed)))

            fig_s = go.Figure()
            fig_s.add_bar(x=d["Day"], y=d["doctors"], name="Doctors", marker_color=PRIMARY)
            fig_s.add_bar(x=d["Day"], y=d["nurses"], name="Nurses", marker_color=WARN)
            fig_s.add_bar(x=d["Day"], y=d["support_staff"], name="Support", marker_color=OK_GREEN)
            fig_s.update_layout(title="Staff on Duty", barmode="group")
            wk_blocks.append(("Staffing", style_fig(fig_s)))

            fig_a = go.Figure()
            fig_a.add_bar(x=d["Day"], y=d["ambulances_available"], name="Available",
                          marker_color=OK_GREEN)
            fig_a.add_scatter(x=d["Day"], y=d["ambulances_total"], name="Fleet total",
                              mode="lines+markers", line=dict(color=DANGER, dash="dot"))
            fig_a.update_layout(title="Ambulances Available vs Fleet")
            wk_blocks.append(("Ambulances", style_fig(fig_a)))

        if not blood.empty:
            latest_b = blood["entry_date"].max()
            b = blood[blood["entry_date"] == latest_b].copy()
            b["low"] = b["units"] < 5
            fig_bb = px.bar(b, x="blood_type", y="units", color="low",
                            color_discrete_map={True: DANGER, False: PRIMARY},
                            title=f"Units by type — {day_label(latest_b)} (red = low, <5)")
            fig_bb.update_layout(showlegend=False)
            wk_blocks.append((f"Blood Bank — as of {day_label(latest_b)}", style_fig(fig_bb)))

        latest_dept_status = None
        if not depts.empty:
            depts["Day"] = depts["entry_date"].map(day_label)
            depts["score"] = depts["status"].map(STATUS_SCORE).fillna(0)
            order = sorted(depts["entry_date"].unique())
            day_order = [day_label(x) for x in order]
            pivot = depts.pivot_table(index="name", columns="Day", values="score",
                                      aggfunc="last").reindex(columns=day_order)
            z = pivot.values
            _lbl = {2: "OK", 1: "Ltd", 0: "Closed"}
            txt = [["" if pd.isna(v) else _lbl[int(v)] for v in row] for row in z]
            fig_h = go.Figure(go.Heatmap(
                z=z, x=list(pivot.columns), y=list(pivot.index), text=txt, texttemplate="%{text}",
                colorscale=[[0, DANGER], [0.5, WARN], [1, OK_GREEN]],
                zmin=0, zmax=2, showscale=False, xgap=3, ygap=3))
            fig_h.update_layout(title="Operational status by day "
                                "(green = Operational, amber = Limited, red = Closed)")
            wk_blocks.append(("Department Status", style_fig(fig_h, h=max(300, 42 * pivot.shape[0]))))
            latest_d = depts[depts["entry_date"] == order[-1]]
            latest_dept_status = (day_label(order[-1]), list(zip(latest_d["name"], latest_d["status"])))

        if not meds.empty:
            latest_m = meds["entry_date"].max()
            m = meds[meds["entry_date"] == latest_m]
            wk_blocks.append((f"Medication Availability — as of {day_label(latest_m)}",
                              med_status_fig(list(zip(m["name"], m["status"])))))

        if tv and not tests.empty:
            latest_t = tests["entry_date"].max()
            t = tests[tests["entry_date"] == latest_t]
            wk_blocks.append(("Tests",
                              tests_fig(t[t["available"] == 1]["name"].tolist(),
                                        t[t["available"] == 0]["name"].tolist())))

        latest_absent = (absent[absent["entry_date"] == absent["entry_date"].max()]
                         if not absent.empty else absent)
        if tv and not latest_absent.empty:
            _abs_fig = absent_fig(latest_absent)
            if _abs_fig is not None:
                _abs_lbl = day_label(absent["entry_date"].max())
                wk_blocks.append((f"Specialist Availability — as of {_abs_lbl}", _abs_fig))

        _mort_fig = mortality_by_ward_fig(mort, "Mortality by ward (this week)")
        if _mort_fig is not None:
            wk_blocks.append(("Mortality by Ward", _mort_fig))
        if tv:
            _mort_tbl = mortality_ward_table_fig(mort, "Mortality by ward (this week)")
            if _mort_tbl is not None:
                wk_blocks.append(("Mortality — ward detail", _mort_tbl))
        for _t, _cr in _custom_rng:
            _blk = custom_table_block(_t, _cr, " (this week)")
            if _blk is not None:
                wk_blocks.append(_blk)
        _wk_trend = trend_fig(daily, [("current_inpatients", "Patients in hospital"),
                                      ("admitted", "New admissions"),
                                      ("deaths", "Mortality")], "Daily trend (this week)")
        if _wk_trend is not None:
            wk_blocks.append(("Daily Trend", _wk_trend))
        show_charts(perf_blocks + wk_blocks, tv,
                    note_scope="week", note_ref=start.isoformat())

        if not tv:
            render_rollup_normal(daily, mort, _custom_rng, tests, latest_dept_status,
                                 latest_absent, "week", "wk_trend_pick")
            render_observations("week", start.isoformat())

    # ──────────────────────────────────────────
    # MONTHLY ROLL-UP DASHBOARD
    # ──────────────────────────────────────────
    else:
        if tv:
            ref_day = _qp_date("d")
        else:
            ref_day = st.date_input("Show month containing", value=date.today(), key="mo_ref")
        start, end = month_bounds(ref_day)
        if not tv:
            st.markdown(f'<div class="sub">Month of <b>{start:%B %Y}</b> '
                        f'&nbsp;•&nbsp; updated {now_local():%H:%M}</div>',
                        unsafe_allow_html=True)

        daily, depts, meds, tests, blood, absent = load_range(start, end)
        mort = load_mortality_range(start, end)
        if all(x.empty for x in (daily, depts, meds, tests, blood, absent)) and mort.empty:
            st.info("No data recorded for this month yet.")
            if not tv:
                render_observations("month", f"{start:%Y-%m}")
            st.stop()

        _custom_rng = [(t, load_custom_rows_range(t["id"], start, end, t["columns"]))
                       for t in order_custom_tables(get_custom_tables())]
        # download the rolled-up month (shown at the top, like the Day view)
        mlabel = lambda s: datetime.fromisoformat(s).strftime("%d")
        d = daily.copy()
        if not d.empty:
            d["Day"] = d["entry_date"].map(mlabel)
        latest = d.iloc[-1] if not d.empty else None
        perf_summary = health_summary(*perf_inputs_range(latest, depts, meds, tests, blood))
        if not tv and st.session_state.get("authed"):
            render_period_downloads(daily, mort, _custom_rng, start, end, "month",
                                    f"monthly_summary_{start:%Y%m}.csv",
                                    f"monthly_report_{start:%Y%m}.pdf", "Monthly Report",
                                    perf_summary=perf_summary, depts=depts, meds=meds,
                                    tests=tests, blood=blood, absent=absent)
        elif not tv:
            st.caption("🔒 Sign in as an analyst (Data Entry) to download this "
                       "month's data and reports.")

        # ── Performance dashboard (always shown first) ──
        perf_figs_list = performance_figs(perf_summary)
        if tv:
            perf_blocks = perf_figs_list
        else:
            render_performance_normal(perf_summary, perf_figs_list)
            perf_blocks = []

        st.markdown('<div class="section">This Month at a Glance</div>', unsafe_allow_html=True)
        g = lambda col: int(d[col].sum()) if not d.empty else 0
        avg = lambda col: round(d[col].mean(), 1) if not d.empty else 0
        occ_txt = "—"
        if latest is not None and latest["beds_total"]:
            occ_txt = f"{(latest['beds_total']-latest['beds_available'])/latest['beds_total']*100:.0f}%"
        days_in_month = (end - start).days + 1
        kpis = [
            ("🤒 Patients in hospital",
             int(latest["current_inpatients"]) if latest is not None else 0),
            ("📥 Admissions (mo)", f"{g('admitted'):,}"),
            ("📤 Discharged (mo)", f"{g('discharged'):,}"),
            ("🚨 ER visits (mo)", f"{g('er_visits'):,}"),
            ("🔪 Surgeries (mo)", f"{g('surgeries'):,}"),
            ("👶 Births (mo)", f"{g('births'):,}"),
            ("🕊️ Stillbirths (mo)", f"{g('stillbirths'):,}"),
            ("⚰️ Mortality (mo)", f"{g('deaths'):,}"),
            ("🛏️ Bed occupancy", occ_txt),
            ("🧑‍⚕️ Avg doctors/day", avg("doctors")),
            ("👩‍⚕️ Avg nurses/day", avg("nurses")),
            ("🚑 Avg ambulances", avg("ambulances_available")),
            ("📅 Days reported",
             f"{d['entry_date'].nunique() if not d.empty else 0}/{days_in_month}"),
        ]
        render_kpis(kpis, tv)

        mo_blocks = []
        if not d.empty:
            fig_p = go.Figure()
            fig_p.add_bar(x=d["Day"], y=d["admitted"], name="Admissions", marker_color=PRIMARY)
            fig_p.add_bar(x=d["Day"], y=d["discharged"], name="Discharges", marker_color=TEAL2)
            fig_p.update_layout(title="Admissions vs Discharges", barmode="group")
            mo_blocks.append(("Patient Flow", style_fig(fig_p)))

            fig_act = go.Figure()
            for col, name, color in [("er_visits", "ER visits", DANGER),
                                     ("surgeries", "Surgeries", PRIMARY),
                                     ("icu_patients", "ICU patients", WARN)]:
                fig_act.add_scatter(x=d["Day"], y=d[col], name=name, mode="lines+markers",
                                    line=dict(color=color))
            fig_act.update_layout(title="ER / Surgeries / ICU")
            mo_blocks.append(("ER / Surgeries / ICU", style_fig(fig_act)))

            fig_bed = go.Figure()
            fig_bed.add_bar(x=d["Day"], y=d["beds_available"], name="Beds available",
                            marker_color=TEAL2)
            fig_bed.add_scatter(x=d["Day"], y=d["beds_total"], name="Total beds",
                                mode="lines+markers", line=dict(color=INK, dash="dot"))
            fig_bed.update_layout(title="Beds Available vs Total")
            mo_blocks.append(("Beds & Occupancy", style_fig(fig_bed)))

            fig_s = go.Figure()
            fig_s.add_bar(x=d["Day"], y=d["doctors"], name="Doctors", marker_color=PRIMARY)
            fig_s.add_bar(x=d["Day"], y=d["nurses"], name="Nurses", marker_color=WARN)
            fig_s.add_bar(x=d["Day"], y=d["support_staff"], name="Support", marker_color=OK_GREEN)
            fig_s.update_layout(title="Staff on Duty", barmode="group")
            mo_blocks.append(("Staffing", style_fig(fig_s)))

            fig_a = go.Figure()
            fig_a.add_bar(x=d["Day"], y=d["ambulances_available"], name="Available",
                          marker_color=OK_GREEN)
            fig_a.add_scatter(x=d["Day"], y=d["ambulances_total"], name="Fleet total",
                              mode="lines+markers", line=dict(color=DANGER, dash="dot"))
            fig_a.update_layout(title="Ambulances Available vs Fleet")
            mo_blocks.append(("Ambulances", style_fig(fig_a)))

        if not blood.empty:
            latest_b = blood["entry_date"].max()
            b = blood[blood["entry_date"] == latest_b].copy()
            b["low"] = b["units"] < 5
            fig_bb = px.bar(b, x="blood_type", y="units", color="low",
                            color_discrete_map={True: DANGER, False: PRIMARY},
                            title=f"Units by type — {day_label(latest_b)} (red = low, <5)")
            fig_bb.update_layout(showlegend=False)
            mo_blocks.append((f"Blood Bank — as of {day_label(latest_b)}", style_fig(fig_bb)))

        latest_dept_status = None
        if not depts.empty:
            depts["Day"] = depts["entry_date"].map(mlabel)
            depts["score"] = depts["status"].map(STATUS_SCORE).fillna(0)
            order = sorted(depts["entry_date"].unique())
            day_order = [mlabel(x) for x in order]
            pivot = depts.pivot_table(index="name", columns="Day", values="score",
                                      aggfunc="last").reindex(columns=day_order)
            z = pivot.values
            _lbl = {2: "OK", 1: "Ltd", 0: "Closed"}
            txt = [["" if pd.isna(v) else _lbl[int(v)] for v in row] for row in z]
            fig_h = go.Figure(go.Heatmap(
                z=z, x=list(pivot.columns), y=list(pivot.index), text=txt, texttemplate="%{text}",
                colorscale=[[0, DANGER], [0.5, WARN], [1, OK_GREEN]],
                zmin=0, zmax=2, showscale=False, xgap=2, ygap=3))
            fig_h.update_layout(title="Operational status by day of month "
                                "(green = Operational, amber = Limited, red = Closed)")
            mo_blocks.append(("Department Status",
                              style_fig(fig_h, h=max(300, 42 * pivot.shape[0]))))
            latest_d = depts[depts["entry_date"] == order[-1]]
            latest_dept_status = (day_label(order[-1]),
                                  list(zip(latest_d["name"], latest_d["status"])))

        if not meds.empty:
            latest_m = meds["entry_date"].max()
            m = meds[meds["entry_date"] == latest_m]
            mo_blocks.append((f"Medication Availability — as of {day_label(latest_m)}",
                              med_status_fig(list(zip(m["name"], m["status"])))))

        if tv and not tests.empty:
            latest_t = tests["entry_date"].max()
            t = tests[tests["entry_date"] == latest_t]
            mo_blocks.append(("Tests",
                              tests_fig(t[t["available"] == 1]["name"].tolist(),
                                        t[t["available"] == 0]["name"].tolist())))

        latest_absent = (absent[absent["entry_date"] == absent["entry_date"].max()]
                         if not absent.empty else absent)
        if tv and not latest_absent.empty:
            _abs_fig = absent_fig(latest_absent)
            if _abs_fig is not None:
                _abs_lbl = day_label(absent["entry_date"].max())
                mo_blocks.append((f"Specialist Availability — as of {_abs_lbl}", _abs_fig))

        _mort_fig = mortality_by_ward_fig(mort, "Mortality by ward (this month)")
        if _mort_fig is not None:
            mo_blocks.append(("Mortality by Ward", _mort_fig))
        if tv:
            _mort_tbl = mortality_ward_table_fig(mort, "Mortality by ward (this month)")
            if _mort_tbl is not None:
                mo_blocks.append(("Mortality — ward detail", _mort_tbl))
        for _t, _cr in _custom_rng:
            _blk = custom_table_block(_t, _cr, " (this month)")
            if _blk is not None:
                mo_blocks.append(_blk)
        _mo_trend = trend_fig(daily, [("current_inpatients", "Patients in hospital"),
                                      ("admitted", "New admissions"),
                                      ("deaths", "Mortality")], "Daily trend (this month)")
        if _mo_trend is not None:
            mo_blocks.append(("Daily Trend", _mo_trend))
        show_charts(perf_blocks + mo_blocks, tv,
                    note_scope="month", note_ref=f"{start:%Y-%m}")

        if not tv:
            render_rollup_normal(daily, mort, _custom_rng, tests, latest_dept_status,
                                 latest_absent, "month", "mo_trend_pick")
            render_observations("month", f"{start:%Y-%m}")
